#!/usr/bin/env python3
"""
Word文档编辑智能体 - 主启动脚本 (Fixed Version)
功能：上传、分析、修改Word文档
"""

from pathlib import Path
import asyncio
import os
import sys

# 添加父目录到路径，以便导入DrSai模块
sys.path.insert(0, str(Path(__file__).parent.parent.parent.parent.parent))

from drsai.modules.components.model_client import HepAIChatCompletionClient, ModelFamily
from drsai.modules.components.model_client.anthropic import (
    HepAIAnthropicChatCompletionClient,
    get_info,
    get_token_limit,
    _MODEL_INFO
)
from drsai.modules.components.model_context import DrSaiChatCompletionContext
from drsai.modules.agents.skills_agent import SkillAgent, DrSaiAssistant
from drsai.modules.managers.database import DatabaseManager
from drsai.modules.managers.messages import (
    TextMessage,
    FileInfo,
    FilesContent,
    FilesEvent,
)
from drsai.utils.utils import upload_to_hepai_filesystem
import base64
from typing import AsyncGenerator, Sequence
from autogen_agentchat.messages import BaseAgentEvent, BaseChatMessage, ToolCallSummaryMessage
from autogen_agentchat.base import Response
from autogen_core import CancellationToken

# Import document processing components
import sys
from pathlib import Path
sys.path.append(str(Path(__file__).parent))
try:
    from document_processor import DocumentProcessor, print_document_summary
    from document_skills.process_document_skill import DocumentProcessingSkill
    from document_skills.docx_template_skill import DocxTemplateSkill
    from document_skills.template_library_skill import TemplateLibrarySkill
    from document_skills.doc_to_docx_skill import DocToDocxSkill
    from document_skills.contract_review_skill import ContractReviewSkill
    DOCUMENT_PROCESSING_AVAILABLE = True
except ImportError as e:
    DOCUMENT_PROCESSING_AVAILABLE = False
    print(f"⚠️ Document processing components not available: {e}")
    print("Install with: pip install python-docx PyPDF2 python-pptx pandas openpyxl")

try:
    import docxtpl  # noqa: F401
    DOCX_TEMPLATE_JINJA_AVAILABLE = True
except ImportError:
    DOCX_TEMPLATE_JINJA_AVAILABLE = False
    print("⚠️ docxtpl not installed — jinja-mode template filling will return an error; bracket mode still works")

from dotenv import load_dotenv
load_dotenv()

# 工作目录设置
HERE = Path(__file__).parent

# ── Auto-install missing Python dependencies ──────────────────────────────
def _ensure_python_deps():
    """Install missing Python packages from requirements.txt at startup."""
    import subprocess
    req_file = HERE / "requirements.txt"
    if not req_file.exists():
        return
    try:
        result = subprocess.run(
            ["pip", "install", "-r", str(req_file), "--quiet"],
            capture_output=True, text=True, timeout=120
        )
        if result.returncode == 0:
            print("✅ Python dependencies installed from requirements.txt")
        else:
            print(f"⚠️ Some Python dependencies could not be installed: {result.stderr[:200]}")
    except Exception as e:
        print(f"⚠️ Could not run pip install: {e}")

_ensure_python_deps()

# ── Monkey-patch fix_and_parse_json to handle unescaped quotes ────────────
# The upstream fix_and_parse_json in drsai.utils.utils fails when LLMs put
# unescaped " (U+0022) inside JSON string values (common with Chinese text
# like 前两句"床前明月光"以月光). We wrap it to fall back to json_repair.
#
# IMPORTANT: Must patch BOTH the module attribute AND all modules that did
# `from drsai.utils.utils import fix_and_parse_json` (which creates a local
# reference that won't see module-level patches).
import drsai.utils.utils as _drsai_utils
_original_fix_and_parse_json = _drsai_utils.fix_and_parse_json

def _patched_fix_and_parse_json(json_str, debug=True):
    result = _original_fix_and_parse_json(json_str, debug=debug)
    if isinstance(result, str) and "[JSON" in result and "解析失败" in result:
        try:
            from json_repair import repair_json
            import json
            repaired = repair_json(json_str, return_objects=False)
            parsed = json.loads(repaired)
            if debug:
                print("[json_repair 修复成功]")
            return parsed
        except Exception as e:
            if debug:
                print(f"[json_repair 也失败] {e}")
    return result

# Patch the module attribute
_drsai_utils.fix_and_parse_json = _patched_fix_and_parse_json
# Patch the local reference in drsai_assistant (imported via `from ... import`)
import drsai.modules.agents.skills_agent.drsai_assistant as _drsai_assistant
_drsai_assistant.fix_and_parse_json = _patched_fix_and_parse_json

WORKSPACE = HERE / "workspace"
WORKSPACE.mkdir(parents=True, exist_ok=True)
WORKDIR = WORKSPACE / "runs"
WORKDIR.mkdir(parents=True, exist_ok=True)


# ── PPT skill paths ───────────────────────────────────────────────────────
# All `ppt_*_tool` functions resolve scripts / references through these
# constants so the LLM never needs to guess the on-disk layout. The directory
# uses kebab-case to match the SKILL.md frontmatter `name: ppt-polished-deck-collab`.
PPT_SKILL_ROOT = (
    HERE / "skills" / "presentation-skills" / "ppt-polished-deck-collab-traditional"
)
PPT_SCRIPTS_DIR = PPT_SKILL_ROOT / "scripts"
PPT_REFERENCES_DIR = PPT_SKILL_ROOT / "references"

# Template library is served by the deployed UI backend so that DocMaster and
# the 模板库 tab see the *same* catalog. Local disk under WORKSPACE/templates is
# no longer used by these four tools — outputs (filled docs, edits) still go
# to WORKDIR on this machine.
TEMPLATE_API_URL = os.environ.get(
    "DOCMASTER_TEMPLATE_API_URL",
    "https://opendrsai.ihep.ac.cn/api",
).rstrip("/")


def _template_api_envelope(resp_json: dict, key: str, default):
    """Server wraps skill output as {status, message, data:{...}}. Unwrap it
    back into the {success, message, <skill-fields>} shape the four tools
    used to return when they called TemplateLibrarySkill directly."""
    data = resp_json.get("data") or {}
    return data.get(key, default)


def _template_api_list(user_id, category=None, query=None):
    import urllib.parse, urllib.request, json as _json
    params = {}
    if user_id: params["user_id"] = user_id
    if category: params["category"] = category
    if query: params["query"] = query
    url = f"{TEMPLATE_API_URL}/docmaster/templates"
    if params:
        url += "?" + urllib.parse.urlencode(params)
    req = urllib.request.Request(url, method="GET")
    try:
        with urllib.request.urlopen(req, timeout=15) as r:
            body = _json.loads(r.read().decode("utf-8"))
    except Exception as exc:
        return {"success": False, "shared": [], "mine": [], "message": f"模板库服务调用失败: {exc}"}
    return {
        "success": bool(body.get("status", True)),
        "shared": _template_api_envelope(body, "shared", []),
        "mine": _template_api_envelope(body, "mine", []),
        "message": body.get("message", ""),
    }


def _template_api_download(template_id, source, user_id):
    """GET /docmaster/templates/file → bytes + filename, written to a temp .docx."""
    import urllib.parse, urllib.request, tempfile as _tf
    params = {"template_id": template_id, "source": source}
    if source == "mine" and user_id:
        params["user_id"] = user_id
    url = f"{TEMPLATE_API_URL}/docmaster/templates/file?" + urllib.parse.urlencode(params)
    req = urllib.request.Request(url, method="GET")
    try:
        with urllib.request.urlopen(req, timeout=30) as r:
            data = r.read()
    except Exception as exc:
        return None, f"下载模板失败: {exc}"
    tmp = _tf.NamedTemporaryFile(prefix="docmaster_tpl_", suffix=".docx", delete=False)
    try:
        tmp.write(data)
    finally:
        tmp.close()
    return tmp.name, None


def _template_api_save(source_path, user_id, name, description, category, tags, aliases, template_id):
    import urllib.request, json as _json, uuid
    # Hand-built multipart so we don't pull in `requests`.
    boundary = f"----docmaster{uuid.uuid4().hex}"
    crlf = b"\r\n"
    parts = []
    def _field(key, value):
        if value is None: return
        parts.append(f"--{boundary}".encode())
        parts.append(f'Content-Disposition: form-data; name="{key}"'.encode())
        parts.append(b"")
        parts.append(str(value).encode("utf-8"))
    _field("user_id", user_id)
    _field("name", name)
    _field("description", description or "")
    if category: _field("category", category)
    if tags: _field("tags", _json.dumps(list(tags), ensure_ascii=False))
    if aliases: _field("aliases", _json.dumps(list(aliases), ensure_ascii=False))
    if template_id: _field("template_id", template_id)
    # The .docx file
    try:
        with open(source_path, "rb") as f:
            file_bytes = f.read()
    except Exception as exc:
        return {"success": False, "message": f"读取模板文件失败: {exc}"}
    fname = os.path.basename(source_path) or "template.docx"
    parts.append(f"--{boundary}".encode())
    parts.append(
        f'Content-Disposition: form-data; name="file"; filename="{fname}"'.encode()
    )
    parts.append(b"Content-Type: application/vnd.openxmlformats-officedocument.wordprocessingml.document")
    parts.append(b"")
    parts.append(file_bytes)
    parts.append(f"--{boundary}--".encode())
    parts.append(b"")
    body = crlf.join(parts)
    req = urllib.request.Request(
        f"{TEMPLATE_API_URL}/docmaster/templates",
        data=body,
        method="POST",
        headers={"Content-Type": f"multipart/form-data; boundary={boundary}"},
    )
    try:
        with urllib.request.urlopen(req, timeout=60) as r:
            resp = _json.loads(r.read().decode("utf-8"))
    except Exception as exc:
        return {"success": False, "message": f"保存模板失败: {exc}"}
    return {
        "success": bool(resp.get("status", True)),
        "template_id": _template_api_envelope(resp, "template_id", None),
        "metadata": _template_api_envelope(resp, "metadata", None),
        "message": resp.get("message", ""),
    }


def _template_api_delete(template_id, user_id):
    import urllib.parse, urllib.request, json as _json
    url = (
        f"{TEMPLATE_API_URL}/docmaster/templates/{urllib.parse.quote(template_id, safe='')}"
        f"?{urllib.parse.urlencode({'user_id': user_id})}"
    )
    req = urllib.request.Request(url, method="DELETE")
    try:
        with urllib.request.urlopen(req, timeout=30) as r:
            resp = _json.loads(r.read().decode("utf-8"))
    except Exception as exc:
        return {"success": False, "message": f"删除模板失败: {exc}"}
    return {
        "success": bool(resp.get("status", True)),
        "removed_id": _template_api_envelope(resp, "removed_id", None),
        "message": resp.get("message", ""),
    }

# 支持的模型配置
llm_mode_config = {
    "deepseek-v4-flash(Fast)": "hepai/deepseek-v4-flash",
    "qwen3_30b": "hepai/qwen3_30b",
    "minimax-m2.7": "hepai/minimax-m2.7",
    "minimax-m2.7-highspeed": "hepai/minimax-m2.7-highspeed",
}

def _build_files_event_data(file_path: str, description: str) -> dict | None:
    """
    Build a serialized FilesEvent payload for a generated/edited file.

    Tries to upload via HepAI filesystem (URL method) first.
    Falls back to base64 encoding if the upload fails.

    Returns a dict (serialized ``FilesContent``) to be appended to the
    pending files-events side-channel, or *None* if both methods fail.
    """
    import mimetypes

    file_path_obj = Path(file_path)
    if not file_path_obj.exists():
        return None

    file_name = file_path_obj.name
    file_size = file_path_obj.stat().st_size
    mime_type = mimetypes.guess_type(file_path)[0] or "application/octet-stream"

    file_info = None

    # --- Primary: upload to HepAI filesystem for a URL ---
    try:
        file_obj = upload_to_hepai_filesystem(file_path=file_path)
        url = file_obj["url"]
        file_info = FileInfo(
            name=file_name,
            url=url,
            description=description,
            download_method="url",
            size=file_size,
            mime_type=mime_type,
            path=file_path,  # Store the file path for tracking in on_messages_stream
        )
        print(f"📤 File uploaded for FilesEvent (URL): {url}")
    except Exception as upload_err:
        print(f"⚠️ HepAI upload failed, falling back to base64: {upload_err}")

    # --- Fallback: base64 encode the file ---
    if file_info is None:
        try:
            with open(file_path, "rb") as f:
                encoded = base64.b64encode(f.read()).decode("utf-8")
            file_info = FileInfo(
                name=file_name,
                base64_content=encoded,
                description=description,
                download_method="base64",
                size=file_size,
                mime_type=mime_type,
                path=file_path,  # Store the file path for tracking in on_messages_stream
            )
            print(f"📦 File encoded for FilesEvent (base64): {file_name}")
        except Exception as b64_err:
            print(f"❌ base64 fallback also failed: {b64_err}")
            return None

    files_content = FilesContent(
        files=[file_info],
        title=file_name,
        description=description,
    )
    return files_content.model_dump()


class DocMasterAgent(DrSaiAssistant):
    """DrSaiAssistant subclass that emits FilesEvent for generated/edited documents.
    
    Uses filesystem-scanning approach: snapshots file mtimes before tool execution,
    then detects new/modified files after — works regardless of HOW files are changed
    (run_write, run_bash, python scripts, etc.).
    """

    # File extensions to track for FilesEvent registration
    # NOTE: 
    # - .xml is excluded because unpacked DOCX/PPTX/XLSX files create many
    #   intermediate XML files that are internal to Office documents
    TRACKED_EXTENSIONS = {
        '.docx', '.pdf', '.pptx', '.xlsx', '.csv',
        '.txt', '.md', '.rtf', '.tex', '.json', '.yaml', '.yml',
        '.ini', '.cfg', '.conf', '.log', '.html', '.htm', '.css', '.js',
    }
    
    # Directories to exclude from workspace scanning
    # NOTE: skills/ and configs/ are excluded because:
    # - skills/ contains copied skill files (SKILL.md etc.) from first_time_setup
    # - configs/ contains system config files (AGENTS.md, TOOLS.md, USER.md etc.)
    # These are internal DrSai files, not user documents
    EXCLUDED_DIRS = {'skills', 'configs', 'document_skills', 'scripts', '__pycache__', '.git'}

    def __init__(self, pending_files_events: list, **kwargs):
        super().__init__(**kwargs)
        self._pending_files_events = pending_files_events
        self._install_filesystem_tool_guards()

    # ---- filesystem-tool guards ------------------------------------------
    # The framework auto-registers run_bash / run_glob / run_read as basic
    # tools. Despite explicit prompt rules forbidding their use for template
    # lookup, the agent regularly reaches for run_glob to "find" a template
    # by name — bypassing the template library and sometimes picking up a
    # stale duplicate. We can't drop the basic tools without forking the
    # framework, so we wrap their underlying callables and intercept calls
    # whose arguments look like template-hunting, returning a directive
    # error that redirects to the right tool.

    # Signals that an argument is about template/DOCX lookup. Keep these
    # narrow on purpose — we do NOT want to break legitimate run_bash /
    # run_glob uses (logs, code search, generic project files).
    _TEMPLATE_HUNT_TOKENS = (
        ".docx",
        ".doc",
        "template",
        "模板",
        "合同",
        "contract",
        "workspace/templates",
        "/templates/",
    )

    @classmethod
    def _looks_like_template_hunt(cls, *args: str) -> bool:
        blob = " ".join(a for a in args if isinstance(a, str)).lower()
        # Explicit allowlist: reads under the PPT skill directory or any deck
        # workspace are legitimate. The template-hunt heuristic exists to stop
        # DOCX template lookups via run_glob, not to block PPT artifacts that
        # happen to live in a path named "template_audit" or "templates/".
        if (
            "skills/presentation-skills" in blob
            or "ppt-polished-deck-collab" in blob
            or "/decks/" in blob              # per-user deck workspaces
            or "validation/template_audit" in blob  # audit reports under decks
        ):
            return False
        return any(tok.lower() in blob for tok in cls._TEMPLATE_HUNT_TOKENS)

    @staticmethod
    def _template_redirect_message(tool_name: str, args_blob: str) -> str:
        return (
            f"{tool_name} blocked: the arguments ({args_blob[:160]}) look like "
            "a template / DOCX file lookup. Filesystem tools are NOT the "
            "template entry point — they bypass the template library and "
            "regularly pick up stale duplicates. Recover with: "
            "(1) if the user named a template, call "
            "get_template_path_tool(template_ref=<user's words>); "
            "(2) to browse what's available, call "
            "list_templates_tool(category=None, query=None); "
            "(3) if the user just uploaded a file, re-read the upload event "
            "for the absolute path. To inspect DOCX content, use "
            "extract_docx_content_tool, NOT run_read."
        )

    def _install_filesystem_tool_guards(self) -> None:
        import functools

        def _wrap(orig_func, tool_name: str, arg_extractor):
            @functools.wraps(orig_func)
            async def guarded(*args, **kwargs):
                hunt_args = arg_extractor(args, kwargs)
                if self._looks_like_template_hunt(*hunt_args):
                    return self._template_redirect_message(
                        tool_name, " | ".join(str(a) for a in hunt_args)
                    )
                return await orig_func(*args, **kwargs)
            return guarded

        # Extract the user-visible arg(s) we want to sniff per tool.
        extractors = {
            "run_glob": lambda args, kwargs: (
                kwargs.get("pattern") or (args[0] if args else ""),
                kwargs.get("search_path") or (args[1] if len(args) > 1 else ""),
            ),
            "run_bash": lambda args, kwargs: (
                kwargs.get("cmd") or kwargs.get("command")
                or (args[0] if args else ""),
            ),
            "run_read": lambda args, kwargs: (
                kwargs.get("path") or kwargs.get("file_path")
                or (args[0] if args else ""),
            ),
        }

        tools_list = getattr(self, "_tools", None) or []
        for tool in tools_list:
            name = getattr(tool, "name", None) or getattr(tool, "_name", None)
            if name in extractors and hasattr(tool, "_func"):
                tool._func = _wrap(tool._func, name, extractors[name])

    @staticmethod
    async def _execute_tool_call(tool_call, workbench, handoff_tools, agent_name, cancellation_token):
        """Override autogen's _execute_tool_call to repair broken JSON before parsing.
        
        Autogen does json.loads(tool_call.arguments) which fails on unescaped
        quotes in Chinese text. We repair the JSON first, then delegate to the
        parent implementation.
        """
        import json as _json
        try:
            _json.loads(tool_call.arguments)
        except _json.JSONDecodeError:
            try:
                from json_repair import repair_json
                tool_call.arguments = repair_json(tool_call.arguments, return_objects=False)
            except Exception:
                pass  # let autogen handle the error as usual
        # Call the parent (autogen) implementation with repaired arguments
        from autogen_agentchat.agents._assistant_agent import AssistantAgent
        return await AssistantAgent._execute_tool_call(
            tool_call, workbench, handoff_tools, agent_name, cancellation_token,
        )

    def _snapshot_workspace(self) -> dict[str, float]:
        """Snapshot all tracked files in the CURRENT USER's workspace only."""
        user_dir = self._user_profile_manager.work_dir
        snapshot: dict[str, float] = {}
        for ext in self.TRACKED_EXTENSIONS:
            for f in user_dir.rglob(f'*{ext}'):
                if f.is_file():
                    # Skip files in excluded directories (skills, configs, etc.)
                    if any(excluded in f.parts for excluded in self.EXCLUDED_DIRS):
                        continue
                    try:
                        snapshot[str(f)] = f.stat().st_mtime
                    except OSError:
                        pass
        return snapshot

    def _detect_changed_files(self, before: dict[str, float]) -> list[str]:
        """Compare current workspace state with a previous snapshot.
        Returns list of new or modified file paths."""
        after = self._snapshot_workspace()
        changed: list[str] = []
        for fpath, mtime in after.items():
            if fpath not in before or mtime > before[fpath]:
                changed.append(fpath)
        return changed

    async def on_messages_stream(
        self,
        messages: Sequence[BaseChatMessage],
        cancellation_token: CancellationToken,
    ) -> AsyncGenerator[BaseAgentEvent | BaseChatMessage | Response, None]:
        """
        Override on_messages_stream to detect file changes via filesystem scanning.
        Takes a snapshot before processing, then emits FilesEvent for any
        new/modified files after each Response.
        """
        # Snapshot workspace BEFORE any tool execution
        snapshot_before = self._snapshot_workspace()
        already_emitted: set[str] = set()

        async for event in super().on_messages_stream(messages, cancellation_token):
            yield event

            # Check for changed files after:
            # - ToolCallSummaryMessage: yielded after each tool execution round
            #   (parent unwraps Response from _process_model_result into chat_message)
            # - Response: yielded for final text response, paused state, etc.
            if isinstance(event, (ToolCallSummaryMessage, Response)):
                # Get files that will be emitted via _pending_files_events
                pending_files = set()
                for pe in self._pending_files_events:
                    files_list = pe.get('files', [{}])
                    for f in files_list:
                        fpath = f.get('path', '')  # Use 'path' field (not 'file_path')
                        if fpath:
                            pending_files.add(fpath)
                
                changed_files = self._detect_changed_files(snapshot_before)
                for fpath in changed_files:
                    # Skip if already in pending events (will be emitted below)
                    if fpath in pending_files:
                        print(f"⏭️ Skipping filesystem scan for pending file: {Path(fpath).name}")
                        continue
                    if fpath not in already_emitted:
                        desc = f"File created/modified: {Path(fpath).name}"
                        fe_data = _build_files_event_data(fpath, desc)
                        if fe_data:
                            print(f"📄 File event emitted (scanner): {Path(fpath).name}")
                            already_emitted.add(fpath)
                            yield FilesEvent(
                                content=FilesContent(**fe_data),
                                source=self.name,
                            )
                # Update snapshot for next round of tool calls
                snapshot_before = self._snapshot_workspace()

            # Drain pending file events from docx tools (edit_docx_tool, add_comment_tool, etc.)
            # Always emit these — each tool call produces a new version of the file,
            # so the UI should always get the latest version even if the same file
            # was emitted by a previous tool call in this stream.
            while self._pending_files_events:
                fe_data = self._pending_files_events.pop(0)
                files_list = fe_data.get('files', [{}])
                for f in files_list:
                    fpath = f.get('path', '')  # Use 'path' field (not 'file_path')
                    if fpath:
                        already_emitted.add(fpath)  # still track to avoid scanner duplicates
                        print(f"📄 File event emitted (pending): {Path(fpath).name}")
                        yield FilesEvent(
                            content=FilesContent(**fe_data),
                            source=self.name,
                        )
                        break  # Only emit once per fe_data


_HALLUCINATED_PATH_HINTS = (
    "/Users/",           # macOS home — agent invents this from training priors
    "C:\\",              # Windows path — same
    "C:/",
    "/Desktop/",         # common desktop/downloads guess regardless of root
    "/Downloads/",
    "/Documents/",
)


def _guard_docx_file_path(file_path, *, tool_label: str) -> dict | None:
    """Reject obviously-hallucinated DOCX paths before tool execution.

    The LLM regularly invents paths like `/Users/jerry/Desktop/<filename>.docx`
    when the user mentions a template by name. The generic "File not found"
    response gives it no recovery target, so it falls back to `run_bash` /
    `run_glob` to "find" the file — which either misses entirely or pulls a
    stale duplicate. This guard returns a directive error that names the
    exact recovery path (get_template_path_tool or the upload event).
    """
    import os
    if not file_path or not isinstance(file_path, str):
        return {
            "success": False,
            "error": "Missing file_path",
            "message": (
                f"{tool_label} requires an absolute file_path. If the user "
                "referred to a template by name, call get_template_path_tool "
                "first; if the user uploaded a file, use the absolute path "
                "from the upload event. Do NOT guess paths, and do NOT use "
                "run_bash / run_glob / run_read to search for the file."
            ),
        }
    looks_hallucinated = any(hint in file_path for hint in _HALLUCINATED_PATH_HINTS)
    if looks_hallucinated and not os.path.exists(file_path):
        return {
            "success": False,
            "error": "Hallucinated file path",
            "message": (
                f"{tool_label}: the path {file_path!r} does not exist on this "
                "system and looks invented (macOS/Windows-style or "
                "Desktop/Downloads/Documents). This server is Linux and user "
                "files live under the docmaster workspace. NEVER guess "
                "filesystem paths. To recover: "
                "(1) if the user named a template (e.g. \"用 X 模板\"), call "
                "get_template_path_tool(template_ref=<user's words>) — it "
                "returns the canonical absolute path; "
                "(2) if the user just uploaded a file, re-read the upload "
                "event in the conversation for the absolute path; "
                "(3) if neither applies, ask the user — do NOT use "
                "run_bash / run_glob / run_read to search the filesystem."
            ),
        }
    if not os.path.exists(file_path):
        return {
            "success": False,
            "error": "File not found",
            "message": (
                f"{tool_label}: no file at {file_path!r}. If you got this "
                "path from get_template_path_tool the catalog may be stale — "
                "call list_templates_tool then get_template_path_tool again. "
                "If from an upload event, re-check the absolute path from "
                "that event. Do NOT use run_bash / run_glob / run_read to "
                "search for the file."
            ),
        }
    return None


def _guard_template_path(template_path) -> dict | None:
    """Validate template_path before calling DocxTemplateSkill.

    Returns an error dict (to be returned directly to the agent) when the
    path is empty / relative / non-existent. Returns None when the path is
    acceptable. The error messages are intentionally directive — the LLM
    has a habit of falling back to run_glob / run_bash when a tool call
    fails, which can silently pick up the wrong copy of a template (e.g.
    a stray export under downloads/). Spelling out the recovery in the
    tool error lands in fresh attention and is followed reliably.
    """
    if not template_path or not isinstance(template_path, str):
        return {
            "success": False,
            "error": "Missing template_path",
            "message": (
                "template_path is required. If the user named a template, call "
                "get_template_path_tool first and use the absolute template_path "
                "it returns. Do NOT use run_glob / run_bash / run_read to find "
                "template files — those can pick the wrong copy of the template."
            ),
        }
    p = Path(template_path)
    if not p.is_absolute():
        return {
            "success": False,
            "error": "Relative template_path not accepted",
            "message": (
                f"template_path={template_path!r} is a relative path. Use the "
                "absolute path returned by get_template_path_tool (the same value "
                "you received earlier in this conversation — re-call "
                "get_template_path_tool with the template name if you lost it). "
                "Do NOT pass bare filenames like 'template.docx', and do NOT use "
                "run_glob / run_bash / run_read to search the filesystem for the "
                "template — it can find a stale duplicate under downloads/ and "
                "bypass the template library."
            ),
        }
    if not p.exists():
        return {
            "success": False,
            "error": "Template file not found",
            "message": (
                f"No file at {template_path}. If you got this path from "
                "get_template_path_tool, the catalog entry may be stale — call "
                "list_templates_tool then get_template_path_tool again. If the "
                "user uploaded the template earlier in the session, re-check the "
                "absolute path from the upload event. Do NOT use run_glob / "
                "run_bash / run_read to search for template files — those bypass "
                "the template library and can pick up an old export."
            ),
        }
    return None


def create_word_editor_agent(
        api_key: str|None = None, 
        thread_id: str|None = None, 
        user_id: str|None = None, 
        db_manager: DatabaseManager|None = None,
        default_config_name: str|None = None,
) -> DocMasterAgent:
    """
    创建Word文档编辑智能体
    
    Args:
        api_key: HepAI API密钥
        thread_id: 对话线程ID
        user_id: 用户ID
        db_manager: 数据库管理器
        default_config_name: 默认模型配置名称
    
    Returns:
        DocMasterAgent实例
    """
    
    def set_model_client(default_config_name: str|None = None):
        """设置模型客户端"""
        # Try different models if the default fails
        if default_config_name is None:
            default_config_name = "deepseek-v4-flash(Fast)"
        
        # List of models to try in order (fastest/lightest first)
        models_to_try = [
            "deepseek-v4-flash(Fast)",   # Default - fast and reliable
            "qwen3_30b",                  # Qwen 30B
            "minimax-m2.7-highspeed",    # Fast minimax
            "minimax-m2.7",              # Standard minimax
        ]
        
        # If specified model is in the list, try it first
        if default_config_name in llm_mode_config:
            models_to_try.insert(0, default_config_name)
        
        # Remove duplicates
        models_to_try = list(dict.fromkeys(models_to_try))
        
        # Try each model until one works
        for model_name in models_to_try:
            if model_name in llm_mode_config:
                llm_model = llm_mode_config[model_name]
                print(f"🔄 Trying model: {model_name} ({llm_model})")
                break
        else:
            # Fallback to default
            llm_model = "hepai/deepseek-v4-flash"
            model_name = "deepseek-v4-flash(Fast)"
            print(f"⚠️ Model not found in config, using default: {llm_model}")
        
        # Create model client with timeout and retry settings
        try:
            # minimax models use Anthropic API
            if llm_model.startswith("hepai/minimax"):
                model_client = HepAIAnthropicChatCompletionClient(
                    model=llm_model,
                    base_url="https://aiapi.ihep.ac.cn/apiv2/anthropic",
                    api_key=api_key or os.environ.get("HEPAI_API_KEY"),
                    model_info=_MODEL_INFO.get("claude-sonnet-4-5", _MODEL_INFO["claude-sonnet-4-5"]),
                    max_tokens=16000,
                    temperature=0.3,
                    timeout=120.0,
                    max_retries=2,
                )
            else:
                # deepseek-v4-flash and qwen3_30b use OpenAI-compatible API
                model_client = HepAIChatCompletionClient(
                    model=llm_model,
                    api_key=api_key or os.environ.get("HEPAI_API_KEY"),
                    base_url="https://aiapi.ihep.ac.cn/apiv2",
                    model_info={
                        "vision": False,
                        "function_calling": True,
                        "json_output": True,
                        "structured_output": False,
                        "family": ModelFamily.UNKNOWN,
                        "multiple_system_messages": True,
                        "token_model": "hepai/deepseek-v4-flash",
                    },
                    temperature=0.3,
                    max_tokens=16000,
                    timeout=120.0,
                    max_retries=2,
                )
            
            print(f"✅ Successfully created model client for {model_name}")
            return model_client
            
        except Exception as e:
            print(f"❌ Failed to create model client for {model_name}: {e}")
            print("🔄 Falling back to deepseek-v4-flash")
            
            # Fallback to deepseek-v4-flash
            return HepAIChatCompletionClient(
                model="hepai/deepseek-v4-flash",
                api_key=api_key or os.environ.get("HEPAI_API_KEY"),
                base_url="https://aiapi.ihep.ac.cn/apiv2",
                model_info={
                    "vision": False,
                    "function_calling": True,
                    "json_output": True,
                    "structured_output": False,
                    "family": ModelFamily.UNKNOWN,
                    "multiple_system_messages": True,
                    "token_model": "hepai/deepseek-v4-flash",
                },
                temperature=0.3,
                max_tokens=16000,
                timeout=120.0,
                max_retries=2,
            )

    # 子智能体配置 - DocMaster的助手
    SUB_AGENTS = {
        "doc_processor": {
            "type": "DrSaiAgent",
            "description": "DocMaster的文档处理助手，专门处理Word文档的分析、编辑和格式化",
            "tools": ["run_bash", "run_read", "run_write", "run_edit", "run_glob"],
            "prompt": """你是DocMaster的文档处理助手。你的任务是：
1. 分析用户上传的Word文档内容
2. 根据DocMaster的指令编辑和修改文档
3. 保持文档格式、样式和结构的一致性
4. 处理文档中的文本、表格、图片、超链接等元素
5. 协助生成高质量的修改结果

请使用专业的文档处理工具和技术来协助DocMaster完成工作。""",
        },
        "code_executor": {
            "type": "CodeExecutorAgent",
            "description": "DocMaster的代码执行助手，用于运行Python脚本处理Word文档",
            "tools": [],
            "prompt": "执行Python代码来协助DocMaster处理Word文档，包括文档解析、内容编辑、格式调整等复杂操作。",
        }
    }

    # 系统提示词 - 专注于Word文档处理
    SYSTEM = """你是 DocMaster，一个以 DOCX 和 PPTX 为核心的文档分析与编辑助手。

你的目标不是夸大能力，而是稳定、准确地理解用户意图，并选择最合适的工具完成任务。


【关键行为准则】
⚠️ 重要：当用户要求对文档进行多项修改（如"扩写"、"重写"、"添加多个章节"等）时，你必须：
1. **先分析**：使用 extract_docx_content_tool 查看文档当前结构
2. **再规划**：在脑子里规划好所有需要的编辑操作
3. **一次性执行**：将所有编辑放在一个 edit_docx_tool 调用中完成，不要分成多次调用
4. **不要中途汇报**：完成所有编辑后再向用户报告结果，不要在编辑过程中停下来询问

⚠️ 重要：工具返回后必须检查 `changes` 数组！
- 如果 `changes: []` 是**空数组**，说明**没有任何编辑被实际执行**！
- 此时绝对不能向用户谎报"已完成XX修改"
- 要如实告诉用户："工具返回成功但 changes 数组为空，可能是编辑条件不满足，请检查工具输出或尝试其他方法"
- 信任 `changes` 数组内容，不信任 `success: true` 单独判断（因为底层工具可能返回 true 但 changes 为空）

❌ 错误做法示例：
- 调用一次 edit_docx_tool，汇报一次，然后再调用第二次 → 这是浪费时间和打断用户
- "好的，我先添加标题..." → 调用工具 → "标题已添加，接下来..." → 再调用 → "现在让我..."
- 分5次调用 edit_docx_tool，每次只做一个修改

✅ 正确做法示例：
- 调用一次 edit_docx_tool，edits=[{所有需要的修改}], 一次性完成所有工作
- "我来为您扩写文档，将一次性完成所有修改..."
- 调用 edit_docx_tool → 完成后直接向用户展示最终结果

【能力边界】
1. 你可以分析多种文档格式：DOCX、PDF、PPTX、XLSX、CSV、TXT、MD。
2. 你主要支持对 DOCX 文件进行编辑。
3. 你最擅长的 DOCX 操作包括：
   - 提取段落和表格内容
   - 创建结构化新文档
   - 添加标题、段落、表格
   - 添加项目符号列表和编号列表（支持嵌套层级）
   - 执行结构化文本修改与替换
   - 修改样式和字体
   - 删除文档内容
   - 添加批注和回复（使用 add_comment_tool，针对文档中的特定文本添加评论）
   - 删除批注（使用 remove_comment_tool，根据批注ID删除指定的批注）
   - 合同审查（使用 review_contract_tool，对 .docx 合同做格式 / 填写 / 一致性 / 法律风险四方面体检，可同时输出带批注的副本）
4. 对图片、超链接、页眉页脚、复杂版式重排等高级 Word 元素，不要假装已经可靠支持；如果用户提出这类需求，可以先说明当前能力更适合文本、标题、段落、表格、批注和字体层面的处理。
5. 你支持以 DOCX 模板填充方式批量生成文档：用户上传一个带占位符的 .docx，你可以读取占位符并按其提供的值生成填充后的新文档。占位符支持两种风格：Jinja 风格（{{ name }}、{% for x in xs %}…{% endfor %}、{%tr for %} 表格行循环）以及方括号风格（[NAME]、[DATE]）。
6. 你支持以 PPTX 模板填充的方式生成演示文稿：用户上传一个带文字占位符的.pptx文件，你可以读取文字占位符，删除原值后按照提供的值填充内容，形成新的演示文档。

【核心工作原则】
1. 先判断任务类型，再选择工具。
2. 如果用户没有提供文件路径、文件内容或明确目标，不要猜，先提出一个简短的澄清问题。
3. 对于非简单替换类编辑请求，优先先检查文档内容或结构，再执行修改。
4. 不要声称已经完成工具未实际执行的操作。
5. 若任务超出当前工具能力，要明确说明限制，并给出最接近的可执行方案。
6. 编辑操作默认会直接覆盖原始 DOCX 文件，必要时应提醒用户这一点。
7. 信任工具返回结果。如果 edit_docx_tool 或 create_docx_with_content_tool 返回了 success: true，任务就已经完成了——不要再用 run_read 去验证，不要再用子智能体去重做，不要再用其他方式重试。直接向用户汇报结果。
8. 不要使用 run_read 来读取 DOCX 文件内容。读取文档内容必须使用 extract_docx_content_tool。run_read 只在 XML 编辑工作流中用于读取已解包的 XML 文件。

【收到用户请求后的标准流程】
第一步：判断任务属于哪一类：
- 文档分析
- DOCX 内容检查
- 新建 DOCX
- 修改现有 DOCX
- 字体调整
- 清空文档内容
- 仅提供建议或说明
- 生成演示文档
- 修改现有演示文档
- 根据用户上传的演示文档模板生成新的演示文档

第二步：判断是否具备执行条件：
- 如果用户提到“这个文档/这份文件”，但没有给出文件路径或可识别文件，就先询问文件。
- 如果用户要求修改现有 DOCX，但没有说明改哪里，先询问目标段落、目标文本，或先读取文档内容。
- 如果用户要求“润色/改写/更专业/更简洁”这类语义编辑，不要直接盲改；应先查看相关内容，再生成修改方案或执行编辑。
- 如果用户要求新建文档但没有给出内容，也要先确认要写入什么。
- 如果用户提到演示文档 / PPT / 汇报 deck，按下方【PPT / 演示文档任务的标准流程】走，**绝不要**自己拼 skills 目录路径。

第三步：选择工具：
- 分析上传或给定文件：使用 process_document
- 检查 DOCX 实际内容：使用 extract_docx_content_tool
- 创建新 DOCX：使用 create_docx_with_content_tool
- 修改现有 DOCX：使用 edit_docx_tool
- 修改字体：使用 modify_docx_fonts_tool
- 删除全部内容：使用 delete_docx_content_tool
- 生成 / 修改 / 验证演示文档（pptx）：使用下方列出的 `ppt_*_tool` 系列工具。

【PPT / 演示文档任务的标准流程】
所有 PPT 相关任务都遵循 `ppt-polished-deck-collab` 这个 skill 的主链路。**所有
skill 内容（references / scripts）都通过 `ppt_*_tool` 工具访问，绝不要用
run_bash / run_read / run_glob / run_write 去执行 `skills/presentation-skills/...`
下任何 .py 脚本或读 .md 文档。**

标准顺序（按这个顺序走，每步停下来确认上一步真的成功）：

0. **澄清** — 如果用户的请求很模糊（没说目标读者、页数、是否有参考 pptx），
   先简短澄清一两个最关键的问题，再开始建 workspace。
1. **环境探测** — `ppt_check_environment_tool(deck_workspace=None)`。
   查 routes 数组：若缺 `editable_pptx`，直接如实告诉用户当前环境不支持，
   不要硬上；缺 `preview_powerpoint` / `preview_libreoffice` 时告知预览会受限。
2. **建 workspace** — `ppt_init_workspace_tool(deck_slug, deck_title)`。
   · deck_slug 用 kebab-case（a-z 0-9 -，例如 "ihep-2026-safety"）。
   · 工具返回 `deck_workspace`、`brief_path`、`narrative_path`。
   · **后续所有 PPT 工具的 deck_workspace 参数都必须用工具返回的这个绝对路径**，
     不允许自己拼字符串。
3. **填两份主文档** — `brief.md` 与 `deck_narrative.md`。这是 .md 不是 .docx，
   直接用 run_write 落盘（run_write 在 skills/presentation-skills 路径下是
   白名单放行的；写到 deck_workspace 下也没事）。
   · 先写 brief.md：目标读者、主使用场景、目标动作、模板/品牌约束、验证要求、
     免责声明 / 风险边界。
   · 再写 deck_narrative.md：保留工具产生的 YAML frontmatter（theme_tokens 已
     按 zh_formal 中文宋体 / 英文 Times New Roman 默认值填好），主体按
     `### S01 | <title>` 分页，每页含一个 ```yaml slide_spec``` 代码块，
     代码块必须包含 8 个字段：title / reader_question / page_task /
     reading_mode / archetype / asset_mode / validation_mode / key_message。

【slide_spec 字段契约（build 工具要用到的字段）】
derive 工具只校验 8 个必填字段；其余字段是 build 工具的"画页所需输入"，
agent 在写 deck_narrative.md 时要按页面类型决定要不要附加，不要漏掉：

  - `archetype` 必须是这 8 个之一：
      hero-statement / decision-logic / board-memo /
      chart-spotlight / comparison-matrix / process-flow /
      research-note / appendix-dense
    （未知值会 fallback 到 hero-statement 并标 warning。）

  - `asset_mode` 决定 chart-spotlight 等页里资产走哪条路：
      office-chart-native / python-figure-image / table-native /
      diagram-connector / diagram-visual / text-layout-native /
      icon-accent / image-hero / mixed

  - `bullets: [str]` — 决策点 / 解释要点 / 步骤名 / takeaway 句。
    · decision-logic 用作判断要点；board-memo 用前 4 项填入 2x2 panel；
      chart-spotlight 用作右侧 Takeaways；process-flow 在没有 diagram
      字段时被当作步骤名直接生成节点链；research-note 用作注释。

  - `chart` （chart-spotlight 且 asset_mode=office-chart-native 时必填）：
    ```yaml
    chart:
      title: "Coverage by phase"
      chart_type: bar           # bar|column|line|stacked_bar|stacked_column
      categories: ["Phase 1", "Phase 2", "Phase 3"]
      series:
        - {name: "Coverage", values: [92, 88, 73]}
      number_format: "0"        # 可选
      show_legend: false        # 可选
    ```

  - `image_path` （chart-spotlight + python-figure-image，或 research-note）：
    绝对路径，或相对于 deck_workspace 的路径（推荐放 assets/charts/ 或
    build/rendered/python_figures/）。build 会自动解析。

  - `table` （comparison-matrix 或 appendix-dense 时填一个就行；
              也可在 comparison-matrix 用 `matrix` 替代，见下）：
    ```yaml
    table:
      headers: ["项", "Q1", "Q2", "Q3"]
      rows:
        - ["营收", 120, 138, 142]
      numeric_columns: [1, 2, 3]   # 0-indexed；这些列右对齐
    ```
    表格自动套用 theme_tokens 中的 table policy（10.5pt、单倍行距、
    上下居中、表头居中、文本列居左、数值列右对齐）。

  - `matrix` （仅 comparison-matrix 的另一种写法，更口语化）：
    ```yaml
    matrix:
      - label: "方案 A"
        attributes: {成本: "低", 速度: "中", 可维护性: "高"}
      - label: "方案 B"
        attributes: {成本: "中", 速度: "高", 可维护性: "中"}
    ```
    build 会把 matrix 折叠成 table，attribute 的 key 合并为列头。

  - `diagram` （process-flow 想精确控制节点位置时填；否则用 bullets）：
    ```yaml
    diagram:
      nodes:
        - {key: "n1", text: "采集", left: 0.7, top: 4.0, width: 2.4, height: 1.2}
        - {key: "n2", text: "清洗", left: 4.0, top: 4.0, width: 2.4, height: 1.2}
      edges:
        - {from: "n1", to: "n2", from_site: "right", to_site: "left"}
    ```
    from_site / to_site 限定在 top / left / bottom / right。
    用 diagram 时配 `validation_mode: diagram_connector`，build 后再跑
    ppt_connectors_check_tool 校验真绑定。

  - `caption: str` — 页脚 caption（可选）。
  - `notes` 或 `narrative_markdown` — research-note / decision-logic
    右侧或左侧的说明文字（可选，超过 1200 字会截断）。

字段缺失时 build 会画占位 panel 并在 per_slide 里返回 warning/error，不会
默默编内容——这是故意的，agent 看到 warning 要么补字段，要么向用户确认。

【在 deck_narrative.md 里写 YAML 时的安全规则（必须遵守，否则 derive 会失败）】
agent 用 run_write 给 deck_narrative.md 写内容时，frontmatter 之外还会有大量
```yaml slide_spec``` 代码块。**避开下面这些 YAML 解析地雷**：

  · **字符串内嵌双引号不要再用双引号包**：`title: "Q2 \"Safety\" Report"`
    会触发 ParserError——内嵌的 `"` 直接关闭了外层 quote。
    正确写法：`title: 'Q2 "Safety" Report'`（外层用单引号），或者
    `title: Q2 "Safety" Report`（不用引号，YAML 会把整行 trim 后当字符串）。

  · **字符串带反斜杠 `\\` 一律用单引号或不用引号**：`title: "Report\\path"`
    会让 PyYAML 报 "unknown escape character 'p'"。
    正确：`title: 'Report\\path'` 或 `title: Report\\path`。
    单引号 scalar 里 `\\` 是普通字符，不当转义；双引号 scalar 里 `\\` 是
    转义引导符（只有 `\\n`、`\\t`、`\\\\`、`\\"` 等有限几种合法）。

  · **字符串带真实换行**直接换成 `|` 或 `>` 块标量，或者把换行替成空格：
    ```yaml
    notes: |
      第一行内容
      第二行内容
    ```
    不要把多行内容塞进单行双引号里。

  · **冒号 `:` 在双引号字符串中没问题，但在不加引号的 scalar 中放在词中间会被当 mapping 起点**：
    `title: Q2: Safety` → 解析成 mapping，会失败。
    `title: 'Q2: Safety'` 或 `title: "Q2: Safety"` 都安全。

  · **`---` 出现在字符串值的开头**（如 `summary: --- placeholder`）会被当成
    新一段文档的开始。要么改用引号，要么避开行首的 `---`。

  · **保守起见：字符串值统一用单引号**。除非你 100% 确认值里不含 `'`，
    遇到 `'` 时改用双引号（同时保证值里没有 `"` 和 `\\`），或退而用块标量 `|`。
    数字 / bool / 列表 / 嵌套 mapping 保持原样。

  · derive 工具失败时 stderr 通常会指出哪一行 YAML 解析失败——把这行文本
    念给用户，问他们是想改 title 文案、还是确认要用块标量。**不要**在没看
    懂错误前盲目重写整段 YAML。

4. **（仅当有参考 pptx）模板取证** — `ppt_audit_template_tool(pptx_path,
   deck_workspace)`。审计结果写入 validation/template_audit/，把字号梯度、
   layout 家族、共享母版元素结论回填到 brief.md。
5. **派生 slide_specs** — `ppt_derive_slide_specs_tool(deck_workspace)`。默认
   读 deck_narrative.md，写 build/generated/slide_specs.yaml。若返回 "missing
   field" 类错误，把 stderr_tail 念给用户，他们的 narrative 缺字段，先修再跑。
6. **workspace 体检** — `ppt_lint_workspace_tool(deck_workspace)`。确认目录、
   两份主文档、派生 specs 都齐了，缺什么补什么。
7. **构建可编辑 pptx** — `ppt_build_pptx_tool(slide_specs_path, output_pptx,
   deck_workspace)`。默认 `output_pptx` 用 `<deck_workspace>/build/pptx/deck_v1.pptx`。
   工具会按每页的 `archetype` 路由到对应 renderer，自动注入 narrative 中的
   `theme_tokens`（中文宋体 / 英文 Times、字号梯度、行距），并把生成的 pptx 通过
   FilesEvent 推给前端。
   · **必须先跑 `ppt_derive_slide_specs_tool`**，不能拿手写的 slide_specs.yaml
     直接喂给 build——derive 会校验 8 个必填字段，build 不再重做。
   · 工具返回里有 `per_slide` 数组。任何 `status != "ok"` 的页要把 `error` 念给
     用户并问"是修 narrative 再 rebuild，还是接受当前结果跑质量 gate？"
   · build 完之后**立刻**继续走第 8-10 步的质量 gate / 预览，不要在这里停下来
     向用户征求 build 是否成功的确认——`success` 字段已经告诉你了。
8. **三道 deck 级质量 gate**（build 出 pptx 之后立刻按顺序跑）：
   a) `ppt_package_preflight_tool(pptx_path, deck_workspace)` — 文件级。
   b) `ppt_structure_precheck_tool(pptx_path, deck_workspace)` — 结构层。
   c) 如有 diagram 页，再跑 `ppt_connectors_check_tool(pptx_path,
      deck_workspace, slides=[...], min_connectors=N)`。
   每道 gate 的报告都自动归档到 validation/<gate>/history/<gate>_<timestamp>.{json,md}，
   并随结果返回 `summary` 让你判断是否要继续。
9. **导出预览** — `ppt_export_previews_tool(pptx_path, deck_workspace,
   backend='auto')`。PNG 落到 build/rendered/ppt_preview/slide_NNN.png。
   页数不一致会直接 failure（不静默降级）；失败时换 backend 重试。
10. **成图层 gate** — `ppt_render_review_tool(pptx_path, deck_workspace)`。
    检查边界触墨与扁平化图像内部文字风险，必须在 export_previews 成功之后跑。
11. **first-draft checkpoint** — 主动停一次。把 validation 报告摘要 + 预览图
    路径给用户，问"要进入详细修订吗？如果不需要就交付当前初稿"。**不要**无限
    自我打磨。

【icon 资产（可选）】
- 用 `ppt_icon_search_tool(query, pack=None)` 查 icon。pack 可选
  "general-layout" 或 "llm-research"。
- 用 `ppt_icon_render_tool(deck_workspace, pack=..., color_mode='auto',
  background_color=..., accent_color=...)` 把 SVG 渲染成 deck-aware 的 PNG，
  落到 deck_workspace/assets/icons/<pack>/<theme>/。
- icon 是补充资产，永远不是主信息载体；当页面核心是趋势 / 比较 / 流程 /
  机制 / 架构 / 证据时不要用 icon 替代图表或语言本体。

【需要查 PPT 方法论时】
调 `ppt_read_skill_reference_tool(name=...)`，name 限定为以下 11 个：
  principles / deck_workflow / technical_support / design_support /
  slide_design_system / quality_gates / build_routes /
  diagram_support / office_chart_support / python_figure_support / icon_system
**不要**用 run_read 去读 references/*.md 文件。

【PPT 任务的 don't】
- 不要 run_bash 跑 `scripts/*.py` 里的任何脚本——一律走 `ppt_*` 工具。
- 不要绕过 init_workspace，自己在 user workdir 下手建 brief.md / deck_narrative.md。
- 不要凭印象拼接 `skills/presentation-skills/...` 路径——全部经由 `ppt_*` 工具。
- 不要在 slide_specs 派生失败时硬继续——先告诉用户 narrative 哪里缺字段。
- 不要在质量 gate 报 error 时直接进 preview——先看 summary 决定是否回头修。
- **不要在 ppt_build_pptx_tool 之后自己写 python-pptx 代码"微调"页面**——
  改 narrative + slide_specs 字段然后 rebuild，比手改 .pptx 更稳。如果某一页
  真的需要 build 工具不支持的版式，告诉用户当前 build 工具的能力边界，让
  用户决定是否手工编辑导出的 pptx。
- 不要无限自我打磨；初稿就绪后停一次，让用户决定是否进入详细修订。



【工具选择规则】
1. 如果用户只是想“了解文档是什么”，优先用 process_document。
2. 如果用户要查看 DOCX 里的实际段落、表格或目标文本，使用 extract_docx_content_tool。
3. 如果用户说“添加项目符号列表”“添加编号列表”“列出以下要点”，使用 add_bullet_list_tool 或 add_numbered_list_tool。
4. 如果用户说“把 A 改成 B”“在末尾增加一段”“插入标题”“添加表格”，统一使用 edit_docx_tool。
5. 如果用户说“把中文改成宋体、英文改成 Times New Roman”，使用 modify_docx_fonts_tool。
6. 如果用户说“重写引言/缩短结论/让措辞更正式”，先用 extract_docx_content_tool 查看内容，再进行后续编辑。
7. 如果用户要新建文档，使用 create_docx_with_content_tool。
8. 如果用户只是咨询写作或格式建议，不必强行调用工具。
9. 如果用户希望"填空"/"按模板生成"新文档：
   - **第零步（先查模板库）**：在要求用户上传文件之前，先看模板库里有没有现成的。
     · 用户说"用 X 模板""用 3-1 合同模板""用我的 XX 模板"——把用户的原话当作 template_ref 传给 `get_template_path_tool(template_ref)`。如果返回 success=True，拿到的 template_path 直接进入第一步，不要再要求上传。
     · 用户问"我有哪些模板""现在能用哪些合同模板"或没有具体指向——调 `list_templates_tool(category=None, query=None)` 把结果（共享 + 我的）念给用户挑。
     · `get_template_path_tool` 返回 ambiguous=True 时，**不要**自己挑——把 candidates 念给用户，让用户从中确认一个，再用确认后的 id 再调一次。
     · 模板库里查无匹配（success=False & ambiguous=False）才回退到要求用户上传文件。已经上传过的 .docx 进入第一步。
   - 第一步：必须先用 inspect_docx_template_tool 检查模板，了解 mode_detected、jinja_variables、bracket_tokens、slots 以及 removals。
   - 第二步（slots / 占位符）：
     · 如果有 jinja_variables 或 bracket_tokens：向用户询问尚未提供的值（用户已给出的字段不要重复问）。
     · 如果只有 slots（模板没有显式占位符）：**逐个**用 slots 里的 label + context 向用户确认每个槽位应填什么。slot 的 kind 可能是 highlighted / underscores / label_blank / empty_cell / angle_bracketed / placeholder_phrase / hint_text / section_body_empty / option_choice——其中 **highlighted（带黄/绿/青等 Word 高亮的文字）是最强的"请修改我"信号**，用户上传带高亮的模板就是希望把高亮处替换并清除高亮；填充时工具会**自动清除高亮**，同时保留字体/字号/加粗/斜体/颜色等其他格式。即使如此也要**逐条向用户确认替换内容**——不要仅凭高亮就自动决定写什么进去。
     · **option_choice（二选一/三选一）槽位**：slot 自带一个 `options` 列表，每项有 `index`、`header`（如"第一种：…"）和 `preview`（该方案正文的开头一段）。把所有选项的 header 念给用户，问"请问选第几种？"。用户答完之后，把 **选项的索引**（1、2、3…）或对应的标签（如 "第一种"、"第二种"）传回 `slot_values[slot_id]`。填充工具会**自动**保留所选方案的正文、删除提示语（"以下两种选择适合的一种…请删除"）、删除未选方案的 header 和全部正文段落——**你不需要**再用 edit_docx_tool / replace_text 去手动删除任何"第N种"标记或未选方案的正文，**也不要**把提示语当成 removal 重复提交（inspect 已经故意不把它作为单独 removal 列出）。如果未选方案中有用户**特别想保留**的某一段话，建议先用 edit_docx_tool 把那段话挪到所选方案下，再让 option_choice 槽位执行删除。
     · 对 highlighted 槽位，确认内容时**必须把整段高亮文字原样念给用户**（用 slot 的 `span_text` 字段，里面是高亮区域的完整原文）。例如高亮文字是「15个工作日」，要问"高亮的『15个工作日』要改成什么？"而不是只问"工作日改几天"。
     · **underscores 槽位的 `replaces` 字段**：inspect 会返回该字段，给出**精确**要被替换的字符（例如 `'     '` 5 个空格，或者 `'_______'` 7 个下划线）。slot 的 `context` 显示槽位前后的句子作为环境信息——但**填值时只替换 `replaces` 那一段**。**绝对不要**把 context 中前后已有的模板正文（比如紧跟在空白后面的括号说明 "（其中合同总额的20%作为定金）"、单位 "元/工作日/%" 等）复制进 slot_values 里——否则会出现 "90%（其中合同总额的20%作为定金）（其中合同总额的20%作为定金）" 这类括号被重复粘贴的 bug，因为模板里的原括号本来就还在。slot_values 里只放**真正要填进空白的值**（如 "90%" 或 "90"），其余 prose 让模板自己保留。
     · 如果 highlighted slot 带有 `scaffold` 字段（说明工具识别出了"变量 + 单位"形式，比如 15+个工作日、¥+850、50+%、2025年5月14日 等），**意味着填充时只换变量部分、保留前后的单位/币符/百分号**：用户回"20"，最终会写成"20个工作日"。即使如此，**你给 slot_values 时也最好直接传完整字符串**（"20个工作日"），不要只传"20"——只把数字作为兜底逻辑，避免歧义。绝对不要把"15个工作日"原样替换成"20"丢掉单位。
     · **`**` 双星号占位符（kind=angle_bracketed, source=asterisk_marker）**：在中文合同模板里，**每一个** `**` 都是一个独立的填空位，**`**` 之间和周围的文字是模板里要保留的原文，不要碰**。例：`项 目 名 称 ：  **项目50台**设备运输`——这里有 **两个** `**` 标记，所以是 **两个独立 slot**；中间的 `项目50台` 是模板作者写的内容（要保留），后面的 `设备运输` 也是模板正文（要保留）。**正确填法**：在第一个 `**` 处填"高能物理"，第二个 `**` 处填"试探"——最终输出 `项 目 名 称 ：  高能物理项目50台试探设备运输`。**绝对不要**把 `项目50台`、`设备运输` 这些已有原文复制到你的 slot_values 里——那样会变成 `高能物理项目50台试探项目50台设备运输` 这种重复。同理 `乙方单位名称（承运人）：上海**物流有限公司` 是一个 `**` slot：在 `**` 处填"顺达"，输出 `上海顺达物流有限公司`，**不要**再写"物流有限公司"。问用户时按 slot 的位置语境提问，例如"项目名称这一行有两个填空位，第一个 `**` 前面是空，后面跟着『项目50台』；要填什么？"。
     · **`is_prefilled: true` 槽位（关键，避免重复填写）**：inspect 返回的 slot 如果带 `is_prefilled: true` 字段，说明该槽位的段落（或 label_blank 下面的正文段落）**已经写好了实质内容**——`existing_text` 字段会显示当前的内容（最多 200/400 字）。这通常发生在：用户在上传模板前已经手动填了某个章节（如"1.5 合同文件的优先顺序"下已经列了 7 条文件清单），或者模板自带示例填充。**绝对不要**把这种槽位当成普通空白槽位往 `slot_values` 里塞值——填了也会被工具默认 skip 掉（防止 2026-05 "（1）变更洽商…（1）变更洽商…" 重复行 bug），返回里会出现 `skipped_prefilled_slot_ids`。正确流程：(1) 把 `existing_text` **原样读给用户**："我看到 1.5 节已经写了这些：……，要保留原样、修改、还是替换？"(2) 用户说**保留**——不用做任何事，槽位会维持原样；(3) 用户说**替换**——把新内容通过 `fill_docx_template_tool` 的 `replace_prefilled={slot_id: 新内容}` 参数（**不是 slot_values**）传入。工具会**清空整段原有内容并写入新值**，保留段落的对齐、缩进、编号样式以及第一个 run 的字体格式。`label_blank` 类槽位的 body 可能跨多个段落（比如"合同文件组成"下的 7 条），新值里用 `\\n` 分隔每一项，工具会按行写入既有段落（多余的清空，不够的克隆最后一段插入），保证编号列表的视觉结构不变。
     · **当章节标题（如"一、甲方委托乙方提供以下维修服务："）带有"以下/如下/下表/following/below"等字样、且后面紧跟一张表格时，要把每条维修服务/物品作为表格的一行来填，而不是把描述文字塞在标题和表格之间的空段落里。**inspect 已经默认不会在这种情况下emit section_body_empty 槽位；如果用户需要新增多条服务/产品行：先调 `edit_docx_tool` 用 `add_table_row` 增行（`{'type': 'add_table_row', 'table_index': N, 'values': ['第一列', '第二列', ...], 'position': 'end'}`，工具会**自动克隆最后一行的格式**——列宽、边框、对齐都会跟着——所以新加的行视觉上和原有行一致），如果还要改原有行用 `set_cell_text` / `replace_in_cell`，一次 `edit_docx_tool` 调用可以把多个 `add_table_row` + 多个 `set_cell_text` 全部放进 `edits` 数组里。**绝对不要**因为某个操作"看似不支持"就回退到 `unpack_docx_tool` + 手动改 XML + `pack_docx_tool`——这条路几乎一定会破坏文档结构（曾经把样式表搞坏）。如果 `add_table_row` 真的失败了，先把错误信息念给用户，再讨论替代方案，不要静默地直接拆包改 XML。
     · 如果多个 slots 共享相同或近似的 label（比如表格里多个"总价"、"大写"、"小写"单元格），**逐个**问用户该填什么，**不要**把同一个数字往所有看似相同的格子里灌。"总价" = 单价×数量的合计；"大写" / "小写" 是同一个金额的中文大写 / 阿拉伯数字两种写法——三者**值不同**，要按语义分别计算/转换后再填。
     · **中文数字的用法（关键，避免误用大写）**：大写数字 "壹/贰/叁/肆/伍/陆/柒/捌/玖/拾/佰/仟/万/亿/元/整" **只用于金额**——而且只在 slot 的 label 或上下文出现"大写"/"in words"/"capital amount"，或与一个"小写"金额槽位配对时才用。**所有其它中文数字填充**——年限（"保修期 一 年"）、月数、周数、天数（"3天内"）、工作日数、合同期限、产品数量、序号、百分比、版次、人数、份数、第几条等——一律用**阿拉伯数字**（"1"、"3"、"15"）或**小写中文数字**（"一/二/三/四/五/六/七/八/九/十/百/千"），**严禁**用"壹/贰/叁..."。错误示范："保修期或质量保证期 壹 年"（壹是大写）；正确："保修期或质量保证期 1 年" 或 "保修期或质量保证期 一 年"。如果用户没明确指定"用大写"，默认用阿拉伯数字；用户说"用中文"再用小写。
   - 第三步（removals 删除候选）：
     · 如果 inspect 返回的 removals 非空，**逐条**把 removal.text 朗读给用户，问"这一段需要从最终文档中删除吗？"
     · 把用户确认要删除的 id 收集进列表。**永远不要自动删除**——红色斜体的备注也可能是用户故意保留的注解。
   - 第四步：用 fill_docx_template_tool 生成新文档：
     · 显式占位符的值放入 context。
     · 用户确认的 slot 值放入 slot_values={"slot_003_签订时间": "...", "slot_005_签订地点": "...", ...}。slot id 形如 `slot_NNN_<标签>`（如 `slot_004_签订时间`），**必须**把 inspect_template 返回的 id **整段**作为 key，包括尾部的描述标签。**绝对禁止**自己改写成 `slot_4` / `slot_004` 这种**无标签**的旧短格式——填充工具会在保存前**直接拒绝**这种调用并返回 `rejected_legacy_slot_ids`，整次 fill 失败、不会写出文件。上次因为 LLM 用 `slot_0..slot_120` 这种裸编号给一个只有 79 个槽位的模板编号，结果值被错误地按数字前缀路由到语义完全不同的槽位（"甲方"被填进了"中华人民共和国合同法"），整份合同作废。**每个 key 都要原样复制 inspect 返回的 `id`**——一个字符都不要省。
     · 用户确认要删除的项放入 removal_ids=["rm_0", "rm_2", ...]。
     · **必须**输出一个**新文件**，文件名在模板原名后加 "_filled" 后缀（例如 contract.docx → contract_filled.docx），放在用户工作目录下。**绝对不要覆盖**用户上传的原模板——用户保留原始模板用于多次填写。fill_docx_template_tool 会以原模板为底直接复制并仅替换占位符所在 run，**保留原模板的全部其他内容、样式、页眉页脚、表格、图片、批注等不变**。
     · **优先一次调用填完所有 slot**——把所有用户确认的 slot 值放进**同一个** fill_docx_template_tool 调用的 slot_values 里。这样模型的 token 预算、JSON 输出和 slot id 编号都在同一个 inspect 上下文里、最稳定。**不要**为了"分批确认"就分多次调用——每次调用都要重新 inspect、重新做编号映射，出错面变大。
     · **如果实在因为槽位太多（>80个）需要分批 fill 才能避免 completion token 截断**：
       1) **第二次及之后**的调用，必须把 `template_path` 设为**上一次的 output_path**（即正在累积的 _filled 文件），而不是用户原始模板。原模板每次都用、会**抹掉**前一次填好的所有 slot——只有最后一次填进去的少数几个 slot 会留下，前面全部丢失。
       2) slot_values 的 key 用**第一次** inspect_template 返回的 canonical id（如 `slot_003_受托方乙方`）。工具检测到 output_path 已存在时会自动把 canonical id 翻译成当前 partial 文档的 id，**不要**自己用第二次 inspect 的新编号——后续 inspect 的编号是基于"还剩下的空槽"重新计数的，跟第一次不一样。
       3) **绝对不要**两次调用都把 `template_path` 设为原模板而 `output_path` 设为同一个文件——这就是去年（2026-05-20）智能仓储合同那次大量空白的根因，前三次填的 slot 全部被第四次调用从原模板复制时覆盖掉。
     · **填错了想从头再来怎么办（关键，避免死循环）**：如果 fill 的结果不对，**绝对不要**用 run_bash 去删除/移动 _filled 文件——workspace guard 会拒绝，并且每次重试都会让 tool 默默地把已有的 _filled 当成"上一次的部分填充"来续填，slot id 偏移越来越严重，agent 就会陷入"重新 inspect → 槽位变了 → 再填一次 → 又偏了"的循环（2026-05 真实发生过）。正确做法只有两条：（A）在 fill_docx_template_tool 上加 `force_fresh=True`，工具会忽略已有的 _filled、直接用原模板**覆盖**写新文件；（B）换一个新的 output_path（例如 `contract_filled_v2.docx`）。如果 fill 的返回里出现 `chunked_continuation: true` 或 `continuation_notice` 字段、而你又不是在做分块 continuation，就是这个陷阱——立刻用 force_fresh=True 重新调用。
   - **填写成功后（仅限新上传模板）**：如果这次填的模板是用户**新上传**的（不是从模板库通过 `get_template_path_tool` 取来的），主动问一句："要把这个模板保存进你的模板库吗？以后可以直接说'用 XX 模板'调用。要起什么名字 / 分类 / 别名？" 用户同意后调 `save_template_tool(template_path=<原模板路径>, name=..., description=..., category=..., tags=..., aliases=...)`。注意 template_path 传**原模板**，不是 _filled 文件。**模板已经在库里的不要重复问**——会重复保存。
   - 第五步（强制约束）：模板相关流程**只能**用以下工具：`list_templates_tool` / `get_template_path_tool` / `save_template_tool` / `delete_template_tool` / `inspect_docx_template_tool` / `fill_docx_template_tool` / `convert_doc_to_docx_tool`。**禁止**用 run_bash、run_glob、run_read、run_write、run_edit 去浏览模板库、定位模板文件、读取或填充模板——即便某个工具返回看起来为空、超时或慢，也要**重试同一个工具**或把情况告诉用户，**不要**回退到 bash/文件系统操作来"找文件""列目录"或"读 XML"。模板路径**只能**通过 `get_template_path_tool` / 用户上传事件取得，不要 glob 模板目录；DOCX 内部结构由 inspect/fill 工具内部处理，外部 bash 操作只会破坏格式或读不到正确字段。
   - **填写后的修正流程（关键，避免胡乱回退到 run_bash）**：`fill_docx_template_tool` 成功之后，**这个 _filled 文件已经没有占位符了**——再次调用 `inspect_docx_template_tool` / `fill_docx_template_tool` 通常什么都不会做（slots 都已消费），**不要重复 fill**。如果检查后发现某些表格单元格仍为空或填错了位置，必须按以下顺序处理：
     · **第一步：诊断**——调 `extract_docx_content_tool(file_path=<_filled 文件路径>)` 把 tables 读出来，记下错位/空白单元格的 `table_index` / `row` / `col` 和当前文本。
     · **第二步：用 `edit_docx_tool` 修复**——把每个修复都写成 `{'type': 'set_cell_text', ...}` 或 `{'type': 'replace_in_cell', ...}`（见规则 14）放进一次 `edit_docx_tool` 调用的 edits 数组里。**绝对不要**用 run_bash / run_write / run_edit 去改 .docx 文件——任何把 .docx 当文本/二进制改的尝试都会**损坏文档结构**（出现重复段落、丢失格式、xml 报错），过去用户就因为这种回退导致 "（价格含税，单位：人民币元）" 被重复粘了一份还把总计行覆盖掉了。
     · **第三步：验证**——再调一次 `extract_docx_content_tool` 把刚改过的那张表读出来，**核对**修复目标的单元格是不是符合预期（金额单元格非空、写对位置）。不要凭印象就说"已修复"。
     · 标准名是 `edit_docx_tool`。`edit_docx_content_tool` / `edit_docx_content` 是兼容别名，参数完全一致，调用任意一个都行——但**绝对**不要因为"找不到工具"就回退到 `run_bash` / `run_write` / `run_edit`。如果某个名字真的报"未注册"，先把所有 `edit_*` 名字都试一遍（`edit_docx_tool`、`edit_docx_content_tool`、`edit_docx_content`），再来汇报"工具都找不到"——绝不直接拿 bash 去改 .docx。
10. fill_docx_template_tool 默认 mode="auto"，会自动检测占位符风格——除非用户明确要求，否则不要强行指定 mode。若模板同时含 {{ }} 与 [TOKEN]，auto 模式会先按 Jinja 渲染再做一次方括号替换；slot_values 总是在最后一步应用，removal_ids 在保存输出文件后执行。
11. 不要把 inspect_docx_template_tool 用在普通文档上——那应该使用 extract_docx_content_tool 来查看内容。但模板里**没有任何**占位符也属于合法用法：inspect 会返回 slots / removals 让你识别可填空和可删除的位置。
12. fill_docx_template_tool 会保留整个文档的字体、颜色、加粗、斜体、对齐、段距等格式——只修改占位符所在的 run，周围的 run 和段落的样式都不动。对于 highlighted 类型的槽位，工具会**只清除该处的高亮**，但保留同一 run 的字体/字号/加粗/斜体/颜色等。因此**不要**在填模板之后再去"统一字体/格式"或调用 modify_docx_fonts_tool，那会覆盖用户模板的样式。
13. 如果用户上传的是 **.doc**（旧版二进制 Word）文件，必须先调用 convert_doc_to_docx_tool 把它转换成 .docx，然后再继续后续流程（模板检测、内容编辑、批注等）。转换后的文件路径在工具返回的 output_path 字段里——之后所有 DOCX 工具都用这个新路径。如果工具返回 success=False 且 error="soffice not found"，把 message 中的安装提示告诉用户并停止——LibreOffice 没装好之前下游工具都用不了 .doc 文件。convert_doc_to_docx_tool 也可以对 .docx 文件无害地调用（会返回 note="already .docx" 并不做任何修改），所以遇到 Word 文件统一先调一次很安全。
14. **表格编辑必须按单元格定位**——表格内容中相同的数字/词常常分布在多个单元格（单价 vs 总价、各行重复值等），用 replace_text 做全文档替换会误改。正确流程：
    - 第一步：用 extract_docx_content_tool 读出 tables，记下要改的 table_index / row / col，并把目标单元格的当前文本完整记下来（用于 replace_in_cell 的 old_text）。
    - 第二步：选其一：
      · 整格重写：{'type': 'set_cell_text', 'table_index': 0, 'row': 3, 'col': 5, 'value': '2550'}。如果单元格里要分多段（例如同一格里 "（大写）..." 和 "（小写）..." 各占一段），用 '\n' 分隔。
      · 在该单元格内精准替换：{'type': 'replace_in_cell', 'table_index': 0, 'row': 3, 'col': 5, 'old_text': '850', 'new_text': '2550'}。replace_in_cell 同时支持单段落 run-aware 替换以及跨段落匹配（针对一格内分两段的"大写/小写"情形）。
    - 一次改多格：把多个 set_cell_text / replace_in_cell 放进同一个 edit_docx_tool 调用的 edits 数组里，一次性提交。
    - 这两种工具都保留单元格原有 run 的字体/字号/加粗/斜体/颜色等格式（set_cell_text 保留第一个 run 的 rPr；replace_in_cell 是 run-aware 替换）。
    - 表格内容**禁止**用 replace_text。仅当用户明确要"全文统一替换"时才用 replace_text。

【优先使用 edit_docx_tool】
大多数编辑任务都应该用 edit_docx_tool 完成，包括：
- 添加或修改页眉页脚：{'type': 'add_header', 'header_type': 'custom', 'text': '标题'}
- {'type': 'add_footer', 'footer_type': 'page_number'}  # "Page X of Y"
- 在指定段落处添加/修改内容（用 position 参数，不要盲目用 replace_text）
- 添加标题、段落、表格、列表

如果需要精确修改某个段落，优先用 position 参数 + add_paragraph/add_heading，而不是 replace_text。
replace_text 会替换文档中所有匹配的文本，如果相同内容出现在多处会导致误替换。

【避免 replace_text 的陷阱】
- 替换"1. xxx"可能同时影响"议程"和"决议事项"等多个章节
- 如果必须用 replace_text，先用 extract_docx_content_tool 查看结构，确认目标文本是唯一的
- 或者使用 add_paragraph/add_heading 在指定位置插入新内容，比替换更安全
- **表格内容禁止用 replace_text**——表格里的同一数字/同一文字通常会在多个单元格出现（单价列 vs 总价列、各行重复值等），用全文替换会误改其他单元格。表格请走 set_cell_text / replace_in_cell（见规则 14）。

【⚠️ replace_text 必须使用完全精确的文本】
- replace_text 要求 old_text 与文档中的文本**100%完全一致**
- 包括：标点符号（全角/半角）、空格数量、引号类型（"" vs ''）都必须完全匹配
- **常见错误**：LLM 生成的"简洁版"内容会导致 old_text 不匹配，从而替换失败
- **正确流程**：
  1. 先用 extract_docx_content_tool 获取文档中**原始的 exact 文本**
  2. 在编辑时直接复制粘贴提取的文本作为 old_text
  3. 用 AI 重新生成 new_text 内容
  4. 这样才能确保 old_text 匹配成功
- 如果工具返回 success: False 和 "not found" 消息，说明 old_text 与文档文本不一致，请重新用 extract_docx_content_tool 获取精确文本

【格式修改的正确做法】
- 添加项目符号（bullet points）：使用 add_bullet_list 编辑类型或 add_bullet_list_tool
  例：{'type': 'add_bullet_list', 'items': ['第一点', '第二点', '第三点'], 'position': 5}
- 如果用户要求"把这几个段落改成列表/加bullet points"，在一次 edit_docx_tool 调用中完成：
  1. 先用 delete_paragraph 删除原有的纯文本段落（从最后一个开始往前删，避免索引偏移）
  2. 然后在原位置用 add_bullet_list 添加列表（包含原文内容）
  所有编辑放在同一个 edits 数组里，一次调用完成。
- 禁止用 replace_text 把内容替换成空字符串 ''！这会留下空白段落，不会真正删除段落。
  要删除段落必须用 delete_paragraph。
- 改变现有段落样式：使用 format_paragraph 类型修改段落格式

【批注操作规则】
批注操作必须使用专用工具，禁止通过XML编辑批注：
- 删除批注：使用 remove_comment_tool（可多次调用逐个删除）
- 添加批注：使用 add_comment_tool（针对文档中的特定文本）
- 绝对不要用 unpack_docx_tool 来手动编辑 comments.xml / commentsExtended.xml / commentsIds.xml
- 不要使用 run_bash 来操作批注相关的 XML 文件

【合同审查（review_contract_tool）】
- 触发场景：用户说"审查这份合同"/"审核合同"/"帮我看看这份合同"/"检查格式问题"/
  "看看有没有空着没填"/"合同有什么问题"等——文件类型是 .docx 且内容像合同。
- 调用方式：`review_contract_tool(file_path=<合同路径>, annotate=True)`。
  · 默认 annotate=True：除了返回结构化报告外，还会在 WORKDIR 下生成
    `<原名>_审查.docx`，把每条问题作为 Word 批注挂到原文相应位置，让用户下载对照。
  · annotate=False：只返回报告（适合用户只想要清单不要批注文件的情形）。
- 工具会自动跑四类检查：格式 / 填写缺失 / 内容一致性 / 法律风险（含 LLM 红线扫描，
  fail-soft——LLM 不可用时只返回前三类，summary 会注明）。**不要**自己再用
  extract_docx_content_tool / inspect_docx_template_tool 重复跑一遍这些检查。
- 收到返回后：把 summary 用作开场（如"共发现 10 处问题：…"），然后按 severity
  分段把 issues 念给用户——高优先级先讲，每条给出 location、message、suggestion。
  如果 annotated_path 非空，在末尾告诉用户"已生成带批注副本：<文件名>，可点击下载查看"。
- **不要**把这个工具用在非合同文档上（论文、说明书、邮件等）——它对那些场景的
  规则会误报。如果用户问的是普通文档校对，走 extract_docx_content_tool +
  edit_docx_tool 的常规路径。

【多步编辑任务的正确做法】
当用户要求同时修改内容和批注时，分步完成：
1. 先用 remove_comment_tool 删除需要删除的批注
2. 再用 edit_docx_tool 修改文档内容（替换文本、更新段落等）
3. 最后用 add_comment_tool 添加新批注
每一步都用专用工具，不要试图一次性用 XML 完成所有操作。

【XML编辑工具（最后手段）】
只有当以下情况出现时，才使用 unpack/edit/pack 工作流：
- 需要添加 tracked changes（修订痕迹）+ 作者归属 + 删除线
- edit_docx_tool / add_comment_tool / remove_comment_tool 明确无法完成的功能

注意：批注操作不属于"XML编辑才能完成的功能"。永远不要为了批注而解包 DOCX。

XML编辑工作流（仅用于 tracked changes）：
1. unpack_docx_tool(file_path, output_dir) → 解包DOCX到可编辑XML目录
2. 使用 run_read/run_edit 直接编辑 document.xml 中的特定段落
3. pack_docx_tool(input_dir, output_file, original_file) → 重新打包并验证
4. 验证通过后，删除 unpacked 目录（保持工作区整洁）

【处理模糊请求的规则】
遇到以下情况时，优先提一个简洁问题，而不是直接行动：
- 不知道要操作哪个文件
- 不知道要修改哪一段内容
- 用户要求“优化一下”“改得更好”但没有说明目标
- 用户要求的操作可能覆盖原文件且风险较高

你的澄清问题应尽量短，例如：
- “请提供要处理的 DOCX 文件路径。”
- “你是想修改内容、格式，还是两者都改？”
- “请指出要改写的段落，或让我先读取文档内容。”

【语义编辑规则】
当用户要求润色、改写、缩写、专业化、通俗化时，按以下方式处理：
1. 先确定目标文件和目标段落/章节。
2. 如果目标内容不明确，先读取文档内容。
3. 先基于原文生成合适的新文本，再执行替换或结构化编辑。
4. 完成后简要说明你改了什么。

【edit_docx_tool 的编辑格式】
使用 edit_docx_tool 时，edits 参数应为列表，每个元素是一个字典。常见格式如下：
- 替换文本：{'type': 'replace_text', 'old_text': '原文本', 'new_text': '新文本'}
- 也可使用等价替换格式：{'type': 'replace', 'target': '原文本', 'replacement': '新文本'}
- 添加段落：{'type': 'add_paragraph', 'content': '段落内容', 'position': 'end'}
- 添加标题：{'type': 'add_heading', 'content': '标题内容', 'level': 1}，其中 level=0 可作为 Title，level=2/3 适合子标题
- 修改样式：{'type': 'modify_style', 'style_name': 'Normal', 'font_name': '宋体', 'font_size': 12, 'bold': True, 'italic': False, 'underline': False, 'color': '1F1F1F', 'alignment': 'justify', 'spacing_before': 6, 'spacing_after': 6}
- 设置局部文字格式：{'type': 'format_text', 'target_text': '关键词', 'bold': True, 'italic': True, 'underline': True, 'font_size': 13, 'color': 'C00000'}
- 设置段落格式：{'type': 'format_paragraph', 'position': 3, 'alignment': 'center', 'spacing_before': 6, 'spacing_after': 6, 'line_spacing': 1.5}
- 插入分页符：{'type': 'add_page_break', 'position': 'end'}
- 添加表格：{'type': 'add_table', 'data': [['A', 'B'], ['1', '2']], 'table_style': 'Table Grid'}
- 设置表格样式：{'type': 'set_table_style', 'table_index': 0, 'table_style': 'Light Grid Accent 1'}
- 添加项目符号列表：{'type': 'add_bullet_list', 'items': ['第一点', '第二点', '第三点'], 'position': 'end'}
  也支持嵌套层级：{'type': 'add_bullet_list', 'items': ['主项', {'text': '子项1', 'level': 1}, {'text': '子项2', 'level': 1}], 'position': 'end'}
- 添加编号列表：{'type': 'add_numbered_list', 'items': ['第一步', '第二步', '第三步'], 'position': 'end'}
  也支持嵌套层级：{'type': 'add_numbered_list', 'items': ['步骤一', {'text': '子步骤A', 'level': 1}], 'position': 'end'}
- 注意：add_paragraph 的 position 支持整数索引；'end' 表示追加到文档末尾；嵌套列表的 level 从 0 开始

【输出风格】
1. 回答要专业、直接、清楚。
2. 执行工具前，内部先判断是否真的需要工具。
3. ⚠️ 对于多步骤编辑任务，**不要在中间步骤汇报**，等所有步骤完成后再统一报告最终结果。
4. 如果失败，明确说明失败原因和下一步建议。

记住：你的重点是正确理解用户对文档的真实意图，并以最小、最可靠的步骤完成任务。"""

    # Define document processing tools - define as actual functions
    # Side-channel for file events: tool functions append here,
    # DocMasterAgent.on_messages_stream drains it.
    _pending_files_events: list = []

    tools = []
    
    if DOCUMENT_PROCESSING_AVAILABLE:
        # Define document processing function - SIMPLIFIED
        def process_document(file_path: str):
            """
            Analyze a document file and return a machine-readable summary.

            Use this when the user wants to understand a file rather than edit it,
            for example: summarize the document, identify its type, inspect its
            structure, preview its contents, or extract high-level metadata.

            Best for:
            - DOCX, PDF, PPTX, XLSX, CSV, TXT, MD analysis
            - first-pass inspection before deciding what to edit
            - requests like "analyze this file", "what is in this document?",
              "summarize this report"

            Not for directly editing DOCX content.

            Args:
                file_path: Path to the input file.
            """
            processor = DocumentProcessor(str(WORKSPACE))
            return processor.process_uploaded_file(file_path)
        
        # Define DOCX editing function
        def edit_docx_tool(file_path: str, edits: list):
            """
            Apply one or more structured edits to an existing DOCX file.

            This is the single general-purpose DOCX editing tool. Use it for
            exact replacements, semantic rewrites after inspection, adding
            headings or paragraphs, inserting tables, and style-related changes.

            Best for:
            - replace one phrase with another
            - append or insert new content
            - add headings or tables
            - run several planned edits in one call
            - execute semantic edits after reading the target content first

            This tool overwrites the original DOCX file.

            Args:
                file_path: Path to the DOCX file.
                edits: List of edit operations. Accepted examples:
                    - {'type': 'replace_text', 'old_text': 'old', 'new_text': 'new'}
                    - {'type': 'replace', 'target': 'old', 'replacement': 'new'}
                    - {'type': 'add_paragraph', 'content': 'text', 'position': 'end', 'alignment': 'justify', 'spacing_after': 6}
                    - {'type': 'add_heading', 'content': 'Section 2', 'level': 2, 'bold': True, 'color': '1F4E79'}
                    - {'type': 'modify_style', 'style_name': 'Heading 2', 'font_name': 'Calibri', 'font_size': 14, 'bold': True, 'spacing_before': 12, 'spacing_after': 6}
                    - {'type': 'format_text', 'target_text': 'important', 'bold': True, 'italic': True, 'underline': True, 'color': 'C00000'}
                    - {'type': 'format_paragraph', 'position': 3, 'alignment': 'center', 'line_spacing': 1.5}
                    - {'type': 'add_page_break', 'position': 'end'}
                    - {'type': 'add_table', 'data': [['A', 'B'], ['1', '2']], 'table_style': 'Table Grid'}
                    - {'type': 'set_table_style', 'table_index': 0, 'table_style': 'Light Grid Accent 1'}
                    - {'type': 'set_cell_text', 'table_index': 0, 'row': 3, 'col': 5, 'value': '2550'}   # overwrite one cell; '\\n' splits paragraphs
                    - {'type': 'replace_in_cell', 'table_index': 0, 'row': 3, 'col': 5, 'old_text': '850', 'new_text': '2550'}   # scoped find/replace inside one cell (cross-paragraph aware)
                    - {'type': 'add_table_row', 'table_index': 0, 'values': ['管理交换机', '型号X', '2', '12000', '24000'], 'position': 'end'}   # append a new row; auto-picks the row with the MOST cells as the format template (avoids merged 总计/header rows). 'position' may be an int (insert BEFORE that row index — use this to insert above a totals row). Optional 'template_row_index': N pins which existing row's formatting to clone.
                    - {'type': 'delete_table_row', 'table_index': 0, 'row': 5}   # remove row 5 from table 0
                    - {'type': 'add_bullet_list', 'items': ['要点一', '要点二'], 'position': 'end'}
                    - {'type': 'add_numbered_list', 'items': ['第一步', '第二步'], 'position': 'end'}
                    - {'type': 'insert_image', 'image_path': '/abs/path/to/image.png', 'position': 0, 'width_inches': 5.0}
                    - {'type': 'add_footer', 'footer_type': 'page_number'}   # adds "Page X of Y" centered footer
                    - {'type': 'add_footer', 'footer_type': 'page_x'}        # adds "Page X" centered footer
                    - {'type': 'add_footer', 'footer_type': 'custom', 'text': 'Confidential'}  # custom text footer
                    - {'type': 'add_header', 'header_type': 'custom', 'text': 'Document Title'}  # custom text header
                    - {'type': 'add_header', 'header_type': 'title'}         # document title from file properties
                    - {'type': 'add_header', 'header_type': 'filename', 'text': 'filename.docx'}  # filename header
            """
            import json

            # Log the incoming request for debugging
            print(f"🔧 edit_docx_tool called with:")
            print(f"   File: {file_path}")
            print(f"   Edits count: {len(edits)}")
            print(f"   First edit sample: {json.dumps(edits[0] if edits else {}, ensure_ascii=False, indent=2)[:200]}...")
            
            # File-path sanity guard — rejects hallucinated paths with strong
            # recovery guidance so the agent doesn't fall back to run_bash.
            guard = _guard_docx_file_path(file_path, tool_label="edit_docx_tool")
            if guard is not None:
                print(f"❌ {guard.get('error')}: {guard.get('message')[:200]}")
                return guard
            
            # Convert LLM format to standard format if needed
            standardized_edits = []
            conversion_stats = {'llm_format': 0, 'standard_format': 0, 'unknown': 0}
            
            for i, edit in enumerate(edits):
                if isinstance(edit, dict):
                    edit_type = edit.get('type', '')
                    
                    # Handle LLM format: 'replace' -> 'replace_text'
                    # Supports both 'target'/'replacement' and 'old_text'/'new_text' variants
                    if edit_type == 'replace':
                        if 'target' in edit and 'replacement' in edit:
                            standardized_edits.append({
                                'type': 'replace_text',
                                'old_text': edit['target'],
                                'new_text': edit['replacement']
                            })
                            conversion_stats['llm_format'] += 1
                            continue
                        elif 'old_text' in edit and 'new_text' in edit:
                            # Model is using old_text/new_text with type 'replace'
                            standardized_edits.append({
                                'type': 'replace_text',
                                'old_text': edit['old_text'],
                                'new_text': edit['new_text']
                            })
                            conversion_stats['llm_format'] += 1
                            continue
                    
                    # Handle 'replace_text' type with old_text/new_text (also common LLM format)
                    if edit_type == 'replace_text' and 'old_text' in edit and 'new_text' in edit:
                        standardized_edits.append(edit)
                        conversion_stats['standard_format'] += 1
                        continue
                    
                    # Handle other potential LLM format variations
                    if edit_type == 'add':
                        if 'content' in edit:
                            standardized_edits.append({
                                'type': 'add_paragraph',
                                'content': edit['content'],
                                'position': edit.get('position', 'end')
                            })
                            conversion_stats['llm_format'] += 1
                            continue
                    
                    # If it's already in standard format
                    if edit_type in ['replace_text', 'add_paragraph', 'add_heading', 'modify_style', 'add_table', 'format_text', 'format_paragraph', 'add_page_break', 'set_table_style', 'insert_image', 'add_footer', 'add_header', 'add_bullet_list', 'add_numbered_list', 'delete_text', 'delete_paragraph']:
                        standardized_edits.append(edit)
                        conversion_stats['standard_format'] += 1
                    else:
                        # Unknown format, try to use as-is
                        standardized_edits.append(edit)
                        conversion_stats['unknown'] += 1
                        print(f"⚠️ Unknown edit format at index {i}: {edit_type}")
                else:
                    # If edit is not a dict, keep it as-is
                    standardized_edits.append(edit)
                    conversion_stats['unknown'] += 1
            
            print(f"📊 Edit conversion stats: {conversion_stats}")
            print(f"📝 Standardized edits count: {len(standardized_edits)}")
            
            processor = DocumentProcessor(str(WORKSPACE))
            result = processor.edit_docx(file_path, standardized_edits, overwrite_original=True)
            
            print(f"✅ edit_docx_tool result: {result.get('success', False)}")
            if not result.get('success', False):
                print(f"❌ Error: {result.get('error', 'Unknown error')}")
                print(f"📝 Message: {result.get('message', 'No message')}")
            else:
                fe_data = _build_files_event_data(file_path, f"Edited DOCX: {Path(file_path).name}")
                if fe_data:
                    _pending_files_events.append(fe_data)
            
            # Trim result to avoid flooding the context window.
            # The raw result repeats change details 3x (document_info.changes_made,
            # changes, debug_info.edits_sample). Keep only what the agent needs.
            def _truncate(s, maxlen=80):
                return s[:maxlen] + '...' if len(s) > maxlen else s
            
            changes_raw = result.get('changes', result.get('document_info', {}).get('changes_made', []))
            trimmed_changes = [_truncate(c) for c in changes_raw] if isinstance(changes_raw, list) else changes_raw
            
            return {
                'success': result.get('success', False),
                'message': result.get('message', ''),
                'file_path': result.get('file_path', file_path),
                'paragraph_count': result.get('document_info', {}).get('paragraph_count'),
                'table_count': result.get('document_info', {}).get('table_count'),
                'changes': trimmed_changes,
            }
        
        # Define DOCX content extraction function
        def extract_docx_content_tool(file_path: str):
            """
            Extract detailed DOCX content for inspection before editing.

            Use this when you need the actual document text or structure, not
            just metadata. This is the preferred inspection tool before nontrivial
            rewrite, polish, shorten, reorganize, or section-level edits.

            Best for:
            - reading paragraphs and tables from a DOCX
            - locating target text before semantic edits
            - understanding document structure before applying changes

            Not for editing by itself.

            ⚠️ 表格里的合并单元格（merged cells）—— 读结果里 `tables[i].data` 是按
            **grid 坐标**返回的二维数组，所以一个合并单元格只在它的**锚点位置**
            出现一次，其它被合并进来的位置会显示为
            `'⟨merged with r{R}c{C}⟩'` 哨兵字符串（指向锚点）。同时
            `tables[i].merges` 会列出每一处合并区域：
                {anchor:[r,c], colspan, rowspan, text}

            合并单元格的使用规则：
            - **绝对不要**对哨兵位置写值（set_cell_text / replace_in_cell 之类）。
              那些位置和锚点共用同一个底层 <w:tc>，写过去会把整个合并块的文字
              一次性覆盖成你写的短字符串——这是上次"甲方开票信息"被改成
              "100049" 的根因。
            - 要修改合并块的内容，**只对锚点 (r, c) 写**。例如 `merges` 里有
              `{anchor:[6,1], colspan:3, ...}`，所有改动都用 row=6, col=1。
            - 合并块里的文本本身就是原模板的固定内容（例如 "甲方开票信息如下…"
              整段、印章占位 "合 同 章…年 月 日"、表头中"甲方/乙方"的标签）——
              通常这些是模板里**有意要保留的整段固定信息**，**不要**当成
              "重复 / 错位 / 需要清理"的脏数据去改。看到合并块里有连续多行
              prose（公司名/纳税人识别号/地址/账号 等），默认保留，**除非用户
              明确要求修改**。
            - 真正空的可填单元格（不是合并）会显示成空字符串 ''，和哨兵
              `'⟨merged with …⟩'` 是不同的；不要混淆。

            Args:
                file_path: Path to the DOCX file.
            """
            guard = _guard_docx_file_path(file_path, tool_label="extract_docx_content_tool")
            if guard is not None:
                return guard
            processor = DocumentProcessor(str(WORKSPACE))
            result = processor.extract_docx_content(file_path)
            return result
        
        # ============ NEW ENHANCED EDITING TOOLS ============
        
        def delete_docx_content_tool(file_path: str):
            """
            Remove all content from an existing DOCX document.

            Use this only when the user clearly wants the document emptied,
            cleared, or reset. This is a destructive operation on the target file.

            Best for:
            - "clear this document"
            - "delete all text/content"
            - preparing an existing DOCX to be rebuilt

            Args:
                file_path: Path to the DOCX file.
            """
            processor = DocumentProcessor(str(WORKSPACE))
            result = processor.delete_docx_content(file_path)
            if result.get('success', False):
                fe_data = _build_files_event_data(file_path, f"Cleared DOCX: {Path(file_path).name}")
                if fe_data:
                    _pending_files_events.append(fe_data)
            return result
        
        def modify_docx_fonts_tool(file_path: str, font_rules: dict = None):
            """
            Change DOCX fonts according to language or content-type rules.

            Use this when the request is specifically about typography rather than
            wording, for example changing Chinese text to 宋体 and English text to
            Times New Roman.

            Best for:
            - document-wide font normalization
            - Chinese/English font separation
            - formatting requests focused on font family rules

            Not for rewriting content.

            Args:
                file_path: Path to the DOCX file.
                font_rules: Mapping of content categories to fonts, e.g.
                    {'chinese': '宋体', 'english': 'Times New Roman'}
            """
            processor = DocumentProcessor(str(WORKSPACE))
            result = processor.modify_docx_fonts(file_path, font_rules)
            if result.get('success', False):
                fe_data = _build_files_event_data(file_path, f"Font-modified DOCX: {Path(file_path).name}")
                if fe_data:
                    _pending_files_events.append(fe_data)
            return result
        
        def create_docx_with_content_tool(output_path: str, content: list):
            """
            Create a new DOCX file from structured content elements.

            This is the single document-creation tool. Use it whenever the user
            wants a new Word document, whether simple or structured. If the user
            only gives plain text, convert it into a reasonable structured content
            list before calling this tool.

            Best for:
            - creating a new report or letter
            - building a document from headings and paragraphs
            - generating a formatted DOCX from an outline or structured data

            Args:
                output_path: Path where the new DOCX should be saved.
                content: List of structured content items, for example:
                    - {'type': 'heading', 'text': 'Document Title', 'level': 0}
                    - {'type': 'heading', 'text': 'Introduction', 'level': 1}
                    - {'type': 'subheading', 'text': 'Background', 'level': 2}
                    - {'type': 'paragraph', 'text': 'Body text', 'font_name': 'Times New Roman', 'font_size': 12, 'alignment': 'justify', 'spacing_after': 6}
                    - {'type': 'page_break'}
                    - {'type': 'table', 'data': [['A', 'B'], ['1', '2']], 'table_style': 'Table Grid'}
                    Optional fields may include font/font_name, bold, italic, underline, font_size, color, alignment, and spacing.
            """
            processor = DocumentProcessor(str(WORKSPACE))
            result = processor.create_docx_with_content(output_path, content)
            if result.get('success', False):
                fe_data = _build_files_event_data(output_path, f"Created DOCX: {Path(output_path).name}")
                if fe_data:
                    _pending_files_events.append(fe_data)
            
            # Trim verbose result for context window
            doc_info = result.get('document_info', {})
            elements = doc_info.get('created_elements', result.get('elements_created', []))
            return {
                'success': result.get('success', False),
                'message': result.get('message', ''),
                'file_path': result.get('file_path', output_path),
                'paragraph_count': doc_info.get('paragraph_count'),
                'table_count': doc_info.get('table_count'),
                'elements_created': len(elements),
            }

        def convert_doc_to_docx_tool(file_path: str):
            """
            Convert a legacy Word .doc file to modern .docx via LibreOffice headless.

            Call this FIRST whenever the user uploads a .doc file — python-docx
            and every other DOCX tool here cannot read the legacy binary format.
            The result includes `output_path`; use that path for any follow-up
            tool calls (inspect_docx_template_tool, edit_docx_tool, etc.).

            If the input is already .docx, returns success with note="already .docx"
            and the same path — safe to call on any uploaded Word file.

            Requires `soffice` (LibreOffice) on PATH. If missing, returns a clear
            error with install instructions; stop and tell the user.

            Args:
                file_path: Path to the uploaded .doc (or .docx) file.

            Returns dict with: success, input_path, output_path, soffice_used,
                message. On failure: success=False, error=<reason>, message=<hint>.
            """
            skill = DocToDocxSkill(str(WORKSPACE))
            result = skill.convert(file_path)
            if (
                result.get('success')
                and result.get('soffice_used')
                and result.get('output_path') != file_path
            ):
                fe_data = _build_files_event_data(
                    result['output_path'],
                    f"Converted to DOCX: {Path(result['output_path']).name}",
                )
                if fe_data:
                    _pending_files_events.append(fe_data)
            return result

        def inspect_docx_template_tool(template_path: str):
            """
            Inspect a DOCX template and discover its placeholders and fillable slots.

            Use this BEFORE asking the user to provide values for a template.
            Returns the detected placeholder style and the list of variables /
            tokens found, plus heuristic "slots" for templates that do not use
            any explicit placeholder syntax (e.g. underscore lines `_____`,
            "Label:" with empty tail, empty table cells under a header).

            Best for:
            - "I uploaded a template, what fields does it have?"
            - planning a conversational fill flow before generating a document
            - filling user-uploaded templates that don't use {{ }} or [TOKEN]

            Args:
                template_path: Path to the .docx template file.

            Returns a dict with:
                mode_detected: 'jinja' (uses {{ name }}), 'bracket' (uses [NAME]),
                    'both', or 'none'
                jinja_variables: list of top-level Jinja variable names
                bracket_tokens: list of bracket tokens (without the brackets)
                has_loops, has_conditionals: True if {% for %} / {% if %} present
                slots: list of {id, kind, label, context, ...} dicts where kind
                    is one of:
                      - 'highlighted'         — run(s) with a Word highlight
                                                (yellow/green/cyan/…); the
                                                strongest fill signal. On fill
                                                the text is replaced AND the
                                                highlight is cleared. Extra
                                                fields on this kind:
                                                  • span_text: full original
                                                    highlighted text — show it
                                                    to the user verbatim when
                                                    asking what to fill.
                                                  • scaffold (optional): when
                                                    the span looks like
                                                    "<variable><unit>"
                                                    (15个工作日, ¥850, 50%,
                                                    2025年5月14日, 10kg, …),
                                                    a {kind, prefix, variable,
                                                    suffix} dict. If the
                                                    slot_values reply is a
                                                    bare number, prefix+suffix
                                                    are auto-reattached on
                                                    fill so "20" becomes
                                                    "20个工作日". The agent
                                                    SHOULD still pass the full
                                                    intended string when it
                                                    can — auto-reattach is a
                                                    safety net.
                      - 'underscores'         — run of 3+ underscores
                      - 'label_blank'         — paragraph ending in "Label:"
                      - 'empty_cell'          — empty table cell under a header
                      - 'angle_bracketed'     — <token> or 《占位符》, OR
                                                **per-marker `**` asterisk
                                                fill positions** common in
                                                CN contract templates.
                                                EACH `**` is one slot that
                                                replaces just the 2-char
                                                marker; the literal text
                                                around and between markers
                                                is template content to
                                                preserve. So
                                                `**项目50台**设备运输` has
                                                TWO slots (flanking the
                                                preserved sample
                                                `项目50台`, with `设备运输`
                                                kept as static text); fill
                                                "高能物理" + "试探" → output
                                                `高能物理项目50台试探设备运输`.
                                                NEVER echo the surrounding
                                                template text in the slot
                                                value — only the words that
                                                replace the `**` itself.
                                                Each such slot carries
                                                `source: "asterisk_marker"`
                                                plus `asterisk_marker_index`
                                                / `asterisk_marker_total`
                                                so you can tell the user
                                                which `**` in the paragraph
                                                you're asking about.
                      - 'placeholder_phrase'  — "your text here", "请填写", "TBD"…
                      - 'hint_text'           — italic/grey instructional run
                      - 'section_body_empty'  — empty body under a Heading
                      - 'option_choice'       — 二选一 / 三选一 pattern: an
                                                instruction paragraph
                                                ("以下两种选择适合的一种…请删除")
                                                followed by two or more
                                                "第N种" option headers, each
                                                with its own body paragraphs.
                                                Extra fields on this kind:
                                                  • options: list of
                                                    {index, header, preview}
                                                    describing each branch.
                                                  • fill_policy: instructions
                                                    on what value to pass.
                                                Ask the user which option
                                                ("请问选第几种？") and pass
                                                the chosen index back as
                                                slot_values[slot_id] = 1
                                                (or "第二种" / "first" /
                                                "2" — all accepted). The
                                                fill tool then keeps only
                                                the chosen option's body
                                                and deletes the prompt +
                                                all other options.
                                                **Do NOT also pass the
                                                instruction prompt as a
                                                removal_id** — its
                                                deletion is owned by the
                                                slot fill, and it is
                                                deliberately not emitted
                                                as a separate removal.
                  label = best-guess field name (may be None for stray cases).
                  context = surrounding snippet for disambiguation.
                  is_prefilled (optional bool) + existing_text: when present
                  and true, the slot's paragraph (or its label_blank body)
                  ALREADY contains substantive content — e.g. the user
                  pre-typed values under a section heading before asking
                  docmaster to fill the rest. By default fill_docx_template_tool
                  SKIPS these slots when given via `slot_values` (passes them
                  through unchanged) to avoid duplicate-content bugs. To
                  overwrite, read `existing_text` to the user, ask whether
                  to keep or replace, and if replace, pass the new value
                  through fill_docx_template_tool's `replace_prefilled`
                  argument (NOT slot_values). The tool wipes the existing
                  paragraph(s) and writes the new value while preserving
                  paragraph-level formatting.
                  Pass the slot ids back via fill_docx_template_tool's
                  slot_values argument once the user confirms each. Each id
                  has the form `slot_NNN_<label>` (e.g.
                  `slot_004_PatientName`); copy it VERBATIM — the trailing
                  label is part of the id, and fill_docx_template_tool
                  REJECTS bare numeric forms like `slot_4` outright.
                removals: list of {id, kind, text, reason} dicts of paragraphs/
                    runs that look like template instructions to be deleted
                    before publishing (e.g. "Delete this before submitting",
                    "Note to author: ...", "请删除本段"). kind is
                    'instruction_paragraph' (whole paragraph removed) or
                    'instruction_run' (just one run blanked). NEVER auto-
                    delete: read each entry to the user, get confirmation,
                    then pass approved ids via fill_docx_template_tool's
                    removal_ids argument.
                warnings: notes (e.g. mixed-mode template, ambiguous slots,
                    removal candidates present)
            """
            guard = _guard_template_path(template_path)
            if guard is not None:
                return guard
            skill = DocxTemplateSkill(str(WORKSPACE))
            return skill.inspect_template(template_path)

        def fill_docx_template_tool(
            template_path: str,
            output_path: str,
            context: dict = None,
            mode: str = "auto",
            slot_values: dict = None,
            removal_ids: list = None,
            force_fresh: bool = False,
            replace_prefilled: dict = None,
        ):
            """
            Fill a DOCX template with values and save to a new DOCX file.

            Supports three input styles, freely combinable:
            - Jinja (docxtpl): {{ name }}, {% for x in xs %}…{% endfor %},
              {%tr for row in rows %}…{%tr endtr %} for table rows.
              `context` is a regular nested dict matching the variables.
            - Bracket: [NAME], [DATE] — literal substitution across paragraphs,
              tables, and headers/footers. `context` keys are the token names
              without brackets, e.g. {"NAME": "Alice"} fills [NAME]. Bracket
              tokens must be uppercase-leading (alnum, underscore, dash).
            - Heuristic slots: for templates without placeholder syntax. Call
              inspect_docx_template_tool first to discover slot ids, confirm
              what each holds with the user, then pass slot_values using the
              ids returned VERBATIM — they look like
              {"slot_000_PatientName": "Alice",
               "slot_001_Diagnosis": "Mild hypertension", ...}.
              The trailing label after the counter is PART OF the id. Bare
              forms like "slot_0" or "slot_001" (no descriptive tail) are
              REJECTED at the fill boundary — the tool refuses to save and
              returns `slot_fill.rejected_legacy_slot_ids` plus a
              `legacy_canonical_map` showing the correct id for each rejected
              key. Past incident (2026-05): the LLM emitted `slot_0..slot_120`
              for a 79-slot contract template, values were routed by numeric
              prefix to semantically wrong slots, and the whole contract had
              to be discarded. Copy each id character-for-character from
              inspect's output.

            Chunked fills (large templates with many slots): prefer ONE call
            covering all slots. If you must split into multiple calls, the
            2nd+ calls MUST set `template_path` to the PREVIOUS call's
            `output_path` so prior fills are preserved. You may pass the
            canonical slot ids from the FIRST inspect_docx_template_tool —
            the tool auto-translates them to the partial document's current
            ids.

            Implicit continuation: if `output_path` already exists from a
            previous call (and you didn't set `template_path` to that file),
            the tool auto-switches the source to the partial doc so prior
            fills survive — the response will carry `chunked_continuation:
            True` and a `continuation_notice` field. If that was NOT what
            you wanted (e.g. the previous fill went sideways and you want
            to start over from the original template), re-call this tool
            with `force_fresh=True`. **Do NOT use run_bash / rm / mv to
            delete or move the output file** — the workspace guards reject
            those, and it makes the loop worse. Either `force_fresh=True`
            or a fresh `output_path` is the right recovery move.

            Run-level formatting (bold/italic/color/font) in the template is
            preserved — only the runs that span the matched placeholder text
            are edited; other runs in the same paragraph are left untouched.

            Use mode='auto' (default) to let the tool detect from the template.
            Force 'jinja' or 'bracket' only when the user explicitly asks.
            For mixed templates ({{ }} AND [TOKEN]) in auto mode, Jinja renders
            first, then a bracket pass runs over the rendered output. Slot
            fills always run last, on the rendered document.

            Best for:
            - "Fill this template with these values"
            - generating one or more documents from an uploaded template
            - filling templates whose authors didn't add placeholders

            Args:
                template_path: Path to the .docx template.
                output_path: Where to write the filled .docx. Use a different
                    path from template_path to preserve the original (a
                    `_filled.docx` suffix is conventional).
                context: Mapping of placeholder values (Jinja vars or bracket
                    tokens). May be omitted/None if you're only filling slots.
                mode: 'auto' (default), 'jinja', or 'bracket'.
                slot_values: Optional {slot_id: value} from
                    inspect_docx_template_tool's `slots` list.
                removal_ids: Optional list of `id`s from inspect's `removals`
                    list. Each id corresponds to a paragraph or run that the
                    user confirmed should be DELETED from the final document
                    (template instructions, "Delete this before submitting"
                    notes, etc.). Removals run last, after all fill passes,
                    on the saved output file. Pass None or [] to keep all
                    template prose intact.
                force_fresh: When True, skip the implicit chunked-continuation
                    auto-detect: even if `output_path` already exists, write
                    a fresh fill from `template_path` and overwrite the
                    output. Use this to recover from a botched previous
                    fill — NEVER use run_bash to delete the file. Default
                    False (preserve prior partial fills when re-running).
                replace_prefilled: Optional {slot_id: new_text} for slots
                    inspect_docx_template_tool flagged with
                    `is_prefilled: true`. A prefilled slot is one whose
                    paragraph (or its body, for `label_blank` kind) already
                    contains substantive content — e.g. the user pre-typed
                    a numbered list under "1.5 合同文件的优先顺序" before
                    asking docmaster to fill the rest of the contract.
                    By default these slots are SKIPPED by the normal
                    `slot_values` pass to avoid the duplication bug
                    (2026-05) where the agent's value was appended next to
                    the existing content. To OVERWRITE a prefilled slot,
                    confirm the change with the user, then pass the new
                    text here. The tool wipes the existing paragraph
                    (or all consecutive body paragraphs for label_blank)
                    and writes the new value, preserving paragraph-level
                    formatting (alignment, style, numbering) and the
                    first run's character formatting. For label_blank
                    bodies with multiple lines, separate items with `\\n`
                    — each line becomes a separate body paragraph (extra
                    paragraphs are inserted by cloning the last existing
                    body paragraph so list numbering / indents survive).
                    Workflow: (1) inspect, (2) note slots with
                    `is_prefilled: true` and their `existing_text`,
                    (3) read `existing_text` to the user and ASK whether
                    to keep / replace, (4) only if the user says replace,
                    pass {slot_id: new_text} here.
            """
            guard = _guard_template_path(template_path)
            if guard is not None:
                return guard
            skill = DocxTemplateSkill(str(WORKSPACE))
            result = skill.fill_template(
                template_path,
                output_path,
                context=context or {},
                mode=mode,
                slot_values=slot_values or {},
                removal_ids=removal_ids or [],
                force_fresh=force_fresh,
                replace_prefilled=replace_prefilled or {},
            )
            if result.get('success', False):
                fe_data = _build_files_event_data(
                    output_path,
                    f"Template-filled DOCX: {Path(output_path).name}",
                )
                if fe_data:
                    _pending_files_events.append(fe_data)
            return result

        # ---------------------------------------------------------------- #
        # Template library — shared + per-user .docx template catalog       #
        # ---------------------------------------------------------------- #

        def list_templates_tool(category: str = None, query: str = None):
            """
            列出当前用户可用的 DOCX 模板（共享库 + 用户自己保存的模板）。

            什么时候用：
            - 用户问"我有哪些模板？""现在能用哪些合同模板？"
            - 用户没有上传文件但提到要用某种模板，**先**调这个工具看看有没有；
              不要直接要求用户上传。
            - 用户要按类别浏览（如"看看采购合同类的模板"）—— 用 category 参数。
            - 用户给的是关键词（如"含'保密'的模板"）—— 用 query 参数。

            Args:
                category: 可选，按类别过滤（如 "合同/采购"）。匹配前缀，
                    所以 "合同" 也会命中 "合同/采购" 与 "合同/技术开发"。
                query: 可选，关键词子串匹配 name / description / id / tags /
                    aliases（大小写不敏感）。

            Returns dict with:
                shared: 共享库里的模板列表，每项含 id / name / description /
                    category / tags / aliases / source="shared"
                mine:   当前用户自己保存的模板，结构同上，source="mine"
                message: 简要中文摘要

            注意：返回的是元数据，不含模板文件本身。要打开/填一个模板，请把
            id（或 alias / name）传给 get_template_path_tool 拿到具体路径，
            再用 inspect_docx_template_tool / fill_docx_template_tool。
            """
            return _template_api_list(user_id=user_id, category=category, query=query)

        def get_template_path_tool(template_ref: str):
            """
            根据 id / alias / 模板名（支持模糊子串）定位模板库里的模板，返回
            它的本地路径，方便接下来用 inspect_docx_template_tool /
            fill_docx_template_tool 进行检查和填写。

            什么时候用：
            - 用户说"用 3-1 模板""用 技术开发合同 模板""用我的采购合同模板"
              —— 把用户说的字符串原样传进 template_ref。
            - 模板填写流程的**第零步**：先尝试在模板库里找，找到就用，
              找不到再回退到要求上传。

            解析顺序：
              1. 用户自己库里的精确 id
              2. 共享库里的精确 id
              3. 别名（aliases）精确匹配（用户优先于共享）
              4. 在 name / id / aliases / tags 上做子串匹配（用户优先于共享）

            Args:
                template_ref: 用户口中的"模板名"——可以是 id（如 "tech-dev-3-1"）、
                    别名（如 "3-1"）、或显示名的一部分（如 "技术开发"）。

            Returns dict with:
                success: True / False
                template_path: 命中时返回模板 .docx 的绝对路径
                source: "mine" 或 "shared"
                metadata: 命中模板的完整元数据
                ambiguous: True 表示匹配到多个候选 → candidates 字段里给出
                    候选列表，让用户挑一个
                message: 中文提示

            如果 ambiguous=True，**不要**自己挑—— 把 candidates 念给用户，
            让用户确认是哪一个，再用确认后的 id 再调一次 get_template_path_tool。
            """
            ref = (template_ref or "").strip()
            if not ref:
                return {"success": False, "ambiguous": False, "message": "template_ref 不能为空。"}
            listing = _template_api_list(user_id=user_id)
            if not listing.get("success"):
                return {"success": False, "ambiguous": False, "message": listing.get("message") or "模板库服务调用失败"}
            mine = listing.get("mine") or []
            shared = listing.get("shared") or []
            ref_lower = ref.lower()

            def _by_id(entries, rid):
                for e in entries:
                    if e.get("id") == rid:
                        return e
                return None

            hit_entry = None
            hit_source = None
            e = _by_id(mine, ref)
            if e:
                hit_entry, hit_source = e, "mine"
            else:
                e = _by_id(shared, ref)
                if e:
                    hit_entry, hit_source = e, "shared"

            if hit_entry is None:
                # Stage 3 — exact alias
                alias_hits = []
                for src_name, entries in (("mine", mine), ("shared", shared)):
                    for e in entries:
                        for a in (e.get("aliases") or []):
                            sa = str(a)
                            if sa == ref or sa.lower() == ref_lower:
                                alias_hits.append((e, src_name))
                                break
                if len(alias_hits) == 1:
                    hit_entry, hit_source = alias_hits[0]
                elif len(alias_hits) > 1:
                    return {
                        "success": False, "ambiguous": True,
                        "candidates": [
                            {"id": e["id"], "name": e.get("name", ""), "source": s, "description": e.get("description", "")}
                            for e, s in alias_hits
                        ],
                        "message": f"匹配到 {len(alias_hits)} 个模板，请让用户从候选中挑选一个。",
                    }

            if hit_entry is None:
                # Stage 4 — substring match on name / id / aliases / tags
                sub_hits = []
                for src_name, entries in (("mine", mine), ("shared", shared)):
                    for e in entries:
                        haystack = " ".join(
                            [str(e.get("name", "")), str(e.get("id", ""))]
                            + [str(a) for a in (e.get("aliases") or [])]
                            + [str(t) for t in (e.get("tags") or [])]
                        ).lower()
                        if ref_lower in haystack:
                            sub_hits.append((e, src_name))
                if len(sub_hits) == 1:
                    hit_entry, hit_source = sub_hits[0]
                elif len(sub_hits) > 1:
                    return {
                        "success": False, "ambiguous": True,
                        "candidates": [
                            {"id": e["id"], "name": e.get("name", ""), "source": s, "description": e.get("description", "")}
                            for e, s in sub_hits
                        ],
                        "message": f"匹配到 {len(sub_hits)} 个模板，请让用户从候选中挑选一个。",
                    }

            if hit_entry is None:
                return {
                    "success": False, "ambiguous": False,
                    "message": (
                        f"模板库里没找到匹配 '{template_ref}' 的模板。"
                        " 请用 list_templates_tool 看一下可用模板，或让用户上传新模板。"
                    ),
                }

            local_path, err = _template_api_download(hit_entry["id"], hit_source, user_id)
            if err:
                return {"success": False, "ambiguous": False, "message": err}
            return {
                "success": True,
                "template_path": local_path,
                "source": hit_source,
                "metadata": hit_entry,
                "message": f"已定位模板 '{hit_entry.get('name', hit_entry['id'])}'。",
            }

        def save_template_tool(
            template_path: str,
            name: str,
            description: str,
            category: str = None,
            tags: list = None,
            aliases: list = None,
        ):
            """
            把一份用户**新上传**的 .docx 模板保存进当前用户的模板库，下次
            可以直接通过 list_templates_tool / get_template_path_tool 调用。

            什么时候用：
            - 用户成功用 fill_docx_template_tool 填完一份**新上传**的模板后，
              询问用户"要把这个模板存进你的模板库吗？要起什么名字、分类、别名？"
              用户同意后调本工具。
            - 用户明说"把这个模板保存起来"。
            - **不要**在用户从模板库里取出的模板上重复调本工具——会出现重复
              条目。

            Args:
                template_path: 用户上传的 .docx 文件路径（**模板原件**，不是填写后的 _filled 文件）。
                name: 中文显示名，用户能一眼认出（如 "技术开发合同（3-1）"）。
                description: 一句话说明用途（如 "科技部印制的技术开发委托合同模板"）。
                category: 可选，分类路径（如 "合同/技术开发"）。建议两级，
                    用 / 分隔。
                tags: 可选，关键词列表（如 ["合同", "技术开发", "科技部"]）。
                aliases: 可选，用户日常口语里的别名（如 ["3-1", "3-1技术开发"]），
                    会被 get_template_path_tool 拿来做精确匹配，**强烈建议
                    填几个常用别名**。

            Returns dict with success / template_id / template_path / metadata /
            message。文件会被复制到用户私有库里（不影响原始上传文件）。
            """
            return _template_api_save(
                source_path=template_path,
                user_id=user_id,
                name=name,
                description=description,
                category=category,
                tags=tags,
                aliases=aliases,
                template_id=None,
            )

        def delete_template_tool(template_id: str):
            """
            从**当前用户自己的**模板库里删除一个模板（同时移除 catalog 条目和磁盘上的 .docx）。

            什么时候用：
            - 用户明说"删掉我的 XX 模板""不要这个模板了""把 XX 模板移除"。
            - 用户先用 list_templates_tool 浏览之后明确指向某个模板要求删除。

            ⚠️ 这是**破坏性**操作，文件会被真正删掉，不可撤销。调用前必须：
            1. **先用 `get_template_path_tool` 或 `list_templates_tool` 把用户口中的
               模板解析成一个明确的 id**（确认 source="mine"）。不要凭印象猜 id。
            2. **必须先得到用户的明确确认**——念出要删除的模板名（和 id），
               让用户回答"是/确定/删"之后再调本工具。如果用户只是问"我有哪些模板"
               或泛泛抱怨，**不要**主动调用。
            3. 如果 get_template_path_tool 返回 ambiguous=True，**不要**自己挑——
               把候选念给用户，让用户确认具体是哪一个。

            限制：
            - 只能删除**当前用户自己的**模板（source="mine"）。
            - **共享模板（source="shared"）不可删除**——本工具会返回失败。
              如果用户想删共享模板，告诉他/她需要联系管理员。

            Args:
                template_id: 模板的精确 id（如 "tech-dev-3-1" 或 "cai-gou-he-tong-2"）。
                    通常由 get_template_path_tool 的返回值或 list_templates_tool
                    列表里的 id 字段得来。

            Returns dict with:
                success: True / False
                removed_id: 成功时返回被删除的模板 id
                message: 中文提示

            删除成功后，可以提示用户"已删除，要看看剩下的模板吗？"，
            如果用户想看再调 list_templates_tool。
            """
            return _template_api_delete(template_id=template_id, user_id=user_id)

        def add_bullet_list_tool(file_path: str, items: list, position: int | str = "end", style: str = "List Bullet"):
            """
            Add a bullet list to an existing DOCX file.

            Use this when the user wants to add a bulleted list — for example a list
            of key points, requirements, features, or any collection of items that
            should be visually grouped.

            Best for:
            - "add a bullet list of ..."
            - "insert a list of items"
            - "list the following points"
            - adding structured list content to a document

            Args:
                file_path: Path to the DOCX file.
                items: List of item strings, or list of dicts with 'text' and optional
                       'level' (0=normal, 1+=nested), e.g.
                       ['Item one', {'text': 'Item two', 'level': 1}, 'Item three']
                position: Integer index to insert before that paragraph, or 'end'
                          to append at the end of the document.
                style: Base list style to use; defaults to 'List Bullet'.
            """
            import os
            import json

            print(f"🔧 add_bullet_list_tool called:")
            print(f"   File: {file_path}")
            print(f"   Items: {json.dumps(items[:3], ensure_ascii=False)}...")

            if not os.path.exists(file_path):
                return {
                    'success': False,
                    'error': 'File not found',
                    'message': f'File not found: {file_path}'
                }

            processor = DocumentProcessor(str(WORKSPACE))
            edits = [{
                'type': 'add_bullet_list',
                'items': items,
                'position': position,
                'style': style,
            }]
            result = processor.edit_docx(file_path, edits, overwrite_original=True)

            if result.get('success', False):
                fe_data = _build_files_event_data(file_path, f"Bullet list added to: {Path(file_path).name}")
                if fe_data:
                    _pending_files_events.append(fe_data)

            return {
                'success': result.get('success', False),
                'message': result.get('message', ''),
                'file_path': file_path,
                'items_added': len(items),
            }
        
        def add_numbered_list_tool(file_path: str, items: list, position: int | str = "end", style: str = "List Number"):
            """
            Add a numbered list to an existing DOCX file.

            Use this when the user wants an ordered, sequential list — for example
            steps in a process, a ranked list, or an enumerated set of items.

            Best for:
            - "add a numbered list of ..."
            - "list the steps in order"
            - "add these items as a numbered sequence"
            - "insert a numbered sequence"

            Args:
                file_path: Path to the DOCX file.
                items: List of item strings, or list of dicts with 'text' and optional
                       'level' (0=normal, 1+=nested), e.g.
                       ['Step 1: Do X', {'text': 'Step 2: Do Y', 'level': 1}, 'Step 3: Do Z']
                position: Integer index to insert before that paragraph, or 'end'
                          to append at the end of the document.
                style: Base list style to use; defaults to 'List Number'.
            """
            import os
            import json

            print(f"🔧 add_numbered_list_tool called:")
            print(f"   File: {file_path}")
            print(f"   Items: {json.dumps(items[:3], ensure_ascii=False)}...")

            if not os.path.exists(file_path):
                return {
                    'success': False,
                    'error': 'File not found',
                    'message': f'File not found: {file_path}'
                }

            processor = DocumentProcessor(str(WORKSPACE))
            edits = [{
                'type': 'add_numbered_list',
                'items': items,
                'position': position,
                'style': style,
            }]
            result = processor.edit_docx(file_path, edits, overwrite_original=True)

            if result.get('success', False):
                fe_data = _build_files_event_data(file_path, f"Numbered list added to: {Path(file_path).name}")
                if fe_data:
                    _pending_files_events.append(fe_data)

            return {
                'success': result.get('success', False),
                'message': result.get('message', ''),
                'file_path': file_path,
                'items_added': len(items),
            }
        
        def add_comment_tool(
            file_path: str,
            comments: str = None,
            # Legacy single-comment parameters (for backward compatibility)
            target_text: str = None,
            comment_text: str = None,
            comment_id: int = 0,
            author: str = "DocMaster",
            initials: str = "DM",
            parent_comment_id: int | None = None,
        ):
            """
            Add one or more comments to a DOCX document, attached to specific text ranges.

            Uses a direct zipfile + lxml approach for reliable XML manipulation.
            ALL comments are processed in a SINGLE pass, and ONE file event is emitted
            at the end containing the complete document with all comments.

            Best for:
            - "add comments to ..."
            - "add multiple comments/annotations"
            - "add feedback to this essay"
            - "comment on all sections"
            - "annotate [text] with [comments]"

            Args:
                file_path: Path to the DOCX file to add comment to.
                comments: JSON string or list of comment dicts. Each dict should have:
                    - target_text: The exact text string in the document to attach to
                    - comment_text: The content of the comment
                    - comment_id: Unique integer ID for this comment (0, 1, 2, ...)
                    - author: (optional) Author name, defaults to "DocMaster"
                    - initials: (optional) Author initials, defaults to "DM"
                    - parent_comment_id: (optional) If set, this is a reply to that comment
                    Example: '[{"target_text": "Introduction", "comment_text": "Great intro!", "comment_id": 0}]'
                target_text: (Legacy) Target text for single comment
                comment_text: (Legacy) Comment text for single comment  
                comment_id: (Legacy) Comment ID for single comment
                author: (Legacy) Author for single comment
                initials: (Legacy) Initials for single comment
                parent_comment_id: (Legacy) Parent comment ID for reply
            """
            import json
            import re
            import zipfile
            from lxml import etree
            from datetime import datetime, timezone
            import random

            # Handle batch comments - support both JSON string and list
            if comments is not None:
                if isinstance(comments, str):
                    # Try to parse as JSON, handling curly quotes
                    try:
                        # First, normalize curly quotes to regular quotes for JSON parsing
                        normalized = comments.replace('\u201c', '"').replace('\u201d', '"').replace('\u2018', "'").replace('\u2019', "'")
                        comment_list = json.loads(normalized)
                    except json.JSONDecodeError as e:
                        return {
                            'success': False,
                            'error': f'Invalid JSON in comments parameter: {e}',
                            'message': 'Failed to parse comments JSON. Make sure to use regular double quotes.'
                        }
                elif isinstance(comments, list):
                    comment_list = comments
                else:
                    return {
                        'success': False,
                        'error': 'Invalid comments format',
                        'message': 'comments must be a JSON string or a list'
                    }
            elif target_text is not None and comment_text is not None:
                # Legacy single comment - convert to list format
                comment_list = [{
                    "target_text": target_text,
                    "comment_text": comment_text,
                    "comment_id": comment_id,
                    "author": author,
                    "initials": initials,
                    "parent_comment_id": parent_comment_id,
                }]
            else:
                return {
                    'success': False,
                    'error': 'Invalid arguments',
                    'message': 'Either provide a "comments" list OR both "target_text" and "comment_text"'
                }

            print(f"🔧 add_comment_tool called:")
            print(f"   File: {file_path}")
            print(f"   Comments to add: {len(comment_list)}")

            if not os.path.exists(file_path):
                return {
                    'success': False,
                    'error': 'File not found',
                    'message': f'File not found: {file_path}'
                }

            # Namespace definitions
            NSMAP = {
                'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main',
                'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
                'w14': 'http://schemas.microsoft.com/office/word/2010/wordml',
                'w15': 'http://schemas.microsoft.com/office/word/2012/wordml',
                'w16cid': 'http://schemas.microsoft.com/office/word/2016/wordml/cid',
                'w16cex': 'http://schemas.microsoft.com/office/word/2018/wordml/cex',
                'mc': 'http://schemas.openxmlformats.org/markup-compatibility/2006',
            }

            def qn(tag):
                """Resolve a qualified name like 'w:body' to '{ns}body'."""
                prefix, local = tag.split(':')
                return f'{{{NSMAP[prefix]}}}{local}'

            def _generate_hex_id():
                """Generate a unique hex ID for paraId/durableId."""
                return f"{random.randint(1, 0x7FFFFFFE):08X}"

            def _sanitize_text(text):
                """Replace curly/smart quotes with regular quotes to prevent JSON parsing issues."""
                if not isinstance(text, str):
                    return text
                # Replace curly quotes with regular quotes
                return text.replace('\u201c', '"').replace('\u201d', '"').replace('\u2018', "'").replace('\u2019', "'")

            try:
                # Step 0: Sanitize all comment texts to remove curly quotes
                for comment_spec in comment_list:
                    if 'target_text' in comment_spec:
                        comment_spec['target_text'] = _sanitize_text(comment_spec['target_text'])
                    if 'comment_text' in comment_spec:
                        comment_spec['comment_text'] = _sanitize_text(comment_spec['comment_text'])

                # Step 1: Read the DOCX as a zip file (once)
                with zipfile.ZipFile(file_path, 'r') as zin:
                    doc_xml = zin.read('word/document.xml')
                    has_comments = 'word/comments.xml' in zin.namelist()
                    comments_xml = zin.read('word/comments.xml') if has_comments else None
                    has_comments_extended = 'word/commentsExtended.xml' in zin.namelist()
                    comments_extended_xml = zin.read('word/commentsExtended.xml') if has_comments_extended else None
                    has_comments_ids = 'word/commentsIds.xml' in zin.namelist()
                    comments_ids_xml = zin.read('word/commentsIds.xml') if has_comments_ids else None
                    has_comments_extensible = 'word/commentsExtensible.xml' in zin.namelist()
                    comments_extensible_xml = zin.read('word/commentsExtensible.xml') if has_comments_extensible else None
                    rels_xml = zin.read('word/_rels/document.xml.rels')
                    content_types_xml = zin.read('[Content_Types].xml')
                    
                    other_files = {}
                    for name in zin.namelist():
                        if name not in ('word/document.xml', 'word/comments.xml', 'word/commentsExtended.xml',
                                       'word/commentsIds.xml', 'word/commentsExtensible.xml',
                                       'word/_rels/document.xml.rels', '[Content_Types].xml'):
                            other_files[name] = zin.read(name)

                # Step 2: Parse document.xml and build paragraph text index
                doc_tree = etree.fromstring(doc_xml)
                body = doc_tree.find(qn('w:body'))
                paragraphs = body.findall(qn('w:p'))
                
                # Build a map of paragraph text -> paragraph element for fast lookup
                para_by_text = {}
                for p in paragraphs:
                    texts = list(p.itertext())
                    full_text = ''.join(texts).strip()
                    if full_text:
                        para_by_text[full_text] = p
                        # Also store first 100 chars for partial matching
                        short_text = full_text[:100]
                        if short_text not in para_by_text:
                            para_by_text[short_text] = p

                # Step 3: Process each comment
                added_comments = []
                errors = []
                comment_para_ids = {}  # Maps comment_id -> para_id for reply support
                
                for comment_spec in comment_list:
                    c_target = comment_spec.get('target_text')
                    c_text = comment_spec.get('comment_text')
                    c_id = comment_spec.get('comment_id', 0)
                    c_author = comment_spec.get('author', author)
                    c_initials = comment_spec.get('initials', initials)
                    c_parent = comment_spec.get('parent_comment_id')
                    
                    print(f"   Processing comment #{c_id}: {c_target[:40]}..." if c_target else f"   Processing comment #{c_id}")
                    
                    # Find the target paragraph
                    target_para = None
                    
                    # Try exact match first
                    if c_target in para_by_text:
                        target_para = para_by_text[c_target]
                    else:
                        # Try partial match - look for target text within any paragraph
                        for p in paragraphs:
                            texts = list(p.itertext())
                            full_text = ''.join(texts)
                            if c_target in full_text:
                                target_para = p
                                break
                    
                    if target_para is None:
                        errors.append(f"Comment #{c_id}: Target text not found: '{c_target[:50]}...'")
                        continue

                    # Get runs for inserting markers
                    runs = target_para.findall(qn('w:r'))
                    if not runs:
                        errors.append(f"Comment #{c_id}: No runs in target paragraph")
                        continue

                    first_run = runs[0]
                    para_id = _generate_hex_id()
                    comment_para_ids[c_id] = para_id
                    
                    # Create comment markers
                    comment_range_start = etree.Element(qn('w:commentRangeStart'))
                    comment_range_start.set(qn('w:id'), str(c_id))

                    comment_range_end = etree.Element(qn('w:commentRangeEnd'))
                    comment_range_end.set(qn('w:id'), str(c_id))

                    comment_ref = etree.Element(qn('w:r'))
                    rPr = etree.SubElement(comment_ref, qn('w:rPr'))
                    rStyle = etree.SubElement(rPr, qn('w:rStyle'))
                    rStyle.set(qn('w:val'), 'CommentReference')
                    comment_ref_elem = etree.SubElement(comment_ref, qn('w:commentReference'))
                    comment_ref_elem.set(qn('w:id'), str(c_id))

                    # Insert commentRangeStart before the first run
                    target_para.insert(list(target_para).index(first_run), comment_range_start)

                    # Append commentRangeEnd and commentReference at the end
                    target_para.append(comment_range_end)
                    target_para.append(comment_ref)
                    
                    added_comments.append({
                        'comment_id': c_id,
                        'target_text': c_target,
                        'comment_text': c_text,
                        'author': c_author,
                        'para_id': para_id,
                    })
                    
                    print(f"      ✓ Added markers for comment #{c_id}")

                # Step 4: Update comments.xml (main comment storage)
                if has_comments and comments_xml:
                    comments_tree = etree.fromstring(comments_xml)
                else:
                    comments_tree = etree.Element(qn('w:comments'))
                    for prefix, uri in NSMAP.items():
                        etree.register_namespace(prefix, uri)

                timestamp = datetime.now(timezone.utc).strftime("%Y-%m-%dT%H:%M:%SZ")
                
                for added in added_comments:
                    # Create comment element
                    new_comment = etree.SubElement(comments_tree, qn('w:comment'))
                    new_comment.set(qn('w:id'), str(added['comment_id']))
                    new_comment.set(qn('w:author'), added['author'])
                    new_comment.set(qn('w:initials'), added.get('initials', initials))
                    new_comment.set(qn('w:date'), timestamp)

                    # Create paragraph inside comment
                    comment_p = etree.SubElement(new_comment, qn('w:p'))
                    comment_p.set(qn('w14:paraId'), added['para_id'])
                    comment_p.set(qn('w14:textId'), '77777777')
                    
                    comment_r = etree.SubElement(comment_p, qn('w:r'))
                    comment_t = etree.SubElement(comment_r, qn('w:t'))
                    comment_t.set('{http://www.w3.org/XML/1998/namespace}space', 'preserve')
                    comment_t.text = added['comment_text']

                new_comments_xml = etree.tostring(comments_tree, xml_declaration=True, encoding='UTF-8', standalone=True)

                # Step 5: Update commentsExtended.xml (for reply threading and done status)
                if has_comments_extended and comments_extended_xml:
                    ext_tree = etree.fromstring(comments_extended_xml)
                else:
                    ext_tree = etree.Element(qn('w15:commentsEx'))
                    etree.register_namespace('w15', NSMAP['w15'])

                for added in added_comments:
                    # Find parent para_id if this is a reply
                    parent_para = None
                    for spec in comment_list:
                        if spec.get('comment_id') == spec.get('parent_comment_id'):
                            parent_para = comment_para_ids.get(spec.get('comment_id'))
                            break
                    
                    comment_ex = etree.SubElement(ext_tree, qn('w15:commentEx'))
                    comment_ex.set(qn('w15:paraId'), added['para_id'])
                    if parent_para:
                        comment_ex.set(qn('w15:paraIdParent'), parent_para)
                    comment_ex.set(qn('w15:done'), '0')

                new_comments_extended_xml = etree.tostring(ext_tree, xml_declaration=True, encoding='UTF-8', standalone=True)

                # Step 6: Update commentsIds.xml (for comment identity tracking)
                if has_comments_ids and comments_ids_xml:
                    ids_tree = etree.fromstring(comments_ids_xml)
                else:
                    ids_tree = etree.Element(qn('w16cid:commentsIds'))
                    etree.register_namespace('w16cid', NSMAP.get('w16cid', 'http://schemas.microsoft.com/office/word/2016/wordml/cid'))

                for added in added_comments:
                    durable_id = _generate_hex_id()
                    comment_id_elem = etree.SubElement(ids_tree, qn('w16cid:commentId'))
                    comment_id_elem.set(qn('w16cid:paraId'), added['para_id'])
                    comment_id_elem.set(qn('w16cid:durableId'), durable_id)

                new_comments_ids_xml = etree.tostring(ids_tree, xml_declaration=True, encoding='UTF-8', standalone=True)

                # Step 7: Update commentsExtensible.xml (for Office 365 collaboration)
                if has_comments_extensible and comments_extensible_xml:
                    ext2_tree = etree.fromstring(comments_extensible_xml)
                else:
                    ext2_tree = etree.Element(qn('w16cex:commentsExtensible'))
                    etree.register_namespace('w16cex', NSMAP.get('w16cex', 'http://schemas.microsoft.com/office/word/2018/wordml/cex'))

                for added in added_comments:
                    comment_ext = etree.SubElement(ext2_tree, qn('w16cex:commentExtensible'))
                    durable_id = _generate_hex_id()
                    comment_ext.set(qn('w16cex:durableId'), durable_id)
                    comment_ext.set(qn('w16cex:dateUtc'), timestamp)

                new_comments_extensible_xml = etree.tostring(ext2_tree, xml_declaration=True, encoding='UTF-8', standalone=True)

                # Step 8: Update relationships
                rels_tree = etree.fromstring(rels_xml)
                
                # Check which comment relationships already exist
                existing_targets = {rel.get('Target') for rel in rels_tree}
                
                if 'comments.xml' not in existing_targets:
                    max_id = 0
                    for rel in rels_tree:
                        rid = rel.get('Id', '')
                        if rid.startswith('rId'):
                            try:
                                max_id = max(max_id, int(rid[3:]))
                            except ValueError:
                                pass
                    
                    comment_rels = [
                        ('comments', 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/comments'),
                        ('commentsExtended', 'http://schemas.microsoft.com/office/2011/relationships/commentsExtended'),
                        ('commentsIds', 'http://schemas.microsoft.com/office/2016/09/relationships/commentsIds'),
                        ('commentsExtensible', 'http://schemas.microsoft.com/office/2018/08/relationships/commentsExtensible'),
                    ]
                    
                    for name, rel_type in comment_rels:
                        new_rid = f'rId{max_id + 1}'
                        new_rel = etree.SubElement(rels_tree, 'Relationship')
                        new_rel.set('Id', new_rid)
                        new_rel.set('Type', rel_type)
                        new_rel.set('Target', f'{name}.xml')
                        max_id += 1

                new_rels_xml = etree.tostring(rels_tree, xml_declaration=True, encoding='UTF-8', standalone=True)

                # Step 9: Update [Content_Types].xml
                ct_tree = etree.fromstring(content_types_xml)
                ns_ct = 'http://schemas.openxmlformats.org/package/2006/content-types'
                existing_parts = {child.get('PartName') for child in ct_tree}
                
                if '/word/comments.xml' not in existing_parts:
                    overrides = [
                        ('/word/comments.xml', 'application/vnd.openxmlformats-officedocument.wordprocessingml.comments+xml'),
                        ('/word/commentsExtended.xml', 'application/vnd.openxmlformats-officedocument.wordprocessingml.commentsExtended+xml'),
                        ('/word/commentsIds.xml', 'application/vnd.openxmlformats-officedocument.wordprocessingml.commentsIds+xml'),
                        ('/word/commentsExtensible.xml', 'application/vnd.openxmlformats-officedocument.wordprocessingml.commentsExtensible+xml'),
                    ]
                    
                    for part_name, content_type in overrides:
                        override = etree.SubElement(ct_tree, f'{{{ns_ct}}}Override')
                        override.set('PartName', part_name)
                        override.set('ContentType', content_type)

                new_content_types_xml = etree.tostring(ct_tree, xml_declaration=True, encoding='UTF-8', standalone=True)

                # Step 10: Write document.xml back
                new_doc_xml = etree.tostring(doc_tree, xml_declaration=True, encoding='UTF-8', standalone=True)

                # Step 11: Write back to docx (ONE file write)
                tmp_path = file_path + '.tmp'
                with zipfile.ZipFile(tmp_path, 'w', zipfile.ZIP_DEFLATED) as zout:
                    zout.writestr('word/document.xml', new_doc_xml)
                    zout.writestr('word/comments.xml', new_comments_xml)
                    zout.writestr('word/commentsExtended.xml', new_comments_extended_xml)
                    zout.writestr('word/commentsIds.xml', new_comments_ids_xml)
                    zout.writestr('word/commentsExtensible.xml', new_comments_extensible_xml)
                    zout.writestr('word/_rels/document.xml.rels', new_rels_xml)
                    zout.writestr('[Content_Types].xml', new_content_types_xml)
                    for name, data in other_files.items():
                        zout.writestr(name, data)

                # Replace original
                os.replace(tmp_path, file_path)
                
                print(f"   ✅ Added {len(added_comments)} comment(s) successfully!")

                # Step 12: Emit SINGLE file event after ALL comments are added
                fe_data = _build_files_event_data(file_path, f"Comments added to: {Path(file_path).name}")
                if fe_data:
                    _pending_files_events.append(fe_data)

                return {
                    'success': len(errors) == 0,
                    'message': f'Added {len(added_comments)} comment(s) to document',
                    'comments_added': len(added_comments),
                    'errors': errors if errors else None,
                    'file_path': file_path
                }

            except Exception as e:
                import traceback
                traceback.print_exc()
                return {
                    'success': False,
                    'error': str(e),
                    'message': f'Failed to add comments: {e}'
                }

        def remove_comment_tool(
            file_path: str,
            comment_ids: list = None,
            # Legacy single comment_id parameter
            comment_id: int = None,
        ):
            """
            Remove one or more comments from a DOCX document.

            This tool removes the comment(s) with the specified ID(s) from the document,
            including:
            - Comment markers from document.xml (commentRangeStart, commentRangeEnd, commentReference)
            - The comment entry from comments.xml, commentsExtended.xml, commentsIds.xml, commentsExtensible.xml
            - Cleanup of relationships and content types if no comments remain

            ALL comments are processed in a SINGLE pass, and ONE file event is emitted
            at the end after all comments are removed.

            Best for:
            - "remove all comments"
            - "clear all annotations"
            - "delete comments #0, #1, #2"
            - "清除所有批注"
            - "删除第N条批注"

            Args:
                file_path: Path to the DOCX file to remove comment from.
                comment_ids: List of comment IDs to remove. Example: [0, 1, 2] removes comments with IDs 0, 1, and 2.
                    Use "all" as a special value to remove all comments.
                comment_id: (Legacy) Single comment ID to remove.
            """
            import zipfile
            from lxml import etree

            # Handle batch comment IDs
            remove_all = False
            ids_to_remove = []
            
            if comment_ids is not None:
                # Normalize to handle various input types
                if isinstance(comment_ids, str):
                    if comment_ids.lower() == "all":
                        remove_all = True
                        print(f"   Mode: Remove ALL comments")
                    else:
                        # Single ID as string
                        ids_to_remove = [comment_ids]
                elif isinstance(comment_ids, list):
                    if len(comment_ids) == 0:
                        return {
                            'success': False,
                            'error': 'Empty list',
                            'message': 'comment_ids list is empty'
                        }
                    ids_to_remove = [str(i) for i in comment_ids]
                else:
                    ids_to_remove = [str(comment_ids)]
            elif comment_id is not None:
                ids_to_remove = [str(comment_id)]
            else:
                return {
                    'success': False,
                    'error': 'Invalid arguments',
                    'message': 'Either provide "comment_ids" list or "comment_id"'
                }

            print(f"🔧 remove_comment_tool called:")
            print(f"   File: {file_path}")
            print(f"   remove_all: {remove_all}")
            print(f"   ids_to_remove: {ids_to_remove}")

            if not os.path.exists(file_path):
                return {
                    'success': False,
                    'error': 'File not found',
                    'message': f'File not found: {file_path}'
                }

            # Namespace definitions
            NSMAP = {
                'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main',
                'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
                'w14': 'http://schemas.microsoft.com/office/word/2010/wordml',
                'w15': 'http://schemas.microsoft.com/office/word/2012/wordml',
                'w16cid': 'http://schemas.microsoft.com/office/word/2016/wordml/cid',
                'mc': 'http://schemas.openxmlformats.org/markup-compatibility/2006',
            }

            def qn(tag):
                """Resolve a qualified name like 'w:body' to '{ns}body'."""
                prefix, local = tag.split(':')
                return f'{{{NSMAP[prefix]}}}{local}'

            try:
                # Step 1: Read the DOCX as a zip file (once)
                with zipfile.ZipFile(file_path, 'r') as zin:
                    has_comments = 'word/comments.xml' in zin.namelist()
                    has_comments_extended = 'word/commentsExtended.xml' in zin.namelist()
                    has_comments_ids = 'word/commentsIds.xml' in zin.namelist()
                    has_comments_extensible = 'word/commentsExtensible.xml' in zin.namelist()
                    
                    if not has_comments:
                        return {
                            'success': False,
                            'error': 'No comments in document',
                            'message': 'The document has no comments to remove'
                        }
                    
                    doc_xml = zin.read('word/document.xml')
                    comments_xml = zin.read('word/comments.xml')
                    comments_extended_xml = zin.read('word/commentsExtended.xml') if has_comments_extended else None
                    comments_ids_xml = zin.read('word/commentsIds.xml') if has_comments_ids else None
                    comments_extensible_xml = zin.read('word/commentsExtensible.xml') if has_comments_extensible else None
                    rels_xml = zin.read('word/_rels/document.xml.rels')
                    content_types_xml = zin.read('[Content_Types].xml')
                    
                    other_files = {}
                    for name in zin.namelist():
                        if name not in ('word/document.xml', 'word/comments.xml', 'word/commentsExtended.xml',
                                       'word/commentsIds.xml', 'word/commentsExtensible.xml',
                                       'word/_rels/document.xml.rels', '[Content_Types].xml'):
                            other_files[name] = zin.read(name)

                # Step 2: If removing all, get all comment IDs first
                if remove_all:
                    comments_tree = etree.fromstring(comments_xml)
                    all_comments = comments_tree.findall(qn('w:comment'))
                    print(f"   Found {len(all_comments)} comments in document")
                    for c in all_comments:
                        cid = c.get(qn('w:id'))
                        print(f"      - Comment ID: {cid}")
                        ids_to_remove.append(cid)
                    print(f"   ids_to_remove populated: {ids_to_remove}")

                # Step 3: Parse and clean document.xml - remove comment markers for ALL IDs
                doc_tree = etree.fromstring(doc_xml)
                
                removed_markers = 0
                for c_id in ids_to_remove:
                    # Remove commentRangeStart elements
                    for elem in doc_tree.findall(f'.//{qn("w:commentRangeStart")}'):
                        if elem.get(qn('w:id')) == c_id:
                            # Find parent paragraph and remove
                            for para in doc_tree.findall(f'.//{qn("w:p")}'):
                                if elem in list(para):
                                    para.remove(elem)
                                    removed_markers += 1
                                    break
                    
                    # Remove commentRangeEnd elements
                    for elem in doc_tree.findall(f'.//{qn("w:commentRangeEnd")}'):
                        if elem.get(qn('w:id')) == c_id:
                            for para in doc_tree.findall(f'.//{qn("w:p")}'):
                                if elem in list(para):
                                    para.remove(elem)
                                    removed_markers += 1
                                    break
                    
                    # Remove commentReference runs (entire w:r containing the reference)
                    for para in doc_tree.findall(f'.//{qn("w:p")}'):
                        for run in list(para.findall(qn('w:r'))):
                            comment_ref = run.find(qn('w:commentReference'))
                            if comment_ref is not None and comment_ref.get(qn('w:id')) == c_id:
                                para.remove(run)
                                removed_markers += 1

                new_doc_xml = etree.tostring(doc_tree, xml_declaration=True, encoding='UTF-8', standalone=True)
                print(f"   Removed {removed_markers} comment marker(s) from document.xml")

                # Step 4: Parse comments.xml and remove all target comments
                comments_tree = etree.fromstring(comments_xml)
                removed_comments = []
                
                print(f"   Step 4: Looking for comment IDs: {ids_to_remove}")
                
                for comment in list(comments_tree.findall(qn('w:comment'))):
                    c_id = comment.get(qn('w:id'))
                    if c_id in ids_to_remove:
                        # Get paraId before removing for later use in extended files
                        para_elem = comment.find(qn('w:p'))
                        para_id = para_elem.get(qn('w14:paraId')) if para_elem is not None else None
                        
                        comments_tree.remove(comment)
                        removed_comments.append({'id': c_id, 'para_id': para_id})
                        print(f"   Removed comment #{c_id} from comments.xml")

                if not removed_comments:
                    return {
                        'success': False,
                        'error': 'Comment not found',
                        'message': f'Comment(s) not found in document'
                    }

                new_comments_xml = etree.tostring(comments_tree, xml_declaration=True, encoding='UTF-8', standalone=True)

                # Step 5: Update commentsExtended.xml
                new_comments_extended_xml = None
                if has_comments_extended and comments_extended_xml:
                    ext_tree = etree.fromstring(comments_extended_xml)
                    for para_id in [c['para_id'] for c in removed_comments if c['para_id']]:
                        for ex in ext_tree.findall(qn('w15:commentEx')):
                            if ex.get(qn('w15:paraId')) == para_id:
                                ext_tree.remove(ex)
                                break
                    new_comments_extended_xml = etree.tostring(ext_tree, xml_declaration=True, encoding='UTF-8', standalone=True)

                # Step 6: Update commentsIds.xml
                new_comments_ids_xml = None
                if has_comments_ids and comments_ids_xml:
                    ids_tree = etree.fromstring(comments_ids_xml)
                    for c in removed_comments:
                        for cid in list(ids_tree.findall(qn('w16cid:commentId'))):
                            if cid.get(qn('w16cid:paraId')) == c['para_id']:
                                ids_tree.remove(cid)
                                break
                    new_comments_ids_xml = etree.tostring(ids_tree, xml_declaration=True, encoding='UTF-8', standalone=True)

                # Step 7: Update commentsExtensible.xml
                new_comments_extensible_xml = None
                if has_comments_extensible and comments_extensible_xml:
                    ext2_tree = etree.fromstring(comments_extensible_xml)
                    # Remove all extensible entries for removed comments (we need durableId matching which is complex)
                    # For simplicity, we'll just regenerate if empty
                    if len(ext2_tree) == 0:
                        new_comments_extensible_xml = None
                        has_comments_extensible = False

                # Step 8: Check if there are any remaining comments
                remaining_comments = comments_tree.findall(qn('w:comment'))
                
                if remaining_comments:
                    # There are still comments, keep all comment XML files
                    keep_comments = True
                else:
                    # No more comments, remove all comment XML files and clean up
                    keep_comments = False
                    print(f"   No more comments, will clean up all comment files")

                # Step 9: Update relationships if removing comments entirely
                rels_tree = etree.fromstring(rels_xml)
                comment_rel_types = [
                    'http://schemas.openxmlformats.org/officeDocument/2006/relationships/comments',
                    'http://schemas.microsoft.com/office/2011/relationships/commentsExtended',
                    'http://schemas.microsoft.com/office/2016/09/relationships/commentsIds',
                    'http://schemas.microsoft.com/office/2018/08/relationships/commentsExtensible',
                ]
                
                if not keep_comments:
                    # Remove all comment relationships
                    for rel in list(rels_tree):
                        if rel.get('Type', '') in comment_rel_types:
                            rels_tree.remove(rel)
                    print(f"   Removed all comment relationships")
                
                new_rels_xml = etree.tostring(rels_tree, xml_declaration=True, encoding='UTF-8', standalone=True)

                # Step 10: Update [Content_Types].xml
                ct_tree = etree.fromstring(content_types_xml)
                ns_ct = 'http://schemas.openxmlformats.org/package/2006/content-types'
                
                comment_parts = [
                    '/word/comments.xml',
                    '/word/commentsExtended.xml',
                    '/word/commentsIds.xml',
                    '/word/commentsExtensible.xml',
                ]
                
                if not keep_comments:
                    # Remove all comment content types
                    for child in list(ct_tree):
                        if child.get('PartName') in comment_parts:
                            ct_tree.remove(child)
                    print(f"   Removed all comment content types")
                
                new_content_types_xml = etree.tostring(ct_tree, xml_declaration=True, encoding='UTF-8', standalone=True)

                # Step 11: Write back to docx (ONE file write)
                tmp_path = file_path + '.tmp'
                with zipfile.ZipFile(tmp_path, 'w', zipfile.ZIP_DEFLATED) as zout:
                    zout.writestr('word/document.xml', new_doc_xml)
                    if keep_comments:
                        zout.writestr('word/comments.xml', new_comments_xml)
                        if new_comments_extended_xml:
                            zout.writestr('word/commentsExtended.xml', new_comments_extended_xml)
                        if new_comments_ids_xml:
                            zout.writestr('word/commentsIds.xml', new_comments_ids_xml)
                    zout.writestr('word/_rels/document.xml.rels', new_rels_xml)
                    zout.writestr('[Content_Types].xml', new_content_types_xml)
                    for name, data in other_files.items():
                        zout.writestr(name, data)

                os.replace(tmp_path, file_path)
                print(f"   ✅ Removed {len(removed_comments)} comment(s) successfully!")

                # Step 12: Emit SINGLE file event after all comments removed
                fe_data = _build_files_event_data(file_path, f"Comments removed from: {Path(file_path).name}")
                if fe_data:
                    _pending_files_events.append(fe_data)

                return {
                    'success': True,
                    'message': f'Removed {len(removed_comments)} comment(s) from document',
                    'comments_removed': len(removed_comments),
                    'file_path': file_path
                }

            except Exception as e:
                import traceback
                traceback.print_exc()
                return {
                    'success': False,
                    'error': str(e),
                    'message': f'Failed to remove comment(s): {e}'
                }

        def review_contract_tool(file_path: str, annotate: bool = True):
            """
            合同审查工具：对一份 .docx 合同做四方面体检并（可选）输出带批注的副本。

            审查维度：
              1. 格式：字体/字号一致性、中英文混排、半角/全角标点、条款编号风格统一。
              2. 填写缺失：未替换的 ____、{{...}}、**XX**、空白字段（甲方:）、空白日期。
              3. 内容一致性：大写/小写金额对账、条款编号连续性、对『第X条』/『附件X』
                 的悬空引用、甲乙方主体名称在全文中是否一致。
              4. 法律风险：用模型做红线扫描（缺失条款、不公平条款、模糊措辞、合规问题）。
                 fail-soft —— 若 LLM 30s 内未返回，仅返回前三类启发式结果，并在
                 summary 中注明。

            产物：
              - 一份结构化的中文报告（issues 列表 + stats + summary）；
              - 当 annotate=True 时，把 issues 转为 Word 批注，写到一份新文件
                `<原名>_审查.docx`（不覆盖原文件），并发出 FilesEvent 让用户下载。

            适用场景：
              - "帮我审查这份合同 / 帮我看看这份合同"
              - "检查格式问题"
              - "合同里有没有空着的字段 / 大小写金额对不对"

            Args:
                file_path: 合同 .docx 路径（必须在工作区内）。
                annotate:  是否同时输出带 Word 批注的副本，默认 True。
            """
            src = Path(file_path)
            if not src.is_file():
                return {"success": False, "message": f"文件不存在: {file_path}"}
            if src.suffix.lower() != ".docx":
                return {"success": False, "message": f"仅支持 .docx 文件（当前 {src.suffix}）。"}

            # Build a sync llm_call around the agent's model_client. Fail-soft: any
            # error in the wrapper (timeout, network, JSON, etc.) is caught inside
            # the skill itself, which then returns the heuristic-only result with a
            # note in summary.
            def _llm_call(prompt: str) -> str:
                import asyncio
                from autogen_core.models import UserMessage

                client = set_model_client(default_config_name)

                async def _run():
                    result = await client.create([UserMessage(content=prompt, source="docmaster")])
                    content = getattr(result, "content", "") or ""
                    if isinstance(content, list):
                        # Some clients return a list of content blocks.
                        parts = []
                        for c in content:
                            if isinstance(c, str):
                                parts.append(c)
                            elif isinstance(c, dict) and "text" in c:
                                parts.append(c["text"])
                        content = "".join(parts)
                    return str(content or "")

                try:
                    return asyncio.run(asyncio.wait_for(_run(), timeout=30.0))
                except RuntimeError:
                    # We are already inside a running event loop — run the
                    # coroutine in a separate thread with its own loop.
                    import concurrent.futures
                    def _bg():
                        loop = asyncio.new_event_loop()
                        try:
                            return loop.run_until_complete(asyncio.wait_for(_run(), timeout=30.0))
                        finally:
                            loop.close()
                    with concurrent.futures.ThreadPoolExecutor(max_workers=1) as ex:
                        return ex.submit(_bg).result(timeout=35.0)

            skill = ContractReviewSkill(str(WORKSPACE))
            result = skill.review(str(src), llm_call=_llm_call)
            if not result.get("success"):
                return result

            issues = result.get("issues") or []
            annotated_path = None
            annotate_note = ""

            if annotate and issues:
                # Copy source → <stem>_审查.docx in WORKDIR, then add comments.
                import shutil
                target = WORKDIR / f"{src.stem}_审查{src.suffix}"
                try:
                    shutil.copyfile(src, target)
                except Exception as exc:
                    annotate_note = f"复制副本失败，未生成带批注文档: {exc}"
                    target = None

                if target is not None:
                    comments_payload = []
                    cid = 0
                    for it in issues:
                        ct = (it.get("comment_target") or "").strip()
                        if not ct:
                            continue
                        # Word can't anchor a comment on whitespace-only ranges;
                        # skip those.
                        if not ct.strip():
                            continue
                        sev_tag = {"high": "高", "medium": "中", "low": "低"}.get(
                            it.get("severity", "medium"), "中"
                        )
                        body = f"[{sev_tag}/{it.get('category', '')}] {it.get('message', '')}"
                        suggestion = (it.get("suggestion") or "").strip()
                        if suggestion:
                            body += f"\n建议：{suggestion}"
                        comments_payload.append({
                            "target_text": ct,
                            "comment_text": body,
                            "comment_id": cid,
                            "author": "DocMaster",
                            "initials": "DM",
                        })
                        cid += 1

                    if comments_payload:
                        ac_result = add_comment_tool(
                            file_path=str(target),
                            comments=comments_payload,
                        )
                        if ac_result.get("success"):
                            annotated_path = str(target)
                        else:
                            annotate_note = (
                                f"批注写入失败：{ac_result.get('message', '未知错误')}"
                            )
                    else:
                        # No anchorable comments — still emit the copy so the user
                        # can read/edit alongside the chat report.
                        fe_data = _build_files_event_data(
                            str(target), f"Contract review copy: {target.name}"
                        )
                        if fe_data:
                            _pending_files_events.append(fe_data)
                        annotated_path = str(target)

            out = {
                "success": True,
                "summary": result.get("summary", ""),
                "stats": result.get("stats", {}),
                "issues": issues,
                "annotated_path": annotated_path,
                "source_path": str(src),
            }
            if annotate_note:
                out["annotate_note"] = annotate_note
            return out

        # ============ DOCX SKILL TOOLS (XML-level, formatting-safe) ============
        # These tools wrap the scripts in skills/docx/scripts/ and provide
        # a formatting-preserving workflow: unpack → edit XML → repack.

        _SKILL_SCRIPTS_DIR = Path(__file__).parent / "skills" / "docx" / "scripts"
        _OFFICE_SCRIPTS_DIR = _SKILL_SCRIPTS_DIR / "office"

        def unpack_docx_tool(file_path: str, output_dir: str, merge_runs: bool = True, simplify_redlines: bool = True):
            """
            Unpack a DOCX file into a directory of editable XML files.

            This is the first step of the formatting-safe editing workflow.
            It extracts the ZIP archive, pretty-prints XML, merges adjacent
            runs with identical formatting, simplifies tracked changes, and
            escapes smart quotes for safe editing.

            Use this when you need to make advanced edits that python-docx
            cannot handle without formatting loss, such as:
            - Complex tracked changes (deletions + insertions with author attribution)
            - Advanced comments via the comment.py script
            - Preserving every formatting detail during text edits

            After unpacking, edit the XML files in output_dir/word/ directly,
            then use pack_docx_tool to reassemble.

            Args:
                file_path: Path to the DOCX file to unpack.
                output_dir: Directory to extract into (will be created).
                merge_runs: Merge adjacent runs with identical formatting (default True).
                simplify_redlines: Simplify adjacent tracked changes from same author (default True).
            """
            import sys as _sys
            _saved = _sys.path[:]
            _sys.path.insert(0, str(_OFFICE_SCRIPTS_DIR))
            try:
                from unpack import unpack
                _, message = unpack(file_path, output_dir, merge_runs=merge_runs, simplify_redlines=simplify_redlines)
                success = "Error" not in message
                return {'success': success, 'message': message, 'output_dir': output_dir}
            except Exception as e:
                return {'success': False, 'error': str(e), 'message': f'Failed to unpack: {e}'}
            finally:
                _sys.path[:] = _saved

        def pack_docx_tool(input_dir: str, output_file: str, original_file: str = None, validate: bool = True):
            """
            Pack an unpacked directory back into a DOCX file.

            This is the final step of the formatting-safe editing workflow.
            It validates the XML with auto-repair, condenses formatting, and
            creates the output DOCX.

            Auto-repair fixes:
            - durableId values that exceed OOXML limits (regenerates valid IDs)
            - Missing xml:space="preserve" on <w:t> elements with whitespace

            Args:
                input_dir: Path to the unpacked directory (from unpack_docx_tool).
                output_file: Path for the output DOCX file.
                original_file: (Optional) Path to the original DOCX for validation comparison.
                validate: Run schema validation with auto-repair (default True).
            """
            import sys as _sys
            _saved = _sys.path[:]
            _sys.path.insert(0, str(_OFFICE_SCRIPTS_DIR))
            try:
                from pack import pack
                _, message = pack(input_dir, output_file, original_file=original_file, validate=validate)
                success = "Error" not in message
                if success:
                    fe_data = _build_files_event_data(output_file, f"Packed DOCX: {Path(output_file).name}")
                    if fe_data:
                        _pending_files_events.append(fe_data)
                return {'success': success, 'message': message, 'output_file': output_file}
            except Exception as e:
                return {'success': False, 'error': str(e), 'message': f'Failed to pack: {e}'}
            finally:
                _sys.path[:] = _saved

        def validate_docx_tool(path: str, original: str = None, auto_repair: bool = True):
            """
            Validate a DOCX file or unpacked directory against OOXML schemas.

            Can validate either a packed .docx file or an unpacked directory.
            Optionally compares against the original file to check for
            tracked-change consistency.

            Args:
                path: Path to the unpacked directory or packed DOCX file.
                original: (Optional) Path to the original DOCX for comparison.
                auto_repair: Automatically repair common issues (default True).
            """
            import sys as _sys
            import tempfile
            import zipfile
            _saved = _sys.path[:]
            _sys.path.insert(0, str(_OFFICE_SCRIPTS_DIR))
            try:
                from validators import DOCXSchemaValidator, RedliningValidator
                unpacked_dir = Path(path)
                original_file = Path(original) if original else None

                # If path is a .docx file, unpack to temp dir
                if unpacked_dir.is_file() and unpacked_dir.suffix.lower() == '.docx':
                    temp_dir = tempfile.mkdtemp()
                    with zipfile.ZipFile(unpacked_dir, 'r') as zf:
                        zf.extractall(temp_dir)
                    unpacked_dir = Path(temp_dir)

                validators = [DOCXSchemaValidator(unpacked_dir, original_file)]
                if original_file and original_file.exists():
                    validators.append(RedliningValidator(unpacked_dir, original_file, author="DocMaster"))

                output_lines = []
                if auto_repair:
                    total_repairs = sum(v.repair() for v in validators)
                    if total_repairs:
                        output_lines.append(f"Auto-repaired {total_repairs} issue(s)")

                success = all(v.validate() for v in validators)
                if success:
                    output_lines.append("All validations PASSED!")

                return {'success': success, 'message': '\n'.join(output_lines)}
            except Exception as e:
                return {'success': False, 'error': str(e), 'message': f'Validation failed: {e}'}
            finally:
                _sys.path[:] = _saved

        def accept_tracked_changes_tool(input_file: str, output_file: str):
            """
            Accept all tracked changes in a DOCX file using LibreOffice.

            Produces a clean document with all insertions accepted and
            deletions removed. Requires LibreOffice to be installed.

            Args:
                input_file: Path to the DOCX file with tracked changes.
                output_file: Path for the clean output DOCX.
            """
            import sys as _sys
            _saved = _sys.path[:]
            _sys.path.insert(0, str(_SKILL_SCRIPTS_DIR))
            _sys.path.insert(0, str(_OFFICE_SCRIPTS_DIR))
            try:
                from accept_changes import accept_changes
                _, message = accept_changes(input_file, output_file)
                success = "Error" not in message
                if success:
                    fe_data = _build_files_event_data(output_file, f"Tracked changes accepted: {Path(output_file).name}")
                    if fe_data:
                        _pending_files_events.append(fe_data)
                return {'success': success, 'message': message, 'output_file': output_file}
            except Exception as e:
                return {'success': False, 'error': str(e), 'message': f'Failed to accept changes: {e}'}
            finally:
                _sys.path[:] = _saved

        def add_xml_comment_tool(unpacked_dir: str, comment_id: int, text: str, author: str = "DocMaster", initials: str = "DM", parent_id: int = None):
            """
            Add a comment to an unpacked DOCX directory (XML-level).

            This handles all the boilerplate across comments.xml,
            commentsExtended.xml, commentsIds.xml, and commentsExtensible.xml.

            After calling this, you still need to add markers to document.xml:
              <w:commentRangeStart w:id="ID"/>
              ...commented runs...
              <w:commentRangeEnd w:id="ID"/>
              <w:r><w:rPr><w:rStyle w:val="CommentReference"/></w:rPr><w:commentReference w:id="ID"/></w:r>

            Args:
                unpacked_dir: Path to the unpacked DOCX directory.
                comment_id: Unique integer ID for this comment.
                text: Comment text (should be pre-escaped XML).
                author: Author name (default "DocMaster").
                initials: Author initials (default "DM").
                parent_id: (Optional) Parent comment ID for replies.
            """
            import sys as _sys
            _saved = _sys.path[:]
            _sys.path.insert(0, str(_SKILL_SCRIPTS_DIR))
            _sys.path.insert(0, str(_OFFICE_SCRIPTS_DIR))
            try:
                from comment import add_comment
                para_id, message = add_comment(unpacked_dir, comment_id, text, author, initials, parent_id)
                success = "Error" not in message
                result = {'success': success, 'message': message, 'para_id': para_id}
                if parent_id is not None:
                    result['marker_hint'] = f'Nest markers: commentRangeStart id="{parent_id}" then id="{comment_id}"'
                else:
                    result['marker_hint'] = f'Add markers: commentRangeStart/End id="{comment_id}" around target text'
                return result
            except Exception as e:
                return {'success': False, 'error': str(e), 'message': f'Failed to add comment: {e}'}
            finally:
                _sys.path[:] = _saved

        # Aliases for names the model frequently hallucinates. Each is a thin
        # wrapper around `edit_docx_tool` so a wrong-name call still succeeds
        # instead of triggering a multi-turn "tool not found" recovery loop
        # (observed 2026-05 — agent typed `edit_docx_content_tool` and
        # `edit_docx_content`, then fell back to forbidden run_bash).
        def edit_docx_content_tool(file_path: str, edits: list):
            """Alias for edit_docx_tool — accepts the same arguments. Prefer
            calling edit_docx_tool directly; this alias exists only so that
            calls written under the wrong name still work."""
            return edit_docx_tool(file_path, edits)

        def edit_docx_content(file_path: str, edits: list):
            """Alias for edit_docx_tool — see edit_docx_content_tool."""
            return edit_docx_tool(file_path, edits)

        # ============ PPT SKILL TOOLS (Phase 1) ============
        # These wrap the `ppt-polished-deck-collab` skill at
        # `skills/presentation-skills/ppt-polished-deck-collab-traditional/`.
        # The agent MUST NOT run scripts/*.py via run_bash — use these tools so
        # script paths, JSON report locations and workspace structure stay
        # consistent.

        _PPT_REFERENCE_NAMES = {
            "principles",
            "deck_workflow",
            "technical_support",
            "design_support",
            "slide_design_system",
            "quality_gates",
            "build_routes",
            "diagram_support",
            "office_chart_support",
            "python_figure_support",
            "icon_system",
        }

        def _ppt_user_workdir() -> Path:
            """Return the current user's workdir (matches DocMasterAgent)."""
            sub = user_id or "_default"
            wd = WORKDIR / sub
            wd.mkdir(parents=True, exist_ok=True)
            return wd

        def _run_ppt_script(
            script_name: str,
            args: list,
            *,
            timeout: int = 300,
            capture_json: Path | None = None,
        ) -> dict:
            """Run a script from PPT_SCRIPTS_DIR with cwd set to that dir.

            cwd is fixed so the scripts' internal sibling imports (e.g.
            `from ppt_quality_helpers import ...`) resolve. stdout/stderr are
            tail-trimmed to 2000 chars each to keep the model's context window
            sane; full reports live on disk (see `capture_json`).
            """
            import subprocess
            import sys as _sys

            script_path = PPT_SCRIPTS_DIR / script_name
            if not script_path.exists():
                return {
                    "success": False,
                    "error": "Script not found",
                    "message": (
                        f"{script_name} does not exist under {PPT_SCRIPTS_DIR}. "
                        "The PPT skill may be missing or named differently."
                    ),
                }
            try:
                proc = subprocess.run(
                    [_sys.executable, str(script_path), *map(str, args)],
                    cwd=str(PPT_SCRIPTS_DIR),
                    capture_output=True,
                    text=True,
                    timeout=timeout,
                )
            except subprocess.TimeoutExpired as exc:
                return {
                    "success": False,
                    "error": "Timeout",
                    "message": f"{script_name} did not finish within {timeout}s",
                    "stderr_tail": (exc.stderr or "")[-2000:] if exc.stderr else "",
                }
            except Exception as exc:
                return {
                    "success": False,
                    "error": str(exc),
                    "message": f"Failed to invoke {script_name}: {exc}",
                }

            result = {
                "success": proc.returncode == 0,
                "returncode": proc.returncode,
                "stdout_tail": (proc.stdout or "")[-2000:],
                "stderr_tail": (proc.stderr or "")[-2000:],
            }
            if capture_json is not None and capture_json.exists():
                try:
                    import json as _json
                    result["report"] = _json.loads(
                        capture_json.read_text(encoding="utf-8")
                    )
                    result["report_path"] = str(capture_json)
                except Exception as exc:
                    result["report_read_error"] = str(exc)
            return result

        def ppt_read_skill_reference_tool(name: str) -> dict:
            """
            Read one of the PPT skill's reference documents and return its text.

            Use this BEFORE planning a deck whenever you need methodology:
            page archetypes, slide design system, quality gate semantics, build
            routes, diagram / chart / icon / python figure rules.

            Args:
                name: One of:
                    - "principles"
                    - "deck_workflow"
                    - "technical_support"
                    - "design_support"
                    - "slide_design_system"
                    - "quality_gates"
                    - "build_routes"
                    - "diagram_support"
                    - "office_chart_support"
                    - "python_figure_support"
                    - "icon_system"

            Returns dict with:
                success / name / path / content (full markdown) / message.
            """
            if name not in _PPT_REFERENCE_NAMES:
                return {
                    "success": False,
                    "error": "Unknown reference",
                    "message": (
                        f"name={name!r} is not a valid PPT skill reference. "
                        f"Valid names: {sorted(_PPT_REFERENCE_NAMES)}"
                    ),
                }
            ref_path = PPT_REFERENCES_DIR / f"{name}.md"
            if not ref_path.exists():
                return {
                    "success": False,
                    "error": "Reference file missing",
                    "message": (
                        f"{ref_path} does not exist. The PPT skill may be "
                        "incomplete; check that "
                        "skills/presentation-skills/ppt-polished-deck-collab-traditional/references/ "
                        "is intact."
                    ),
                }
            try:
                content = ref_path.read_text(encoding="utf-8")
            except Exception as exc:
                return {
                    "success": False,
                    "error": str(exc),
                    "message": f"Failed to read {ref_path}: {exc}",
                }
            return {
                "success": True,
                "name": name,
                "path": str(ref_path),
                "content": content,
                "message": f"Loaded reference {name} ({len(content)} chars)",
            }

        def ppt_check_environment_tool(deck_workspace: str | None = None) -> dict:
            """
            Probe the local environment for the PPT skill's required tooling
            and return the set of available build / preview routes.

            Use this as the FIRST step of every PPT task — it tells you which
            preview backend (PowerPoint vs LibreOffice) you can actually use,
            and which optional capabilities (Python figure, Mermaid) are
            present. The agent should branch on the returned `routes` list
            instead of assuming a backend is available.

            Args:
                deck_workspace: Optional. When provided, the JSON env report is
                    written to <deck_workspace>/validation/env_check.json so the
                    deck has a durable record of which routes were available
                    when it was built.

            Returns dict with:
                success / routes (list) / report (parsed JSON) / stdout_tail /
                stderr_tail / message.
            """
            json_out: Path
            if deck_workspace:
                ws = Path(deck_workspace).resolve()
                target_dir = ws / "validation"
                target_dir.mkdir(parents=True, exist_ok=True)
                json_out = target_dir / "env_check.json"
            else:
                # Fall back to a per-user scratch path so we can still read the
                # structured report even when the agent has not picked a deck
                # workspace yet.
                scratch = _ppt_user_workdir() / "_ppt_env_check.json"
                json_out = scratch

            result = _run_ppt_script(
                "check_environment.py",
                ["--json-out", str(json_out)],
                timeout=120,
                capture_json=json_out,
            )
            routes = []
            report = result.get("report") or {}
            if isinstance(report, dict):
                routes = list(report.get("routes") or [])
            result["routes"] = routes
            result["message"] = (
                f"Detected {len(routes)} available route(s): "
                + (", ".join(routes) if routes else "(none)")
            )
            return result

        _PPT_SLUG_RE = __import__("re").compile(r"^[a-z0-9][a-z0-9\-]{0,62}$")

        def _ppt_brief_md(deck_title: str) -> str:
            """Return a filled-in brief.md template (matches deck_workflow.md)."""
            return (
                f"# {deck_title}\n\n"
                "## 任务定义\n"
                "- 目标读者：\n"
                "- 主使用场景：\n"
                "- 目标动作：\n"
                "- 参考模板文件：\n"
                "- 模板 / 品牌约束：\n"
                "- 交付物要求：\n"
                "- 验证要求：\n\n"
                "## 模板取证\n"
                "- 页面系统判断：\n"
                "- 关键母版 / layout 元素：\n"
                "- 字号系统：\n"
                "- 计划采用的构建路线：\n"
                "- 最小 PoC 结论：\n\n"
                "## 风格与边界\n"
                "- 风格参考：\n"
                "- typography_profile：zh_formal\n"
                "- domain_profile：\n"
                "- 允许使用的素材：\n"
                "- 禁止使用的品牌元素：\n"
                "- 免责声明 / 风险边界：\n"
                "- 不允许发生的错误：\n"
            )

        def _ppt_narrative_md(deck_title: str) -> str:
            """Return a starter deck_narrative.md (zh_formal theme tokens).

            The YAML frontmatter is built via ``yaml.safe_dump`` rather than
            string concatenation: hand-rolled f-string injection broke on
            titles containing ``"`` (which closed the double-quoted scalar
            prematurely) or ``\\`` (PyYAML treats it as an escape lead-in
            and raises ``ScannerError`` on ``\\p`` / ``\\n`` etc.). Routing
            through ``safe_dump`` lets PyYAML pick the right quoting style.
            """
            import yaml as _yaml

            frontmatter = {
                "deck": {
                    "title": deck_title,
                    "audience": "<target audience>",
                    "scenario": "<primary scenario>",
                    "objective": "<primary decision or action>",
                    "theme_tokens": {
                        "typography_profile": "zh_formal",
                        "domain_profile": None,
                        "hero_title_font_pt": 24,
                        "section_title_font_pt": 20,
                        "page_title_font_pt": 24,
                        "subtitle_font_pt": 16,
                        "minor_title_font_pt": 14,
                        "body_font_pt": 12,
                        "label_font_pt": 10.5,
                        "caption_font_pt": 9,
                        "title_line_spacing_multiple": 1.0,
                        "body_line_spacing_multiple": 1.5,
                        "title_paragraph_space_lines": 0.5,
                        "body_first_line_indent_chars": 2,
                        "body_paragraph_space_lines": 0.5,
                        "latin_font_name": "Times New Roman",
                        "east_asia_font_name": "宋体",
                        "table_font_pt": 10.5,
                        "table_line_spacing_multiple": 1.0,
                        "table_paragraph_space_lines": 0,
                        "table_first_line_indent_chars": 0,
                        "table_vertical_anchor": "middle",
                        "table_header_alignment": "center",
                        "table_index_alignment": "left",
                        "table_text_alignment": "left",
                        "table_numeric_alignment": "right",
                    },
                }
            }
            yaml_body = _yaml.safe_dump(
                frontmatter,
                allow_unicode=True,
                sort_keys=False,
                default_flow_style=False,
            )
            # Also defang the H1 line: a deck_title containing a real
            # newline would split the heading into two lines and confuse
            # tools that treat the first line as the H1.
            h1_safe = " ".join((deck_title or "").splitlines()).strip() or "Deck"
            return (
                "---\n"
                f"{yaml_body}"
                "---\n\n"
                f"# {h1_safe}\n\n"
                "## Global Narrative\n"
                "- 这套 deck 的主判断：\n"
                "- 这套 deck 的论证主线：\n"
                "- 这套 deck 的主题词和禁区：\n\n"
                "### S01 | <slide title>\n"
                "```yaml slide_spec\n"
                "title: '<slide title>'\n"
                "reader_question: '<what this page should answer>'\n"
                "page_task: 'persuade'\n"
                "reading_mode: 'decision'\n"
                "archetype: 'hero-statement'\n"
                "asset_mode: 'text-layout-native'\n"
                "validation_mode: 'preview_only'\n"
                "key_message: '<single core message>'\n"
                "required_assets: []\n"
                "```\n\n"
                "**Narrative Role.** 这页为什么存在、要帮助读者完成什么判断。\n\n"
                "**Content Notes.** 这页准备放什么内容、什么判断句、什么证据。\n\n"
                "**Layout Notes.** 这页倾向使用什么版式、什么 icon 或图表策略。\n"
            )

        def ppt_init_workspace_tool(deck_slug: str, deck_title: str) -> dict:
            """
            Create a deck workspace under the current user's work dir, with
            `brief.md` + `deck_narrative.md` (zh_formal theme_tokens already
            filled in) and the six standard sub-directories required by the
            `ppt-polished-deck-collab` skill.

            Call this as the FIRST PPT tool of every deck task. All subsequent
            PPT tools should pass the returned `deck_workspace` value as their
            `--workspace-dir` — never assemble the path yourself.

            Args:
                deck_slug: kebab-case identifier for this deck (a–z, 0–9, '-';
                    max 63 chars). Used as the directory name. Example:
                    "ihep-2026-q2-safety".
                deck_title: Human-readable deck title — appears in brief.md
                    and the narrative document's YAML frontmatter `deck.title`.

            Returns dict with:
                success / deck_workspace (abs path) / brief_path /
                narrative_path / created (list of created paths) /
                already_exists (bool) / message.
            """
            slug = (deck_slug or "").strip().lower()
            if not _PPT_SLUG_RE.match(slug):
                return {
                    "success": False,
                    "error": "Invalid deck_slug",
                    "message": (
                        f"deck_slug={deck_slug!r} must be kebab-case "
                        "(a-z, 0-9, '-', start with alnum, max 63 chars). "
                        "Examples: 'ihep-2026-safety', 'q2-product-review'."
                    ),
                }
            if not deck_title or not deck_title.strip():
                return {
                    "success": False,
                    "error": "Missing deck_title",
                    "message": "deck_title is required and must be non-empty.",
                }

            base = _ppt_user_workdir() / "decks" / slug
            already_exists = base.exists()
            base.mkdir(parents=True, exist_ok=True)

            sub_dirs = [
                "data",
                "assets/diagrams",
                "assets/charts",
                "assets/icons",
                "assets/images",
                "assets/tables",
                "build/generated",
                "build/pptx",
                "build/rendered/ppt_preview",
                "build/rendered/python_figures",
                "validation/template_audit",
                "validation/package_preflight/history",
                "validation/structure_precheck/history",
                "validation/render_review/history",
                "validation/visual",
                "final",
            ]
            created: list[str] = []
            for rel in sub_dirs:
                target = base / rel
                if not target.exists():
                    target.mkdir(parents=True, exist_ok=True)
                    created.append(str(target))

            brief_path = base / "brief.md"
            narrative_path = base / "deck_narrative.md"
            if not brief_path.exists():
                brief_path.write_text(
                    _ppt_brief_md(deck_title.strip()), encoding="utf-8"
                )
                created.append(str(brief_path))
            if not narrative_path.exists():
                narrative_path.write_text(
                    _ppt_narrative_md(deck_title.strip()), encoding="utf-8"
                )
                created.append(str(narrative_path))

            return {
                "success": True,
                "deck_workspace": str(base),
                "brief_path": str(brief_path),
                "narrative_path": str(narrative_path),
                "created": created,
                "already_exists": already_exists,
                "message": (
                    f"Deck workspace ready at {base}. "
                    f"{'Re-used existing structure.' if already_exists else 'Created fresh.'} "
                    "Next: edit brief.md and deck_narrative.md, then call "
                    "ppt_derive_slide_specs_tool."
                ),
            }

        # ============ PPT SKILL TOOLS (Phase 2) ============
        # Audit / derive / lint / 3 quality gates / preview export /
        # connector check / icon search & render. All Phase 2 tools share the
        # same conventions:
        #   - deck_workspace MUST come from ppt_init_workspace_tool; the agent
        #     never invents the path.
        #   - Reports auto-archive under <deck_workspace>/validation/<gate>/
        #     history/<gate>_<timestamp>.{json,md}.
        #   - Result dicts include `report` (parsed JSON when available) so
        #     the agent can branch without re-reading files.

        def _resolve_deck_workspace(deck_workspace: str) -> tuple[Path | None, dict | None]:
            """Validate a deck_workspace string.

            Returns (path, None) on success or (None, error_dict) when the
            input is missing, not absolute, or does not exist. Phase 2 tools
            short-circuit on the error dict so the agent gets a directive
            recovery hint instead of a generic OS error.
            """
            if not deck_workspace or not isinstance(deck_workspace, str):
                return None, {
                    "success": False,
                    "error": "Missing deck_workspace",
                    "message": (
                        "deck_workspace is required. Call ppt_init_workspace_tool "
                        "first and pass the returned `deck_workspace` value here."
                    ),
                }
            p = Path(deck_workspace)
            if not p.is_absolute():
                return None, {
                    "success": False,
                    "error": "Relative deck_workspace not accepted",
                    "message": (
                        f"deck_workspace={deck_workspace!r} is a relative path. "
                        "Use the absolute path returned by ppt_init_workspace_tool."
                    ),
                }
            if not p.exists() or not p.is_dir():
                return None, {
                    "success": False,
                    "error": "deck_workspace not found",
                    "message": (
                        f"No directory at {deck_workspace}. Re-run "
                        "ppt_init_workspace_tool to create it, or check the "
                        "value you received from that tool earlier in the "
                        "conversation."
                    ),
                }
            return p, None

        def _resolve_pptx_path(pptx_path: str, label: str = "pptx_path") -> tuple[Path | None, dict | None]:
            """Validate a .pptx path argument."""
            if not pptx_path or not isinstance(pptx_path, str):
                return None, {
                    "success": False,
                    "error": f"Missing {label}",
                    "message": f"{label} is required.",
                }
            p = Path(pptx_path)
            if not p.is_absolute():
                return None, {
                    "success": False,
                    "error": f"Relative {label} not accepted",
                    "message": (
                        f"{label}={pptx_path!r} is relative. Use the absolute "
                        "path returned by your previous build step."
                    ),
                }
            if not p.exists():
                return None, {
                    "success": False,
                    "error": "File not found",
                    "message": f"{label}: no file at {pptx_path}.",
                }
            if p.suffix.lower() != ".pptx":
                return None, {
                    "success": False,
                    "error": "Not a .pptx file",
                    "message": f"{label}={pptx_path!r} is not a .pptx file.",
                }
            return p, None

        def ppt_audit_template_tool(
            pptx_path: str,
            deck_workspace: str,
            sample_limit: int = 3,
            text_preview_limit: int = 90,
        ) -> dict:
            """
            Audit a reference .pptx template — discover its layout family,
            master/layout/slide text inventory and font-size distribution.

            Use this when the user provides an existing .pptx and wants the
            new deck to inherit its page system. Run BEFORE writing
            deck_narrative.md so the narrative can be anchored to the
            template's real font sizes and layout names, not a guess.

            Outputs are written to:
              <deck_workspace>/validation/template_audit/template_audit.json
              <deck_workspace>/validation/template_audit/template_audit.md

            After running, fold the key findings (font-size ladder, layout
            family, shared master elements, build-route choice) back into
            brief.md so subsequent steps treat them as deck-level facts.

            Args:
                pptx_path: Absolute path to the reference .pptx.
                deck_workspace: Value returned by ppt_init_workspace_tool.
                sample_limit: How many sample text strings to retain per
                    font-size bucket (default 3).
                text_preview_limit: Max characters per retained sample
                    (default 90).

            Returns dict with success, json_path, md_path, report (parsed
            JSON), stdout_tail, stderr_tail, message.
            """
            pptx, err = _resolve_pptx_path(pptx_path)
            if err is not None:
                return err
            ws, err = _resolve_deck_workspace(deck_workspace)
            if err is not None:
                return err

            target_dir = ws / "validation" / "template_audit"
            target_dir.mkdir(parents=True, exist_ok=True)
            json_out = target_dir / "template_audit.json"
            md_out = target_dir / "template_audit.md"

            result = _run_ppt_script(
                "audit_pptx_template.py",
                [
                    "--pptx", str(pptx),
                    "--json-out", str(json_out),
                    "--md-out", str(md_out),
                    "--sample-limit", str(sample_limit),
                    "--text-preview-limit", str(text_preview_limit),
                ],
                timeout=180,
                capture_json=json_out,
            )
            result["json_path"] = str(json_out) if json_out.exists() else None
            result["md_path"] = str(md_out) if md_out.exists() else None
            if result.get("success"):
                summary = (result.get("report") or {}).get("summary") or {}
                result["message"] = (
                    f"Template audit OK. slides={summary.get('slide_count', '?')}, "
                    f"masters={summary.get('master_count', '?')}, "
                    f"layouts={summary.get('default_slide_layout_count', '?')}. "
                    f"Findings written to {target_dir}."
                )
            else:
                result["message"] = (
                    "Template audit failed. Check stderr_tail for details; "
                    "if soffice / pptx parsing complains, confirm the file is "
                    "a real .pptx (not .ppt — convert via Office or LibreOffice first)."
                )
            return result

        def ppt_derive_slide_specs_tool(
            deck_workspace: str,
            narrative_path: str | None = None,
            out_yaml: str | None = None,
        ) -> dict:
            """
            Parse deck_narrative.md and write a structured slide_specs.yaml
            ready for the build step.

            What the script does:
              - Reads YAML frontmatter as deck-level metadata.
              - Splits the body by `### Sxx | <title>` headings.
              - Pulls the first ```yaml slide_spec``` block from each section.
              - Validates the eight required fields per slide: title,
                reader_question, page_task, reading_mode, archetype,
                asset_mode, validation_mode, key_message.
              - Carries the remaining markdown as narrative_markdown.

            Defaults:
              narrative_path = <deck_workspace>/deck_narrative.md
              out_yaml       = <deck_workspace>/build/generated/slide_specs.yaml

            If the script fails with "missing field" errors, read the stderr
            tail to the user — those are authoring problems in
            deck_narrative.md (a slide section forgot its yaml block, the
            yaml block lacks a required field, etc.). Fix in the narrative,
            then re-run this tool.

            Args:
                deck_workspace: Value returned by ppt_init_workspace_tool.
                narrative_path: Optional override; defaults to deck_workspace/deck_narrative.md.
                out_yaml: Optional override; defaults to
                    deck_workspace/build/generated/slide_specs.yaml.

            Returns dict with success / slide_specs_path / slide_count /
            stdout_tail / stderr_tail / message.
            """
            ws, err = _resolve_deck_workspace(deck_workspace)
            if err is not None:
                return err

            narr = Path(narrative_path) if narrative_path else (ws / "deck_narrative.md")
            if not narr.exists():
                return {
                    "success": False,
                    "error": "Narrative not found",
                    "message": (
                        f"deck_narrative.md not found at {narr}. Edit the "
                        "narrative file produced by ppt_init_workspace_tool "
                        "before deriving slide_specs."
                    ),
                }
            out = Path(out_yaml) if out_yaml else (ws / "build" / "generated" / "slide_specs.yaml")
            out.parent.mkdir(parents=True, exist_ok=True)

            # The derive script doesn't take --json-out for the structured
            # result, but we can scrape the slide count from stdout.
            result = _run_ppt_script(
                "derive_slide_specs_from_narrative.py",
                ["--narrative", str(narr), "--out-yaml", str(out)],
                timeout=60,
            )
            slide_count = None
            if result.get("success"):
                import re as _re
                m = _re.search(r"slides=(\d+)", result.get("stdout_tail") or "")
                if m:
                    slide_count = int(m.group(1))
                result["slide_specs_path"] = str(out)
                result["slide_count"] = slide_count
                result["message"] = (
                    f"Derived {slide_count if slide_count is not None else '?'} "
                    f"slide spec(s) → {out}."
                )
            else:
                result["message"] = (
                    "Derive failed. Common causes: deck_narrative.md is "
                    "missing YAML frontmatter, a `### Sxx | <title>` heading, "
                    "a ```yaml slide_spec``` block, or one of the required "
                    "slide fields (title, reader_question, page_task, "
                    "reading_mode, archetype, asset_mode, validation_mode, "
                    "key_message). Read stderr_tail and fix the narrative, "
                    "then call again."
                )
            return result

        def ppt_lint_workspace_tool(deck_workspace: str) -> dict:
            """
            Check that a deck workspace has the required directories, the two
            human-authored markdown files, and a derived slide_specs.yaml.

            Use this as a pre-flight before build to catch missing inputs
            without trying to compile a half-finished deck. The script also
            reports asset-folder occupancy (diagrams / charts / icons /
            images / tables) so the agent can spot under-supplied assets.

            Args:
                deck_workspace: Value returned by ppt_init_workspace_tool.

            Returns dict with success, report (parsed JSON), errors,
            warnings, message.
            """
            ws, err = _resolve_deck_workspace(deck_workspace)
            if err is not None:
                return err

            json_out = ws / "validation" / "workspace_lint.json"
            result = _run_ppt_script(
                "lint_deck_assets.py",
                ["--workspace-dir", str(ws), "--json-out", str(json_out)],
                timeout=30,
                capture_json=json_out,
            )
            report = result.get("report") or {}
            result["errors"] = list(report.get("errors") or [])
            result["warnings"] = list(report.get("warnings") or [])
            if result.get("success"):
                result["message"] = (
                    "Workspace lint passed."
                    + (f" Warnings: {len(result['warnings'])}."
                       if result["warnings"] else "")
                )
            else:
                result["message"] = (
                    f"Workspace lint failed with {len(result['errors'])} "
                    "error(s). Fix the missing directories/files, then re-run."
                )
            return result

        def ppt_package_preflight_tool(
            pptx_path: str,
            deck_workspace: str,
            fail_on: str = "error",
        ) -> dict:
            """
            File-level quality gate: zip integrity, slide-count consistency
            (presentation.xml vs docProps/app.xml vs actual slide files),
            stale section_lst references, missing slide relationships, and
            embedded-object mobile-compatibility risk.

            This is the FIRST gate after build. Run BEFORE structure_precheck
            and BEFORE preview export — if the deck can't be opened by a
            fragile parser (WeChat / mobile WPS), there is no point checking
            its layout.

            Output is auto-archived to:
              <deck_workspace>/validation/package_preflight/history/
                package_preflight_<YYYYMMDD_HHMMSS>.{json,md}

            Args:
                pptx_path: Absolute path to the deck .pptx.
                deck_workspace: Value returned by ppt_init_workspace_tool.
                fail_on: 'error' (default), 'warning', or 'never'. Controls
                    the script's exit code, NOT what is reported — issues are
                    always returned in the result dict.

            Returns dict with success, returncode, report (parsed JSON
            including `summary` counts and `issues` list), stdout_tail,
            stderr_tail, message.
            """
            pptx, err = _resolve_pptx_path(pptx_path)
            if err is not None:
                return err
            ws, err = _resolve_deck_workspace(deck_workspace)
            if err is not None:
                return err
            if fail_on not in {"error", "warning", "never"}:
                return {
                    "success": False,
                    "error": "Invalid fail_on",
                    "message": "fail_on must be one of: error, warning, never.",
                }

            result = _run_ppt_script(
                "check_pptx_package_preflight.py",
                [
                    "--pptx", str(pptx),
                    "--workspace-dir", str(ws),
                    "--fail-on", fail_on,
                ],
                timeout=120,
            )
            report = result.get("report") or {}
            summary = report.get("summary") or {}
            # The script writes its own timestamped report — locate the
            # newest one so the agent can point the user at it.
            hist = ws / "validation" / "package_preflight" / "history"
            newest = None
            if hist.exists():
                jsons = sorted(hist.glob("package_preflight_*.json"))
                if jsons:
                    newest = jsons[-1]
                    try:
                        import json as _json
                        report = _json.loads(newest.read_text(encoding="utf-8"))
                        result["report"] = report
                        summary = report.get("summary") or {}
                    except Exception:
                        pass
            result["report_path"] = str(newest) if newest else None
            result["summary"] = summary
            result["message"] = (
                f"package_preflight: errors={summary.get('error', 0)}, "
                f"warnings={summary.get('warning', 0)}, "
                f"not_checked={summary.get('not_checked', 0)}. "
                + (f"Report: {newest.name}." if newest else "(no report on disk)")
            )
            return result

        def ppt_structure_precheck_tool(
            pptx_path: str,
            deck_workspace: str,
            fail_on: str = "error",
        ) -> dict:
            """
            Structure-layer quality gate: textbox fit / near-overflow,
            compact-width pressure on short labels, text occluded by higher
            z-order shapes, critical content (table/chart/picture) covered
            by overlay shapes, and explicit `not_checked` records for
            structured chart labels and flattened pictures.

            Run AFTER package_preflight and BEFORE preview export. The
            issues here have shape/slide-level locations so they are easy
            to drive into targeted fixes.

            Output is auto-archived to:
              <deck_workspace>/validation/structure_precheck/history/
                structure_precheck_<YYYYMMDD_HHMMSS>.{json,md}
              <deck_workspace>/validation/structure_precheck/shape_inventory.json

            Args:
                pptx_path: Absolute path to the deck .pptx.
                deck_workspace: Value returned by ppt_init_workspace_tool.
                fail_on: 'error' (default), 'warning', or 'never'.

            Returns dict with success, returncode, report (parsed JSON),
            summary, inventory_path, report_path, stdout_tail, stderr_tail,
            message.
            """
            pptx, err = _resolve_pptx_path(pptx_path)
            if err is not None:
                return err
            ws, err = _resolve_deck_workspace(deck_workspace)
            if err is not None:
                return err
            if fail_on not in {"error", "warning", "never"}:
                return {
                    "success": False,
                    "error": "Invalid fail_on",
                    "message": "fail_on must be one of: error, warning, never.",
                }

            inventory_out = ws / "validation" / "structure_precheck" / "shape_inventory.json"
            inventory_out.parent.mkdir(parents=True, exist_ok=True)
            result = _run_ppt_script(
                "check_pptx_structure_precheck.py",
                [
                    "--pptx", str(pptx),
                    "--workspace-dir", str(ws),
                    "--inventory-out", str(inventory_out),
                    "--fail-on", fail_on,
                ],
                timeout=180,
            )
            hist = ws / "validation" / "structure_precheck" / "history"
            newest = None
            summary = {}
            if hist.exists():
                jsons = sorted(hist.glob("structure_precheck_*.json"))
                if jsons:
                    newest = jsons[-1]
                    try:
                        import json as _json
                        report = _json.loads(newest.read_text(encoding="utf-8"))
                        result["report"] = report
                        summary = report.get("summary") or {}
                    except Exception:
                        pass
            result["report_path"] = str(newest) if newest else None
            result["inventory_path"] = str(inventory_out) if inventory_out.exists() else None
            result["summary"] = summary
            result["message"] = (
                f"structure_precheck: errors={summary.get('error', 0)}, "
                f"warnings={summary.get('warning', 0)}, "
                f"not_checked={summary.get('not_checked', 0)}. "
                + (f"Report: {newest.name}." if newest else "(no report on disk)")
            )
            return result

        def ppt_export_previews_tool(
            pptx_path: str,
            deck_workspace: str,
            backend: str = "auto",
            render_backend: str = "auto",
            prefix: str = "slide_",
            keep_pdf: bool = False,
        ) -> dict:
            """
            Render the deck as per-slide PNG previews via PowerPoint or
            LibreOffice, then pdftoppm or PyMuPDF for PDF→PNG.

            Output directory defaults to:
              <deck_workspace>/build/rendered/ppt_preview/slide_001.png ...

            Manifest written to:
              <deck_workspace>/validation/preview_manifest.json

            Page-count mismatch (e.g. LibreOffice silently dropping a slide)
            is treated as failure, not a warning — re-export with the other
            backend if it happens.

            Args:
                pptx_path: Absolute path to the deck .pptx.
                deck_workspace: Value returned by ppt_init_workspace_tool.
                backend: 'auto' (default), 'powerpoint', or 'libreoffice'.
                render_backend: 'auto' (default), 'pdftoppm', or 'fitz'.
                prefix: Output PNG filename prefix (default 'slide_').
                keep_pdf: When true, the intermediate PDF is moved into the
                    preview directory; default false.

            Returns dict with success, preview_dir, manifest_path,
            generated_pages, pdf_backend, render_backend, stdout_tail,
            stderr_tail, message.
            """
            pptx, err = _resolve_pptx_path(pptx_path)
            if err is not None:
                return err
            ws, err = _resolve_deck_workspace(deck_workspace)
            if err is not None:
                return err
            if backend not in {"auto", "powerpoint", "libreoffice"}:
                return {
                    "success": False,
                    "error": "Invalid backend",
                    "message": "backend must be one of: auto, powerpoint, libreoffice.",
                }
            if render_backend not in {"auto", "pdftoppm", "fitz"}:
                return {
                    "success": False,
                    "error": "Invalid render_backend",
                    "message": "render_backend must be one of: auto, pdftoppm, fitz.",
                }

            out_dir = ws / "build" / "rendered" / "ppt_preview"
            out_dir.mkdir(parents=True, exist_ok=True)
            manifest = ws / "validation" / "preview_manifest.json"
            manifest.parent.mkdir(parents=True, exist_ok=True)

            args = [
                "--pptx", str(pptx),
                "--out-dir", str(out_dir),
                "--backend", backend,
                "--render-backend", render_backend,
                "--prefix", prefix,
                "--json-out", str(manifest),
            ]
            if keep_pdf:
                args.append("--keep-pdf")
            result = _run_ppt_script(
                "export_pptx_previews.py",
                args,
                timeout=300,
                capture_json=manifest,
            )
            report = result.get("report") or {}
            result["preview_dir"] = str(out_dir)
            result["manifest_path"] = str(manifest) if manifest.exists() else None
            result["generated_pages"] = report.get("generated_pages")
            result["pdf_backend"] = report.get("pdf_backend")
            result["render_backend"] = report.get("render_backend")
            if result.get("success"):
                result["message"] = (
                    f"Exported {result['generated_pages']} preview(s) via "
                    f"{result['pdf_backend']} → {result['render_backend']}. "
                    f"PNGs in {out_dir}."
                )
            else:
                result["message"] = (
                    "Preview export failed. If the error mentions PowerPoint "
                    "automation, check macOS Privacy & Security → Automation. "
                    "If LibreOffice converted but pages count is off, retry "
                    "with backend='powerpoint' or vice versa."
                )
            return result

        def ppt_render_review_tool(
            pptx_path: str,
            deck_workspace: str,
            preview_dir: str | None = None,
            fail_on: str = "error",
        ) -> dict:
            """
            Render-layer quality gate: boundary-touch-ink at bottom/right of
            text frames (font strokes within ~3px of the inner edge in the
            PNG) and flattened-graphic internal-text `not_checked` entries.

            Run AFTER ppt_export_previews_tool — this gate consumes the
            preview PNGs. It complements structure_precheck by catching
            issues only visible after rasterization (e.g. last-line clipped
            by 1-2 px when the structure-level math says it just barely
            fits).

            Output auto-archived to:
              <deck_workspace>/validation/render_review/history/
                render_review_<YYYYMMDD_HHMMSS>.{json,md}

            Args:
                pptx_path: Absolute path to the deck .pptx.
                deck_workspace: Value returned by ppt_init_workspace_tool.
                preview_dir: Optional override; defaults to
                    deck_workspace/build/rendered/ppt_preview.
                fail_on: 'error' (default), 'warning', or 'never'.

            Returns dict with success, summary, report, report_path,
            stdout_tail, stderr_tail, message.
            """
            pptx, err = _resolve_pptx_path(pptx_path)
            if err is not None:
                return err
            ws, err = _resolve_deck_workspace(deck_workspace)
            if err is not None:
                return err
            if fail_on not in {"error", "warning", "never"}:
                return {
                    "success": False,
                    "error": "Invalid fail_on",
                    "message": "fail_on must be one of: error, warning, never.",
                }

            pv = Path(preview_dir) if preview_dir else (ws / "build" / "rendered" / "ppt_preview")
            if not pv.exists():
                return {
                    "success": False,
                    "error": "Preview directory not found",
                    "message": (
                        f"No preview directory at {pv}. Run "
                        "ppt_export_previews_tool first."
                    ),
                }

            result = _run_ppt_script(
                "check_pptx_render_review.py",
                [
                    "--pptx", str(pptx),
                    "--preview-dir", str(pv),
                    "--workspace-dir", str(ws),
                    "--fail-on", fail_on,
                ],
                timeout=180,
            )
            hist = ws / "validation" / "render_review" / "history"
            newest = None
            summary = {}
            if hist.exists():
                jsons = sorted(hist.glob("render_review_*.json"))
                if jsons:
                    newest = jsons[-1]
                    try:
                        import json as _json
                        report = _json.loads(newest.read_text(encoding="utf-8"))
                        result["report"] = report
                        summary = report.get("summary") or {}
                    except Exception:
                        pass
            result["report_path"] = str(newest) if newest else None
            result["summary"] = summary
            result["message"] = (
                f"render_review: errors={summary.get('error', 0)}, "
                f"warnings={summary.get('warning', 0)}, "
                f"not_checked={summary.get('not_checked', 0)}. "
                + (f"Report: {newest.name}." if newest else "(no report on disk)")
            )
            return result

        def ppt_connectors_check_tool(
            pptx_path: str,
            deck_workspace: str,
            slides: list | None = None,
            min_connectors: int = 0,
            forbid_prefixes: list | None = None,
        ) -> dict:
            """
            Module-level gate for diagram pages: verify each connector is
            REALLY glued to two shapes (stCxn + endCxn present, target shape
            ids resolve, no connections to forbidden parent shapes such as
            lane / cluster outer frames).

            Run after a diagram page with asset_mode=diagram-connector is
            built. A passing report is the evidence the user can rely on
            that dragging a node will not break the diagram.

            Output:
              <deck_workspace>/validation/connectors/connector_report.json

            Args:
                pptx_path: Absolute path to the deck .pptx.
                deck_workspace: Value returned by ppt_init_workspace_tool.
                slides: Optional list of slide numbers (1-based) to limit
                    the check. Defaults to all slides.
                min_connectors: Optional. Total connector count must reach
                    this value or the check fails. Use it on dedicated
                    diagram pages where you know connectors must exist.
                forbid_prefixes: Optional list of forbidden prefixes for
                    connector endpoints. Default is `["Lane "]` — connector
                    endpoints starting with these strings are flagged as
                    illegal (they typically mean the line is glued to a
                    swimlane outer frame instead of a business node).

            Returns dict with success, total_connectors, report (parsed
            JSON, mapping slide → records), report_path, stdout_tail,
            stderr_tail, message.
            """
            pptx, err = _resolve_pptx_path(pptx_path)
            if err is not None:
                return err
            ws, err = _resolve_deck_workspace(deck_workspace)
            if err is not None:
                return err

            json_out = ws / "validation" / "connectors" / "connector_report.json"
            json_out.parent.mkdir(parents=True, exist_ok=True)
            args = ["--pptx", str(pptx), "--json-out", str(json_out)]
            if isinstance(slides, list):
                for s in slides:
                    args.extend(["--slide", str(int(s))])
            if min_connectors and min_connectors > 0:
                args.extend(["--min-connectors", str(int(min_connectors))])
            for prefix in (forbid_prefixes or []):
                args.extend(["--forbid-prefix", str(prefix)])

            result = _run_ppt_script(
                "check_pptx_connectors.py",
                args,
                timeout=60,
                capture_json=json_out,
            )
            report = result.get("report") or {}
            # report maps slide_num -> list of connector records
            total = 0
            if isinstance(report, dict):
                for v in report.values():
                    if isinstance(v, list):
                        total += len(v)
            result["report_path"] = str(json_out) if json_out.exists() else None
            result["total_connectors"] = total
            if result.get("success"):
                result["message"] = (
                    f"connector check passed: {total} connector(s) verified."
                )
            else:
                result["message"] = (
                    f"connector check FAILED. {total} connector(s) seen. "
                    "Read stdout_tail for the specific issues — usually one "
                    "of: stCxn/endCxn missing (line drawn but not glued), "
                    "endpoint id unresolved (target shape deleted), or "
                    "connector glued to a lane/cluster outer frame instead "
                    "of a business node."
                )
            return result

        def ppt_icon_search_tool(query: str, pack: str | None = None) -> dict:
            """
            Search the PPT skill's Tabler-Outline icon registry.

            Use this when planning an icon-accent page (asset_mode=icon-accent)
            or when looking for a section header icon. Icons are SUPPORTING
            assets — they never carry primary information.

            Args:
                query: Space-separated keywords (English or Chinese aliases
                    both work). Example: "risk safety" or "趋势 增长".
                pack: Optional pack id. One of:
                    - "general-layout" (default scope — titles, cards, sections)
                    - "llm-research" (ACL/EMNLP/LLM/Agent/RAG topics)
                    Omit to search across all packs.

            Returns dict with success, matches (list of {score, id,
            source_name, packs, aliases, usage_note}), stdout_tail, message.
            """
            if not query or not isinstance(query, str) or not query.strip():
                return {
                    "success": False,
                    "error": "Missing query",
                    "message": "query is required (space-separated keywords).",
                }
            args = ["search", "--query", query]
            if pack:
                args.extend(["--pack", str(pack)])
            result = _run_ppt_script("icon_registry.py", args, timeout=30)
            # Parse stdout's [MATCH ...] lines into a structured list.
            matches = []
            current = None
            for line in (result.get("stdout_tail") or "").splitlines():
                line = line.rstrip()
                if line.startswith("[MATCH]"):
                    if current:
                        matches.append(current)
                    parts = line[len("[MATCH]"):].strip().split()
                    rec = {"score": None, "id": None, "source_name": None, "packs": []}
                    for p in parts:
                        if "=" in p:
                            k, v = p.split("=", 1)
                            if k == "score":
                                try:
                                    rec["score"] = int(v)
                                except ValueError:
                                    rec["score"] = v
                            elif k == "packs":
                                rec["packs"] = [x for x in v.split(",") if x]
                            elif k in {"id", "source"}:
                                rec["id" if k == "id" else "source_name"] = v
                    current = rec
                elif current and "aliases=" in line:
                    current["aliases"] = [
                        x for x in line.split("aliases=", 1)[1].split(",") if x
                    ]
                elif current and "usage=" in line:
                    current["usage_note"] = line.split("usage=", 1)[1]
            if current:
                matches.append(current)
            result["matches"] = matches
            result["message"] = (
                f"Found {len(matches)} icon match(es) for query={query!r}."
            )
            return result

        def ppt_icon_render_tool(
            deck_workspace: str,
            pack: str | None = None,
            size: int = 128,
            color_mode: str = "auto",
            background_color: str = "#F8FAFC",
            accent_color: str = "#2563EB",
            theme_name: str = "default",
            icon_color: str | None = None,
        ) -> dict:
            """
            Render icon PNGs into the deck workspace, with deck-aware
            recoloring (auto mode picks colors from the icon's role + the
            slide background + the accent color, then enforces WCAG ≥3.0
            contrast).

            Output goes to:
              <deck_workspace>/assets/icons/<pack or 'all'>/<theme_name>/

            so a build can later reference these PNGs by relative path
            without polluting the skill directory.

            Args:
                deck_workspace: Value returned by ppt_init_workspace_tool.
                pack: Optional. 'general-layout' / 'llm-research'. Omit to
                    render every pack.
                size: Square PNG side in pixels (default 128).
                color_mode: 'auto' (default — recommend per icon), 'original'
                    (keep SVG default colors), or 'fixed' (use icon_color).
                background_color: Slide background hex, used by 'auto' mode.
                    Default '#F8FAFC'.
                accent_color: Deck accent hex, used by 'auto' mode. Default
                    '#2563EB'.
                theme_name: Sub-directory name under assets/icons/<pack>/
                    (so multiple light/dark variants can coexist).
                icon_color: Required when color_mode='fixed'; ignored
                    otherwise.

            Returns dict with success, out_dir, stdout_tail, stderr_tail,
            message.
            """
            ws, err = _resolve_deck_workspace(deck_workspace)
            if err is not None:
                return err
            if color_mode not in {"auto", "original", "fixed"}:
                return {
                    "success": False,
                    "error": "Invalid color_mode",
                    "message": "color_mode must be one of: auto, original, fixed.",
                }
            if color_mode == "fixed" and not icon_color:
                return {
                    "success": False,
                    "error": "Missing icon_color",
                    "message": "icon_color is required when color_mode='fixed'.",
                }

            scope = pack or "all"
            out_dir = ws / "assets" / "icons" / scope / theme_name
            out_dir.mkdir(parents=True, exist_ok=True)

            args = [
                "render",
                "--size", str(int(size)),
                "--color-mode", color_mode,
                "--background-color", background_color,
                "--accent-color", accent_color,
                "--theme-name", theme_name,
                "--out-dir", str(out_dir),
            ]
            if pack:
                args.extend(["--pack", pack])
            if icon_color:
                args.extend(["--icon-color", icon_color])

            result = _run_ppt_script("icon_registry.py", args, timeout=120)
            result["out_dir"] = str(out_dir)
            # The script may exit non-zero if a single SVG is missing — be
            # explicit so the agent can advise running `icon_registry.py sync`
            # (the agent can't run sync directly; this is a skill-maintenance
            # operation handled out-of-band).
            if result.get("success"):
                pngs = sorted(out_dir.glob("*.png"))
                result["icon_count"] = len(pngs)
                result["message"] = (
                    f"Rendered {len(pngs)} icon(s) under {out_dir}."
                )
            else:
                result["message"] = (
                    "Icon render failed. If stderr mentions a missing .svg, "
                    "the icon registry needs `icon_registry.py sync` first — "
                    "this is a one-off skill-maintenance task (it downloads "
                    "SVGs from the Tabler GitHub repo). Tell the user; this "
                    "tool does not auto-sync."
                )
            return result

        # ============ PPT SKILL TOOLS (Phase 3 — build) ============
        # `ppt_build_pptx_tool` reads a derived slide_specs.yaml and produces
        # a real editable pptx by routing each slide through an archetype
        # renderer. The renderers call into skills/.../scripts/ppt_asset_helpers
        # so theme tokens (fonts, panels, palette) stay consistent with the
        # rest of the PPT skill.
        #
        # Field contract per slide (in addition to the 8 fields derive enforces):
        #   - archetype: one of {hero-statement, decision-logic, board-memo,
        #     chart-spotlight, comparison-matrix, process-flow, research-note,
        #     appendix-dense}; unknown values fall back to hero-statement.
        #   - asset_mode: one of {text-layout-native, office-chart-native,
        #     python-figure-image, table-native, diagram-connector,
        #     diagram-visual, icon-accent, image-hero, mixed}. Drives which
        #     renderer block is used inside the archetype.
        #   - key_message: short conclusion sentence shown as subtitle/answer.
        #   - bullets: list[str] of supporting points (optional).
        #   - chart: optional dict
        #       {chart_type: bar|column|line|stacked_bar, categories: [...],
        #        series: [{name: str, values: [num]}], number_format?: str,
        #        show_legend?: bool}
        #   - image_path: optional absolute or workspace-relative image path
        #     (used by asset_mode in {python-figure-image, image-hero}).
        #   - table: optional dict {headers: [...], rows: [[...], ...],
        #            numeric_columns?: [int]}
        #   - matrix: optional list of {label, attributes: {col: val}} for
        #     comparison-matrix.
        #   - diagram: optional dict
        #       {nodes: [{key, text, left, top, width, height, fill?, line?}],
        #        edges: [{from, to, from_site: top|left|bottom|right,
        #                 to_site: ..., line_rgb?}]}
        #   - caption: optional footer text per slide.

        _ARCHETYPES = {
            "hero-statement",
            "decision-logic",
            "board-memo",
            "chart-spotlight",
            "comparison-matrix",
            "process-flow",
            "research-note",
            "appendix-dense",
        }

        _ASSET_MODES = {
            "text-layout-native",
            "office-chart-native",
            "python-figure-image",
            "table-native",
            "diagram-connector",
            "diagram-visual",
            "icon-accent",
            "image-hero",
            "mixed",
        }

        def _ppt_load_yaml(path: Path) -> tuple[dict | None, str | None]:
            try:
                import yaml  # python-pptx already pulls it in transitively
            except ImportError as exc:
                return None, f"PyYAML not installed: {exc}"
            try:
                data = yaml.safe_load(path.read_text(encoding="utf-8"))
            except Exception as exc:
                return None, f"YAML parse error: {exc}"
            if not isinstance(data, dict):
                return None, "slide_specs.yaml top-level must be a mapping"
            return data, None

        def _ppt_chart_type(name: str):
            """Map a friendly name to pptx XL_CHART_TYPE."""
            from pptx.enum.chart import XL_CHART_TYPE
            return {
                "bar": XL_CHART_TYPE.BAR_CLUSTERED,
                "bar_clustered": XL_CHART_TYPE.BAR_CLUSTERED,
                "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
                "column_clustered": XL_CHART_TYPE.COLUMN_CLUSTERED,
                "line": XL_CHART_TYPE.LINE,
                "stacked_bar": XL_CHART_TYPE.BAR_STACKED,
                "stacked_column": XL_CHART_TYPE.COLUMN_STACKED,
            }.get((name or "bar").lower(), XL_CHART_TYPE.BAR_CLUSTERED)

        def ppt_build_pptx_tool(
            slide_specs_path: str,
            output_pptx: str,
            deck_workspace: str,
        ) -> dict:
            """
            Build an editable .pptx from a derived slide_specs.yaml.

            Each slide is routed to an archetype renderer (hero-statement /
            decision-logic / board-memo / chart-spotlight / comparison-matrix
            / process-flow / research-note / appendix-dense). Theme tokens
            from deck.theme_tokens are injected into the underlying
            ppt_asset_helpers module so fonts, sizes, panel colors and CJK
            (东亚) font slots match the rest of the skill.

            **Do NOT call this tool with hand-crafted YAML.** Always run
            ppt_derive_slide_specs_tool first so the structural fields are
            validated against deck_narrative.md.

            What the tool does NOT do:
              - It does not invent content. Bullets, chart data, table data
                and diagram structure must come from the slide_spec itself.
              - It does not run quality gates. After build, call
                ppt_package_preflight_tool → ppt_structure_precheck_tool →
                ppt_export_previews_tool → ppt_render_review_tool.

            Args:
                slide_specs_path: Absolute path, typically
                    <deck_workspace>/build/generated/slide_specs.yaml.
                output_pptx: Absolute path to write, typically
                    <deck_workspace>/build/pptx/deck_v1.pptx.
                deck_workspace: Value returned by ppt_init_workspace_tool.
                    Used to resolve relative image / chart paths inside the
                    spec.

            Returns dict with:
                success / output_pptx / slide_count / per_slide (list of
                {slide_id, archetype, status, error?}) / message.
            """
            ws, err = _resolve_deck_workspace(deck_workspace)
            if err is not None:
                return err

            specs_path = Path(slide_specs_path)
            if not specs_path.is_absolute():
                return {
                    "success": False,
                    "error": "Relative slide_specs_path",
                    "message": (
                        f"slide_specs_path={slide_specs_path!r} must be absolute. "
                        "Use the slide_specs_path returned by "
                        "ppt_derive_slide_specs_tool."
                    ),
                }
            if not specs_path.exists():
                return {
                    "success": False,
                    "error": "slide_specs not found",
                    "message": (
                        f"No file at {specs_path}. Run "
                        "ppt_derive_slide_specs_tool first."
                    ),
                }

            out_path = Path(output_pptx)
            if not out_path.is_absolute():
                return {
                    "success": False,
                    "error": "Relative output_pptx",
                    "message": "output_pptx must be absolute.",
                }
            out_path.parent.mkdir(parents=True, exist_ok=True)

            data, parse_err = _ppt_load_yaml(specs_path)
            if parse_err:
                return {
                    "success": False,
                    "error": "Spec parse failed",
                    "message": parse_err,
                }
            deck = data.get("deck") or {}
            slides = data.get("slides") or []
            if not isinstance(slides, list) or not slides:
                return {
                    "success": False,
                    "error": "No slides",
                    "message": (
                        "slide_specs.yaml has no `slides` list. Re-run "
                        "ppt_derive_slide_specs_tool after editing the "
                        "narrative."
                    ),
                }

            # Make scripts dir importable for ppt_asset_helpers.
            import sys as _sys
            _saved_path = _sys.path[:]
            _sys.path.insert(0, str(PPT_SCRIPTS_DIR))
            try:
                import ppt_asset_helpers as pah  # type: ignore
                from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_ALIGN
                from pptx.dml.color import RGBColor
                from pptx.util import Inches, Pt
            except Exception as exc:
                _sys.path[:] = _saved_path
                return {
                    "success": False,
                    "error": "Helper import failed",
                    "message": (
                        f"Could not import ppt_asset_helpers: {exc}. "
                        "Confirm python-pptx is installed and the PPT skill "
                        "scripts directory is intact."
                    ),
                }

            # ---- theme tokens injection (restored in finally) ----
            theme_tokens = (deck.get("theme_tokens") or {}) if isinstance(deck, dict) else {}
            saved_tokens = dict(pah.DEFAULT_TYPOGRAPHY_TOKENS)
            saved_latin = pah.DEFAULT_LATIN_FONT_NAME
            saved_ea = pah.DEFAULT_EAST_ASIA_FONT_NAME
            saved_font = pah.DEFAULT_FONT_NAME
            saved_line = pah.DEFAULT_LINE_SPACING_MULTIPLE
            try:
                for token_key in (
                    "hero_title_font_pt", "section_title_font_pt",
                    "page_title_font_pt", "subtitle_font_pt",
                    "minor_title_font_pt", "body_font_pt",
                    "label_font_pt", "caption_font_pt",
                    "title_line_spacing_multiple",
                    "body_line_spacing_multiple",
                    "title_paragraph_space_lines",
                ):
                    if token_key in theme_tokens:
                        try:
                            pah.DEFAULT_TYPOGRAPHY_TOKENS[token_key] = float(theme_tokens[token_key])
                        except (TypeError, ValueError):
                            pass
                latin_font = theme_tokens.get("latin_font_name") or saved_latin
                ea_font = theme_tokens.get("east_asia_font_name") or saved_ea
                pah.DEFAULT_LATIN_FONT_NAME = latin_font
                pah.DEFAULT_EAST_ASIA_FONT_NAME = ea_font
                pah.DEFAULT_FONT_NAME = latin_font
                body_lsm = theme_tokens.get("body_line_spacing_multiple")
                if isinstance(body_lsm, (int, float)):
                    pah.DEFAULT_LINE_SPACING_MULTIPLE = float(body_lsm)

                # CJK font helper: walk every run on a shape and ensure the
                # east-asia font slot is set. python-pptx by default only
                # sets the latin slot, so Chinese characters fall back to a
                # system default in PowerPoint.
                from docx.oxml.ns import qn  # python-docx ships with python-pptx envs
                from docx.oxml import OxmlElement  # noqa: F401

                def _set_ea_font(run, ea_name: str):
                    try:
                        rPr = run._r.get_or_add_rPr()
                        rFonts = rPr.find(qn("w:rFonts"))
                        # pptx uses the drawingml namespace, not w:; fall back
                        # to direct latin attribute writes.
                        from pptx.oxml.ns import qn as pqn
                        rPr2 = run._r.get_or_add_rPr()
                        # Try drawingml-style east-asia slot.
                        ea_elem = rPr2.find(pqn("a:ea"))
                        if ea_elem is None:
                            ea_elem = OxmlElement("a:ea")
                            rPr2.append(ea_elem)
                        ea_elem.set("typeface", ea_name)
                    except Exception:
                        pass

                def _apply_fonts(shape, latin: str = latin_font, ea: str = ea_font):
                    if not getattr(shape, "has_text_frame", False):
                        return
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            try:
                                run.font.name = latin
                            except Exception:
                                pass
                            _set_ea_font(run, ea)

                palette = pah.default_palette()
                tokens = pah.default_typography_tokens()

                def _resolve_asset_path(maybe_path: str | None) -> Path | None:
                    if not maybe_path:
                        return None
                    p = Path(maybe_path)
                    if not p.is_absolute():
                        p = ws / maybe_path
                    return p if p.exists() else None

                # ====== Internal helpers ======================================

                def _add_native_table(
                    slide,
                    *,
                    headers: list,
                    rows: list,
                    left: float,
                    top: float,
                    width: float,
                    height: float,
                    accent_rgb: tuple,
                    numeric_columns: set,
                    table_tokens: dict,
                ):
                    """Add a python-pptx native table that follows the deck's table policy."""
                    rows_count = len(rows) + 1
                    cols_count = max(len(headers), max((len(r) for r in rows), default=0))
                    if cols_count == 0:
                        return
                    table_shape = slide.shapes.add_table(
                        rows_count, cols_count,
                        Inches(left), Inches(top),
                        Inches(width), Inches(height),
                    ).table

                    # Header row.
                    for col_idx in range(cols_count):
                        cell = table_shape.cell(0, col_idx)
                        cell.text = str(headers[col_idx] if col_idx < len(headers) else "")
                        cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(*pah.tint(accent_rgb, 0.18))
                        for para in cell.text_frame.paragraphs:
                            para.alignment = PP_ALIGN.CENTER
                            for run in para.runs:
                                run.font.bold = True
                                run.font.size = Pt(float(table_tokens.get("table_font_pt", 10.5)))
                                run.font.color.rgb = RGBColor(*palette["title"])
                        _apply_fonts(cell)

                    # Body rows.
                    text_align = table_tokens.get("table_text_alignment", "left")
                    numeric_align = table_tokens.get("table_numeric_alignment", "right")
                    index_align = table_tokens.get("table_index_alignment", "left")
                    for r_idx, row in enumerate(rows, start=1):
                        for c_idx in range(cols_count):
                            cell = table_shape.cell(r_idx, c_idx)
                            value = row[c_idx] if c_idx < len(row) else ""
                            cell.text = str(value)
                            cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
                            if c_idx == 0:
                                align = index_align
                            elif c_idx in numeric_columns:
                                align = numeric_align
                            else:
                                align = text_align
                            pp_align = {
                                "left": PP_ALIGN.LEFT,
                                "right": PP_ALIGN.RIGHT,
                                "center": PP_ALIGN.CENTER,
                            }.get(align, PP_ALIGN.LEFT)
                            for para in cell.text_frame.paragraphs:
                                para.alignment = pp_align
                                for run in para.runs:
                                    run.font.size = Pt(float(table_tokens.get("table_font_pt", 10.5)))
                                    run.font.color.rgb = RGBColor(*palette["subtitle"])
                            _apply_fonts(cell)

                def _add_bullets(
                    slide,
                    bullets: list,
                    *,
                    left: float,
                    top: float,
                    width: float,
                    height: float,
                ):
                    if not bullets:
                        return None
                    box = slide.shapes.add_textbox(
                        Inches(left), Inches(top), Inches(width), Inches(height),
                    )
                    tf = box.text_frame
                    tf.word_wrap = True
                    body_size = float(tokens["body_font_pt"])
                    body_lsm = float(tokens["body_line_spacing_multiple"])
                    for i, item in enumerate(bullets):
                        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                        para.text = f"• {item}"
                        para.font.size = Pt(body_size)
                        para.line_spacing = body_lsm
                        para.font.color.rgb = RGBColor(*palette["subtitle"])
                    _apply_fonts(box)
                    return box

                def _add_slide_caption(slide, text: str):
                    """Footer caption with enough height for the configured
                    caption font (pah.add_caption's default 0.22-inch box
                    triggers a textbox_fit_failure in structure_precheck).
                    Also leaves the bottom-right 1.1 inch clear so the page
                    figure_tag set by add_slide_header is not occluded."""
                    if not text:
                        return None
                    caption_pt = float(tokens["caption_font_pt"])
                    # Box height: at least 2x font size in inches (72 pt = 1 in).
                    box_h = max(0.32, caption_pt * 2.4 / 72.0)
                    box_top = 9.0 - box_h - 0.32
                    # Stop before the figure_tag (added at left=14.85 by
                    # add_slide_header) so structure_precheck does not flag
                    # an occlusion between caption and tag.
                    box = slide.shapes.add_textbox(
                        Inches(0.7), Inches(box_top),
                        Inches(13.7), Inches(box_h),
                    )
                    box.text_frame.word_wrap = True
                    box.text_frame.margin_left = Inches(0.05)
                    box.text_frame.margin_right = Inches(0.05)
                    box.text_frame.margin_top = Inches(0.02)
                    box.text_frame.margin_bottom = Inches(0.02)
                    para = box.text_frame.paragraphs[0]
                    para.text = text
                    para.font.size = Pt(caption_pt)
                    para.line_spacing = 1.0
                    para.font.color.rgb = RGBColor(*palette["subtitle"])
                    _apply_fonts(box)
                    return box

                # ====== Archetype renderers ===================================

                def _render_hero_statement(slide, spec: dict):
                    title = spec.get("title") or "Untitled"
                    key_message = spec.get("key_message") or ""
                    slide.background.fill.solid()
                    slide.background.fill.fore_color.rgb = RGBColor(*palette["bg"])

                    title_box = slide.shapes.add_textbox(
                        Inches(0.9), Inches(2.6), Inches(14.2), Inches(1.6),
                    )
                    p = title_box.text_frame.paragraphs[0]
                    p.text = title
                    p.font.bold = True
                    p.font.size = Pt(float(tokens["hero_title_font_pt"]) * 1.4)
                    p.font.color.rgb = RGBColor(*palette["title"])
                    p.line_spacing = float(tokens["title_line_spacing_multiple"])
                    _apply_fonts(title_box)

                    if key_message:
                        sub = slide.shapes.add_textbox(
                            Inches(0.9), Inches(4.6), Inches(14.2), Inches(1.6),
                        )
                        sp = sub.text_frame.paragraphs[0]
                        sp.text = key_message
                        sp.font.size = Pt(float(tokens["subtitle_font_pt"]) * 1.2)
                        sp.font.color.rgb = RGBColor(*palette["subtitle"])
                        sp.line_spacing = float(tokens["body_line_spacing_multiple"])
                        _apply_fonts(sub)

                    return ["hero_layout"]

                def _render_standard_header(slide, spec: dict, figure_tag: str):
                    pah.add_slide_header(
                        slide,
                        figure_tag=figure_tag,
                        title=spec.get("title") or "",
                        subtitle=spec.get("key_message") or "",
                    )
                    # Re-apply CJK font on the freshly added header shapes.
                    for shape in slide.shapes:
                        _apply_fonts(shape)

                def _render_decision_logic(slide, spec: dict, idx: int):
                    _render_standard_header(slide, spec, f"S{idx:02d}")
                    bullets = spec.get("bullets") or []
                    accent = palette["blue"]
                    pah.add_panel(
                        slide, title="Key judgement",
                        left=0.7, top=1.5, width=8.4, height=6.4,
                        accent_rgb=accent,
                    )
                    box = _add_bullets(slide, bullets, left=1.0, top=2.1,
                                       width=7.8, height=5.5)
                    pah.add_panel(
                        slide, title="Why this matters",
                        left=9.4, top=1.5, width=5.9, height=6.4,
                        accent_rgb=palette["emerald"],
                    )
                    notes = spec.get("notes") or spec.get("narrative_markdown") or ""
                    if notes:
                        nb = slide.shapes.add_textbox(
                            Inches(9.7), Inches(2.1), Inches(5.3), Inches(5.5),
                        )
                        np_para = nb.text_frame.paragraphs[0]
                        np_para.text = str(notes)[:600]
                        np_para.font.size = Pt(float(tokens["body_font_pt"]))
                        np_para.font.color.rgb = RGBColor(*palette["subtitle"])
                        np_para.line_spacing = float(tokens["body_line_spacing_multiple"])
                        _apply_fonts(nb)
                    if spec.get("caption"):
                        _add_slide_caption(slide, spec["caption"])
                    return ["panel_left", "panel_right"]

                def _render_board_memo(slide, spec: dict, idx: int):
                    _render_standard_header(slide, spec, f"S{idx:02d}")
                    bullets = spec.get("bullets") or []
                    panels = ["Summary", "Progress", "Risks", "Next steps"]
                    # 2x2 grid of panels.
                    positions = [
                        (0.7, 1.5, 7.3, 3.0),
                        (8.3, 1.5, 7.0, 3.0),
                        (0.7, 4.7, 7.3, 3.2),
                        (8.3, 4.7, 7.0, 3.2),
                    ]
                    palette_colors = [palette["blue"], palette["emerald"],
                                      palette["amber"], palette["violet"]]
                    for i, ((l, t, w, h), color, label) in enumerate(
                        zip(positions, palette_colors, panels)
                    ):
                        pah.add_panel(slide, label, l, t, w, h, color)
                        text = ""
                        if i < len(bullets):
                            text = str(bullets[i])
                        if text:
                            tb = slide.shapes.add_textbox(
                                Inches(l + 0.15), Inches(t + 0.6),
                                Inches(w - 0.3), Inches(h - 0.8),
                            )
                            para = tb.text_frame.paragraphs[0]
                            para.text = text
                            para.font.size = Pt(float(tokens["body_font_pt"]))
                            para.font.color.rgb = RGBColor(*palette["subtitle"])
                            para.line_spacing = float(tokens["body_line_spacing_multiple"])
                            tb.text_frame.word_wrap = True
                            _apply_fonts(tb)
                    if spec.get("caption"):
                        _add_slide_caption(slide, spec["caption"])
                    return ["2x2_panels"]

                def _render_chart_spotlight(slide, spec: dict, idx: int):
                    _render_standard_header(slide, spec, f"S{idx:02d}")
                    asset_mode = spec.get("asset_mode") or "office-chart-native"
                    chart = spec.get("chart") or {}
                    image_path = _resolve_asset_path(spec.get("image_path"))
                    accent = palette["blue"]

                    if asset_mode == "office-chart-native" and chart.get("categories") and chart.get("series"):
                        series_list = [
                            (s.get("name", "Series"), list(s.get("values") or []))
                            for s in (chart.get("series") or [])
                            if isinstance(s, dict)
                        ]
                        pah.add_native_chart_card(
                            slide,
                            title=chart.get("title") or "Chart",
                            left=0.7, top=1.5, width=9.6, height=6.4,
                            accent_rgb=accent,
                            categories=list(chart["categories"]),
                            series_list=series_list,
                            chart_type=_ppt_chart_type(chart.get("chart_type")),
                            number_format=chart.get("number_format") or "0",
                            show_legend=bool(chart.get("show_legend", False)),
                        )
                    elif asset_mode == "python-figure-image" and image_path:
                        pah.add_picture_card(
                            slide,
                            title=chart.get("title") or spec.get("title") or "Figure",
                            image_path=image_path,
                            left=0.7, top=1.5, width=9.6, height=6.4,
                            accent_rgb=accent,
                            caption=chart.get("caption"),
                        )
                    else:
                        # Fallback: panel + note that chart data is missing.
                        pah.add_panel(slide, "Chart pending", 0.7, 1.5, 9.6, 6.4, accent)
                        nb = slide.shapes.add_textbox(
                            Inches(1.0), Inches(2.0), Inches(8.8), Inches(5.0),
                        )
                        nb_para = nb.text_frame.paragraphs[0]
                        nb_para.text = (
                            "Chart data not provided in slide_spec. Add a "
                            "`chart:` block with categories + series, or an "
                            "`image_path:` for python-figure-image."
                        )
                        nb_para.font.size = Pt(float(tokens["body_font_pt"]))
                        nb_para.font.color.rgb = RGBColor(*palette["muted"])
                        _apply_fonts(nb)

                    # Right column: takeaways.
                    pah.add_panel(
                        slide, "Takeaways",
                        left=10.5, top=1.5, width=4.8, height=6.4,
                        accent_rgb=palette["emerald"],
                    )
                    _add_bullets(slide, spec.get("bullets") or [],
                                 left=10.8, top=2.1, width=4.3, height=5.5)
                    if spec.get("caption"):
                        _add_slide_caption(slide, spec["caption"])
                    return ["chart_or_image", "takeaways_panel"]

                def _render_comparison_matrix(slide, spec: dict, idx: int):
                    _render_standard_header(slide, spec, f"S{idx:02d}")
                    table_spec = spec.get("table") or {}
                    matrix = spec.get("matrix") or []
                    table_tokens = {
                        k: theme_tokens.get(k)
                        for k in (
                            "table_font_pt", "table_text_alignment",
                            "table_numeric_alignment", "table_index_alignment",
                        )
                        if theme_tokens.get(k) is not None
                    }
                    numeric_columns = set(table_spec.get("numeric_columns") or [])

                    if table_spec.get("headers") and table_spec.get("rows"):
                        _add_native_table(
                            slide,
                            headers=list(table_spec["headers"]),
                            rows=[list(r) for r in table_spec["rows"]],
                            left=0.7, top=1.5, width=14.6, height=6.4,
                            accent_rgb=palette["blue"],
                            numeric_columns=numeric_columns,
                            table_tokens=table_tokens,
                        )
                    elif matrix:
                        # Convert matrix list-of-dicts into a table.
                        all_attrs: list = []
                        for row in matrix:
                            for k in (row.get("attributes") or {}).keys():
                                if k not in all_attrs:
                                    all_attrs.append(k)
                        headers = ["选项"] + all_attrs
                        rows = []
                        for row in matrix:
                            line = [row.get("label", "")]
                            attrs = row.get("attributes") or {}
                            for k in all_attrs:
                                line.append(str(attrs.get(k, "")))
                            rows.append(line)
                        _add_native_table(
                            slide,
                            headers=headers,
                            rows=rows,
                            left=0.7, top=1.5, width=14.6, height=6.4,
                            accent_rgb=palette["blue"],
                            numeric_columns=numeric_columns,
                            table_tokens=table_tokens,
                        )
                    else:
                        pah.add_panel(slide, "Comparison data missing",
                                      0.7, 1.5, 14.6, 6.4, palette["amber"])
                    if spec.get("caption"):
                        _add_slide_caption(slide, spec["caption"])
                    return ["matrix_table"]

                def _render_process_flow(slide, spec: dict, idx: int):
                    _render_standard_header(slide, spec, f"S{idx:02d}")
                    diagram = spec.get("diagram") or {}
                    nodes_spec = diagram.get("nodes") or []
                    edges_spec = diagram.get("edges") or []
                    bullets = spec.get("bullets") or []

                    if nodes_spec:
                        # Draw user-provided nodes.
                        node_objs = {}
                        for nd in nodes_spec:
                            key = nd.get("key") or f"n{len(node_objs)}"
                            style = pah.NodeStyle(
                                fill_rgb=tuple(nd.get("fill") or palette["blue"]),
                                line_rgb=tuple(nd.get("line") or palette["line"]),
                            )
                            shape = pah.add_node(
                                slide, key=key, text=str(nd.get("text") or key),
                                left=float(nd.get("left", 1.0)),
                                top=float(nd.get("top", 4.0)),
                                width=float(nd.get("width", 2.4)),
                                height=float(nd.get("height", 1.2)),
                                style=style,
                            )
                            _apply_fonts(shape)
                            node_objs[key] = shape
                        for ed in edges_spec:
                            src = node_objs.get(ed.get("from"))
                            dst = node_objs.get(ed.get("to"))
                            if src is None or dst is None:
                                continue
                            pah.add_glued_connector(
                                slide,
                                from_shape=src, to_shape=dst,
                                from_site=ed.get("from_site") or "right",
                                to_site=ed.get("to_site") or "left",
                                line_rgb=tuple(ed.get("line_rgb") or palette["line"]),
                            )
                    elif bullets:
                        # Visual fallback: evenly-spaced node row from bullets.
                        n = min(len(bullets), 6)
                        margin = 0.7
                        width = 2.2
                        gap = (16 - 2 * margin - n * width) / max(1, n - 1) if n > 1 else 0
                        accents = [palette["blue"], palette["emerald"],
                                   palette["amber"], palette["violet"],
                                   palette["teal"], palette["rose"]]
                        prev_shape = None
                        for i in range(n):
                            left = margin + i * (width + gap)
                            style = pah.NodeStyle(
                                fill_rgb=accents[i % len(accents)],
                                line_rgb=palette["line"],
                            )
                            shape = pah.add_node(
                                slide, key=f"step_{i}", text=str(bullets[i]),
                                left=left, top=4.0, width=width, height=1.4,
                                style=style,
                            )
                            _apply_fonts(shape)
                            if prev_shape is not None:
                                pah.add_glued_connector(
                                    slide,
                                    from_shape=prev_shape, to_shape=shape,
                                    from_site="right", to_site="left",
                                    line_rgb=palette["line"],
                                )
                            prev_shape = shape
                    else:
                        pah.add_panel(slide, "Process steps missing",
                                      0.7, 1.5, 14.6, 6.4, palette["amber"])
                    if spec.get("caption"):
                        _add_slide_caption(slide, spec["caption"])
                    return ["process_nodes"]

                def _render_research_note(slide, spec: dict, idx: int):
                    _render_standard_header(slide, spec, f"S{idx:02d}")
                    image_path = _resolve_asset_path(spec.get("image_path"))
                    pah.add_panel(slide, "Mechanism / structure",
                                  0.7, 1.5, 10.2, 6.4, palette["blue"])
                    if image_path:
                        pah.add_picture_card(
                            slide,
                            title=spec.get("image_title") or "Figure",
                            image_path=image_path,
                            left=0.85, top=1.85, width=9.9, height=5.7,
                            accent_rgb=palette["blue"],
                            caption=spec.get("image_caption"),
                        )
                    else:
                        nb = slide.shapes.add_textbox(
                            Inches(1.0), Inches(2.1), Inches(9.7), Inches(5.5),
                        )
                        para = nb.text_frame.paragraphs[0]
                        para.text = str(spec.get("notes") or spec.get("narrative_markdown") or "")[:1200]
                        para.font.size = Pt(float(tokens["body_font_pt"]))
                        para.font.color.rgb = RGBColor(*palette["subtitle"])
                        para.line_spacing = float(tokens["body_line_spacing_multiple"])
                        nb.text_frame.word_wrap = True
                        _apply_fonts(nb)
                    pah.add_panel(slide, "Annotations",
                                  11.1, 1.5, 4.2, 6.4, palette["emerald"])
                    _add_bullets(slide, spec.get("bullets") or [],
                                 left=11.35, top=2.1, width=3.85, height=5.5)
                    if spec.get("caption"):
                        _add_slide_caption(slide, spec["caption"])
                    return ["panel_left", "panel_right"]

                def _render_appendix_dense(slide, spec: dict, idx: int):
                    _render_standard_header(slide, spec, f"S{idx:02d}")
                    table_spec = spec.get("table") or {}
                    table_tokens = {
                        k: theme_tokens.get(k)
                        for k in (
                            "table_font_pt", "table_text_alignment",
                            "table_numeric_alignment", "table_index_alignment",
                        )
                        if theme_tokens.get(k) is not None
                    }
                    if table_spec.get("headers") and table_spec.get("rows"):
                        _add_native_table(
                            slide,
                            headers=list(table_spec["headers"]),
                            rows=[list(r) for r in table_spec["rows"]],
                            left=0.7, top=1.5, width=14.6, height=6.4,
                            accent_rgb=palette["slate"],
                            numeric_columns=set(table_spec.get("numeric_columns") or []),
                            table_tokens=table_tokens,
                        )
                    else:
                        pah.add_panel(slide, "Appendix data",
                                      0.7, 1.5, 14.6, 6.4, palette["slate"])
                    if spec.get("caption"):
                        _add_slide_caption(slide, spec["caption"])
                    return ["appendix_table"]

                # Dispatcher.
                renderers = {
                    "hero-statement":
                        lambda slide, spec, idx: _render_hero_statement(slide, spec),
                    "decision-logic": _render_decision_logic,
                    "board-memo": _render_board_memo,
                    "chart-spotlight": _render_chart_spotlight,
                    "comparison-matrix": _render_comparison_matrix,
                    "process-flow": _render_process_flow,
                    "research-note": _render_research_note,
                    "appendix-dense": _render_appendix_dense,
                }

                # ====== Build loop ===========================================
                prs = pah.new_presentation()
                blank_layout = prs.slide_layouts[6]  # 'Blank'
                per_slide = []

                for idx, spec in enumerate(slides, start=1):
                    if not isinstance(spec, dict):
                        per_slide.append({
                            "slide_id": f"S{idx:02d}",
                            "archetype": None,
                            "status": "error",
                            "error": "spec is not a mapping",
                        })
                        continue
                    archetype = spec.get("archetype") or "hero-statement"
                    if archetype not in _ARCHETYPES:
                        # Unknown archetype → fall back, but record it.
                        per_slide.append({
                            "slide_id": spec.get("slide_id") or f"S{idx:02d}",
                            "archetype": archetype,
                            "status": "warning",
                            "error": (
                                f"unknown archetype {archetype!r}, "
                                "falling back to hero-statement"
                            ),
                        })
                        archetype = "hero-statement"

                    slide = prs.slides.add_slide(blank_layout)
                    renderer = renderers.get(archetype, renderers["hero-statement"])
                    try:
                        if archetype == "hero-statement":
                            renderer(slide, spec, idx)
                        else:
                            renderer(slide, spec, idx)
                    except Exception as exc:
                        import traceback
                        per_slide.append({
                            "slide_id": spec.get("slide_id") or f"S{idx:02d}",
                            "archetype": archetype,
                            "status": "error",
                            "error": f"{type(exc).__name__}: {exc}",
                            "traceback_tail": traceback.format_exc()[-1500:],
                        })
                        continue

                    per_slide.append({
                        "slide_id": spec.get("slide_id") or f"S{idx:02d}",
                        "archetype": archetype,
                        "status": "ok",
                    })

                pah.save_presentation(prs, out_path)

                # Post-process: python-pptx does not update docProps/app.xml's
                # <Slides> counter on save. ppt_package_preflight_tool flags
                # this as a hard error (`docprops_slide_count_mismatch`) since
                # mobile parsers reject it. Patch the counter in-place so a
                # freshly built deck always passes preflight on its first run.
                try:
                    import zipfile as _zipfile
                    import re as _re

                    actual_pages = len(prs.slides)
                    tmp_path = out_path.with_suffix(out_path.suffix + ".tmp")
                    with _zipfile.ZipFile(str(out_path), "r") as zin:
                        names = zin.namelist()
                        with _zipfile.ZipFile(
                            str(tmp_path), "w", _zipfile.ZIP_DEFLATED
                        ) as zout:
                            for name in names:
                                data = zin.read(name)
                                if name == "docProps/app.xml":
                                    text = data.decode("utf-8", "ignore")
                                    if "<Slides>" in text:
                                        text = _re.sub(
                                            r"<Slides>\s*\d+\s*</Slides>",
                                            f"<Slides>{actual_pages}</Slides>",
                                            text,
                                            count=1,
                                        )
                                    else:
                                        # Insert a <Slides> element just before
                                        # the closing </Properties> tag.
                                        text = text.replace(
                                            "</Properties>",
                                            f"<Slides>{actual_pages}</Slides></Properties>",
                                            1,
                                        )
                                    data = text.encode("utf-8")
                                zout.writestr(name, data)
                    out_path.unlink()
                    tmp_path.rename(out_path)
                except Exception as _exc:  # noqa: BLE001
                    # The build itself succeeded; surface the post-process
                    # issue but do not fail the tool — preflight will pick it
                    # up if the counter is still wrong.
                    print(f"⚠️ docProps post-process skipped: {_exc}")

                fe_data = _build_files_event_data(
                    str(out_path), f"Built PPTX: {out_path.name}"
                )
                if fe_data:
                    _pending_files_events.append(fe_data)

                ok_count = sum(1 for s in per_slide if s["status"] == "ok")
                err_count = sum(1 for s in per_slide if s["status"] == "error")
                warn_count = sum(1 for s in per_slide if s["status"] == "warning")
                return {
                    "success": err_count == 0,
                    "output_pptx": str(out_path),
                    "slide_count": len(per_slide),
                    "ok_count": ok_count,
                    "warning_count": warn_count,
                    "error_count": err_count,
                    "per_slide": per_slide,
                    "message": (
                        f"Built {ok_count}/{len(per_slide)} slide(s) into {out_path.name}. "
                        + (f"{warn_count} warning(s). " if warn_count else "")
                        + (f"{err_count} error(s) — check per_slide. " if err_count else "")
                        + "Next: run ppt_package_preflight_tool → "
                        "ppt_structure_precheck_tool → ppt_export_previews_tool "
                        "→ ppt_render_review_tool."
                    ),
                }
            except Exception as exc:
                import traceback
                return {
                    "success": False,
                    "error": f"Build aborted: {type(exc).__name__}: {exc}",
                    "traceback_tail": traceback.format_exc()[-2000:],
                    "message": (
                        "ppt_build_pptx_tool crashed before finishing. "
                        "Check traceback_tail; the most common cause is a "
                        "malformed chart/table/diagram block in slide_specs.yaml."
                    ),
                }
            finally:
                pah.DEFAULT_TYPOGRAPHY_TOKENS.clear()
                pah.DEFAULT_TYPOGRAPHY_TOKENS.update(saved_tokens)
                pah.DEFAULT_LATIN_FONT_NAME = saved_latin
                pah.DEFAULT_EAST_ASIA_FONT_NAME = saved_ea
                pah.DEFAULT_FONT_NAME = saved_font
                pah.DEFAULT_LINE_SPACING_MULTIPLE = saved_line
                _sys.path[:] = _saved_path

        tools = [
            process_document,
            edit_docx_tool,
            edit_docx_content_tool,   # alias — DO NOT remove (anti-loop guard)
            edit_docx_content,        # alias — DO NOT remove (anti-loop guard)
            extract_docx_content_tool,
            # New enhanced editing tools
            delete_docx_content_tool,
            modify_docx_fonts_tool,
            create_docx_with_content_tool,
            # Legacy .doc → .docx conversion (call first for any uploaded .doc)
            convert_doc_to_docx_tool,
            # Template-fill tools (uploaded .docx as template)
            inspect_docx_template_tool,
            fill_docx_template_tool,
            # Template library — shared + per-user persistent catalog
            list_templates_tool,
            get_template_path_tool,
            save_template_tool,
            delete_template_tool,
            add_bullet_list_tool,
            add_numbered_list_tool,
            add_comment_tool,
            remove_comment_tool,
            # Contract review (format + fill + consistency + LLM legal red-flags)
            review_contract_tool,
            # XML-level formatting-safe tools (from docx skill)
            unpack_docx_tool,
            pack_docx_tool,
            validate_docx_tool,
            accept_tracked_changes_tool,
            add_xml_comment_tool,
            # PPT skill tools (Phase 1 — environment probe, reference loader,
            # workspace initializer).
            ppt_read_skill_reference_tool,
            ppt_check_environment_tool,
            ppt_init_workspace_tool,
            # PPT skill tools (Phase 2 — template audit, derive specs,
            # workspace lint, 3 quality gates, preview export, connector
            # check, icon search/render).
            ppt_audit_template_tool,
            ppt_derive_slide_specs_tool,
            ppt_lint_workspace_tool,
            ppt_package_preflight_tool,
            ppt_structure_precheck_tool,
            ppt_export_previews_tool,
            ppt_render_review_tool,
            ppt_connectors_check_tool,
            ppt_icon_search_tool,
            ppt_icon_render_tool,
            # PPT skill tools (Phase 3 — build editable pptx from
            # derived slide_specs.yaml via 8 archetype renderers).
            ppt_build_pptx_tool,
        ]
    
    return DocMasterAgent(
        pending_files_events=_pending_files_events,
        name="DocMaster",
        model_client=set_model_client(default_config_name),
        system_message=SYSTEM,
        reflect_on_tool_use=False,  # Disable reflection to simplify
        model_client_stream=True,  # Disable streaming to avoid timeout issues
        
        # DrSaiAgent特定配置
        thread_id=thread_id,
        db_manager=db_manager,
        user_id=user_id,
        set_model_client=set_model_client,
        llm_mode_config=llm_mode_config,
        
        # 技能和工作目录
        skills_dir=[
            str(Path(__file__).parent / "document_skills") if DOCUMENT_PROCESSING_AVAILABLE else os.getenv("SYSTEM_SKILLS_DIR"),
            str(Path(__file__).parent / "skills"),
        ],
        work_dir=WORKDIR,
        only_in_workspace=True,
        
        # Tools configuration
        tools=tools,
        
        # 子智能体配置
        sub_agent_config=SUB_AGENTS,
        
        # 资源限制
        token_limit=60000,  # deepseek-v4-flash supports 128k context
        
        # RAG集成（可选）
        rag_flow_url=os.getenv('RAGFLOW_URL'),
        rag_flow_token=os.getenv('RAGFLOW_TOKEN'),
        memory_dataset_id=os.getenv('MEMORY_DATASET_ID'),
        
        # 额外配置
        max_turn_count=30,
    )

def main():
    """主函数：启动Word文档编辑智能体"""
    from drsai.backend import run_worker, run_console
    
    # 方式1：作为Worker服务运行（注册到HepAI平台）
    asyncio.run(
        run_worker(
            # 智能体注册信息
            agent_name="DocMaster",
            author="haiuser01@ihep.ac.cn",  # 改成你的邮箱
            description="专业的Word文档处理大师，支持上传、分析、编辑、格式化Word文档，支持添加和删除批注和评论",
            version="1.0.0",
            logo="docmaster_logo.png",  # 需要提供logo URL

            permission='groups: drsai; users: admin, haiuser01@ihep.ac.cn, ddf_free, yqsun@ihep.ac.cn; owner: haiuser01@ihep.ac.cn',
            
            # 示例对话
            examples=[
                "DocMaster，请帮我分析这份文档的主要内容",
                "先读取这份 DOCX 的内容，再帮我润色引言部分",
                "把文档中的技术术语替换为更通俗的表达",
                "在这份 DOCX 末尾新增一个总结段落",
                "新建一份 DOCX，包含标题、正文和一个简单表格",
                "把这份 DOCX 的中文设为宋体、英文设为 Times New Roman",
                "这是一份合同模板，请按以下信息填充并生成新文档：甲方=张三，乙方=李四，日期=2026-05-13",
                "我上传了一份带 {{ name }}、{{ date }} 占位符的模板，请帮我填充",
                "帮我做一份 6 页的 PPT，主题是『2026 Q2 安全合规季报』，目标读者是所领导",
                "做一份 4 页的产品周报 deck，包含一页趋势图、一页方案对比矩阵、一页结论页",
                "我有一份参考 pptx 模板，请按它的页面系统做一份汇报",
            ],
            
            # 模型配置
            agent_config=llm_mode_config,
            default_config_name="deepseek-v4-flash(Fast)",
            
            # 智能体工厂
            agent_factory=create_word_editor_agent,
            
            # 服务配置
            port=42819,  # 选择一个未使用的端口
            no_register=False,  # 注册到HepAI平台
            enable_openwebui_pipeline=True,
            history_mode="backend",
            
            # 其他配置
            join_topics=["document-processing", "office-tools"],
            metadata={
                "category": "文档处理大师",
                "tags": ["docmaster", "word", "文档编辑", "办公自动化", "专业文档"],
                "capabilities": ["文档分析", "内容编辑", "格式优化", "结构重组", "专业排版"],
                "dependencies": {
                    "python": ["python-docx", "PyPDF2", "python-pptx", "pandas", "openpyxl"],
                    "system": ["pandoc", "libreoffice", "poppler-utils (pdftoppm)"],
                    "npm": ["docx"],
                    "install_system": "sudo apt install pandoc libreoffice poppler-utils && npm install -g docx",
                }
            },
        )
    )
    
    # 方式2：控制台测试（取消注释以使用）
    # asyncio.run(
    #     run_console(
    #         agent_factory=create_word_editor_agent,
    #         task="请帮我分析这个Word文档"
    #     )
    # )

if __name__ == "__main__":
    main()
# __DRSAI_CWD__:/aifs/user/home/haiuser01/drsai_code/workspace/runs
