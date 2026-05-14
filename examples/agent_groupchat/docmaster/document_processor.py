"""
Document Processor for DocMaster
Handles document upload, analysis, and basic processing for various file types.
"""

import os
import json
import mimetypes
from pathlib import Path
from typing import Dict, Any, Optional, List
import hashlib
from datetime import datetime
import logging

logger = logging.getLogger(__name__)


# ==================== Shared CJK Font Helpers ====================

_CJK_FONT_KEYWORDS = (
    '宋', '黑', '楷', '仿宋', '隶书', '幼圆', '华文', '微软雅黑',
    'SimSun', 'NSimSun', 'SimHei', 'KaiTi', 'FangSong',
    'STSong', 'STHeiti', 'STKaiti', 'STFangsong', 'STZhongsong',
    'MingLiU', 'PMingLiU', 'MS Mincho', 'MS Gothic',
    'Malgun', 'Batang', 'Gulim', 'DengXian', '等线',
)


def _contains_chinese(text: str) -> bool:
    """Return True if *text* contains any CJK Unified Ideograph."""
    import re
    return bool(re.search(r'[\u4e00-\u9fff]', text))


def _is_cjk_font(font_name: str) -> bool:
    """Return True if *font_name* looks like a CJK (Chinese/Japanese/Korean) font."""
    return any(kw in font_name for kw in _CJK_FONT_KEYWORDS)


def _set_east_asia_font(run_element, font_name: str):
    """Set the ``w:eastAsia`` attribute on a run's ``<w:rFonts>`` element.

    This is required for Chinese/Japanese/Korean characters to render in the
    correct font inside Word.  ``run.font.name`` in *python-docx* only writes
    the ``w:ascii`` / ``w:hAnsi`` attributes.
    """
    from docx.oxml.ns import qn

    rPr = run_element.get_or_add_rPr()
    rFonts = rPr.find(qn('w:rFonts'))
    if rFonts is None:
        from docx.oxml import OxmlElement
        rFonts = OxmlElement('w:rFonts')
        rPr.insert(0, rFonts)
    rFonts.set(qn('w:eastAsia'), font_name)


def _apply_font_name_to_run(run, font_name: str):
    """Set *font_name* on a run, including ``w:eastAsia`` for CJK fonts."""
    if not font_name:
        return
    run.font.name = font_name
    if _is_cjk_font(font_name):
        _set_east_asia_font(run._element, font_name)


def _apply_font_name_to_style(style, font_name: str):
    """Set *font_name* on a style object, including ``w:eastAsia`` for CJK fonts."""
    if not font_name:
        return
    from docx.oxml.ns import qn

    style.font.name = font_name
    if _is_cjk_font(font_name):
        rPr = style.element.get_or_add_rPr()
        rFonts = rPr.find(qn('w:rFonts'))
        if rFonts is None:
            from docx.oxml import OxmlElement
            rFonts = OxmlElement('w:rFonts')
            rPr.insert(0, rFonts)
        rFonts.set(qn('w:eastAsia'), font_name)


class DocumentProcessor:
    """
    Main document processor for handling various file types.
    Supports: .docx, .pdf, .pptx, .xlsx, .txt, .md
    """
    
    def __init__(self, workspace_dir: str):
        """
        Initialize document processor with workspace directory.
        
        Args:
            workspace_dir: Directory where uploaded documents are stored
        """
        self.workspace_dir = Path(workspace_dir)
        self.workspace_dir.mkdir(parents=True, exist_ok=True)
        
        # Check for available libraries
        self._check_dependencies()
    
    def _check_dependencies(self):
        """Check for required document processing libraries."""
        self.dependencies = {
            'docx': self._try_import('docx', 'python-docx'),
            'pdf': self._try_import('PyPDF2', 'PyPDF2'),
            'pptx': self._try_import('pptx', 'python-pptx'),
            'excel': self._try_import('pandas', 'pandas') and self._try_import('openpyxl', 'openpyxl'),
        }
    
    def _try_import(self, module_name: str, pip_name: str) -> bool:
        """Try to import a module and return availability."""
        try:
            __import__(module_name)
            return True
        except ImportError:
            logger.warning(f"Module {module_name} not available. Install with: pip install {pip_name}")
            return False
    
    def process_uploaded_file(self, file_path: str, original_filename: str = None) -> Dict[str, Any]:
        """
        Process an uploaded file and extract information.
        
        Args:
            file_path: Path to the uploaded file
            original_filename: Original filename from user upload (optional)
            
        Returns:
            Dictionary with document analysis results
        """
        file_path = Path(file_path)
        
        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")
        
        # Get file information
        file_info = self._get_file_info(file_path, original_filename)
        
        # Determine file type and process accordingly
        ext = file_path.suffix.lower()
        
        try:
            content_info = self._process_by_extension(file_path, ext)
            file_info.update(content_info)
            file_info['processing_success'] = True
            file_info['processing_error'] = None
        except Exception as e:
            logger.error(f"Error processing file {file_path}: {e}")
            file_info['processing_success'] = False
            file_info['processing_error'] = str(e)
            file_info.update(self._process_generic(file_path))
        
        # Generate a summary
        file_info['summary'] = self._generate_summary(file_info)
        
        return file_info
    
    def _get_file_info(self, file_path: Path, original_filename: str = None) -> Dict[str, Any]:
        """Extract basic file information."""
        stat = file_path.stat()
        
        # Calculate file hash
        file_hash = self._calculate_file_hash(file_path)
        
        return {
            'file_path': str(file_path),
            'original_filename': original_filename or file_path.name,
            'file_name': file_path.name,
            'file_size': stat.st_size,
            'file_size_human': self._human_readable_size(stat.st_size),
            'file_extension': file_path.suffix.lower(),
            'created_time': datetime.fromtimestamp(stat.st_ctime).isoformat(),
            'modified_time': datetime.fromtimestamp(stat.st_mtime).isoformat(),
            'mime_type': mimetypes.guess_type(file_path)[0] or 'application/octet-stream',
            'file_hash': file_hash,
        }
    
    def _process_by_extension(self, file_path: Path, ext: str) -> Dict[str, Any]:
        """Process file based on extension."""
        if ext == '.docx' and self.dependencies['docx']:
            return self._process_docx(file_path)
        elif ext == '.pdf' and self.dependencies['pdf']:
            return self._process_pdf(file_path)
        elif ext == '.pptx' and self.dependencies['pptx']:
            return self._process_pptx(file_path)
        elif ext in ['.xlsx', '.xls'] and self.dependencies['excel']:
            return self._process_excel(file_path)
        elif ext == '.csv' and self.dependencies['excel']:
            return self._process_csv(file_path)
        elif ext in ['.txt', '.md']:
            return self._process_text(file_path)
        else:
            return self._process_generic(file_path)
    
    def _process_docx(self, file_path: Path) -> Dict[str, Any]:
        """Process Word document (.docx)."""
        import docx
        
        doc = docx.Document(file_path)
        
        # Extract paragraphs
        paragraphs = [para.text.strip() for para in doc.paragraphs if para.text.strip()]
        
        # Count tables
        table_count = len(doc.tables)
        
        return {
            'content_type': 'docx',
            'paragraph_count': len(paragraphs),
            'table_count': table_count,
            'content_preview': ' '.join(paragraphs[:3])[:500] if paragraphs else '',
            'has_content': len(paragraphs) > 0,
        }
    
    def _process_pdf(self, file_path: Path) -> Dict[str, Any]:
        """Process PDF document."""
        from PyPDF2 import PdfReader
        
        reader = PdfReader(file_path)
        
        # Extract text from first few pages
        total_text = ""
        for page in reader.pages[:3]:  # First 3 pages only
            text = page.extract_text()
            if text:
                total_text += text + "\n"
        
        return {
            'content_type': 'pdf',
            'page_count': len(reader.pages),
            'has_text': len(total_text.strip()) > 0,
            'content_preview': total_text[:500] if total_text else '',
            'is_encrypted': reader.is_encrypted,
        }
    
    def _process_pptx(self, file_path: Path) -> Dict[str, Any]:
        """Process PowerPoint presentation."""
        import pptx
        
        prs = pptx.Presentation(file_path)
        
        # Count slides
        slide_count = len(prs.slides)
        
        return {
            'content_type': 'pptx',
            'slide_count': slide_count,
            'has_content': slide_count > 0,
        }
    
    def _process_excel(self, file_path: Path) -> Dict[str, Any]:
        """Process Excel spreadsheet."""
        import pandas as pd
        
        # Get sheet names
        xl_file = pd.ExcelFile(file_path)
        sheet_names = xl_file.sheet_names
        
        return {
            'content_type': 'excel',
            'sheet_count': len(sheet_names),
            'sheet_names': sheet_names,
            'has_data': len(sheet_names) > 0,
        }
    
    def _process_csv(self, file_path: Path) -> Dict[str, Any]:
        """Process CSV file."""
        import pandas as pd
        
        # Read first few rows
        df = pd.read_csv(file_path, nrows=5)
        
        return {
            'content_type': 'csv',
            'row_count': len(df),
            'column_count': len(df.columns),
            'column_names': list(df.columns),
            'has_data': len(df) > 0,
        }
    
    def _process_text(self, file_path: Path) -> Dict[str, Any]:
        """Process plain text file."""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read(1000)  # Read first 1000 chars
            
            lines = content.split('\n')
            
            return {
                'content_type': 'text',
                'total_length': len(content),
                'line_count': len(lines),
                'content_preview': content[:500],
                'word_count': len(content.split()),
            }
        except Exception as e:
            return {
                'content_type': 'text',
                'error': f'Could not read text file: {e}'
            }
    
    def _process_generic(self, file_path: Path) -> Dict[str, Any]:
        """Generic processing for unsupported file types."""
        return {
            'content_type': 'generic',
            'can_extract_text': False,
            'note': 'File type not fully supported for content extraction',
        }
    
    # ==================== DOCX Editing and Writing Methods ====================
    
    def create_docx(self, content: str = "", title: str = "Untitled Document") -> Dict[str, Any]:
        """
        Create a new Word document with optional content.
        
        Args:
            content: Initial content to add to the document
            title: Document title
            
        Returns:
            Dictionary with document information
        """
        if not self.dependencies['docx']:
            return {
                'success': False,
                'error': 'python-docx library not available',
                'message': 'Cannot create DOCX files without python-docx library'
            }
        
        try:
            import docx
            from docx.shared import Pt, Inches
            from docx.enum.text import WD_ALIGN_PARAGRAPH
            
            # Create a new document
            doc = docx.Document()
            
            # Add title
            title_para = doc.add_heading(title, 0)
            title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            
            # Add content if provided
            if content:
                doc.add_paragraph(content)
            
            # Save to a temporary file
            temp_file = self.workspace_dir / f"{title.replace(' ', '_')}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"
            doc.save(temp_file)
            
            # Get file info
            file_info = self._get_file_info(temp_file, f"{title}.docx")
            file_info.update({
                'content_type': 'docx',
                'paragraph_count': len(doc.paragraphs),
                'table_count': len(doc.tables),
                'created': True,
                'temp_path': str(temp_file),
                'title': title
            })
            
            return {
                'success': True,
                'document_info': file_info,
                'message': f'Successfully created Word document: {title}.docx',
                'file_path': str(temp_file)
            }
            
        except Exception as e:
            return {
                'success': False,
                'error': str(e),
                'message': f'Failed to create Word document: {e}'
            }
    
    def _add_footer_with_page_numbers(self, doc, footer_type: str, text: str, font_name: str, font_size: float, alignment: str):
        """
        Add a footer with optional page numbers to a document section.

        footer_type options:
            - 'page_number'        → "Page X of Y" centered
            - 'page_x'             → "Page X" centered
            - 'custom'             → just the 'text' string, no fields
        """
        from docx.oxml.ns import qn
        from docx.oxml import OxmlElement

        section = doc.sections[-1]
        footer = section.footer

        # Clear existing footer paragraphs
        for para in footer.paragraphs:
            for run in para.runs:
                run.text = ''
        # Ensure at least one paragraph
        if not footer.paragraphs:
            footer.add_paragraph()

        para = footer.paragraphs[0]
        para.clear()

        # Map alignment string → WD_ALIGN_PARAGRAPH constant
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        alignment_map = {
            'left': WD_ALIGN_PARAGRAPH.LEFT,
            'right': WD_ALIGN_PARAGRAPH.RIGHT,
            'center': WD_ALIGN_PARAGRAPH.CENTER,
            'centre': WD_ALIGN_PARAGRAPH.CENTER,
            'justify': WD_ALIGN_PARAGRAPH.JUSTIFY,
            'distribute': WD_ALIGN_PARAGRAPH.DISTRIBUTE,
        }
        para.alignment = alignment_map.get(alignment, WD_ALIGN_PARAGRAPH.CENTER)

        # Helper: apply font name + size to a Run (not a Paragraph)
        def apply_run_font(run, name, size):
            from docx.shared import Pt
            if name:
                run.font.name = name
                if _is_cjk_font(name):
                    _set_east_asia_font(run._r, name)
            if size:
                run.font.size = Pt(size)  # Pt() creates a Length object in points

        # Helper: build run properties element for field runs
        def make_rPr():
            rPr = OxmlElement('w:rPr')
            if font_name:
                rFonts = OxmlElement('w:rFonts')
                rFonts.set(qn('w:ascii'), font_name)
                rFonts.set(qn('w:hAnsi'), font_name)
                if _is_cjk_font(font_name):
                    rFonts.set(qn('w:eastAsia'), font_name)
                rPr.append(rFonts)
            if font_size:
                sz = OxmlElement('w:sz')
                sz.set(qn('w:val'), str(int(font_size * 2)))
                rPr.append(sz)
                szCs = OxmlElement('w:szCs')
                szCs.set(qn('w:val'), str(int(font_size * 2)))
                rPr.append(szCs)
            return rPr

        # Helper: append a complex field (fldChar/instrText) to the paragraph.
        # Complex fields are rendered reliably across Word, LibreOffice, and web viewers.
        def append_field_run(field_code: str):
            # BEGIN run
            r_begin = OxmlElement('w:r')
            r_begin.append(make_rPr())
            fldChar_begin = OxmlElement('w:fldChar')
            fldChar_begin.set(qn('w:fldCharType'), 'begin')
            r_begin.append(fldChar_begin)
            para._p.append(r_begin)

            # INSTRUCTION run
            r_instr = OxmlElement('w:r')
            r_instr.append(make_rPr())
            instrText = OxmlElement('w:instrText')
            instrText.set(qn('xml:space'), 'preserve')
            instrText.text = f' {field_code} '
            r_instr.append(instrText)
            para._p.append(r_instr)

            # SEPARATE run
            r_sep = OxmlElement('w:r')
            r_sep.append(make_rPr())
            fldChar_sep = OxmlElement('w:fldChar')
            fldChar_sep.set(qn('w:fldCharType'), 'separate')
            r_sep.append(fldChar_sep)
            para._p.append(r_sep)

            # DISPLAY VALUE run (placeholder — recalculated when document opens)
            r_val = OxmlElement('w:r')
            r_val.append(make_rPr())
            t = OxmlElement('w:t')
            t.set(qn('xml:space'), 'preserve')
            t.text = '1'
            r_val.append(t)
            para._p.append(r_val)

            # END run
            r_end = OxmlElement('w:r')
            r_end.append(make_rPr())
            fldChar_end = OxmlElement('w:fldChar')
            fldChar_end.set(qn('w:fldCharType'), 'end')
            r_end.append(fldChar_end)
            para._p.append(r_end)

        if footer_type == 'custom':
            run = para.add_run(text)
            apply_run_font(run, font_name, font_size)
            return

        if footer_type == 'page_number':
            # "Page X of Y"  →  run("Page ") + field(PAGE) + run(" of ") + field(NUMPAGES)
            run0 = para.add_run('Page ')
            apply_run_font(run0, font_name, font_size)
            append_field_run('PAGE')
            run2 = para.add_run(' of ')
            apply_run_font(run2, font_name, font_size)
            append_field_run('NUMPAGES')

        elif footer_type == 'page_x':
            run0 = para.add_run('Page ')
            apply_run_font(run0, font_name, font_size)
            append_field_run('PAGE')

    def _add_header_with_page_numbers(self, doc, header_type: str, text: str, font_name: str, font_size: float, alignment: str):
        """
        Add a header with optional content to a document section.

        header_type options:
            - 'page_number'   → "Page X of Y" centered
            - 'page_x'        → "Page X" centered
            - 'custom'        → just the 'text' string, no fields
            - 'title'         → document title (docx title property or 'text' arg)
            - 'filename'      → use 'text' as filename content
        """
        from docx.oxml.ns import qn
        from docx.oxml import OxmlElement

        section = doc.sections[-1]
        header = section.header

        # Clear existing header paragraphs
        for para in header.paragraphs:
            for run in para.runs:
                run.text = ''
        if not header.paragraphs:
            header.add_paragraph()

        para = header.paragraphs[0]
        para.clear()

        from docx.enum.text import WD_ALIGN_PARAGRAPH
        alignment_map = {
            'left': WD_ALIGN_PARAGRAPH.LEFT,
            'right': WD_ALIGN_PARAGRAPH.RIGHT,
            'center': WD_ALIGN_PARAGRAPH.CENTER,
            'centre': WD_ALIGN_PARAGRAPH.CENTER,
            'justify': WD_ALIGN_PARAGRAPH.JUSTIFY,
            'distribute': WD_ALIGN_PARAGRAPH.DISTRIBUTE,
        }
        para.alignment = alignment_map.get(alignment, WD_ALIGN_PARAGRAPH.CENTER)

        # Helper: apply font name + size to a Run
        def apply_run_font(run, name, size):
            from docx.shared import Pt
            if name:
                run.font.name = name
                if _is_cjk_font(name):
                    _set_east_asia_font(run._r, name)
            if size:
                run.font.size = Pt(size)  # Pt() creates a Length object in points

        # Helper: build run properties element for field runs
        def make_rPr():
            rPr = OxmlElement('w:rPr')
            if font_name:
                rFonts = OxmlElement('w:rFonts')
                rFonts.set(qn('w:ascii'), font_name)
                rFonts.set(qn('w:hAnsi'), font_name)
                if _is_cjk_font(font_name):
                    rFonts.set(qn('w:eastAsia'), font_name)
                rPr.append(rFonts)
            if font_size:
                sz = OxmlElement('w:sz')
                sz.set(qn('w:val'), str(int(font_size * 2)))
                rPr.append(sz)
                szCs = OxmlElement('w:szCs')
                szCs.set(qn('w:val'), str(int(font_size * 2)))
                rPr.append(szCs)
            return rPr

        # Helper: append a complex field (fldChar/instrText) to the paragraph.
        # Complex fields are rendered reliably across Word, LibreOffice, and web viewers.
        def append_field_run(field_code: str):
            # BEGIN run
            r_begin = OxmlElement('w:r')
            r_begin.append(make_rPr())
            fldChar_begin = OxmlElement('w:fldChar')
            fldChar_begin.set(qn('w:fldCharType'), 'begin')
            r_begin.append(fldChar_begin)
            para._p.append(r_begin)

            # INSTRUCTION run
            r_instr = OxmlElement('w:r')
            r_instr.append(make_rPr())
            instrText = OxmlElement('w:instrText')
            instrText.set(qn('xml:space'), 'preserve')
            instrText.text = f' {field_code} '
            r_instr.append(instrText)
            para._p.append(r_instr)

            # SEPARATE run
            r_sep = OxmlElement('w:r')
            r_sep.append(make_rPr())
            fldChar_sep = OxmlElement('w:fldChar')
            fldChar_sep.set(qn('w:fldCharType'), 'separate')
            r_sep.append(fldChar_sep)
            para._p.append(r_sep)

            # DISPLAY VALUE run (placeholder — recalculated when document opens)
            r_val = OxmlElement('w:r')
            r_val.append(make_rPr())
            t = OxmlElement('w:t')
            t.set(qn('xml:space'), 'preserve')
            t.text = '1'
            r_val.append(t)
            para._p.append(r_val)

            # END run
            r_end = OxmlElement('w:r')
            r_end.append(make_rPr())
            fldChar_end = OxmlElement('w:fldChar')
            fldChar_end.set(qn('w:fldCharType'), 'end')
            r_end.append(fldChar_end)
            para._p.append(r_end)

        if header_type == 'custom':
            run = para.add_run(text)
            apply_run_font(run, font_name, font_size)
            return

        if header_type == 'page_number':
            run0 = para.add_run('Page ')
            apply_run_font(run0, font_name, font_size)
            append_field_run('PAGE')
            run2 = para.add_run(' of ')
            apply_run_font(run2, font_name, font_size)
            append_field_run('NUMPAGES')

        elif header_type == 'page_x':
            run0 = para.add_run('Page ')
            apply_run_font(run0, font_name, font_size)
            append_field_run('PAGE')

        elif header_type == 'title':
            title_text = ''
            if hasattr(doc, 'core_properties') and doc.core_properties.title:
                title_text = doc.core_properties.title
            if not title_text:
                title_text = text or 'Document'
            run = para.add_run(title_text)
            apply_run_font(run, font_name, font_size)

        elif header_type == 'filename':
            run = para.add_run(text or 'Document')
            apply_run_font(run, font_name, font_size)

    def edit_docx(self, file_path: str, edits: List[Dict[str, Any]], overwrite_original: bool = True) -> Dict[str, Any]:
        """
        Edit an existing Word document.
        
        Args:
            file_path: Path to the DOCX file
            edits: List of edit operations. Each operation should be a dict with:
                  - 'type': 'add_paragraph', 'add_heading', 'add_table', 'replace_text', 'modify_style', 'insert_image'
                  - 'content': Content to add or text to replace
                  - 'position': Position index (for paragraphs)
                  - 'style': Style information
                  - 'old_text': Text to replace (for replace_text type)
                  - 'new_text': New text (for replace_text type)
                  - 'image_path': Path to image file (for insert_image type)
                  - 'width_inches': Width in inches (for insert_image type, optional)
            overwrite_original: If True, overwrite the original file. If False, create a new file.
        
        Returns:
            Dictionary with edit results
        """
        if not self.dependencies['docx']:
            return {
                'success': False,
                'error': 'python-docx library not available',
                'message': 'Cannot edit DOCX files without python-docx library'
            }
        
        try:
            import docx
            from docx.shared import Pt, RGBColor
            from docx.enum.text import WD_PARAGRAPH_ALIGNMENT, WD_BREAK
            from docx.oxml.ns import qn
            
            file_path = Path(file_path)
            if not file_path.exists():
                return {
                    'success': False,
                    'error': 'File not found',
                    'message': f'File not found: {file_path}'
                }
            
            doc = docx.Document(file_path)
            changes_made = []
            
            def get_alignment_value(alignment):
                if alignment == 'left':
                    return WD_PARAGRAPH_ALIGNMENT.LEFT
                if alignment == 'center':
                    return WD_PARAGRAPH_ALIGNMENT.CENTER
                if alignment == 'right':
                    return WD_PARAGRAPH_ALIGNMENT.RIGHT
                if alignment == 'justify':
                    return WD_PARAGRAPH_ALIGNMENT.JUSTIFY
                return None
            
            def normalize_color(color_value):
                if color_value is None:
                    return None
                color_text = str(color_value).strip().lstrip('#')
                if len(color_text) == 6:
                    return color_text.upper()
                return None
            
            def set_run_font_name(run, font_name):
                _apply_font_name_to_run(run, font_name)
            
            def apply_run_formatting(run, formatting):
                if not formatting:
                    return
                font_name = formatting.get('font_name') or formatting.get('font')
                font_size = formatting.get('font_size')
                bold = formatting.get('bold')
                italic = formatting.get('italic')
                underline = formatting.get('underline')
                color = normalize_color(formatting.get('color'))
                
                if font_name:
                    set_run_font_name(run, font_name)
                    # If text contains Chinese but font is not CJK, set w:eastAsia separately
                    # to ensure CJK characters render correctly
                    run_text = run.text if run.text else ''
                    if _contains_chinese(run_text) and not _is_cjk_font(font_name):
                        _set_east_asia_font(run._element, '宋体')
                if font_size is not None:
                    run.font.size = Pt(font_size)
                if bold is not None:
                    run.font.bold = bold
                if italic is not None:
                    run.font.italic = italic
                if underline is not None:
                    run.font.underline = underline
                if color:
                    run.font.color.rgb = RGBColor.from_string(color)
            
            def apply_run_formatting_obj(run_elem, formatting):
                """Apply formatting to an XML run element (w:r).
                
                Args:
                    run_elem: The XML element representing a run
                    formatting: Dict with font_name, font_size, bold, italic, underline, 
                               color, strikethrough, subscript, superscript, highlight
                """
                if not formatting:
                    return
                
                from docx.oxml import OxmlElement
                from docx.oxml.ns import qn
                from docx.shared import Pt, RGBColor
                
                # Get or create rPr element
                rPr = run_elem.find(qn('w:rPr'))
                if rPr is None:
                    rPr = OxmlElement('w:rPr')
                    run_elem.insert(0, rPr)
                
                font_name = formatting.get('font_name') or formatting.get('font')
                font_size = formatting.get('font_size')
                bold = formatting.get('bold')
                italic = formatting.get('italic')
                underline = formatting.get('underline')
                strikethrough = formatting.get('strikethrough')
                subscript = formatting.get('subscript')
                superscript = formatting.get('superscript')
                color = normalize_color(formatting.get('color'))
                highlight = formatting.get('highlight')
                style_name = formatting.get('style_name')
                
                # Set font name (w:rFonts)
                if font_name:
                    rFonts = rPr.find(qn('w:rFonts'))
                    if rFonts is None:
                        rFonts = OxmlElement('w:rFonts')
                        rPr.insert(0, rFonts)
                    rFonts.set(qn('w:ascii'), font_name)
                    rFonts.set(qn('w:hAnsi'), font_name)
                    rFonts.set(qn('w:eastAsia'), font_name)
                    rFonts.set(qn('w:hint'), 'eastAsia')
                
                # Set font size (w:sz)
                if font_size is not None:
                    sz = rPr.find(qn('w:sz'))
                    if sz is None:
                        sz = OxmlElement('w:sz')
                        rPr.append(sz)
                    # Font size in OOXML is in half-points
                    sz.set(qn('w:val'), str(int(font_size * 2)))
                    szCs = rPr.find(qn('w:szCs'))
                    if szCs is None:
                        szCs = OxmlElement('w:szCs')
                        rPr.append(szCs)
                    szCs.set(qn('w:val'), str(int(font_size * 2)))
                
                # Set bold (w:b)
                if bold is not None:
                    b = rPr.find(qn('w:b'))
                    if bold:
                        if b is None:
                            b = OxmlElement('w:b')
                            rPr.append(b)
                    else:
                        if b is not None:
                            rPr.remove(b)
                
                # Set italic (w:i)
                if italic is not None:
                    i = rPr.find(qn('w:i'))
                    if italic:
                        if i is None:
                            i = OxmlElement('w:i')
                            rPr.append(i)
                    else:
                        if i is not None:
                            rPr.remove(i)
                
                # Set underline (w:u)
                if underline is not None:
                    u = rPr.find(qn('w:u'))
                    if underline:
                        if u is None:
                            u = OxmlElement('w:u')
                            rPr.append(u)
                        u.set(qn('w:val'), 'single')
                    else:
                        if u is not None:
                            rPr.remove(u)
                
                # Set strikethrough (w:strike)
                if strikethrough is not None:
                    strike = rPr.find(qn('w:strike'))
                    if strikethrough:
                        if strike is None:
                            strike = OxmlElement('w:strike')
                            rPr.append(strike)
                    else:
                        if strike is not None:
                            rPr.remove(strike)
                
                # Set subscript (w:vertAlign with val="subscript")
                if subscript is not None:
                    vertAlign = rPr.find(qn('w:vertAlign'))
                    if subscript:
                        if vertAlign is None:
                            vertAlign = OxmlElement('w:vertAlign')
                            rPr.append(vertAlign)
                        vertAlign.set(qn('w:val'), 'subscript')
                    else:
                        if vertAlign is not None and vertAlign.get(qn('w:val')) == 'subscript':
                            rPr.remove(vertAlign)
                
                # Set superscript (w:vertAlign with val="superscript")
                if superscript is not None:
                    vertAlign = rPr.find(qn('w:vertAlign'))
                    if superscript:
                        if vertAlign is None:
                            vertAlign = OxmlElement('w:vertAlign')
                            rPr.append(vertAlign)
                        vertAlign.set(qn('w:val'), 'superscript')
                    else:
                        if vertAlign is not None and vertAlign.get(qn('w:val')) == 'superscript':
                            rPr.remove(vertAlign)
                
                # Set color (w:color)
                if color:
                    c = rPr.find(qn('w:color'))
                    if c is None:
                        c = OxmlElement('w:color')
                        rPr.append(c)
                    c.set(qn('w:val'), color)
                
                # Set highlight (w:highlight)
                if highlight is not None:
                    hl = rPr.find(qn('w:highlight'))
                    if highlight:
                        if hl is None:
                            hl = OxmlElement('w:highlight')
                            rPr.append(hl)
                        hl.set(qn('w:val'), str(highlight))
                    else:
                        if hl is not None:
                            rPr.remove(hl)
            
            def apply_paragraph_formatting(paragraph, formatting):
                if not formatting:
                    return
                alignment_value = get_alignment_value(formatting.get('alignment'))
                if alignment_value is not None:
                    paragraph.alignment = alignment_value
                spacing_before = formatting.get('spacing_before')
                spacing_after = formatting.get('spacing_after')
                line_spacing = formatting.get('line_spacing')
                if spacing_before is not None:
                    paragraph.paragraph_format.space_before = Pt(spacing_before)
                if spacing_after is not None:
                    paragraph.paragraph_format.space_after = Pt(spacing_after)
                if line_spacing is not None:
                    paragraph.paragraph_format.line_spacing = line_spacing
            
            def apply_paragraph_and_run_formatting(paragraph, formatting):
                if not formatting:
                    return
                apply_paragraph_formatting(paragraph, formatting)
                for run in paragraph.runs:
                    apply_run_formatting(run, formatting)
            
            def ensure_paragraph_has_run(paragraph):
                if paragraph.runs:
                    return paragraph.runs[0]
                return paragraph.add_run()
            
            def insert_paragraph_at_position(content, position, style=None):
                # Support multiple field names and formats for position
                if position is None:
                    pos_int = None
                elif isinstance(position, int):
                    pos_int = position
                elif isinstance(position, str):
                    import re
                    numbers = re.findall(r'\d+', str(position))
                    if numbers:
                        try:
                            pos_int = int(numbers[0])
                        except (ValueError, TypeError):
                            pos_int = None
                    elif position.lower() == 'end':
                        pos_int = len(doc.paragraphs) - 1 if doc.paragraphs else 0
                    elif position.lower() == 'start':
                        pos_int = 0
                    else:
                        pos_int = None
                else:
                    pos_int = None
                
                if pos_int is not None and 0 <= pos_int < len(doc.paragraphs):
                    paragraph = doc.paragraphs[pos_int].insert_paragraph_before(content)
                else:
                    paragraph = doc.add_paragraph(content)
                if style:
                    paragraph.style = style
                return paragraph
            
            def move_table_after(table, anchor_paragraph):
                tbl = table._tbl
                anchor = anchor_paragraph._p
                anchor.addnext(tbl)
            
            def iter_table_paragraphs():
                for table in doc.tables:
                    for row in table.rows:
                        for cell in row.cells:
                            for paragraph in cell.paragraphs:
                                yield paragraph
            
            def _delete_text_from_paragraph(paragraph, target_text):
                """Delete target_text from a paragraph, handling cross-run spans.
                
                Uses the same run-aware char-map approach as
                apply_text_replacement_to_paragraph to correctly handle text
                that is split across multiple runs.
                
                Args:
                    paragraph: The paragraph to modify
                    target_text: The text to delete
                
                Returns:
                    Number of deletions made (0 or 1)
                """
                if not target_text or target_text not in paragraph.text:
                    return 0
                
                # Delegate to the replacement function with empty replacement.
                # This preserves formatting on untouched runs correctly.
                count = apply_text_replacement_to_paragraph(
                    paragraph, target_text, '', case_sensitive=True
                )
                
                # Clean up runs that became empty
                _cleanup_empty_runs(paragraph)
                
                return 1 if count > 0 else 0
            
            def _cleanup_empty_runs(paragraph):
                """Remove runs that have empty text content."""
                from docx.oxml.ns import qn
                
                runs_to_remove = []
                for run in paragraph.runs:
                    if not run.text or run.text == '':
                        runs_to_remove.append(run)
                
                for run in runs_to_remove:
                    run_elem = run._element
                    p_elem = run_elem.getparent()
                    if p_elem is not None:
                        p_elem.remove(run_elem)
            
            def get_run_font_info(run):
                """Extract font formatting from a run.
                
                Captures: font_name, font_size, bold, italic, underline, color,
                highlight, strikethrough, subscript, superscript, style_name
                """
                font_info = {}
                
                # Basic font properties
                try:
                    if run.font.name:
                        font_info['font_name'] = run.font.name
                except:
                    pass
                
                try:
                    if run.font.size:
                        font_info['font_size'] = run.font.size.pt
                except:
                    pass
                
                # Style properties
                try:
                    if run.font.bold is not None:
                        font_info['bold'] = run.font.bold
                except:
                    pass
                
                try:
                    if run.font.italic is not None:
                        font_info['italic'] = run.font.italic
                except:
                    pass
                
                try:
                    if run.font.underline is not None:
                        font_info['underline'] = run.font.underline
                except:
                    pass
                
                try:
                    if run.font.strike is not None:
                        font_info['strikethrough'] = run.font.strike
                except:
                    pass
                
                try:
                    if run.font.subscript:
                        font_info['subscript'] = True
                except:
                    pass
                
                try:
                    if run.font.superscript:
                        font_info['superscript'] = True
                except:
                    pass
                
                # Color - handle both RGB and theme colors
                try:
                    if run.font.color.rgb is not None:
                        font_info['color'] = str(run.font.color.rgb)
                    elif run.font.color.theme_color is not None:
                        # Store theme color info if available
                        font_info['color_theme'] = run.font.color.theme_color
                        font_info['color_theme_tint'] = run.font.color.theme_color_tint
                except:
                    pass
                
                # Highlight
                try:
                    if run.font.highlight_color is not None:
                        font_info['highlight'] = run.font.highlight_color
                except:
                    pass
                
                # Character style
                try:
                    if run.style and run.style.name:
                        font_info['style_name'] = run.style.name
                except:
                    pass
                
                return font_info

            def detect_context_formatting(position, doc):
                """Detect formatting from surrounding context when inserting a new paragraph.
                
                Args:
                    position: Integer index for insertion point, or 'end', 'start'
                    doc: The document object
                
                Returns:
                    Dictionary with formatting info to apply to new content
                """
                context_format = {}
                
                # Get the reference paragraph for formatting
                ref_para = None
                has_paragraphs = len(doc.paragraphs) > 0
                
                if isinstance(position, int):
                    # Try to get the paragraph at the position for reference
                    if 0 <= position < len(doc.paragraphs):
                        ref_para = doc.paragraphs[position]
                    elif position >= 0 and has_paragraphs:
                        # position is non-negative but beyond doc length - use last paragraph
                        ref_para = doc.paragraphs[-1]
                    elif position < 0 and has_paragraphs:
                        # Negative index - use first paragraph (or last if out of range)
                        idx = len(doc.paragraphs) + position
                        ref_para = doc.paragraphs[max(0, idx)]
                    # else: position is int but doc is empty -> ref_para stays None
                elif position == 'end':
                    # Use the last paragraph as reference
                    ref_para = doc.paragraphs[-1] if has_paragraphs else None
                elif position == 'start':
                    # Use the first paragraph as reference
                    ref_para = doc.paragraphs[0] if has_paragraphs else None
                else:
                    # Unknown string value - default to 'end' behavior for safety
                    ref_para = doc.paragraphs[-1] if has_paragraphs else None
                
                if ref_para:
                    # Detect paragraph-level formatting
                    try:
                        # Use 'is not None' to properly detect LEFT (0) alignment
                        if ref_para.alignment is not None:
                            alignment_map = {
                                0: 'left',  # LEFT
                                1: 'center',  # CENTER
                                2: 'right',  # RIGHT
                                3: 'justify',  # JUSTIFY
                                4: 'distribute',  # DISTRIBUTE
                                5: 'justified_high',  # JUSTIFIED_H
                                6: 'justified_low',  # JUSTIFIED_L
                            }
                            al = ref_para.alignment
                            if al in alignment_map:
                                context_format['alignment'] = alignment_map[al]
                    except:
                        pass
                    
                    try:
                        # Check if space_before is explicitly set (not None or 0pt)
                        sb = ref_para.paragraph_format.space_before
                        if sb is not None and sb.pt != 0:
                            context_format['spacing_before'] = sb.pt
                    except:
                        pass
                    
                    try:
                        # Check if space_after is explicitly set (not None or 0pt)
                        sa = ref_para.paragraph_format.space_after
                        if sa is not None and sa.pt != 0:
                            context_format['spacing_after'] = sa.pt
                    except:
                        pass
                    
                    try:
                        # Check if line_spacing is explicitly set
                        ls = ref_para.paragraph_format.line_spacing
                        if ls is not None:
                            if isinstance(ls, float) and ls > 0:
                                context_format['line_spacing'] = ls
                            elif hasattr(ls, 'pt') and ls.pt != 0:
                                context_format['line_spacing'] = ls.pt
                    except:
                        pass
                    
                    # Detect run-level formatting from the first run with text
                    for run in ref_para.runs:
                        if run.text.strip():
                            run_info = get_run_font_info(run)
                            # Only inherit primary font properties
                            if run_info:
                                # Prefer font_name over font (for LLM compatibility)
                                if 'font_name' in run_info:
                                    context_format['font'] = run_info['font_name']
                                if 'font_size' in run_info:
                                    context_format['font_size'] = run_info['font_size']
                                if 'bold' in run_info:
                                    context_format['bold'] = run_info['bold']
                                if 'italic' in run_info:
                                    context_format['italic'] = run_info['italic']
                                if 'underline' in run_info:
                                    context_format['underline'] = run_info['underline']
                                if 'color' in run_info:
                                    context_format['color'] = run_info['color']
                                break
                
                return context_format

            def apply_text_replacement_to_paragraph(paragraph, old_text, new_text, case_sensitive: bool = False):
                """Replace text in a paragraph while preserving the original formatting.
                
                Uses a run-aware algorithm that only modifies the runs that contain
                the match, preserving every other run's formatting untouched.
                
                Returns the number of replacements made.
                """
                if not old_text:
                    return 0
                
                # Skip no-op replacements
                if old_text == new_text:
                    return 0
                
                import re
                import copy
                from docx.oxml import OxmlElement
                from docx.oxml.ns import qn
                
                p_elem = paragraph._element
                
                # Collect only w:r elements (runs) — skip pPr, bookmarks, comment markers, etc.
                run_elems = list(p_elem.findall(qn('w:r')))
                if not run_elems:
                    return 0
                
                # Build a character-map: for each char in the concatenated run text,
                # record which run element it belongs to and the offset inside that run's <w:t>.
                char_map = []  # list of (run_elem, char_index_within_run)
                run_texts = []
                for r_elem in run_elems:
                    t_elem = r_elem.find(qn('w:t'))
                    txt = t_elem.text if (t_elem is not None and t_elem.text) else ''
                    run_texts.append(txt)
                    for ci in range(len(txt)):
                        char_map.append((r_elem, ci))
                
                combined = ''.join(run_texts)
                
                # Find all match spans
                flags = 0 if case_sensitive else re.IGNORECASE
                pattern = re.compile(re.escape(old_text), flags)
                matches = list(pattern.finditer(combined))
                if not matches:
                    return 0
                
                # Process matches in reverse order so earlier indices stay valid
                for m in reversed(matches):
                    start, end = m.start(), m.end()
                    
                    # Identify which runs are touched by this match
                    first_run = char_map[start][0]
                    last_run = char_map[end - 1][0]
                    first_offset = char_map[start][1]
                    last_offset = char_map[end - 1][1]
                    
                    if first_run is last_run:
                        # Simple case: match is entirely within one run — just patch its <w:t>
                        t_elem = first_run.find(qn('w:t'))
                        txt = t_elem.text or ''
                        t_elem.text = txt[:first_offset] + new_text + txt[last_offset + 1:]
                        t_elem.set(qn('xml:space'), 'preserve')
                    else:
                        # Complex case: match spans multiple runs.
                        # Strategy:
                        #   - Put replacement text into the first run (preserving its rPr)
                        #   - Trim the last run's consumed prefix
                        #   - Remove any fully-consumed runs in between
                        
                        # 1) Patch first run: keep text before the match, append replacement
                        t_first = first_run.find(qn('w:t'))
                        txt_first = t_first.text or ''
                        t_first.text = txt_first[:first_offset] + new_text
                        t_first.set(qn('xml:space'), 'preserve')
                        
                        # 2) Patch last run: keep text after the match
                        t_last = last_run.find(qn('w:t'))
                        txt_last = t_last.text or ''
                        t_last.text = txt_last[last_offset + 1:]
                        t_last.set(qn('xml:space'), 'preserve')
                        # If last run is now empty, mark for removal
                        remove_last = (not t_last.text)
                        
                        # 3) Remove fully-consumed intermediate runs
                        in_span = False
                        to_remove = []
                        for r_elem in run_elems:
                            if r_elem is first_run:
                                in_span = True
                                continue
                            if r_elem is last_run:
                                if remove_last:
                                    to_remove.append(r_elem)
                                break
                            if in_span:
                                to_remove.append(r_elem)
                        
                        for r_elem in to_remove:
                            p_elem.remove(r_elem)
                    
                    # Rebuild char_map and run_elems for the next (earlier) match
                    run_elems = list(p_elem.findall(qn('w:r')))
                    char_map = []
                    run_texts = []
                    for r_elem in run_elems:
                        t_elem = r_elem.find(qn('w:t'))
                        txt = t_elem.text if (t_elem is not None and t_elem.text) else ''
                        run_texts.append(txt)
                        for ci in range(len(txt)):
                            char_map.append((r_elem, ci))
                    combined = ''.join(run_texts)
                
                return len(matches)

            def _set_cell_text(cell, value: str) -> None:
                """Overwrite a cell's text, preserving the first run's formatting.

                Splits `value` on '\\n' into N paragraph strings:
                  - Cell's existing first paragraph is reused (preserves pPr).
                  - That paragraph's first run is reused (preserves rPr — font,
                    size, bold/italic/color). Its text is set to parts[0].
                  - All other runs in that paragraph are blanked.
                  - All other paragraphs in the cell are removed.
                  - Extra parts become new paragraphs appended via cell.add_paragraph
                    (those inherit cell default styling).
                """
                from docx.oxml.ns import qn
                value = "" if value is None else str(value)
                parts = value.split("\n")

                paragraphs = list(cell.paragraphs)
                if not paragraphs:
                    cell.add_paragraph(parts[0])
                else:
                    first = paragraphs[0]
                    runs = first.runs
                    if runs:
                        runs[0].text = parts[0]
                        for extra in runs[1:]:
                            extra.text = ""
                        # Drop XML elements of blanked runs to keep doc tidy
                        for r in list(first._element.findall(qn('w:r'))):
                            t_elem = r.find(qn('w:t'))
                            if t_elem is not None and not (t_elem.text or "") and r is not first.runs[0]._element:
                                first._element.remove(r)
                    else:
                        first.add_run(parts[0])
                    # Remove any other paragraphs in the cell
                    for extra_p in paragraphs[1:]:
                        p_elem = extra_p._element
                        parent = p_elem.getparent()
                        if parent is not None:
                            parent.remove(p_elem)
                # Append additional paragraphs for multi-line values
                for part in parts[1:]:
                    cell.add_paragraph(part)

            def _replace_in_cell(cell, old_text: str, new_text: str):
                """Run-aware replace scoped to one table cell.

                Tries each paragraph in isolation via apply_text_replacement_to_paragraph
                (handles the common single-paragraph, multi-run case with full
                formatting preservation). If no paragraph matches, falls back to
                a cross-paragraph join over the cell's paragraphs.

                Returns (replaced_count, joined_cell_text_for_hint).
                """
                if not old_text:
                    return 0, ""
                if old_text == new_text:
                    return 0, ""

                paragraphs = list(cell.paragraphs)
                # Per-paragraph pass (case-sensitive — table cells often hold
                # exact strings the agent copied from extract_docx_content_tool)
                total = 0
                for p in paragraphs:
                    total += apply_text_replacement_to_paragraph(
                        p, old_text, new_text, case_sensitive=True
                    )
                if total > 0:
                    return total, ""

                # Cross-paragraph fallback — match spans more than one paragraph
                # inside this cell (e.g. "（大写）...（小写）..." split across two
                # paragraphs in a "total amount" cell).
                if len(paragraphs) < 2:
                    return 0, "\n".join(p.text for p in paragraphs)

                para_texts = [p.text for p in paragraphs]
                joined = "\n".join(para_texts)
                if old_text not in joined:
                    return 0, joined

                match_start = joined.index(old_text)
                match_end = match_start + len(old_text)

                offsets = [0]
                for t in para_texts:
                    offsets.append(offsets[-1] + len(t) + 1)  # +1 for the join '\n'

                def pos_to_para(pos):
                    for i in range(len(paragraphs)):
                        if offsets[i] <= pos < offsets[i + 1]:
                            return i, pos - offsets[i]
                    return len(paragraphs) - 1, len(para_texts[-1])

                p_start_idx, start_off = pos_to_para(match_start)
                # match_end is exclusive; treat it as the position of the next char
                p_end_idx, end_off = pos_to_para(max(match_end - 1, match_start))
                # convert end_off back to "exclusive" coord
                end_off = end_off + 1 if match_end > match_start else end_off

                if p_start_idx == p_end_idx:
                    # Single paragraph after all — re-run on just that paragraph
                    count = apply_text_replacement_to_paragraph(
                        paragraphs[p_start_idx], old_text, new_text, case_sensitive=True
                    )
                    return count, joined

                # True cross-paragraph match.
                # Split new_text on '\n' to map it across paragraphs. parts[0]
                # goes into the start paragraph's tail, parts[-1] into the end
                # paragraph's head, and parts[1:-1] become new paragraphs
                # inserted between them.
                new_parts = new_text.split("\n")
                start_p = paragraphs[p_start_idx]
                end_p = paragraphs[p_end_idx]
                start_tail = para_texts[p_start_idx][start_off:]
                end_head = para_texts[p_end_idx][:end_off]

                # Step 1: rewrite the start paragraph's tail.
                # If only one new part, the whole new_text replaces the tail
                # AND the end paragraph's matched head; nothing remains for the
                # end paragraph to inherit.
                if start_tail:
                    apply_text_replacement_to_paragraph(
                        start_p, start_tail, new_parts[0], case_sensitive=True
                    )
                elif new_parts[0]:
                    start_p.add_run(new_parts[0])

                # Step 2: drop fully-consumed intermediate paragraphs
                for i in range(p_start_idx + 1, p_end_idx):
                    p_elem = paragraphs[i]._element
                    parent = p_elem.getparent()
                    if parent is not None:
                        parent.remove(p_elem)

                # Step 3: handle the end paragraph.
                if len(new_parts) == 1:
                    # No paragraph break in the replacement. The unmatched
                    # suffix of the end paragraph (chars after end_off) needs
                    # to be merged into the start paragraph's run-stream so
                    # the end paragraph can be dropped entirely.
                    suffix = para_texts[p_end_idx][end_off:]
                    if suffix:
                        # Append a fresh run carrying the suffix text to the
                        # start paragraph — formatting of that suffix is lost
                        # by necessity (it lived in a separate paragraph), but
                        # the user's intent for a single-line replacement is
                        # that the cell ends up with one paragraph.
                        start_p.add_run(suffix)
                    end_elem = end_p._element
                    parent = end_elem.getparent()
                    if parent is not None:
                        parent.remove(end_elem)
                else:
                    # parts[-1] replaces the end paragraph's head; suffix stays.
                    if end_head:
                        apply_text_replacement_to_paragraph(
                            end_p, end_head, new_parts[-1], case_sensitive=True
                        )
                    elif new_parts[-1]:
                        # Match ends at the very start of end_p; prepend a run.
                        # python-docx doesn't expose "insert run at index 0", so
                        # we just add a run — it lands at the end, which is fine
                        # because end_head is empty (no leading text to push).
                        end_p.add_run(new_parts[-1])
                    # Step 3b: insert parts[1:-1] as fresh paragraphs between
                    # start_p and end_p (XML insertBefore).
                    middle_parts = new_parts[1:-1]
                    if middle_parts:
                        end_elem = end_p._element
                        parent = end_elem.getparent()
                        if parent is not None:
                            from copy import deepcopy
                            for part_text in middle_parts:
                                # Clone start_p's element shape so list/style
                                # context is inherited from a same-cell sibling.
                                new_p = deepcopy(start_p._element)
                                # Strip runs from the clone
                                for r in list(new_p.findall(qn('w:r'))):
                                    new_p.remove(r)
                                # Insert before end_p; then add a run with text
                                parent.insert(list(parent).index(end_elem), new_p)
                                # Wrap the inserted XML in a Paragraph object so
                                # we can use .add_run, then write the text.
                                from docx.text.paragraph import Paragraph
                                Paragraph(new_p, start_p._parent).add_run(part_text)
                return 1, joined

            def fill_table(table, data, rows, cols):
                for row_idx in range(rows):
                    row_data = data[row_idx] if row_idx < len(data) else []
                    if not isinstance(row_data, (list, tuple)):
                        row_data = [row_data]
                    for col_idx in range(cols):
                        cell_value = row_data[col_idx] if col_idx < len(row_data) else ''
                        table.cell(row_idx, col_idx).text = str(cell_value)
            
            # Debug: log all edits being processed
            print(f"🔍 Processing {len(edits)} edits:")
            for i, edit in enumerate(edits):
                edit_type = edit.get('type', 'add_paragraph')
                position = edit.get('position') or edit.get('after_paragraph') or edit.get('insert_after')
                if edit_type == 'add_table':
                    print(f"   Edit {i}: type={edit_type}, position={position}, data_rows={len(edit.get('data', []))}")
                else:
                    content_preview = str(edit.get('content', ''))[:30] if edit.get('content') else str(edit.get('old_text', ''))[:30]
                    print(f"   Edit {i}: type={edit_type}, position={position}, content='{content_preview}...'")
            
            for i, edit in enumerate(edits):
                edit_type = edit.get('type', 'add_paragraph')
                
                try:
                    if edit_type == 'add_paragraph':
                        content = edit.get('content', '')
                        # Support multiple field names for position
                        position = edit.get('position') or edit.get('after_paragraph') or edit.get('insert_after') or edit.get('after')
                        paragraph = insert_paragraph_at_position(content, position)
                        
                        # Detect and apply context formatting if no explicit formatting provided
                        context_format = detect_context_formatting(position, doc)
                        
                        # Merge context formatting with explicit formatting (explicit takes precedence)
                        merged_formatting = context_format.copy()
                        # Explicit edit formatting overrides context
                        for key in ['font_name', 'font', 'font_size', 'bold', 'italic', 
                                    'underline', 'color', 'alignment', 'spacing_before', 
                                    'spacing_after', 'line_spacing']:
                            if key in edit and edit[key] is not None:
                                merged_formatting[key] = edit[key]
                        
                        # Apply merged formatting
                        apply_paragraph_and_run_formatting(paragraph, merged_formatting)
                        changes_made.append(f"Added paragraph: {content[:50]}...")
                        
                    elif edit_type == 'add_heading':
                        content = edit.get('content', '')
                        level = edit.get('level', 1)
                        # Support multiple field names for position
                        position = edit.get('position') or edit.get('after_paragraph') or edit.get('insert_after') or edit.get('after')
                        if not isinstance(level, int):
                            level = 1
                        if level <= 0:
                            heading_style = 'Title'
                            level_label = 0
                        else:
                            if level > 9:
                                level = 9
                            heading_style = f'Heading {level}'
                            level_label = level
                        heading = insert_paragraph_at_position(content, position, heading_style)
                        
                        # Detect and apply context formatting if no explicit formatting provided
                        context_format = detect_context_formatting(position, doc)
                        
                        # Merge context formatting with explicit formatting (explicit takes precedence)
                        merged_formatting = context_format.copy()
                        for key in ['font_name', 'font', 'font_size', 'bold', 'italic', 
                                    'underline', 'color', 'alignment', 'spacing_before', 
                                    'spacing_after', 'line_spacing']:
                            if key in edit and edit[key] is not None:
                                merged_formatting[key] = edit[key]
                        
                        apply_paragraph_and_run_formatting(heading, merged_formatting)
                        changes_made.append(f"Added heading (level {level_label}): {content}")
                    
                    elif edit_type == 'replace':
                        # 'replace' is an alternative format for text replacement (used by some LLMs)
                        old_text = edit.get('old_text', '') or edit.get('target', '')
                        new_text = edit.get('new_text', '') or edit.get('replacement', '')
                        replaced_count = 0
                        for paragraph in doc.paragraphs:
                            replaced_count += apply_text_replacement_to_paragraph(paragraph, old_text, new_text)
                        for paragraph in iter_table_paragraphs():
                            replaced_count += apply_text_replacement_to_paragraph(paragraph, old_text, new_text)
                        if replaced_count > 0:
                            changes_made.append(f"Replaced '{old_text}' with '{new_text}' in {replaced_count} places")
                        else:
                            changes_made.append(f"Text '{old_text}' not found for replacement")
                    
                    elif edit_type == 'replace_text':
                        old_text = edit.get('old_text', '')
                        new_text = edit.get('new_text', '')
                        replaced_count = 0
                        not_found_info = None  # Track what we actually found
                        for paragraph in doc.paragraphs:
                            count = apply_text_replacement_to_paragraph(paragraph, old_text, new_text)
                            if count > 0:
                                replaced_count += count
                            elif not_found_info is None and old_text and paragraph.text.strip():
                                # Check if this paragraph contains similar text
                                if len(old_text) > 10 and (old_text[:20] in paragraph.text or paragraph.text[:20] in old_text):
                                    not_found_info = {
                                        'actual_text': paragraph.text,
                                        'reason': 'similar_text_found_but_not_exact_match'
                                    }
                        for paragraph in iter_table_paragraphs():
                            count = apply_text_replacement_to_paragraph(paragraph, old_text, new_text)
                            if count > 0:
                                replaced_count += count
                        if replaced_count > 0:
                            changes_made.append(f"Replaced '{old_text}' with '{new_text}' in {replaced_count} places")
                        else:
                            changes_made.append(f"Text '{old_text}' not found for replacement")
                            if not_found_info:
                                changes_made.append(f"HINT: Found similar text: '{not_found_info['actual_text'][:100]}...'")

                    elif edit_type == 'set_cell_text':
                        table_index = edit.get('table_index', 0)
                        row = edit.get('row')
                        col = edit.get('col')
                        value = edit.get('value', '')
                        if not isinstance(table_index, int) or not (0 <= table_index < len(doc.tables)):
                            raise IndexError(f"table_index {table_index} out of range (have {len(doc.tables)} table(s))")
                        table = doc.tables[table_index]
                        if not isinstance(row, int) or not (0 <= row < len(table.rows)):
                            raise IndexError(f"row {row} out of range for table {table_index} (have {len(table.rows)} row(s))")
                        row_cells = table.rows[row].cells
                        if not isinstance(col, int) or not (0 <= col < len(row_cells)):
                            raise IndexError(f"col {col} out of range for row {row} (have {len(row_cells)} cell(s))")
                        cell = row_cells[col]
                        _set_cell_text(cell, value)
                        preview = (str(value)[:40] + '…') if len(str(value)) > 40 else str(value)
                        changes_made.append(f"Set cell ({table_index},{row},{col}) to '{preview}'")

                    elif edit_type == 'replace_in_cell':
                        table_index = edit.get('table_index', 0)
                        row = edit.get('row')
                        col = edit.get('col')
                        old_text = edit.get('old_text', '')
                        new_text = edit.get('new_text', '')
                        if not isinstance(table_index, int) or not (0 <= table_index < len(doc.tables)):
                            raise IndexError(f"table_index {table_index} out of range (have {len(doc.tables)} table(s))")
                        table = doc.tables[table_index]
                        if not isinstance(row, int) or not (0 <= row < len(table.rows)):
                            raise IndexError(f"row {row} out of range for table {table_index} (have {len(table.rows)} row(s))")
                        row_cells = table.rows[row].cells
                        if not isinstance(col, int) or not (0 <= col < len(row_cells)):
                            raise IndexError(f"col {col} out of range for row {row} (have {len(row_cells)} cell(s))")
                        cell = row_cells[col]
                        count, joined_hint = _replace_in_cell(cell, old_text, new_text)
                        if count > 0:
                            old_preview = (old_text[:30] + '…') if len(old_text) > 30 else old_text
                            new_preview = (new_text[:30] + '…') if len(new_text) > 30 else new_text
                            changes_made.append(
                                f"Replaced '{old_preview}' with '{new_preview}' in cell ({table_index},{row},{col})"
                            )
                        else:
                            changes_made.append(
                                f"Text '{old_text[:30]}...' not found in cell ({table_index},{row},{col})"
                            )
                            if joined_hint:
                                changes_made.append(
                                    f"HINT: actual cell text is '{joined_hint[:120]}{'…' if len(joined_hint) > 120 else ''}'"
                                )

                    elif edit_type == 'delete_text':
                        # Properly delete text instead of replacing with empty string
                        target_text = edit.get('old_text', '')
                        if not target_text:
                            changes_made.append("Delete operation requires 'old_text' parameter")
                        else:
                            deleted_count = 0
                            for paragraph in list(doc.paragraphs):
                                deleted_count += _delete_text_from_paragraph(paragraph, target_text)
                                # Remove empty paragraphs after deletion
                                if not paragraph.text.strip():
                                    p_elem = paragraph._element
                                    p_elem.getparent().remove(p_elem)
                            for paragraph in list(iter_table_paragraphs()):
                                deleted_count += _delete_text_from_paragraph(paragraph, target_text)
                            if deleted_count > 0:
                                changes_made.append(f"Deleted '{target_text}' from {deleted_count} places")
                            else:
                                changes_made.append(f"Text '{target_text}' not found for deletion")
                    
                    elif edit_type == 'delete_paragraph':
                        # Delete a specific paragraph by position
                        position = edit.get('position')
                        count = edit.get('count', 1)  # Number of paragraphs to delete
                        deleted_count = 0
                        
                        if isinstance(position, int):
                            paragraphs_to_delete = []
                            for i in range(count):
                                idx = position + i
                                if 0 <= idx < len(doc.paragraphs):
                                    paragraphs_to_delete.append(doc.paragraphs[idx])
                            
                            for para in paragraphs_to_delete:
                                p_elem = para._element
                                p_elem.getparent().remove(p_elem)
                                deleted_count += 1
                        elif position == 'end':
                            # Delete last N paragraphs
                            while count > 0 and doc.paragraphs:
                                last_para = doc.paragraphs[-1]
                                p_elem = last_para._element
                                p_elem.getparent().remove(p_elem)
                                deleted_count += 1
                                count -= 1
                        elif position == 'start':
                            # Delete first N paragraphs
                            while count > 0 and doc.paragraphs:
                                first_para = doc.paragraphs[0]
                                p_elem = first_para._element
                                p_elem.getparent().remove(p_elem)
                                deleted_count += 1
                                count -= 1
                        
                        if deleted_count > 0:
                            changes_made.append(f"Deleted {deleted_count} paragraph(s)")
                        else:
                            changes_made.append(f"No paragraphs found at position {position}")
                            
                    elif edit_type == 'modify_style':
                        style_name = edit.get('style_name', 'Normal')
                        font_name = edit.get('font_name') or edit.get('font')
                        font_size = edit.get('font_size')
                        bold = edit.get('bold')
                        italic = edit.get('italic')
                        underline = edit.get('underline')
                        color = normalize_color(edit.get('color'))
                        alignment = edit.get('alignment')
                        spacing_before = edit.get('spacing_before')
                        spacing_after = edit.get('spacing_after')
                        line_spacing = edit.get('line_spacing')
                        apply_to_existing = edit.get('apply_to_existing', True)
                        
                        try:
                            style = doc.styles[style_name]
                        except KeyError:
                            style = doc.styles.add_style(style_name, 1)
                        
                        style_font = style.font
                        if font_name:
                            _apply_font_name_to_style(style, font_name)
                        if font_size is not None:
                            style_font.size = Pt(font_size)
                        if bold is not None:
                            style_font.bold = bold
                        if italic is not None:
                            style_font.italic = italic
                        if underline is not None:
                            style_font.underline = underline
                        if color:
                            style_font.color.rgb = RGBColor.from_string(color)
                        
                        alignment_value = get_alignment_value(alignment)
                        if alignment_value is not None:
                            style.paragraph_format.alignment = alignment_value
                        if spacing_before is not None:
                            style.paragraph_format.space_before = Pt(spacing_before)
                        if spacing_after is not None:
                            style.paragraph_format.space_after = Pt(spacing_after)
                        if line_spacing is not None:
                            style.paragraph_format.line_spacing = line_spacing
                        
                        affected_paragraphs = 0
                        affected_runs = 0
                        if apply_to_existing:
                            for paragraph in doc.paragraphs:
                                if paragraph.style and paragraph.style.name == style_name:
                                    affected_paragraphs += 1
                                    apply_paragraph_formatting(paragraph, edit)
                                    run_targets = paragraph.runs or [ensure_paragraph_has_run(paragraph)]
                                    for run in run_targets:
                                        apply_run_formatting(run, edit)
                                        affected_runs += 1
                        
                        changes_made.append(f"Modified style '{style_name}' and updated {affected_paragraphs} paragraphs / {affected_runs} runs")
                        
                    elif edit_type == 'add_table':
                        data = edit.get('data', [])
                        # Support multiple field names for position
                        position = edit.get('position') or edit.get('after_paragraph') or edit.get('insert_after') or edit.get('after')
                        # Also support string positions like "after_paragraph_16" or just "16"
                        if isinstance(position, str):
                            # Try to extract number from string like "after_paragraph_16" or "16"
                            import re
                            numbers = re.findall(r'\d+', str(position))
                            if numbers:
                                try:
                                    position = int(numbers[0])
                                except (ValueError, TypeError):
                                    position = None
                            elif position.lower() == 'end':
                                position = len(doc.paragraphs) - 1 if doc.paragraphs else 0
                            elif position.lower() == 'start':
                                position = 0
                            else:
                                position = None
                        table_style = edit.get('table_style') or edit.get('style')
                        if data:
                            rows = len(data)
                            cols = max(len(row) if isinstance(row, (list, tuple)) else 1 for row in data)
                        else:
                            rows = edit.get('rows', 1)
                            cols = edit.get('cols', 1)
                        table = doc.add_table(rows=rows, cols=cols)
                        fill_table(table, data, rows, cols)
                        if table_style:
                            try:
                                table.style = table_style
                            except Exception:
                                changes_made.append(f"Warning: Could not apply table style '{table_style}', using default")
                        if isinstance(position, int) and 0 <= position < len(doc.paragraphs):
                            move_table_after(table, doc.paragraphs[position])
                            changes_made.append(f"Added table with {rows}x{cols} dimensions after paragraph {position}")
                        else:
                            # Default: add at end of document
                            changes_made.append(f"Added table with {rows}x{cols} dimensions (position {position} invalid, appended to end)")
                            if doc.paragraphs:
                                move_table_after(table, doc.paragraphs[-1])
                    
                    elif edit_type == 'format_text':
                        target_text = edit.get('target_text')
                        formatted_count = 0
                        for paragraph in list(doc.paragraphs) + list(iter_table_paragraphs()):
                            if target_text and target_text in paragraph.text:
                                run_targets = paragraph.runs or [ensure_paragraph_has_run(paragraph)]
                                for run in run_targets:
                                    if target_text in run.text:
                                        apply_run_formatting(run, edit)
                                formatted_count += 1
                        changes_made.append(f"Formatted text in {formatted_count} paragraphs")
                    
                    elif edit_type == 'format_paragraph':
                        position = edit.get('position')
                        formatted_count = 0
                        if isinstance(position, int) and 0 <= position < len(doc.paragraphs):
                            apply_paragraph_formatting(doc.paragraphs[position], edit)
                            formatted_count = 1
                        else:
                            target_style = edit.get('style_name')
                            for paragraph in doc.paragraphs:
                                if not target_style or (paragraph.style and paragraph.style.name == target_style):
                                    apply_paragraph_formatting(paragraph, edit)
                                    formatted_count += 1
                        changes_made.append(f"Formatted {formatted_count} paragraphs")
                    
                    elif edit_type == 'add_page_break':
                        position = edit.get('position')
                        if isinstance(position, int) and 0 <= position < len(doc.paragraphs):
                            paragraph = doc.paragraphs[position].insert_paragraph_before('')
                        else:
                            paragraph = doc.add_paragraph()
                        paragraph.add_run().add_break(WD_BREAK.PAGE)
                        changes_made.append("Added page break")
                    
                    elif edit_type == 'set_table_style':
                        table_index = edit.get('table_index', 0)
                        table_style = edit.get('table_style') or edit.get('style')
                        if table_style is None:
                            raise ValueError('table_style is required for set_table_style')
                        if not isinstance(table_index, int) or table_index < 0 or table_index >= len(doc.tables):
                            raise IndexError('table_index out of range')
                        doc.tables[table_index].style = table_style
                        changes_made.append(f"Applied table style '{table_style}' to table {table_index}")

                    elif edit_type == 'insert_image':
                        image_path = edit.get('image_path')
                        if not image_path:
                            changes_made.append("Error in insert_image: image_path is required")
                        elif not Path(image_path).is_file():
                            changes_made.append(f"Error in insert_image: file not found: {image_path}")
                        else:
                            try:
                                from docx.shared import Inches
                                # Support both 'position' (agent convention) and 'insert_after_paragraph' (explicit)
                                raw_pos = edit.get('position', edit.get('insert_after_paragraph', -1))
                                try:
                                    insert_after = int(raw_pos)
                                except (TypeError, ValueError):
                                    insert_after = -1
                                width_inches = edit.get('width_inches')
                                total_paras = len(doc.paragraphs)

                                if insert_after == -1 or insert_after >= total_paras:
                                    target_para = doc.add_paragraph()
                                    log_pos = f"end (of {total_paras} paragraphs)"
                                else:
                                    target_para = doc.paragraphs[insert_after].insert_paragraph_before("")
                                    log_pos = f"after paragraph {insert_after} (of {total_paras})"

                                run = target_para.add_run()
                                if width_inches:
                                    run.add_picture(image_path, width=Inches(width_inches))
                                else:
                                    run.add_picture(image_path)

                                changes_made.append(f"Inserted image at {log_pos}: {Path(image_path).name}")
                            except Exception as img_err:
                                changes_made.append(f"Error in insert_image: {img_err}")

                    elif edit_type == 'add_bullet_list':
                        items = edit.get('items', [])
                        position = edit.get('position')
                        style_override = edit.get('style') or 'List Bullet'
                        start_index = 0

                        added_items = []
                        for idx, item in enumerate(items):
                            if isinstance(item, str):
                                item_text = item
                                item_level = 0
                            elif isinstance(item, dict):
                                item_text = item.get('text', '')
                                item_level = item.get('level', 0)
                            else:
                                continue
                            if not item_text:
                                continue

                            if isinstance(position, int) and 0 <= position < len(doc.paragraphs):
                                para = doc.paragraphs[position].insert_paragraph_before(item_text)
                            else:
                                para = doc.add_paragraph(item_text)

                            # Apply list style based on level
                            list_style = f'List Bullet {item_level + 1}' if item_level > 0 else style_override
                            try:
                                para.style = doc.styles[list_style]
                            except KeyError:
                                try:
                                    para.style = doc.styles['List Bullet']
                                except KeyError:
                                    para.style = doc.styles['Normal']

                            apply_paragraph_and_run_formatting(para, edit)
                            added_items.append(item_text[:50])
                            if isinstance(position, int) and 0 <= position < len(doc.paragraphs):
                                position += 1

                        changes_made.append(f"Added bullet list with {len(added_items)} items")

                    elif edit_type == 'add_numbered_list':
                        items = edit.get('items', [])
                        position = edit.get('position')
                        style_override = edit.get('style') or 'List Number'
                        start_index = 0

                        added_items = []
                        for idx, item in enumerate(items):
                            if isinstance(item, str):
                                item_text = item
                                item_level = 0
                            elif isinstance(item, dict):
                                item_text = item.get('text', '')
                                item_level = item.get('level', 0)
                            else:
                                continue
                            if not item_text:
                                continue

                            if isinstance(position, int) and 0 <= position < len(doc.paragraphs):
                                para = doc.paragraphs[position].insert_paragraph_before(item_text)
                            else:
                                para = doc.add_paragraph(item_text)

                            # Apply list style based on level
                            list_style = f'List Number {item_level + 1}' if item_level > 0 else style_override
                            try:
                                para.style = doc.styles[list_style]
                            except KeyError:
                                try:
                                    para.style = doc.styles['List Number']
                                except KeyError:
                                    para.style = doc.styles['Normal']

                            apply_paragraph_and_run_formatting(para, edit)
                            added_items.append(item_text[:50])
                            if isinstance(position, int) and 0 <= position < len(doc.paragraphs):
                                position += 1

                        changes_made.append(f"Added numbered list with {len(added_items)} items")

                    elif edit_type == 'add_footer':
                        footer_type = edit.get('footer_type', 'page_number')  # 'page_number' | 'page_x' | 'custom'
                        text = edit.get('text', '')
                        font_name = edit.get('font_name', 'Calibri')
                        font_size = edit.get('font_size', 9)
                        alignment = edit.get('alignment', 'center')

                        self._add_footer_with_page_numbers(
                            doc, footer_type, text, font_name, font_size, alignment
                        )
                        changes_made.append(f"Added footer (type={footer_type})")

                    elif edit_type == 'add_header':
                        header_type = edit.get('header_type', 'custom')  # 'page_number' | 'page_x' | 'custom' | 'title' | 'filename'
                        text = edit.get('text', '')
                        font_name = edit.get('font_name', 'Calibri')
                        font_size = edit.get('font_size', 9)
                        alignment = edit.get('alignment', 'center')

                        self._add_header_with_page_numbers(
                            doc, header_type, text, font_name, font_size, alignment
                        )
                        changes_made.append(f"Added header (type={header_type})")

                except Exception as e:
                    changes_made.append(f"Error applying edit {i+1} ({edit_type}): {str(e)}")
            
            # Save the file
            if overwrite_original:
                # Overwrite the original file
                edited_file = Path(file_path)
                doc.save(edited_file)
            else:
                # Save to a new file (don't overwrite original)
                timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
                edited_file = self.workspace_dir / f"{file_path.stem}_edited_{timestamp}.docx"
                doc.save(edited_file)
            
            # Get updated file info
            file_info = self._get_file_info(edited_file, f"{file_path.stem}_edited.docx")
            file_info.update({
                'content_type': 'docx',
                'paragraph_count': len(doc.paragraphs),
                'table_count': len(doc.tables),
                'changes_made': changes_made,
                'original_file': str(file_path),
                'edited_file': str(edited_file)
            })
            
            # Check if any actual changes were made (vs just "not found" messages)
            actual_changes = [c for c in changes_made if 'not found' not in c.lower() and 'error' not in c.lower() and 'hint' not in c.lower()]
            
            # If no actual changes were made, return failure with clear reason
            if len(actual_changes) == 0:
                # Extract hints for retry
                hints = [c for c in changes_made if 'hint' in c.lower()]
                hint_message = '\n'.join(hints) if hints else ''
                return {
                    'success': False,
                    'document_info': file_info,
                    'message': 'No changes were applied. Target text not found - likely due to whitespace/punctuation differences. ' + ('Try using the EXACT text from extract_docx_content tool.' if not hints else ''),
                    'file_path': str(edited_file),
                    'changes': [],
                    'retry_hints': hints if hints else None
                }
            
            return {
                'success': True,
                'document_info': file_info,
                'message': f'Successfully edited document. Changes: {len(actual_changes)}',
                'file_path': str(edited_file),
                'changes': changes_made
            }
            
        except Exception as e:
            return {
                'success': False,
                'error': str(e),
                'message': f'Failed to edit Word document: {e}'
            }
    
    # ==================== ENHANCED DOCX EDITING METHODS ====================
    
    def delete_docx_content(self, file_path: str) -> Dict[str, Any]:
        """
        Delete all content from a Word document while preserving its
        styles, themes, numbering definitions, settings, and other
        document-level metadata.
        
        Args:
            file_path: Path to the DOCX file
            
        Returns:
            Dictionary with deletion results
        """
        if not self.dependencies['docx']:
            return {
                'success': False,
                'error': 'python-docx library not available',
                'message': 'Cannot edit DOCX files without python-docx library'
            }
        
        try:
            import docx
            from docx.oxml.ns import qn
            
            file_path_obj = Path(file_path)
            if not file_path_obj.exists():
                return {
                    'success': False,
                    'error': 'File not found',
                    'message': f'File not found: {file_path}'
                }
            
            # Open the EXISTING document so we keep its styles, themes,
            # numbering definitions, settings, headers/footers, etc.
            doc = docx.Document(file_path)
            
            body = doc.element.body
            
            # Remove all content children (paragraphs, tables, sdt blocks)
            # but preserve <w:sectPr> which defines page layout.
            for child in list(body):
                tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
                if tag != 'sectPr':
                    body.remove(child)
            
            # Ensure at least one empty paragraph exists (Word requires it)
            from docx.oxml import OxmlElement
            new_p = OxmlElement('w:p')
            # Insert before sectPr if it exists, otherwise just append
            sect_pr = body.find(qn('w:sectPr'))
            if sect_pr is not None:
                sect_pr.addprevious(new_p)
            else:
                body.append(new_p)
            
            # Save back to the same file
            doc.save(file_path)
            
            # Verify
            doc_check = docx.Document(file_path)
            paragraph_count = len([p for p in doc_check.paragraphs if p.text.strip()])
            
            return {
                'success': True,
                'message': 'All content successfully deleted from document (styles and settings preserved)',
                'file_path': file_path,
                'is_empty': paragraph_count == 0,
                'remaining_paragraphs': paragraph_count
            }
            
        except Exception as e:
            return {
                'success': False,
                'error': str(e),
                'message': f'Failed to delete document content: {e}'
            }
    
    def modify_docx_fonts(self, file_path: str, font_rules: Dict[str, str] = None) -> Dict[str, Any]:
        """
        Modify fonts in a Word document based on content type.
        
        Args:
            file_path: Path to the DOCX file
            font_rules: Dictionary with font rules, e.g., 
                       {'chinese': '宋体', 'english': 'Times New Roman'}
            
        Returns:
            Dictionary with font modification results
        """
        if not self.dependencies['docx']:
            return {
                'success': False,
                'error': 'python-docx library not available',
                'message': 'Cannot edit DOCX files without python-docx library'
            }
        
        try:
            import docx
            import re
            from docx.oxml.ns import qn
            
            file_path_obj = Path(file_path)
            if not file_path_obj.exists():
                return {
                    'success': False,
                    'error': 'File not found',
                    'message': f'File not found: {file_path}'
                }
            
            # Default font rules
            if font_rules is None:
                font_rules = {
                    'chinese': '宋体',
                    'english': 'Times New Roman',
                    'default': '宋体'
                }
            
            # Helper functions to detect text type
            def contains_chinese(text):
                """Check if text contains Chinese characters"""
                return bool(re.search(r'[\u4e00-\u9fff]', text))
            
            def contains_english(text):
                """Check if text contains English letters"""
                return bool(re.search(r'[a-zA-Z]', text))
            
            # Load the document
            doc = docx.Document(file_path)
            
            # Track changes
            chinese_changes = 0
            english_changes = 0
            other_changes = 0
            
            # Process each paragraph
            for paragraph in doc.paragraphs:
                for run in paragraph.runs:
                    if not run.text.strip():
                        continue
                    
                    run_text = run.text
                    
                    # Determine font based on content
                    if contains_chinese(run_text):
                        # Apply Chinese font
                        chinese_font = font_rules.get('chinese', '宋体')
                        _apply_font_name_to_run(run, chinese_font)
                        # Always ensure w:eastAsia is set for the Chinese branch,
                        # even if the font name isn't in the CJK keyword list.
                        _set_east_asia_font(run._element, chinese_font)
                        chinese_changes += 1
                        
                    elif contains_english(run_text):
                        # Apply English font
                        english_font = font_rules.get('english', 'Times New Roman')
                        _apply_font_name_to_run(run, english_font)
                        english_changes += 1
                        
                    else:
                        # Apply default font (for numbers, punctuation, etc.)
                        default_font = font_rules.get('default', 'Times New Roman')
                        _apply_font_name_to_run(run, default_font)
                        other_changes += 1
            
            # Save the modified document
            doc.save(file_path)
            
            return {
                'success': True,
                'message': f'Fonts modified successfully: {chinese_changes} Chinese, {english_changes} English, {other_changes} other',
                'file_path': file_path,
                'changes': {
                    'chinese_text_modified': chinese_changes,
                    'english_text_modified': english_changes,
                    'other_text_modified': other_changes,
                    'total_runs_modified': chinese_changes + english_changes + other_changes
                },
                'font_rules_applied': font_rules
            }
            
        except Exception as e:
            return {
                'success': False,
                'error': str(e),
                'message': f'Failed to modify document fonts: {e}'
            }
    
    def create_docx_with_content(self, output_path: str, content: List[Dict[str, Any]]) -> Dict[str, Any]:
        """
        Create a new Word document with structured content.
        
        Args:
            output_path: Path where to save the new document
            content: List of content elements, each with:
                    - 'type': 'heading', 'paragraph', 'poem', 'table'
                    - 'text': The text content
                    - 'level': For headings (1-9)
                    - 'font': Font name (optional)
                    - 'alignment': 'left', 'center', 'right', 'justify' (optional)
            
        Returns:
            Dictionary with creation results
        """
        if not self.dependencies['docx']:
            return {
                'success': False,
                'error': 'python-docx library not available',
                'message': 'Cannot create DOCX files without python-docx library'
            }
        
        try:
            import docx
            from docx.shared import Pt, RGBColor
            from docx.enum.text import WD_ALIGN_PARAGRAPH, WD_BREAK
            from docx.oxml.ns import qn
            
            doc = docx.Document()
            created_elements = []
            
            def normalize_color(color_value):
                if color_value is None:
                    return None
                color_text = str(color_value).strip().lstrip('#')
                if len(color_text) == 6:
                    return color_text.upper()
                return None
            
            def apply_font_to_paragraph(paragraph, font_name=None, font_size=None, bold=None, italic=None, underline=None, color=None):
                paragraph_text = paragraph.text if hasattr(paragraph, 'text') else str(paragraph)
                has_cjk = _contains_chinese(paragraph_text)
                
                for run in paragraph.runs:
                    if font_name:
                        _apply_font_name_to_run(run, font_name)
                        # If text contains Chinese but font is not CJK, set w:eastAsia separately
                        # to ensure CJK characters render correctly
                        if has_cjk and not _is_cjk_font(font_name):
                            _set_east_asia_font(run._element, '宋体')
                    if font_size is not None:
                        run.font.size = Pt(font_size)
                    if bold is not None:
                        run.font.bold = bold
                    if italic is not None:
                        run.font.italic = italic
                    if underline is not None:
                        run.font.underline = underline
                    normalized_color = normalize_color(color)
                    if normalized_color:
                        run.font.color.rgb = RGBColor.from_string(normalized_color)
            
            def apply_alignment_to_paragraph(paragraph, alignment_value):
                if not alignment_value:
                    return
                if alignment_value == 'center':
                    paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
                elif alignment_value == 'right':
                    paragraph.alignment = WD_ALIGN_PARAGRAPH.RIGHT
                elif alignment_value == 'justify':
                    paragraph.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
                else:
                    paragraph.alignment = WD_ALIGN_PARAGRAPH.LEFT
            
            def apply_spacing_to_paragraph(paragraph, spacing_before=None, spacing_after=None, line_spacing=None):
                if spacing_before is not None:
                    paragraph.paragraph_format.space_before = Pt(spacing_before)
                if spacing_after is not None:
                    paragraph.paragraph_format.space_after = Pt(spacing_after)
                if line_spacing is not None:
                    paragraph.paragraph_format.line_spacing = line_spacing
            
            for i, element in enumerate(content):
                element_type = element.get('type', 'paragraph')
                text = element.get('text', '')
                font = element.get('font') or element.get('font_name')
                # No default font - user/agent must specify. If using Latin fonts with
                # Chinese text, w:eastAsia will need to be set separately for proper rendering.
                if not font and _contains_chinese(text):
                    font = '宋体'
                font_size = element.get('font_size')
                bold = element.get('bold')
                italic = element.get('italic')
                underline = element.get('underline')
                color = element.get('color')
                alignment = element.get('alignment', 'left')
                spacing_before = element.get('spacing_before')
                spacing_after = element.get('spacing_after')
                line_spacing = element.get('line_spacing')
                current_paragraph = None
                
                if element_type == 'heading':
                    level = element.get('level', 1)
                    if not isinstance(level, int):
                        level = 1
                    if level <= 0:
                        current_paragraph = doc.add_heading(text, 0)
                        created_elements.append(f"Title: {text[:50]}...")
                    else:
                        if level > 9:
                            level = 9
                        current_paragraph = doc.add_heading(text, level)
                        created_elements.append(f"Heading (level {level}): {text[:50]}...")
                elif element_type == 'subheading':
                    level = element.get('level', 2)
                    if not isinstance(level, int) or level < 2:
                        level = 2
                    if level > 9:
                        level = 9
                    current_paragraph = doc.add_heading(text, level)
                    created_elements.append(f"Subheading (level {level}): {text[:50]}...")
                elif element_type == 'paragraph':
                    current_paragraph = doc.add_paragraph(text)
                    created_elements.append(f"Paragraph: {text[:50]}...")
                elif element_type == 'poem':
                    current_paragraph = doc.add_paragraph(text)
                    current_paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
                    created_elements.append(f"Poem: {text[:50]}...")
                elif element_type == 'page_break':
                    current_paragraph = doc.add_paragraph()
                    current_paragraph.add_run().add_break(WD_BREAK.PAGE)
                    created_elements.append("Page break")
                elif element_type == 'table':
                    table_data = element.get('data') or element.get('rows') or []
                    table_style = element.get('table_style') or element.get('style')
                    if table_data and isinstance(table_data, list):
                        normalized_rows = []
                        max_cols = 0
                        for row in table_data:
                            if isinstance(row, (list, tuple)):
                                normalized_row = [str(cell) if cell is not None else '' for cell in row]
                            else:
                                normalized_row = [str(row) if row is not None else '']
                            normalized_rows.append(normalized_row)
                            if len(normalized_row) > max_cols:
                                max_cols = len(normalized_row)
                        if max_cols == 0:
                            max_cols = 1
                        table = doc.add_table(rows=len(normalized_rows), cols=max_cols)
                        if table_style:
                            table.style = table_style
                        for row_idx, row_data in enumerate(normalized_rows):
                            padded_row = row_data + [''] * (max_cols - len(row_data))
                            for col_idx, cell_value in enumerate(padded_row):
                                cell = table.cell(row_idx, col_idx)
                                cell.text = cell_value
                                for cell_paragraph in cell.paragraphs:
                                    apply_font_to_paragraph(cell_paragraph, font, font_size, bold, italic, underline, color)
                                    apply_alignment_to_paragraph(cell_paragraph, alignment)
                                    apply_spacing_to_paragraph(cell_paragraph, spacing_before, spacing_after, line_spacing)
                        created_elements.append(f"Table: {len(normalized_rows)} rows x {max_cols} cols")
                    else:
                        created_elements.append("Table: empty or invalid data")
                else:
                    current_paragraph = doc.add_paragraph(text)
                    created_elements.append(f"Paragraph: {text[:50]}...")
                
                if current_paragraph is not None:
                    apply_font_to_paragraph(current_paragraph, font, font_size, bold, italic, underline, color)
                    if alignment:
                        apply_alignment_to_paragraph(current_paragraph, alignment)
                    apply_spacing_to_paragraph(current_paragraph, spacing_before, spacing_after, line_spacing)
            
            # Ensure output directory exists
            output_path_obj = Path(output_path)
            output_path_obj.parent.mkdir(parents=True, exist_ok=True)
            
            # Save the document
            doc.save(output_path)
            
            # Get file info
            file_info = self._get_file_info(output_path_obj, output_path_obj.name)
            file_info.update({
                'content_type': 'docx',
                'paragraph_count': len(doc.paragraphs),
                'table_count': len(doc.tables),
                'created_elements': created_elements
            })
            
            return {
                'success': True,
                'message': f'Successfully created Word document with {len(content)} elements',
                'file_path': str(output_path),
                'document_info': file_info,
                'elements_created': created_elements
            }
            
        except Exception as e:
            return {
                'success': False,
                'error': str(e),
                'message': f'Failed to create Word document: {e}'
            }
    
    def replace_docx_text(self, file_path: str, find_text: str, replace_text: str, 
                         case_sensitive: bool = False) -> Dict[str, Any]:
        """
        Find and replace text in a Word document.
        
        Args:
            file_path: Path to the DOCX file
            find_text: Text to find
            replace_text: Text to replace with
            case_sensitive: Whether the search should be case sensitive
            
        Returns:
            Dictionary with replacement results
        """
        if not self.dependencies['docx']:
            return {
                'success': False,
                'error': 'python-docx library not available',
                'message': 'Cannot edit DOCX files without python-docx library'
            }
        
        try:
            import docx
            
            file_path_obj = Path(file_path)
            if not file_path_obj.exists():
                return {
                    'success': False,
                    'error': 'File not found',
                    'message': f'File not found: {file_path}'
                }
            
            # Load the document
            doc = docx.Document(file_path)
            
            # Track replacements
            replacements = 0
            
            # Search and replace in paragraphs
            for paragraph in doc.paragraphs:
                original_text = paragraph.text
                
                if case_sensitive:
                    if find_text in original_text:
                        paragraph.text = original_text.replace(find_text, replace_text)
                        replacements += 1
                else:
                    # Case-insensitive replacement
                    import re
                    pattern = re.compile(re.escape(find_text), re.IGNORECASE)
                    if pattern.search(original_text):
                        # This is a simplified approach - for more complex cases,
                        # we'd need to preserve formatting within runs
                        paragraph.text = pattern.sub(replace_text, original_text)
                        replacements += 1
            
            # Save the modified document
            doc.save(file_path)
            
            return {
                'success': True,
                'message': f'Replaced "{find_text}" with "{replace_text}" in {replacements} places',
                'file_path': file_path,
                'replacements': replacements,
                'case_sensitive': case_sensitive
            }
            
        except Exception as e:
            return {
                'success': False,
                'error': str(e),
                'message': f'Failed to replace text in document: {e}'
            }
    
    def get_docx_info(self, file_path: str) -> Dict[str, Any]:
        """
        Get detailed information about a Word document.
        
        Args:
            file_path: Path to the DOCX file
            
        Returns:
            Dictionary with document information
        """
        if not self.dependencies['docx']:
            return {
                'success': False,
                'error': 'python-docx library not available',
                'message': 'Cannot analyze DOCX files without python-docx library'
            }
        
        try:
            import docx
            
            file_path_obj = Path(file_path)
            if not file_path_obj.exists():
                return {
                    'success': False,
                    'error': 'File not found',
                    'message': f'File not found: {file_path}'
                }
            
            # Load the document
            doc = docx.Document(file_path)
            
            # Get basic document info
            paragraphs = []
            for i, para in enumerate(doc.paragraphs):
                if para.text.strip():
                    # Get font information from first run if available
                    font_info = {}
                    if para.runs:
                        first_run = para.runs[0]
                        font_info = {
                            'font_name': first_run.font.name,
                            'font_size': str(first_run.font.size) if first_run.font.size else None,
                            'bold': first_run.font.bold,
                            'italic': first_run.font.italic
                        }
                    
                    paragraphs.append({
                        'index': i,
                        'text': para.text,
                        'style': para.style.name if para.style else 'Normal',
                        'run_count': len(para.runs),
                        'font_info': font_info
                    })
            
            # Get tables
            tables = []
            for i, table in enumerate(doc.tables):
                table_data = []
                for row in table.rows:
                    row_data = [cell.text for cell in row.cells]
                    table_data.append(row_data)
                
                tables.append({
                    'index': i,
                    'rows': len(table.rows),
                    'cols': len(table.columns) if table.columns else 0,
                    'data_preview': table_data[:3]  # First 3 rows only
                })
            
            # Get document properties
            core_properties = {}
            try:
                core_properties = {
                    'title': doc.core_properties.title,
                    'author': doc.core_properties.author,
                    'created': str(doc.core_properties.created),
                    'modified': str(doc.core_properties.modified),
                    'last_modified_by': doc.core_properties.last_modified_by
                }
            except:
                pass
            
            return {
                'success': True,
                'file_path': file_path,
                'paragraph_count': len([p for p in doc.paragraphs if p.text.strip()]),
                'table_count': len(doc.tables),
                'total_runs': sum(len(p.runs) for p in doc.paragraphs),
                'paragraphs_preview': paragraphs[:5],  # First 5 non-empty paragraphs
                'tables_preview': tables,
                'core_properties': core_properties,
                'message': f'Document analysis complete: {len(paragraphs)} paragraphs, {len(tables)} tables'
            }
            
        except Exception as e:
            return {
                'success': False,
                'error': str(e),
                'message': f'Failed to analyze document: {e}'
            }
    
    def save_docx(self, doc_object, file_path: str) -> Dict[str, Any]:
        """
        Save a docx.Document object to a file.
        
        Args:
            doc_object: A docx.Document object
            file_path: Path where to save the file
            
        Returns:
            Dictionary with save results
        """
        if not self.dependencies['docx']:
            return {
                'success': False,
                'error': 'python-docx library not available',
                'message': 'Cannot save DOCX files without python-docx library'
            }
        
        try:
            # Ensure the directory exists
            save_path = Path(file_path)
            save_path.parent.mkdir(parents=True, exist_ok=True)
            
            # Save the document
            doc_object.save(save_path)
            
            # Get file info
            file_info = self._get_file_info(save_path, save_path.name)
            file_info.update({
                'content_type': 'docx',
                'paragraph_count': len(doc_object.paragraphs),
                'table_count': len(doc_object.tables)
            })
            
            return {
                'success': True,
                'document_info': file_info,
                'message': f'Successfully saved Word document to: {file_path}',
                'file_path': str(save_path)
            }
            
        except Exception as e:
            return {
                'success': False,
                'error': str(e),
                'message': f'Failed to save Word document: {e}'
            }
    
    def extract_docx_content(self, file_path: str) -> Dict[str, Any]:
        """
        Extract detailed content from a DOCX file.
        
        Args:
            file_path: Path to the DOCX file
            
        Returns:
            Dictionary with extracted content
        """
        if not self.dependencies['docx']:
            return {
                'success': False,
                'error': 'python-docx library not available',
                'message': 'Cannot extract DOCX content without python-docx library'
            }
        
        try:
            import docx
            
            file_path = Path(file_path)
            if not file_path.exists():
                return {
                    'success': False,
                    'error': 'File not found',
                    'message': f'File not found: {file_path}'
                }
            
            doc = docx.Document(file_path)
            
            # Extract paragraphs with their styles
            paragraphs = []
            for i, para in enumerate(doc.paragraphs):
                if para.text.strip():  # Skip empty paragraphs
                    style = para.style.name if para.style else 'Normal'
                    paragraphs.append({
                        'index': i,
                        'text': para.text,
                        'style': style,
                        'runs': len(para.runs)
                    })
            
            # Extract tables
            tables = []
            for i, table in enumerate(doc.tables):
                table_data = []
                for row in table.rows:
                    row_data = [cell.text for cell in row.cells]
                    table_data.append(row_data)
                
                tables.append({
                    'index': i,
                    'rows': len(table.rows),
                    'cols': len(table.columns) if table.columns else 0,
                    'data': table_data
                })
            
            # Extract styles
            styles = []
            try:
                from docx.enum.style import WD_STYLE_TYPE
                for style in doc.styles:
                    if style.type == WD_STYLE_TYPE.PARAGRAPH:
                        styles.append({
                            'name': style.name,
                            'type': 'paragraph'
                        })
            except ImportError:
                # If WD_STYLE_TYPE is not available, list all styles
                for style in doc.styles:
                    styles.append({
                        'name': style.name,
                        'type': str(style.type)
                    })
            
            return {
                'success': True,
                'paragraphs': paragraphs,
                'tables': tables,
                'styles': styles,
                'paragraph_count': len(paragraphs),
                'table_count': len(tables),
                'style_count': len(styles),
                'message': f'Successfully extracted content from {file_path.name}'
            }
            
        except Exception as e:
            return {
                'success': False,
                'error': str(e),
                'message': f'Failed to extract DOCX content: {e}'
            }
    
    def _generate_summary(self, file_info: Dict[str, Any]) -> str:
        """Generate a human-readable summary of the document."""
        parts = []
        
        # Basic file info
        parts.append(f"Document: {file_info['original_filename']}")
        parts.append(f"Type: {file_info['content_type'].upper()}")
        parts.append(f"Size: {file_info['file_size_human']}")
        
        # Content-specific summary
        if file_info.get('processing_success', False):
            content_type = file_info.get('content_type', '')
            
            if content_type == 'docx':
                parts.append(f"Paragraphs: {file_info.get('paragraph_count', 0)}")
                parts.append(f"Tables: {file_info.get('table_count', 0)}")
                
            elif content_type == 'pdf':
                parts.append(f"Pages: {file_info.get('page_count', 0)}")
                parts.append(f"Has text: {'Yes' if file_info.get('has_text', False) else 'No'}")
                if file_info.get('is_encrypted', False):
                    parts.append("⚠️ Document is encrypted")
                    
            elif content_type == 'pptx':
                parts.append(f"Slides: {file_info.get('slide_count', 0)}")
                
            elif content_type == 'excel':
                parts.append(f"Sheets: {file_info.get('sheet_count', 0)}")
                if file_info.get('sheet_names'):
                    parts.append(f"Sheet names: {', '.join(file_info['sheet_names'][:3])}")
                    
            elif content_type == 'csv':
                parts.append(f"Rows: {file_info.get('row_count', 'Unknown')}")
                parts.append(f"Columns: {file_info.get('column_count', 'Unknown')}")
                
            elif content_type == 'text':
                parts.append(f"Lines: {file_info.get('line_count', 0)}")
                parts.append(f"Words: {file_info.get('word_count', 0)}")
        
        else:
            parts.append("⚠️ Limited processing - file type may not be fully supported")
            if file_info.get('processing_error'):
                parts.append(f"Error: {file_info['processing_error']}")
        
        # Add preview if available
        if file_info.get('content_preview'):
            preview = file_info['content_preview']
            if len(preview) > 100:
                preview = preview[:100] + "..."
            parts.append(f"Preview: {preview}")
        
        return "\n".join(parts)
    
    def _calculate_file_hash(self, file_path: Path) -> str:
        """Calculate MD5 hash of file."""
        hash_md5 = hashlib.md5()
        with open(file_path, "rb") as f:
            for chunk in iter(lambda: f.read(4096), b""):
                hash_md5.update(chunk)
        return hash_md5.hexdigest()
    
    def _human_readable_size(self, size_bytes: int) -> str:
        """Convert bytes to human readable format."""
        if size_bytes == 0:
            return "0B"
        
        units = ["B", "KB", "MB", "GB", "TB"]
        import math
        i = int(math.floor(math.log(size_bytes, 1024)))
        p = math.pow(1024, i)
        s = round(size_bytes / p, 2)
        return f"{s} {units[i]}"


def print_document_summary(document_info: Dict[str, Any]):
    """
    Print a formatted summary of the document to terminal.
    
    Args:
        document_info: Document analysis results
    """
    print("\n" + "="*60)
    print("DOCUMENT ANALYSIS SUMMARY")
    print("="*60)
    
    # Basic info
    print(f"\n📄 File: {document_info.get('original_filename', 'Unknown')}")
    print(f"📁 Type: {document_info.get('content_type', 'Unknown').upper()}")
    print(f"📊 Size: {document_info.get('file_size_human', 'Unknown')}")
    print(f"🆔 Hash: {document_info.get('file_hash', 'Unknown')[:12]}...")
    
    # Status
    if document_info.get('processing_success', False):
        print("✅ Processing: Successful")
    else:
        print("⚠️ Processing: Limited")
        if document_info.get('processing_error'):
            print(f"   Error: {document_info['processing_error']}")
    
    # Content details
    content_type = document_info.get('content_type', '')
    
    if content_type == 'docx':
        print(f"\n📝 Paragraphs: {document_info.get('paragraph_count', 0)}")
        print(f"📊 Tables: {document_info.get('table_count', 0)}")
        
    elif content_type == 'pdf':
        print(f"\n📄 Pages: {document_info.get('page_count', 0)}")
        print(f"📝 Has extractable text: {'Yes' if document_info.get('has_text', False) else 'No'}")
        if document_info.get('is_encrypted', False):
            print("🔒 Document is encrypted")
            
    elif content_type == 'pptx':
        print(f"\n🎯 Slides: {document_info.get('slide_count', 0)}")
        
    elif content_type == 'excel':
        print(f"\n📊 Sheets: {document_info.get('sheet_count', 0)}")
        if document_info.get('sheet_names'):
            print(f"   Names: {', '.join(document_info['sheet_names'])}")
            
    elif content_type == 'csv':
        print(f"\n📊 Rows: {document_info.get('row_count', 'Unknown')}")
        print(f"📈 Columns: {document_info.get('column_count', 'Unknown')}")
        if document_info.get('column_names'):
            print(f"   Columns: {', '.join(document_info['column_names'])}")
            
    elif content_type == 'text':
        print(f"\n📝 Lines: {document_info.get('line_count', 0)}")
        print(f"🔤 Words: {document_info.get('word_count', 0)}")
        print(f"📏 Length: {document_info.get('total_length', 0)} characters")
    
    # Preview
    if document_info.get('content_preview'):
        print(f"\n🔍 Preview:")
        print("-"*40)
        print(document_info['content_preview'])
        if len(document_info['content_preview']) >= 500:
            print("... (truncated)")
        print("-"*40)
    
    print("\n" + "="*60)


# Example usage
if __name__ == "__main__":
    # Create processor with workspace
    processor = DocumentProcessor("/tmp/docmaster_workspace")
    
    # Example: Process a test file (you would replace this with actual file upload)
    test_file = Path(__file__).parent / "test_document.txt"
    
    # Create a test file if it doesn't exist
    if not test_file.exists():
        test_file.write_text("This is a test document.\nIt has multiple lines.\nFor testing document processing.\n" * 10)
    
    try:
        # Process the file
        result = processor.process_uploaded_file(str(test_file))
        
        # Print summary
        print_document_summary(result)
        
        # Also print the summary from the processor
        print("\n📋 Generated Summary:")
        print("-"*40)
        print(result['summary'])
        
    except Exception as e:
        print(f"Error: {e}")