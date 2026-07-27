#!/usr/bin/env python3
"""Claude 能力介绍页 — 单页 PPT 构建脚本。

Archetype: research-note
Page task: explain
风格: 空白页直生，无模板约束，参考 ppt-polished-deck-collab-traditional 规范
"""

from __future__ import annotations

import sys
from pathlib import Path

# 把 skill scripts 目录加入路径，复用 helper
SKILL_SCRIPTS = Path(__file__).resolve().parents[3] / "skills" / "presentation-skills" / "ppt-polished-deck-collab-traditional" / "scripts"
sys.path.insert(0, str(SKILL_SCRIPTS))

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from ppt_asset_helpers import (
    add_panel,
    add_slide_header,
    default_palette,
    default_typography_tokens,
    new_presentation,
    panel_content_box,
    save_presentation,
    tint,
    DEFAULT_FONT_NAME,
    DEFAULT_LINE_SPACING_MULTIPLE,
)

OUTPUT_DIR = Path(__file__).parent / "build"
OUTPUT_PATH = OUTPUT_DIR / "claude_capabilities.pptx"


def _add_capability_card(
    slide,
    title: str,
    bullets: list[str],
    left: float,
    top: float,
    width: float,
    height: float,
    accent_rgb: tuple[int, int, int],
) -> None:
    """在指定区域内绘制一个能力卡片（panel + 正文 bullet）。"""
    tokens = default_typography_tokens()
    palette = default_palette()

    add_panel(slide, title, left, top, width, height, accent_rgb)
    cl, ct, cw, ch = panel_content_box(left, top, width, height)

    box = slide.shapes.add_textbox(Inches(cl), Inches(ct), Inches(cw), Inches(ch))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP

    for i, bullet in enumerate(bullets):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = f"• {bullet}"
        para.font.size = Pt(tokens["body_font_pt"])
        para.font.name = DEFAULT_FONT_NAME
        para.font.color.rgb = RGBColor(*palette["title"])
        para.line_spacing = tokens["body_line_spacing_multiple"]
        para.space_before = Pt(0)
        para.space_after = Pt(2)


def build() -> Path:
    prs = new_presentation()
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)

    palette = default_palette()
    tokens = default_typography_tokens()

    # ── 标题头 ──────────────────────────────────────────────────────────────
    add_slide_header(
        slide,
        figure_tag="01",
        title="Claude 能力概览",
        subtitle="理解、推理、编码、创作 — 您的全栈 AI 协作伙伴",
    )

    # ── 分隔线 ──────────────────────────────────────────────────────────────
    divider = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.LINE = 9 but rectangle works; use connector workaround via textbox bg
        Inches(0.72), Inches(1.18), Inches(14.56), Inches(0.03),
    )
    divider.fill.solid()
    divider.fill.fore_color.rgb = RGBColor(*palette["muted"])
    divider.line.color.rgb = RGBColor(*palette["muted"])

    # ── 6 能力卡片（2行 × 3列）────────────────────────────────────────────
    cards = [
        {
            "title": "理解与分析",
            "accent": palette["blue"],
            "bullets": [
                "深度理解长文档、代码与多模态输入",
                "抽取关键信息，生成结构化摘要",
                "支持 200K token 超长上下文窗口",
            ],
        },
        {
            "title": "推理与规划",
            "accent": palette["violet"],
            "bullets": [
                "逐步拆解复杂问题，多步骤推导",
                "评估方案权衡，给出明确建议",
                "构建 agent 工具链与多轮任务流",
            ],
        },
        {
            "title": "代码与工程",
            "accent": palette["teal"],
            "bullets": [
                "生成、审查、重构 30+ 编程语言代码",
                "调试、解释，并给出测试策略",
                "直接操作文件系统与 shell 命令",
            ],
        },
        {
            "title": "文档与写作",
            "accent": palette["emerald"],
            "bullets": [
                "起草报告、PPT 叙事、邮件与文案",
                "多语言翻译，保留专业领域风格",
                "制作高质量可编辑 pptx 交付物",
            ],
        },
        {
            "title": "多模态理解",
            "accent": palette["amber"],
            "bullets": [
                "分析图表、截图、表格与 PDF 内容",
                "从图像提取数据，辅助视觉说明",
                "结合文字与图像进行综合推断",
            ],
        },
        {
            "title": "工具调用与协作",
            "accent": palette["rose"],
            "bullets": [
                "调用外部 API、搜索、数据库等工具",
                "在 agent 编排框架中担任核心模型",
                "支持 MCP 协议扩展自定义工具集",
            ],
        },
    ]

    # 布局参数
    left_margin = 0.72
    card_width = 4.60
    card_height = 2.95
    col_gap = 0.25
    row_gap = 0.22
    top_start = 1.30

    for idx, card in enumerate(cards):
        row = idx // 3
        col = idx % 3
        left = left_margin + col * (card_width + col_gap)
        top = top_start + row * (card_height + row_gap)
        _add_capability_card(
            slide,
            title=card["title"],
            bullets=card["bullets"],
            left=left,
            top=top,
            width=card_width,
            height=card_height,
            accent_rgb=card["accent"],
        )

    # ── 页脚说明 ─────────────────────────────────────────────────────────────
    footer = slide.shapes.add_textbox(Inches(0.72), Inches(8.50), Inches(14.56), Inches(0.26))
    fp = footer.text_frame.paragraphs[0]
    fp.text = "Claude · Anthropic · claude.ai/claude-code  —  能力以实际部署版本为准"
    fp.font.size = Pt(tokens["caption_font_pt"])
    fp.font.name = DEFAULT_FONT_NAME
    fp.font.color.rgb = RGBColor(*palette["muted"])
    fp.line_spacing = DEFAULT_LINE_SPACING_MULTIPLE

    return save_presentation(prs, OUTPUT_PATH)


if __name__ == "__main__":
    out = build()
    print(f"✓ 已生成: {out}")
