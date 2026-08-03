"""
PPTX Editing Skill for DocMaster.

Lets DocMaster open an **existing** .pptx and make targeted modifications
without rebuilding the whole deck from scratch.

Public API
----------
    skill = PptxProcessorSkill()
    skill.extract_content(pptx_path)          -> dict
    skill.edit_pptx(pptx_path, edits)         -> dict
"""

from __future__ import annotations

import copy
import logging
import re
from pathlib import Path
from typing import Any

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _pptx_import():
    try:
        import pptx  # noqa: F401
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.util import Inches, Pt, Emu
        from pptx.enum.text import PP_ALIGN
        return Presentation, RGBColor, Inches, Pt, Emu, PP_ALIGN
    except ImportError as exc:
        raise ImportError(
            f"python-pptx is required for PPTX editing: {exc}"
        ) from exc


def _parse_hex_color(hex_str: str) -> tuple[int, int, int] | None:
    """Parse '#RRGGBB' or 'RRGGBB' (3- or 6-char) → (R,G,B). None on failure."""
    if not hex_str or not isinstance(hex_str, str):
        return None
    h = hex_str.strip().lstrip("#")
    if len(h) == 3:
        h = h[0]*2 + h[1]*2 + h[2]*2
    if len(h) != 6:
        return None
    try:
        return (int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
    except ValueError:
        return None


def _slide_texts(slide) -> list[dict]:
    """Extract all text runs from a slide for content inspection."""
    items = []
    for idx, shape in enumerate(slide.shapes):
        if not getattr(shape, "has_text_frame", False):
            continue
        for p_idx, para in enumerate(shape.text_frame.paragraphs):
            text = para.text.strip()
            if text:
                sizes = [r.font.size.pt for r in para.runs if r.font.size]
                items.append({
                    "shape_id": shape.shape_id,
                    "shape_name": shape.name,
                    "shape_index": idx,
                    "paragraph_index": p_idx,
                    "text": text,
                    "font_pt": max(sizes) if sizes else None,
                })
    return items


def _find_shapes(slide, *, shape_id: int | None, shape_name: str | None,
                 shape_index: int | None) -> list:
    """Return matching shapes from slide. At least one selector required."""
    results = []
    shapes = list(slide.shapes)
    for i, s in enumerate(shapes):
        if shape_id is not None and s.shape_id == shape_id:
            results.append(s)
        elif shape_name is not None and (
            s.name == shape_name or re.search(shape_name, s.name, re.IGNORECASE)
        ):
            results.append(s)
        elif shape_index is not None and i == shape_index:
            results.append(s)
    return results


def _set_text_in_shape(shape, new_text: str,
                       paragraph_index: int | None = None) -> bool:
    """Replace text in a shape's text frame.

    If paragraph_index is given, replace only that paragraph.
    Otherwise replace the full text frame (preserves first run's formatting).
    """
    if not getattr(shape, "has_text_frame", False):
        return False
    tf = shape.text_frame
    if paragraph_index is not None:
        paras = tf.paragraphs
        if paragraph_index >= len(paras):
            return False
        para = paras[paragraph_index]
        _replace_para_text(para, new_text)
        return True
    # Replace all paragraphs: keep first paragraph's style, clear the rest.
    lines = new_text.split("\n")
    paras = list(tf.paragraphs)
    _replace_para_text(paras[0], lines[0])
    # Remove extra paragraphs beyond what we need
    while len(paras) > len(lines):
        p_elem = paras[-1]._p
        p_elem.getparent().remove(p_elem)
        paras.pop()
    # Add extra paragraphs if needed
    for line in lines[1:]:
        from pptx.oxml.ns import qn
        from lxml import etree
        new_p = copy.deepcopy(paras[0]._p)
        # Clear runs from copy
        for r in new_p.findall(qn("a:r")):
            new_p.remove(r)
        # Add a single run with the text
        r_elem = copy.deepcopy(paras[0]._p.findall(qn("a:r"))[0]) if paras[0]._p.findall(qn("a:r")) else etree.SubElement(new_p, qn("a:r"))
        t_elem = r_elem.find(qn("a:t"))
        if t_elem is None:
            t_elem = etree.SubElement(r_elem, qn("a:t"))
        t_elem.text = line
        new_p.append(r_elem)
        paras[-1]._p.addnext(new_p)
    return True


def _replace_para_text(para, new_text: str):
    """Replace a paragraph's text while preserving the first run's formatting."""
    from pptx.oxml.ns import qn
    runs = para._p.findall(qn("a:r"))
    if not runs:
        # No runs — just set via .text (loses formatting but better than nothing)
        para.text = new_text
        return
    # Keep first run, set its text, remove the rest
    first_run = runs[0]
    t_elem = first_run.find(qn("a:t"))
    if t_elem is None:
        from lxml import etree
        t_elem = etree.SubElement(first_run, qn("a:t"))
    t_elem.text = new_text
    for r in runs[1:]:
        para._p.remove(r)
    # Also remove any <a:br> line-break elements
    for br in para._p.findall(qn("a:br")):
        para._p.remove(br)


# ---------------------------------------------------------------------------
# Main class
# ---------------------------------------------------------------------------

class PptxProcessorSkill:
    """Edit an existing .pptx without rebuilding it from scratch."""

    def extract_content(self, pptx_path: str) -> dict:
        """Return a structured summary of every slide's shapes and texts.

        Use this before editing to discover shape names / indices so you can
        target edits accurately.

        Returns dict with success, slide_count, slides (list of
        {slide_number, layout_name, shapes (list of {shape_id, shape_name,
        shape_index, shape_type, has_chart, has_table, texts})}).
        """
        path = Path(pptx_path)
        if not path.exists():
            return {"success": False, "error": "File not found",
                    "message": f"No file at {pptx_path}"}
        try:
            Presentation, *_ = _pptx_import()
            prs = Presentation(str(path))
        except Exception as exc:
            return {"success": False, "error": str(exc),
                    "message": f"Failed to open {pptx_path}: {exc}"}

        slides_info = []
        for s_idx, slide in enumerate(prs.slides, 1):
            shapes_info = []
            for sh_idx, shape in enumerate(slide.shapes):
                texts = []
                if getattr(shape, "has_text_frame", False):
                    for p in shape.text_frame.paragraphs:
                        t = p.text.strip()
                        if t:
                            sizes = [r.font.size.pt for r in p.runs if r.font.size]
                            texts.append({
                                "paragraph_index": p._p.getparent().index(p._p),
                                "text": t,
                                "font_pt": round(max(sizes), 1) if sizes else None,
                            })
                table_data = None
                if getattr(shape, "has_table", False):
                    tbl = shape.table
                    table_data = {
                        "rows": len(tbl.rows),
                        "cols": len(tbl.columns),
                        "cells": [
                            [tbl.cell(r, c).text
                             for c in range(len(tbl.columns))]
                            for r in range(len(tbl.rows))
                        ],
                    }
                shapes_info.append({
                    "shape_id": shape.shape_id,
                    "shape_name": shape.name,
                    "shape_index": sh_idx,
                    "shape_type": str(getattr(shape, "shape_type", "?")),
                    "has_chart": getattr(shape, "has_chart", False),
                    "has_table": bool(table_data),
                    "texts": texts,
                    "table": table_data,
                })
            slides_info.append({
                "slide_number": s_idx,
                "layout_name": slide.slide_layout.name,
                "shapes": shapes_info,
            })

        return {
            "success": True,
            "pptx_path": str(path),
            "slide_count": len(prs.slides),
            "slide_size_in": {
                "width": round(prs.slide_width / 914400, 2),
                "height": round(prs.slide_height / 914400, 2),
            },
            "slides": slides_info,
        }

    def edit_pptx(self, pptx_path: str, edits: list[dict]) -> dict:
        """Apply a list of structured edits to an existing .pptx in place.

        Supported edit types
        --------------------
        replace_text
            Replace every occurrence of `old_text` with `new_text` across
            the whole presentation (all slides).
            Required: old_text, new_text
            Optional: slide_number (limit to one slide)

        set_shape_text
            Set the text of a specific shape on a specific slide.
            Required: slide_number (1-based), new_text, and at least one of:
              shape_id | shape_name | shape_index
            Optional: paragraph_index (target a specific paragraph)

        set_shape_fill
            Set the fill color of a shape.
            Required: slide_number, color (hex), and one of:
              shape_id | shape_name | shape_index

        set_slide_background
            Set the background fill color of a slide.
            Required: slide_number (or "all"), color (hex)

        add_textbox
            Add a new text box to a slide.
            Required: slide_number, text
            Optional: left, top, width, height (inches, default 1/1/4/1),
              font_pt, bold, color (hex)

        delete_shape
            Remove a shape from a slide.
            Required: slide_number, and one of:
              shape_id | shape_name | shape_index

        delete_slide
            Remove a slide.
            Required: slide_number (1-based)

        duplicate_slide
            Duplicate a slide and insert the copy after it.
            Required: slide_number

        reorder_slides
            Reorder all slides to match the given order list.
            Required: order (list of 1-based slide numbers, e.g. [3,1,2])

        set_table_cell
            Set the text of a specific table cell.
            Required: slide_number, shape_index (or shape_id/shape_name),
              row, col, new_text

        Returns dict with success, applied (list of applied edit descriptions),
        errors (list of {edit_index, edit_type, error}), message.
        """
        path = Path(pptx_path)
        if not path.exists():
            return {"success": False, "error": "File not found",
                    "message": f"No file at {pptx_path}"}
        if not edits:
            return {"success": False, "error": "No edits provided",
                    "message": "edits list is empty."}

        try:
            Presentation, RGBColor, Inches, Pt, Emu, PP_ALIGN = _pptx_import()
            prs = Presentation(str(path))
        except Exception as exc:
            return {"success": False, "error": str(exc),
                    "message": f"Failed to open {pptx_path}: {exc}"}

        applied: list[str] = []
        errors: list[dict] = []
        n_slides = len(prs.slides)

        def _slide(slide_number: int):
            idx = int(slide_number) - 1
            if idx < 0 or idx >= len(prs.slides):
                raise ValueError(
                    f"slide_number={slide_number} out of range "
                    f"(deck has {len(prs.slides)} slides)"
                )
            return prs.slides[idx]

        for i, edit in enumerate(edits):
            if not isinstance(edit, dict):
                errors.append({"edit_index": i, "edit_type": "?",
                                "error": "edit is not a dict"})
                continue
            etype = edit.get("type") or edit.get("edit_type") or ""
            try:
                # ── replace_text ─────────────────────────────────────────
                if etype == "replace_text":
                    old_text = str(edit["old_text"])
                    new_text = str(edit["new_text"])
                    limit_slide = edit.get("slide_number")
                    count = 0
                    slides_iter = ([_slide(limit_slide)] if limit_slide
                                   else list(prs.slides))
                    for sl in slides_iter:
                        for shape in sl.shapes:
                            if not getattr(shape, "has_text_frame", False):
                                continue
                            for para in shape.text_frame.paragraphs:
                                full = para.text
                                if old_text not in full:
                                    continue
                                from pptx.oxml.ns import qn as _qn
                                # Build the full paragraph text from runs,
                                # replace, then re-distribute across runs.
                                run_texts = [r.text or "" for r in para.runs]
                                joined = "".join(run_texts)
                                if old_text not in joined:
                                    continue
                                replaced = joined.replace(old_text, new_text)
                                # Put everything in the first run, clear others
                                runs = para._p.findall(_qn("a:r"))
                                if runs:
                                    t_elem = runs[0].find(_qn("a:t"))
                                    if t_elem is None:
                                        from lxml import etree
                                        t_elem = etree.SubElement(runs[0], _qn("a:t"))
                                    t_elem.text = replaced
                                    for r in runs[1:]:
                                        para._p.remove(r)
                                    count += 1
                    applied.append(
                        f"replace_text: {count} paragraph(s) updated "
                        f"({old_text!r} → {new_text!r})"
                    )

                # ── set_shape_text ────────────────────────────────────────
                elif etype == "set_shape_text":
                    sl = _slide(edit["slide_number"])
                    matches = _find_shapes(
                        sl,
                        shape_id=edit.get("shape_id"),
                        shape_name=edit.get("shape_name"),
                        shape_index=edit.get("shape_index"),
                    )
                    if not matches:
                        raise ValueError(
                            "No shape found with "
                            f"shape_id={edit.get('shape_id')} / "
                            f"shape_name={edit.get('shape_name')!r} / "
                            f"shape_index={edit.get('shape_index')} "
                            f"on slide {edit['slide_number']}"
                        )
                    p_idx = edit.get("paragraph_index")
                    new_text = str(edit["new_text"])
                    ok_shapes = []
                    for shape in matches:
                        if _set_text_in_shape(shape, new_text, p_idx):
                            ok_shapes.append(shape.name)
                    applied.append(
                        f"set_shape_text: slide {edit['slide_number']} "
                        f"shape(s) {ok_shapes} → {new_text[:60]!r}"
                    )

                # ── set_shape_fill ────────────────────────────────────────
                elif etype == "set_shape_fill":
                    sl = _slide(edit["slide_number"])
                    rgb = _parse_hex_color(str(edit["color"]))
                    if rgb is None:
                        raise ValueError(
                            f"Invalid color {edit['color']!r}. "
                            "Use 6-digit hex, e.g. '2563EB'."
                        )
                    matches = _find_shapes(
                        sl,
                        shape_id=edit.get("shape_id"),
                        shape_name=edit.get("shape_name"),
                        shape_index=edit.get("shape_index"),
                    )
                    if not matches:
                        raise ValueError(
                            f"No shape found on slide {edit['slide_number']}"
                        )
                    for shape in matches:
                        shape.fill.solid()
                        shape.fill.fore_color.rgb = RGBColor(*rgb)
                    applied.append(
                        f"set_shape_fill: slide {edit['slide_number']} "
                        f"{len(matches)} shape(s) → #{edit['color']}"
                    )

                # ── set_slide_background ──────────────────────────────────
                elif etype == "set_slide_background":
                    rgb = _parse_hex_color(str(edit["color"]))
                    if rgb is None:
                        raise ValueError(f"Invalid color {edit['color']!r}")
                    hex_val = f"{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}"
                    slide_num = edit.get("slide_number")
                    if slide_num == "all" or slide_num is None:
                        target_slides = list(prs.slides)
                    else:
                        target_slides = [_slide(slide_num)]

                    from pptx.oxml.ns import qn as _qn
                    from lxml import etree as _et

                    for sl in target_slides:
                        # python-pptx's .background.fill.solid() does NOT
                        # remove an existing <a:noFill/> inside <p:bgPr> —
                        # noFill takes precedence and the color is silently
                        # ignored.  Write directly to the XML so we can:
                        #   1. Remove any noFill / existing fill child.
                        #   2. Add a proper solidFill.
                        #   3. Remove any bgRef override from the slide master.
                        cSld = sl._element.find(_qn('p:cSld'))
                        bg = cSld.find(_qn('p:bg')) if cSld is not None else None
                        if bg is None:
                            bg = _et.SubElement(cSld, _qn('p:bg'))
                            cSld.insert(0, bg)

                        bgPr = bg.find(_qn('p:bgPr'))
                        if bgPr is None:
                            bgPr = _et.SubElement(bg, _qn('p:bgPr'))

                        # Remove every existing fill element (noFill, solidFill,
                        # gradFill, blipFill, pattFill, grpFill).
                        _FILL_TAGS = {
                            _qn('a:noFill'), _qn('a:solidFill'),
                            _qn('a:gradFill'), _qn('a:blipFill'),
                            _qn('a:pattFill'), _qn('a:grpFill'),
                        }
                        for child in list(bgPr):
                            if child.tag in _FILL_TAGS:
                                bgPr.remove(child)

                        # Insert solidFill as the first child of bgPr.
                        solidFill = _et.Element(_qn('a:solidFill'))
                        srgbClr   = _et.SubElement(solidFill, _qn('a:srgbClr'))
                        srgbClr.set('val', hex_val)
                        bgPr.insert(0, solidFill)

                        # A <p:bgRef> element ties the slide background to a
                        # theme entry and overrides everything in bgPr in some
                        # renderers.  Remove it so our solidFill wins.
                        bgRef = bg.find(_qn('p:bgRef'))
                        if bgRef is not None:
                            bg.remove(bgRef)

                    applied.append(
                        f"set_slide_background: {len(target_slides)} slide(s) "
                        f"→ #{hex_val}"
                    )

                # ── add_textbox ───────────────────────────────────────────
                elif etype == "add_textbox":
                    sl = _slide(edit["slide_number"])
                    left   = float(edit.get("left",  1.0))
                    top    = float(edit.get("top",   1.0))
                    width  = float(edit.get("width", 4.0))
                    height = float(edit.get("height",1.0))
                    text   = str(edit.get("text", ""))
                    font_pt = float(edit.get("font_pt", 14.0))
                    bold   = bool(edit.get("bold", False))
                    rgb    = _parse_hex_color(str(edit["color"])) if edit.get("color") else None
                    box = sl.shapes.add_textbox(
                        Inches(left), Inches(top), Inches(width), Inches(height)
                    )
                    tf = box.text_frame
                    tf.word_wrap = True
                    for line_idx, line in enumerate(text.split("\n")):
                        para = tf.paragraphs[0] if line_idx == 0 else tf.add_paragraph()
                        para.text = line
                        para.font.size = Pt(font_pt)
                        para.font.bold = bold
                        if rgb:
                            para.font.color.rgb = RGBColor(*rgb)
                    applied.append(
                        f"add_textbox: slide {edit['slide_number']} "
                        f"@ ({left}\", {top}\") text={text[:40]!r}"
                    )

                # ── delete_shape ──────────────────────────────────────────
                elif etype == "delete_shape":
                    sl = _slide(edit["slide_number"])
                    matches = _find_shapes(
                        sl,
                        shape_id=edit.get("shape_id"),
                        shape_name=edit.get("shape_name"),
                        shape_index=edit.get("shape_index"),
                    )
                    if not matches:
                        raise ValueError(
                            f"No shape found on slide {edit['slide_number']}"
                        )
                    removed = []
                    for shape in matches:
                        sp = shape._element
                        sp.getparent().remove(sp)
                        removed.append(shape.name)
                    applied.append(
                        f"delete_shape: slide {edit['slide_number']} "
                        f"removed {removed}"
                    )

                # ── delete_slide ──────────────────────────────────────────
                elif etype == "delete_slide":
                    slide_num = int(edit["slide_number"])
                    if slide_num < 1 or slide_num > len(prs.slides):
                        raise ValueError(
                            f"slide_number={slide_num} out of range"
                        )
                    xml_slides = prs.slides._sldIdLst
                    slide_elem = xml_slides[slide_num - 1]
                    xml_slides.remove(slide_elem)
                    applied.append(f"delete_slide: removed slide {slide_num}")

                # ── reorder_slides ────────────────────────────────────────
                elif etype == "reorder_slides":
                    order = [int(x) for x in edit["order"]]
                    n = len(prs.slides)
                    if sorted(order) != list(range(1, n + 1)):
                        raise ValueError(
                            f"order must be a permutation of 1..{n}, got {order}"
                        )
                    xml_slides = prs.slides._sldIdLst
                    # Detach all slide id elements
                    elems = [xml_slides[i - 1] for i in range(1, n + 1)]
                    for e in elems:
                        xml_slides.remove(e)
                    # Re-attach in new order
                    for pos in order:
                        xml_slides.append(elems[pos - 1])
                    applied.append(f"reorder_slides: new order {order}")

                # ── set_table_cell ────────────────────────────────────────
                elif etype == "set_table_cell":
                    sl = _slide(edit["slide_number"])
                    row = int(edit["row"])
                    col = int(edit["col"])
                    new_text = str(edit["new_text"])
                    matches = _find_shapes(
                        sl,
                        shape_id=edit.get("shape_id"),
                        shape_name=edit.get("shape_name"),
                        shape_index=edit.get("shape_index"),
                    )
                    tbl_shapes = [s for s in (matches or list(sl.shapes))
                                  if getattr(s, "has_table", False)]
                    if not tbl_shapes:
                        raise ValueError(
                            f"No table found on slide {edit['slide_number']}"
                        )
                    table = tbl_shapes[0].table
                    cell = table.cell(row, col)
                    # Clear existing paragraphs; preserve first run formatting
                    tf = cell.text_frame
                    _replace_para_text(tf.paragraphs[0], new_text)
                    for extra_para in list(tf.paragraphs)[1:]:
                        extra_para._p.getparent().remove(extra_para._p)
                    applied.append(
                        f"set_table_cell: slide {edit['slide_number']} "
                        f"[{row},{col}] → {new_text!r}"
                    )

                else:
                    errors.append({
                        "edit_index": i,
                        "edit_type": etype,
                        "error": (
                            f"Unknown edit type {etype!r}. Supported: "
                            "replace_text, set_shape_text, set_shape_fill, "
                            "set_slide_background, add_textbox, delete_shape, "
                            "delete_slide, reorder_slides, set_table_cell"
                        ),
                    })

            except Exception as exc:  # noqa: BLE001
                errors.append({
                    "edit_index": i,
                    "edit_type": etype,
                    "error": str(exc),
                })

        if not applied and errors:
            return {
                "success": False,
                "applied": applied,
                "errors": errors,
                "message": (
                    f"All {len(errors)} edit(s) failed. "
                    "Check `errors` for details."
                ),
            }

        # Save back to the same file
        try:
            prs.save(str(path))
        except Exception as exc:
            return {
                "success": False,
                "error": str(exc),
                "applied": applied,
                "errors": errors,
                "message": f"Edits applied in memory but save failed: {exc}",
            }

        return {
            "success": True,
            "pptx_path": str(path),
            "applied_count": len(applied),
            "error_count": len(errors),
            "applied": applied,
            "errors": errors if errors else None,
            "message": (
                f"Applied {len(applied)} edit(s) to {path.name}."
                + (f" {len(errors)} edit(s) failed." if errors else "")
            ),
        }
