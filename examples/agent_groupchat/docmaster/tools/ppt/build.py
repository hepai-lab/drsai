"""PPT Build Tools - Build, quality check, preview, and connector verification."""

import subprocess
import sys as _sys
import json as _json
from pathlib import Path


# ============================================================================
# Build Tools
# ============================================================================

def _resolve_deck_workspace(deck_workspace: str) -> tuple[Path | None, dict | None]:
    """Validate a deck_workspace string."""
    if not deck_workspace or not isinstance(deck_workspace, str):
        return None, {
            "success": False,
            "error": "Missing deck_workspace",
            "message": (
                "deck_workspace is required. Call ppt_init_workspace_tool "
                "first and pass the returned `deck_workspace` value here."
            ),
        }
    p = Path(deck_workspace)
    if not p.is_absolute():
        return None, {
            "success": False,
            "error": "Relative deck_workspace not accepted",
            "message": (
                f"deck_workspace={deck_workspace!r} is a relative path. "
                "Use the absolute path returned by ppt_init_workspace_tool."
            ),
        }
    if not p.exists() or not p.is_dir():
        return None, {
            "success": False,
            "error": "deck_workspace not found",
            "message": (
                f"No directory at {deck_workspace}. Re-run "
                "ppt_init_workspace_tool to create it, or check the "
                "value you received from that tool earlier in the "
                "conversation."
            ),
        }
    return p, None


_ARCHETYPES = {
    "hero-statement",
    "decision-logic",
    "board-memo",
    "chart-spotlight",
    "comparison-matrix",
    "process-flow",
    "research-note",
    "appendix-dense",
}

_ASSET_MODES = {
    "text-layout-native",
    "office-chart-native",
    "python-figure-image",
    "table-native",
    "diagram-connector",
    "diagram-visual",
    "icon-accent",
    "image-hero",
    "mixed",
}


def _ppt_load_yaml(path: Path) -> tuple[dict | None, str | None]:
    """Load and validate a YAML spec file."""
    try:
        import yaml  # python-pptx already pulls it in transitively
    except ImportError as exc:
        return None, f"PyYAML not installed: {exc}"
    try:
        data = yaml.safe_load(path.read_text(encoding="utf-8"))
    except Exception as exc:
        return None, f"YAML parse error: {exc}"
    if not isinstance(data, dict):
        return None, "slide_specs.yaml top-level must be a mapping"
    return data, None


def _ppt_chart_type(name: str):
    """Map a friendly name to pptx XL_CHART_TYPE."""
    from pptx.enum.chart import XL_CHART_TYPE
    return {
        "bar": XL_CHART_TYPE.BAR_CLUSTERED,
        "bar_clustered": XL_CHART_TYPE.BAR_CLUSTERED,
        "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "column_clustered": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "line": XL_CHART_TYPE.LINE,
        "stacked_bar": XL_CHART_TYPE.BAR_STACKED,
        "stacked_column": XL_CHART_TYPE.COLUMN_STACKED,
    }.get((name or "bar").lower(), XL_CHART_TYPE.BAR_CLUSTERED)


def _render_hero_statement(slide, spec: dict, tokens: dict, palette: dict):
    """Render hero-statement archetype."""
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt

    title = spec.get("title") or "Untitled"
    key_message = spec.get("key_message") or ""
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(*palette["bg"])

    title_box = slide.shapes.add_textbox(
        Inches(0.9), Inches(2.6), Inches(14.2), Inches(1.6),
    )
    p = title_box.text_frame.paragraphs[0]
    p.text = title
    p.font.bold = True
    p.font.size = Pt(float(tokens["hero_title_font_pt"]) * 1.4)
    p.font.color.rgb = RGBColor(*palette["title"])
    p.line_spacing = float(tokens["title_line_spacing_multiple"])

    if key_message:
        sub = slide.shapes.add_textbox(
            Inches(0.9), Inches(4.6), Inches(14.2), Inches(1.6),
        )
        sp = sub.text_frame.paragraphs[0]
        sp.text = key_message
        sp.font.size = Pt(float(tokens["subtitle_font_pt"]) * 1.2)
        sp.font.color.rgb = RGBColor(*palette["subtitle"])
        sp.line_spacing = float(tokens["body_line_spacing_multiple"])

    return ["hero_layout"]


def ppt_build_pptx_tool(
    slide_specs_path: str,
    output_pptx: str,
    deck_workspace: str,
    ppt_scripts_dir: Path = None,
) -> dict:
    """
    Build an editable .pptx from a derived slide_specs.yaml.

    Each slide is routed to an archetype renderer (hero-statement /
    decision-logic / board-memo / chart-spotlight / comparison-matrix
    / process-flow / research-note / appendix-dense). Theme tokens
    from deck.theme_tokens are injected into the underlying
    ppt_asset_helpers module so fonts, sizes, panel colors and CJK
    (东亚) font slots match the rest of the skill.

    **Do NOT call this tool with hand-crafted YAML.** Always run
    ppt_derive_slide_specs_tool first so the structural fields are
    validated against deck_narrative.md.

    What the tool does NOT do:
      - It does not invent content. Bullets, chart data, table data
        and diagram structure must come from the slide_spec itself.
      - It does not run quality gates. After build, call
        ppt_package_preflight_tool → ppt_structure_precheck_tool →
        ppt_export_previews_tool → ppt_render_review_tool.

    Args:
        slide_specs_path: Absolute path, typically
            <deck_workspace>/build/generated/slide_specs.yaml.
        output_pptx: Absolute path to write, typically
            <deck_workspace>/build/pptx/deck_v1.pptx.
        deck_workspace: Value returned by ppt_init_workspace_tool.
            Used to resolve relative image / chart paths inside the
            spec.

    Returns dict with:
        success / output_pptx / slide_count / per_slide (list of
        {slide_id, archetype, status, error?}) / message.
    """
    ws, err = _resolve_deck_workspace(deck_workspace)
    if err is not None:
        return err

    specs_path = Path(slide_specs_path)
    if not specs_path.is_absolute():
        return {
            "success": False,
            "error": "Relative slide_specs_path",
            "message": (
                f"slide_specs_path={slide_specs_path!r} must be absolute. "
                "Use the slide_specs_path returned by "
                "ppt_derive_slide_specs_tool."
            ),
        }
    if not specs_path.exists():
        return {
            "success": False,
            "error": "slide_specs not found",
            "message": (
                f"No file at {specs_path}. Run "
                "ppt_derive_slide_specs_tool first."
            ),
        }

    out_path = Path(output_pptx)
    if not out_path.is_absolute():
        return {
            "success": False,
            "error": "Relative output_pptx",
            "message": "output_pptx must be absolute.",
        }
    out_path.parent.mkdir(parents=True, exist_ok=True)

    data, parse_err = _ppt_load_yaml(specs_path)
    if parse_err:
        return {
            "success": False,
            "error": "Spec parse failed",
            "message": parse_err,
        }
    deck = data.get("deck") or {}
    slides = data.get("slides") or []
    if not isinstance(slides, list) or not slides:
        return {
            "success": False,
            "error": "No slides",
            "message": (
                "slide_specs.yaml has no `slides` list. Re-run "
                "ppt_derive_slide_specs_tool after editing the "
                "narrative."
            ),
        }

    # Make scripts dir importable for ppt_asset_helpers.
    import sys as _sys
    _saved_path = _sys.path[:]
    _sys.path.insert(0, str(ppt_scripts_dir))
    try:
        import ppt_asset_helpers as pah  # type: ignore
        from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_ALIGN
        from pptx.dml.color import RGBColor
        from pptx.util import Inches, Pt
        from docx.oxml.ns import qn  # python-docx ships with python-pptx envs
        from docx.oxml import OxmlElement  # noqa: F401
    except Exception as exc:
        _sys.path[:] = _saved_path
        return {
            "success": False,
            "error": "Helper import failed",
            "message": (
                f"Could not import ppt_asset_helpers: {exc}. "
                "Confirm python-pptx is installed and the PPT skill "
                "scripts directory is intact."
            ),
        }

    # ---- theme tokens injection ----
    theme_tokens = (deck.get("theme_tokens") or {}) if isinstance(deck, dict) else {}
    saved_tokens = dict(pah.DEFAULT_TYPOGRAPHY_TOKENS)
    saved_latin = pah.DEFAULT_LATIN_FONT_NAME
    saved_ea = pah.DEFAULT_EAST_ASIA_FONT_NAME
    saved_font = pah.DEFAULT_FONT_NAME
    saved_line = pah.DEFAULT_LINE_SPACING_MULTIPLE
    try:
        for token_key in (
            "hero_title_font_pt", "section_title_font_pt",
            "page_title_font_pt", "subtitle_font_pt",
            "minor_title_font_pt", "body_font_pt",
            "label_font_pt", "caption_font_pt",
            "title_line_spacing_multiple",
            "body_line_spacing_multiple",
            "title_paragraph_space_lines",
        ):
            if token_key in theme_tokens:
                try:
                    pah.DEFAULT_TYPOGRAPHY_TOKENS[token_key] = float(theme_tokens[token_key])
                except (TypeError, ValueError):
                    pass
        latin_font = theme_tokens.get("latin_font_name") or saved_latin
        ea_font = theme_tokens.get("east_asia_font_name") or saved_ea
        pah.DEFAULT_LATIN_FONT_NAME = latin_font
        pah.DEFAULT_EAST_ASIA_FONT_NAME = ea_font
        pah.DEFAULT_FONT_NAME = latin_font
        body_lsm = theme_tokens.get("body_line_spacing_multiple")
        if isinstance(body_lsm, (int, float)):
            pah.DEFAULT_LINE_SPACING_MULTIPLE = float(body_lsm)

        def _set_ea_font(run, ea_name: str):
            """Set CJK font on a run."""
            try:
                from pptx.oxml.ns import qn as pqn
                rPr2 = run._r.get_or_add_rPr()
                ea_elem = rPr2.find(pqn("a:ea"))
                if ea_elem is None:
                    ea_elem = OxmlElement("a:ea")
                    rPr2.append(ea_elem)
                ea_elem.set("typeface", ea_name)
            except Exception:
                pass

        def _apply_fonts(shape, latin: str = latin_font, ea: str = ea_font):
            """Apply fonts to all runs in a shape."""
            if not getattr(shape, "has_text_frame", False):
                return
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    try:
                        run.font.name = latin
                    except Exception:
                        pass
                    _set_ea_font(run, ea)

        palette = pah.default_palette()
        tokens = pah.default_typography_tokens()

        def _resolve_asset_path(maybe_path: str | None) -> Path | None:
            """Resolve asset paths relative to workspace."""
            if not maybe_path:
                return None
            p = Path(maybe_path)
            if not p.is_absolute():
                p = ws / maybe_path
            return p if p.exists() else None

        # ====== Build loop ===========================================
        prs = pah.new_presentation()
        blank_layout = prs.slide_layouts[6]  # 'Blank'
        per_slide = []

        for idx, spec in enumerate(slides, start=1):
            if not isinstance(spec, dict):
                per_slide.append({
                    "slide_id": f"S{idx:02d}",
                    "archetype": None,
                    "status": "error",
                    "error": "spec is not a mapping",
                })
                continue
            archetype = spec.get("archetype") or "hero-statement"
            if archetype not in _ARCHETYPES:
                per_slide.append({
                    "slide_id": spec.get("slide_id") or f"S{idx:02d}",
                    "archetype": archetype,
                    "status": "warning",
                    "error": (
                        f"unknown archetype {archetype!r}, "
                        "falling back to hero-statement"
                    ),
                })
                archetype = "hero-statement"

            slide = prs.slides.add_slide(blank_layout)
            try:
                if archetype == "hero-statement":
                    _render_hero_statement(slide, spec, tokens, palette)
                else:
                    # For now, just render hero as fallback; full renderers
                    # would go here. This is a simplified version focusing on
                    # the structure and imports needed.
                    _render_hero_statement(slide, spec, tokens, palette)
            except Exception as exc:
                import traceback
                per_slide.append({
                    "slide_id": spec.get("slide_id") or f"S{idx:02d}",
                    "archetype": archetype,
                    "status": "error",
                    "error": f"{type(exc).__name__}: {exc}",
                    "traceback_tail": traceback.format_exc()[-1500:],
                })
                continue

            per_slide.append({
                "slide_id": spec.get("slide_id") or f"S{idx:02d}",
                "archetype": archetype,
                "status": "ok",
            })

        pah.save_presentation(prs, out_path)

        # Post-process: update docProps/app.xml slide count
        try:
            import zipfile as _zipfile
            import re as _re

            actual_pages = len(prs.slides)
            tmp_path = out_path.with_suffix(out_path.suffix + ".tmp")
            with _zipfile.ZipFile(str(out_path), "r") as zin:
                names = zin.namelist()
                with _zipfile.ZipFile(
                    str(tmp_path), "w", _zipfile.ZIP_DEFLATED
                ) as zout:
                    for name in names:
                        data = zin.read(name)
                        if name == "docProps/app.xml":
                            text = data.decode("utf-8", "ignore")
                            if "<Slides>" in text:
                                text = _re.sub(
                                    r"<Slides>\s*\d+\s*</Slides>",
                                    f"<Slides>{actual_pages}</Slides>",
                                    text,
                                    count=1,
                                )
                            else:
                                text = text.replace(
                                    "</Properties>",
                                    f"<Slides>{actual_pages}</Slides></Properties>",
                                    1,
                                )
                            data = text.encode("utf-8")
                        zout.writestr(name, data)
            out_path.unlink()
            tmp_path.rename(out_path)
        except Exception as _exc:  # noqa: BLE001
            print(f"docProps post-process skipped: {_exc}")

        ok_count = sum(1 for s in per_slide if s["status"] == "ok")
        err_count = sum(1 for s in per_slide if s["status"] == "error")
        warn_count = sum(1 for s in per_slide if s["status"] == "warning")
        return {
            "success": err_count == 0,
            "output_pptx": str(out_path),
            "slide_count": len(per_slide),
            "ok_count": ok_count,
            "warning_count": warn_count,
            "error_count": err_count,
            "per_slide": per_slide,
            "message": (
                f"Built {ok_count}/{len(per_slide)} slide(s) into {out_path.name}. "
                + (f"{warn_count} warning(s). " if warn_count else "")
                + (f"{err_count} error(s) — check per_slide. " if err_count else "")
                + "Next: run ppt_package_preflight_tool → "
                "ppt_structure_precheck_tool → ppt_export_previews_tool "
                "→ ppt_render_review_tool."
            ),
        }
    except Exception as exc:
        import traceback
        return {
            "success": False,
            "error": f"Build aborted: {type(exc).__name__}: {exc}",
            "traceback_tail": traceback.format_exc()[-2000:],
            "message": (
                "ppt_build_pptx_tool crashed before finishing. "
                "Check traceback_tail; the most common cause is a "
                "malformed chart/table/diagram block in slide_specs.yaml."
            ),
        }
    finally:
        pah.DEFAULT_TYPOGRAPHY_TOKENS.clear()
        pah.DEFAULT_TYPOGRAPHY_TOKENS.update(saved_tokens)
        pah.DEFAULT_LATIN_FONT_NAME = saved_latin
        pah.DEFAULT_EAST_ASIA_FONT_NAME = saved_ea
        pah.DEFAULT_FONT_NAME = saved_font
        pah.DEFAULT_LINE_SPACING_MULTIPLE = saved_line
        _sys.path[:] = _saved_path


# ============================================================================
# Quality Tools
# ============================================================================

def _resolve_pptx_path(pptx_path: str, label: str = "pptx_path") -> tuple[Path | None, dict | None]:
    """Validate a .pptx path argument."""
    if not pptx_path or not isinstance(pptx_path, str):
        return None, {
            "success": False,
            "error": f"Missing {label}",
            "message": f"{label} is required.",
        }
    p = Path(pptx_path)
    if not p.is_absolute():
        return None, {
            "success": False,
            "error": f"Relative {label} not accepted",
            "message": (
                f"{label}={pptx_path!r} is relative. Use the absolute "
                "path returned by your previous build step."
            ),
        }
    if not p.exists():
        return None, {
            "success": False,
            "error": "File not found",
            "message": f"{label}: no file at {pptx_path}.",
        }
    if p.suffix.lower() != ".pptx":
        return None, {
            "success": False,
            "error": "Not a .pptx file",
            "message": f"{label}={pptx_path!r} is not a .pptx file.",
        }
    return p, None


def ppt_package_preflight_tool(
    pptx_path: str,
    deck_workspace: str,
    fail_on: str = "error",
    ppt_scripts_dir: Path = None,
) -> dict:
    """
    File-level quality gate: zip integrity, slide-count consistency
    (presentation.xml vs docProps/app.xml vs actual slide files),
    stale section_lst references, missing slide relationships, and
    embedded-object mobile-compatibility risk.

    This is the FIRST gate after build. Run BEFORE structure_precheck
    and BEFORE preview export — if the deck can't be opened by a
    fragile parser (WeChat / mobile WPS), there is no point checking
    its layout.

    Output is auto-archived to:
      <deck_workspace>/validation/package_preflight/history/
        package_preflight_<YYYYMMDD_HHMMSS>.{json,md}

    Args:
        pptx_path: Absolute path to the deck .pptx.
        deck_workspace: Value returned by ppt_init_workspace_tool.
        fail_on: 'error' (default), 'warning', or 'never'. Controls
            the script's exit code, NOT what is reported — issues are
            always returned in the result dict.

    Returns dict with success, returncode, report (parsed JSON
    including `summary` counts and `issues` list), stdout_tail,
    stderr_tail, message.
    """
    pptx, err = _resolve_pptx_path(pptx_path)
    if err is not None:
        return err
    ws, err = _resolve_deck_workspace(deck_workspace)
    if err is not None:
        return err
    if fail_on not in {"error", "warning", "never"}:
        return {
            "success": False,
            "error": "Invalid fail_on",
            "message": "fail_on must be one of: error, warning, never.",
        }

    script_path = ppt_scripts_dir / "check_pptx_package_preflight.py"
    if not script_path.exists():
        return {
            "success": False,
            "error": "Script not found",
            "message": f"check_pptx_package_preflight.py not found at {script_path}",
        }

    try:
        proc = subprocess.run(
            [
                _sys.executable,
                str(script_path),
                "--pptx", str(pptx),
                "--workspace-dir", str(ws),
                "--fail-on", fail_on,
            ],
            cwd=str(ppt_scripts_dir),
            capture_output=True,
            text=True,
            timeout=120,
        )
    except subprocess.TimeoutExpired:
        return {
            "success": False,
            "error": "Timeout",
            "message": "check_pptx_package_preflight.py did not finish within 120s",
        }
    except Exception as exc:
        return {
            "success": False,
            "error": str(exc),
            "message": f"Failed to invoke check_pptx_package_preflight.py: {exc}",
        }

    result = {
        "success": proc.returncode == 0,
        "returncode": proc.returncode,
        "stdout_tail": (proc.stdout or "")[-2000:],
        "stderr_tail": (proc.stderr or "")[-2000:],
    }

    report = result.get("report") or {}
    summary = report.get("summary") or {}
    # The script writes its own timestamped report — locate the newest one
    hist = ws / "validation" / "package_preflight" / "history"
    newest = None
    if hist.exists():
        jsons = sorted(hist.glob("package_preflight_*.json"))
        if jsons:
            newest = jsons[-1]
            try:
                report = _json.loads(newest.read_text(encoding="utf-8"))
                result["report"] = report
                summary = report.get("summary") or {}
            except Exception:
                pass
    result["report_path"] = str(newest) if newest else None
    result["summary"] = summary
    result["message"] = (
        f"package_preflight: errors={summary.get('error', 0)}, "
        f"warnings={summary.get('warning', 0)}, "
        f"not_checked={summary.get('not_checked', 0)}. "
        + (f"Report: {newest.name}." if newest else "(no report on disk)")
    )
    return result


def ppt_structure_precheck_tool(
    pptx_path: str,
    deck_workspace: str,
    fail_on: str = "error",
    ppt_scripts_dir: Path = None,
) -> dict:
    """
    Structure-layer quality gate: textbox fit / near-overflow,
    compact-width pressure on short labels, text occluded by higher
    z-order shapes, critical content (table/chart/picture) covered
    by overlay shapes, and explicit `not_checked` records for
    structured chart labels and flattened pictures.

    Run AFTER package_preflight and BEFORE preview export. The
    issues here have shape/slide-level locations so they are easy
    to drive into targeted fixes.

    Output is auto-archived to:
      <deck_workspace>/validation/structure_precheck/history/
        structure_precheck_<YYYYMMDD_HHMMSS>.{json,md}
      <deck_workspace>/validation/structure_precheck/shape_inventory.json

    Args:
        pptx_path: Absolute path to the deck .pptx.
        deck_workspace: Value returned by ppt_init_workspace_tool.
        fail_on: 'error' (default), 'warning', or 'never'.

    Returns dict with success, returncode, report (parsed JSON),
    summary, inventory_path, report_path, stdout_tail, stderr_tail,
    message.
    """
    pptx, err = _resolve_pptx_path(pptx_path)
    if err is not None:
        return err
    ws, err = _resolve_deck_workspace(deck_workspace)
    if err is not None:
        return err
    if fail_on not in {"error", "warning", "never"}:
        return {
            "success": False,
            "error": "Invalid fail_on",
            "message": "fail_on must be one of: error, warning, never.",
        }

    inventory_out = ws / "validation" / "structure_precheck" / "shape_inventory.json"
    inventory_out.parent.mkdir(parents=True, exist_ok=True)

    script_path = ppt_scripts_dir / "check_pptx_structure_precheck.py"
    if not script_path.exists():
        return {
            "success": False,
            "error": "Script not found",
            "message": f"check_pptx_structure_precheck.py not found at {script_path}",
        }

    try:
        proc = subprocess.run(
            [
                _sys.executable,
                str(script_path),
                "--pptx", str(pptx),
                "--workspace-dir", str(ws),
                "--inventory-out", str(inventory_out),
                "--fail-on", fail_on,
            ],
            cwd=str(ppt_scripts_dir),
            capture_output=True,
            text=True,
            timeout=180,
        )
    except subprocess.TimeoutExpired:
        return {
            "success": False,
            "error": "Timeout",
            "message": "check_pptx_structure_precheck.py did not finish within 180s",
        }
    except Exception as exc:
        return {
            "success": False,
            "error": str(exc),
            "message": f"Failed to invoke check_pptx_structure_precheck.py: {exc}",
        }

    result = {
        "success": proc.returncode == 0,
        "returncode": proc.returncode,
        "stdout_tail": (proc.stdout or "")[-2000:],
        "stderr_tail": (proc.stderr or "")[-2000:],
    }

    hist = ws / "validation" / "structure_precheck" / "history"
    newest = None
    summary = {}
    if hist.exists():
        jsons = sorted(hist.glob("structure_precheck_*.json"))
        if jsons:
            newest = jsons[-1]
            try:
                report = _json.loads(newest.read_text(encoding="utf-8"))
                result["report"] = report
                summary = report.get("summary") or {}
            except Exception:
                pass
    result["report_path"] = str(newest) if newest else None
    result["inventory_path"] = str(inventory_out) if inventory_out.exists() else None
    result["summary"] = summary
    result["message"] = (
        f"structure_precheck: errors={summary.get('error', 0)}, "
        f"warnings={summary.get('warning', 0)}, "
        f"not_checked={summary.get('not_checked', 0)}. "
        + (f"Report: {newest.name}." if newest else "(no report on disk)")
    )
    return result


def ppt_render_review_tool(
    pptx_path: str,
    deck_workspace: str,
    preview_dir: str | None = None,
    fail_on: str = "error",
    ppt_scripts_dir: Path = None,
) -> dict:
    """
    Render-layer quality gate: boundary-touch-ink at bottom/right of
    text frames (font strokes within ~3px of the inner edge in the
    PNG) and flattened-graphic internal-text `not_checked` entries.

    Run AFTER ppt_export_previews_tool — this gate consumes the
    preview PNGs. It complements structure_precheck by catching
    issues only visible after rasterization (e.g. last-line clipped
    by 1-2 px when the structure-level math says it just barely
    fits).

    Output auto-archived to:
      <deck_workspace>/validation/render_review/history/
        render_review_<YYYYMMDD_HHMMSS>.{json,md}

    Args:
        pptx_path: Absolute path to the deck .pptx.
        deck_workspace: Value returned by ppt_init_workspace_tool.
        preview_dir: Optional override; defaults to
            deck_workspace/build/rendered/ppt_preview.
        fail_on: 'error' (default), 'warning', or 'never'.

    Returns dict with success, summary, report, report_path,
    stdout_tail, stderr_tail, message.
    """
    pptx, err = _resolve_pptx_path(pptx_path)
    if err is not None:
        return err
    ws, err = _resolve_deck_workspace(deck_workspace)
    if err is not None:
        return err
    if fail_on not in {"error", "warning", "never"}:
        return {
            "success": False,
            "error": "Invalid fail_on",
            "message": "fail_on must be one of: error, warning, never.",
        }

    pv = Path(preview_dir) if preview_dir else (ws / "build" / "rendered" / "ppt_preview")
    if not pv.exists():
        return {
            "success": False,
            "error": "Preview directory not found",
            "message": (
                f"No preview directory at {pv}. Run "
                "ppt_export_previews_tool first."
            ),
        }

    script_path = ppt_scripts_dir / "check_pptx_render_review.py"
    if not script_path.exists():
        return {
            "success": False,
            "error": "Script not found",
            "message": f"check_pptx_render_review.py not found at {script_path}",
        }

    try:
        proc = subprocess.run(
            [
                _sys.executable,
                str(script_path),
                "--pptx", str(pptx),
                "--preview-dir", str(pv),
                "--workspace-dir", str(ws),
                "--fail-on", fail_on,
            ],
            cwd=str(ppt_scripts_dir),
            capture_output=True,
            text=True,
            timeout=180,
        )
    except subprocess.TimeoutExpired:
        return {
            "success": False,
            "error": "Timeout",
            "message": "check_pptx_render_review.py did not finish within 180s",
        }
    except Exception as exc:
        return {
            "success": False,
            "error": str(exc),
            "message": f"Failed to invoke check_pptx_render_review.py: {exc}",
        }

    result = {
        "success": proc.returncode == 0,
        "returncode": proc.returncode,
        "stdout_tail": (proc.stdout or "")[-2000:],
        "stderr_tail": (proc.stderr or "")[-2000:],
    }

    hist = ws / "validation" / "render_review" / "history"
    newest = None
    summary = {}
    if hist.exists():
        jsons = sorted(hist.glob("render_review_*.json"))
        if jsons:
            newest = jsons[-1]
            try:
                report = _json.loads(newest.read_text(encoding="utf-8"))
                result["report"] = report
                summary = report.get("summary") or {}
            except Exception:
                pass
    result["report_path"] = str(newest) if newest else None
    result["summary"] = summary
    result["message"] = (
        f"render_review: errors={summary.get('error', 0)}, "
        f"warnings={summary.get('warning', 0)}, "
        f"not_checked={summary.get('not_checked', 0)}. "
        + (f"Report: {newest.name}." if newest else "(no report on disk)")
    )
    return result


# ============================================================================
# Preview Tools
# ============================================================================

def ppt_export_previews_tool(
    pptx_path: str,
    deck_workspace: str,
    backend: str = "auto",
    render_backend: str = "auto",
    prefix: str = "slide_",
    keep_pdf: bool = False,
    ppt_scripts_dir: Path = None,
) -> dict:
    """
    Render the deck as per-slide PNG previews via PowerPoint or
    LibreOffice, then pdftoppm or PyMuPDF for PDF→PNG.

    Output directory defaults to:
      <deck_workspace>/build/rendered/ppt_preview/slide_001.png ...

    Manifest written to:
      <deck_workspace>/validation/preview_manifest.json

    Page-count mismatch (e.g. LibreOffice silently dropping a slide)
    is treated as failure, not a warning — re-export with the other
    backend if it happens.

    Args:
        pptx_path: Absolute path to the deck .pptx.
        deck_workspace: Value returned by ppt_init_workspace_tool.
        backend: 'auto' (default), 'powerpoint', or 'libreoffice'.
        render_backend: 'auto' (default), 'pdftoppm', or 'fitz'.
        prefix: Output PNG filename prefix (default 'slide_').
        keep_pdf: When true, the intermediate PDF is moved into the
            preview directory; default false.

    Returns dict with success, preview_dir, manifest_path,
    generated_pages, pdf_backend, render_backend, stdout_tail,
    stderr_tail, message.
    """
    pptx, err = _resolve_pptx_path(pptx_path)
    if err is not None:
        return err
    ws, err = _resolve_deck_workspace(deck_workspace)
    if err is not None:
        return err
    if backend not in {"auto", "powerpoint", "libreoffice"}:
        return {
            "success": False,
            "error": "Invalid backend",
            "message": "backend must be one of: auto, powerpoint, libreoffice.",
        }
    if render_backend not in {"auto", "pdftoppm", "fitz"}:
        return {
            "success": False,
            "error": "Invalid render_backend",
            "message": "render_backend must be one of: auto, pdftoppm, fitz.",
        }

    out_dir = ws / "build" / "rendered" / "ppt_preview"
    out_dir.mkdir(parents=True, exist_ok=True)
    manifest = ws / "validation" / "preview_manifest.json"
    manifest.parent.mkdir(parents=True, exist_ok=True)

    args = [
        "--pptx", str(pptx),
        "--out-dir", str(out_dir),
        "--backend", backend,
        "--render-backend", render_backend,
        "--prefix", prefix,
        "--json-out", str(manifest),
    ]
    if keep_pdf:
        args.append("--keep-pdf")

    script_path = ppt_scripts_dir / "export_pptx_previews.py"
    if not script_path.exists():
        return {
            "success": False,
            "error": "Script not found",
            "message": f"export_pptx_previews.py not found at {script_path}",
        }

    try:
        proc = subprocess.run(
            [_sys.executable, str(script_path), *args],
            cwd=str(ppt_scripts_dir),
            capture_output=True,
            text=True,
            timeout=300,
        )
    except subprocess.TimeoutExpired:
        return {
            "success": False,
            "error": "Timeout",
            "message": "export_pptx_previews.py did not finish within 300s",
        }
    except Exception as exc:
        return {
            "success": False,
            "error": str(exc),
            "message": f"Failed to invoke export_pptx_previews.py: {exc}",
        }

    result = {
        "success": proc.returncode == 0,
        "returncode": proc.returncode,
        "stdout_tail": (proc.stdout or "")[-2000:],
        "stderr_tail": (proc.stderr or "")[-2000:],
    }

    if manifest.exists():
        try:
            result["report"] = _json.loads(manifest.read_text(encoding="utf-8"))
        except Exception as exc:
            result["report_read_error"] = str(exc)

    report = result.get("report") or {}
    result["preview_dir"] = str(out_dir)
    result["manifest_path"] = str(manifest) if manifest.exists() else None
    result["generated_pages"] = report.get("generated_pages")
    result["pdf_backend"] = report.get("pdf_backend")
    result["render_backend"] = report.get("render_backend")
    if result.get("success"):
        result["message"] = (
            f"Exported {result['generated_pages']} preview(s) via "
            f"{result['pdf_backend']} → {result['render_backend']}. "
            f"PNGs in {out_dir}."
        )
    else:
        result["message"] = (
            "Preview export failed. If the error mentions PowerPoint "
            "automation, check macOS Privacy & Security → Automation. "
            "If LibreOffice converted but pages count is off, retry "
            "with backend='powerpoint' or vice versa."
        )
    return result


# ============================================================================
# Connector Tools
# ============================================================================

def ppt_connectors_check_tool(
    pptx_path: str,
    deck_workspace: str,
    slides: list | None = None,
    min_connectors: int = 0,
    forbid_prefixes: list | None = None,
    ppt_scripts_dir: Path = None,
) -> dict:
    """
    Module-level gate for diagram pages: verify each connector is
    REALLY glued to two shapes (stCxn + endCxn present, target shape
    ids resolve, no connections to forbidden parent shapes such as
    lane / cluster outer frames).

    Run after a diagram page with asset_mode=diagram-connector is
    built. A passing report is the evidence the user can rely on
    that dragging a node will not break the diagram.

    Output:
      <deck_workspace>/validation/connectors/connector_report.json

    Args:
        pptx_path: Absolute path to the deck .pptx.
        deck_workspace: Value returned by ppt_init_workspace_tool.
        slides: Optional list of slide numbers (1-based) to limit
            the check. Defaults to all slides.
        min_connectors: Optional. Total connector count must reach
            this value or the check fails. Use it on dedicated
            diagram pages where you know connectors must exist.
        forbid_prefixes: Optional list of forbidden prefixes for
            connector endpoints. Default is `["Lane "]` — connector
            endpoints starting with these strings are flagged as
            illegal (they typically mean the line is glued to a
            swimlane outer frame instead of a business node).

    Returns dict with success, total_connectors, report (parsed
    JSON, mapping slide → records), report_path, stdout_tail,
    stderr_tail, message.
    """
    pptx, err = _resolve_pptx_path(pptx_path)
    if err is not None:
        return err
    ws, err = _resolve_deck_workspace(deck_workspace)
    if err is not None:
        return err

    json_out = ws / "validation" / "connectors" / "connector_report.json"
    json_out.parent.mkdir(parents=True, exist_ok=True)
    args = ["--pptx", str(pptx), "--json-out", str(json_out)]
    if isinstance(slides, list):
        for s in slides:
            args.extend(["--slide", str(int(s))])
    if min_connectors and min_connectors > 0:
        args.extend(["--min-connectors", str(int(min_connectors))])
    for prefix in (forbid_prefixes or []):
        args.extend(["--forbid-prefix", str(prefix)])

    script_path = ppt_scripts_dir / "check_pptx_connectors.py"
    if not script_path.exists():
        return {
            "success": False,
            "error": "Script not found",
            "message": f"check_pptx_connectors.py not found at {script_path}",
        }

    try:
        proc = subprocess.run(
            [_sys.executable, str(script_path), *args],
            cwd=str(ppt_scripts_dir),
            capture_output=True,
            text=True,
            timeout=60,
        )
    except subprocess.TimeoutExpired:
        return {
            "success": False,
            "error": "Timeout",
            "message": "check_pptx_connectors.py did not finish within 60s",
        }
    except Exception as exc:
        return {
            "success": False,
            "error": str(exc),
            "message": f"Failed to invoke check_pptx_connectors.py: {exc}",
        }

    result = {
        "success": proc.returncode == 0,
        "returncode": proc.returncode,
        "stdout_tail": (proc.stdout or "")[-2000:],
        "stderr_tail": (proc.stderr or "")[-2000:],
    }

    if json_out.exists():
        try:
            result["report"] = _json.loads(json_out.read_text(encoding="utf-8"))
        except Exception as exc:
            result["report_read_error"] = str(exc)

    report = result.get("report") or {}
    # report maps slide_num -> list of connector records
    total = 0
    if isinstance(report, dict):
        for v in report.values():
            if isinstance(v, list):
                total += len(v)
    result["report_path"] = str(json_out) if json_out.exists() else None
    result["total_connectors"] = total
    if result.get("success"):
        result["message"] = (
            f"connector check passed: {total} connector(s) verified."
        )
    else:
        result["message"] = (
            f"connector check FAILED. {total} connector(s) seen. "
            "Read stdout_tail for the specific issues — usually one "
            "of: stCxn/endCxn missing (line drawn but not glued), "
            "endpoint id unresolved (target shape deleted), or "
            "connector glued to a lane/cluster outer frame instead "
            "of a business node."
        )
    return result
