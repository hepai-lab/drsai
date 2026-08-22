"""Structure-preserving document parsing for auditable retrieval.

Retrieval evidence has to point back at a place in the original file — a page,
a line range, a heading. Extraction that returns one flat string makes
citations file-level at best, which is not enough to open the source at the
cited spot, so every parser here emits ordered units that carry their own
position.

The second reason for this module is negative space: a file that cannot be
read must say so. Returning an empty string for a scanned PDF or an
unsupported format silently shrinks the corpus, and an answer of "the material
does not mention it" then becomes indistinguishable from "that material was
never read". `ParsedDocument.status` keeps that difference visible so the
caller can refuse to claim a complete corpus.
"""

from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
import re
from typing import Callable, Literal

LocatorKind = Literal["line", "page", "heading", "slide", "sheet"]

ParseStatus = Literal[
    "ok",
    "empty",
    "truncated",
    "no_text_layer",
    "parser_unavailable",
    "unsupported_format",
    "failed",
]

TEXT_SUFFIXES = frozenset({
    ".txt", ".rst", ".log", ".ini", ".cfg", ".env",
    ".py", ".js", ".mjs", ".cjs", ".ts", ".tsx", ".jsx", ".java", ".go", ".rs",
    ".c", ".h", ".cc", ".cpp", ".hpp", ".cs", ".rb", ".php", ".swift", ".kt",
    ".sh", ".bash", ".ps1", ".sql", ".r", ".m",
    ".json", ".toml", ".yaml", ".yml", ".xml", ".html", ".htm", ".css",
})
MARKDOWN_SUFFIXES = frozenset({".md", ".markdown", ".mdx"})
TABULAR_SUFFIXES = frozenset({".csv", ".tsv"})

SUPPORTED_SUFFIXES = (
    TEXT_SUFFIXES
    | MARKDOWN_SUFFIXES
    | TABULAR_SUFFIXES
    | {".pdf", ".docx", ".pptx", ".xlsx"}
)

_MARKDOWN_HEADING = re.compile(r"^(#{1,6})\s+(.*\S)\s*$")
_MAX_UNITS = 200_000


@dataclass(frozen=True)
class DocumentLocator:
    """Where a unit sits in its source document.

    ``kind`` is the document's addressing scheme and stays constant for one
    file, so a caller can render "page 4" or "lines 12-18" without inspecting
    every field. Line numbers are additionally carried whenever they are known.
    """

    kind: LocatorKind
    line_start: int | None = None
    line_end: int | None = None
    page: int | None = None
    slide: int | None = None
    sheet: str | None = None
    heading_path: tuple[str, ...] = ()

    def label(self) -> str:
        """Compact, language-neutral form. Surfaces localize from `payload`."""
        if self.kind == "page":
            return f"p.{self.page}"
        if self.kind == "slide":
            return f"slide {self.slide}"
        if self.kind == "sheet":
            return self.sheet or "sheet"
        if self.kind == "heading":
            heading = " > ".join(self.heading_path)
            lines = self._line_label()
            if heading and lines:
                return f"{heading} ({lines})"
            return heading or lines or "document"
        return self._line_label() or "document"

    def _line_label(self) -> str:
        if self.line_start is None:
            return ""
        if self.line_end is None or self.line_end == self.line_start:
            return f"L{self.line_start}"
        return f"L{self.line_start}-{self.line_end}"

    def payload(self) -> dict[str, object]:
        data: dict[str, object] = {"kind": self.kind, "label": self.label()}
        if self.line_start is not None:
            data["line_start"] = self.line_start
            data["line_end"] = self.line_end if self.line_end is not None else self.line_start
        if self.page is not None:
            data["page"] = self.page
        if self.slide is not None:
            data["slide"] = self.slide
        if self.sheet is not None:
            data["sheet"] = self.sheet
        if self.heading_path:
            data["heading_path"] = list(self.heading_path)
        return data

    def merged_with(self, other: "DocumentLocator") -> "DocumentLocator":
        """Span two adjacent units, used when the indexer packs them together."""
        if other.kind != self.kind:
            return self
        starts = [value for value in (self.line_start, other.line_start) if value is not None]
        ends = [value for value in (self.line_end, other.line_end, self.line_start, other.line_start) if value is not None]
        return DocumentLocator(
            kind=self.kind,
            line_start=min(starts) if starts else None,
            line_end=max(ends) if ends else None,
            page=self.page,
            slide=self.slide,
            sheet=self.sheet,
            heading_path=self.heading_path,
        )


@dataclass(frozen=True)
class DocumentUnit:
    order: int
    text: str
    locator: DocumentLocator


@dataclass(frozen=True)
class ParsedDocument:
    source: str
    status: ParseStatus
    units: tuple[DocumentUnit, ...] = ()
    detail: str = ""

    @property
    def ok(self) -> bool:
        return self.status == "ok"

    @property
    def text(self) -> str:
        return "\n".join(unit.text for unit in self.units)


def parse_document(path: str | Path, *, source: str | None = None) -> ParsedDocument:
    """Parse one file into ordered, position-carrying units.

    Never raises for a bad or unreadable document: the failure is reported as a
    status so the caller can record it against corpus completeness.
    """

    target = Path(path)
    name = source if source is not None else target.name
    suffix = target.suffix.lower()
    parser = _PARSERS.get(suffix)
    if parser is None:
        return ParsedDocument(name, "unsupported_format", detail=suffix or "no-suffix")
    try:
        return parser(target, name)
    except _ParserUnavailable as exc:
        return ParsedDocument(name, "parser_unavailable", detail=str(exc))
    except Exception as exc:  # noqa: BLE001 - a bad document must not abort indexing
        return ParsedDocument(name, "failed", detail=type(exc).__name__)


class _ParserUnavailable(RuntimeError):
    """The optional library backing one format is not installed."""


def _finish(source: str, units: list[DocumentUnit], *, empty_status: ParseStatus = "empty") -> ParsedDocument:
    if not units:
        return ParsedDocument(source, empty_status)
    if len(units) > _MAX_UNITS:
        # Keeping the units we read but reporting "ok" would be the same defect
        # this module exists to prevent: the tail is missing from the index
        # while the corpus still claims to be complete, so a question about the
        # tail gets answered "not in the material".
        return ParsedDocument(
            source, "truncated", tuple(units[:_MAX_UNITS]), detail=f"{len(units)}_units_exceeds_{_MAX_UNITS}",
        )
    return ParsedDocument(source, "ok", tuple(units))


def _read_text(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="replace")


def _parse_lines(path: Path, source: str) -> ParsedDocument:
    units: list[DocumentUnit] = []
    for index, raw in enumerate(_read_text(path).splitlines(), start=1):
        text = raw.rstrip()
        if not text.strip():
            continue
        units.append(DocumentUnit(
            order=len(units),
            text=text,
            locator=DocumentLocator("line", line_start=index, line_end=index),
        ))
    return _finish(source, units)


def _parse_markdown(path: Path, source: str) -> ParsedDocument:
    units: list[DocumentUnit] = []
    heading_stack: list[tuple[int, str]] = []
    for index, raw in enumerate(_read_text(path).splitlines(), start=1):
        text = raw.rstrip()
        if not text.strip():
            continue
        heading = _MARKDOWN_HEADING.match(text)
        if heading:
            level = len(heading.group(1))
            while heading_stack and heading_stack[-1][0] >= level:
                heading_stack.pop()
            heading_stack.append((level, heading.group(2)))
        units.append(DocumentUnit(
            order=len(units),
            text=text,
            locator=DocumentLocator(
                "heading",
                line_start=index,
                line_end=index,
                heading_path=tuple(title for _level, title in heading_stack),
            ),
        ))
    return _finish(source, units)


def _parse_tabular(path: Path, source: str) -> ParsedDocument:
    units: list[DocumentUnit] = []
    for index, raw in enumerate(_read_text(path).splitlines(), start=1):
        text = raw.strip()
        if not text:
            continue
        units.append(DocumentUnit(
            order=len(units),
            text=text,
            locator=DocumentLocator("line", line_start=index, line_end=index),
        ))
    return _finish(source, units)


def _parse_pdf(path: Path, source: str) -> ParsedDocument:
    try:
        from pypdf import PdfReader
    except ImportError as exc:
        raise _ParserUnavailable("pypdf") from exc
    reader = PdfReader(str(path))
    units: list[DocumentUnit] = []
    for number, page in enumerate(reader.pages, start=1):
        text = (page.extract_text() or "").strip()
        if not text:
            continue
        units.append(DocumentUnit(
            order=len(units),
            text=text,
            locator=DocumentLocator("page", page=number),
        ))
    # A PDF with pages but no extractable text is a scan, not an empty
    # document. Indexing it as empty would quietly drop it from the corpus.
    empty_status: ParseStatus = "no_text_layer" if len(reader.pages) else "empty"
    return _finish(source, units, empty_status=empty_status)


def _parse_docx(path: Path, source: str) -> ParsedDocument:
    try:
        from docx import Document
    except ImportError as exc:
        raise _ParserUnavailable("python-docx") from exc
    units: list[DocumentUnit] = []
    heading_stack: list[tuple[int, str]] = []
    for paragraph in Document(str(path)).paragraphs:
        text = (paragraph.text or "").strip()
        if not text:
            continue
        level = _docx_heading_level(getattr(paragraph.style, "name", "") or "")
        if level:
            while heading_stack and heading_stack[-1][0] >= level:
                heading_stack.pop()
            heading_stack.append((level, text))
        units.append(DocumentUnit(
            order=len(units),
            text=text,
            locator=DocumentLocator(
                "heading",
                heading_path=tuple(title for _level, title in heading_stack),
            ),
        ))
    return _finish(source, units)


def _docx_heading_level(style_name: str) -> int:
    match = re.match(r"^Heading\s+(\d)$", style_name.strip(), re.IGNORECASE)
    if match:
        return int(match.group(1))
    return 1 if style_name.strip().casefold() == "title" else 0


def _parse_pptx(path: Path, source: str) -> ParsedDocument:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise _ParserUnavailable("python-pptx") from exc
    units: list[DocumentUnit] = []
    for number, slide in enumerate(Presentation(str(path)).slides, start=1):
        lines = [
            (shape.text_frame.text or "").strip()
            for shape in slide.shapes
            if getattr(shape, "has_text_frame", False)
        ]
        text = "\n".join(line for line in lines if line).strip()
        if not text:
            continue
        units.append(DocumentUnit(
            order=len(units),
            text=text,
            locator=DocumentLocator("slide", slide=number),
        ))
    return _finish(source, units)


def _parse_xlsx(path: Path, source: str) -> ParsedDocument:
    try:
        from openpyxl import load_workbook
    except ImportError as exc:
        raise _ParserUnavailable("openpyxl") from exc
    workbook = load_workbook(str(path), read_only=True, data_only=True)
    units: list[DocumentUnit] = []
    try:
        for worksheet in workbook.worksheets:
            for index, row in enumerate(worksheet.iter_rows(values_only=True), start=1):
                cells = [str(value).strip() for value in row if value is not None and str(value).strip()]
                if not cells:
                    continue
                units.append(DocumentUnit(
                    order=len(units),
                    text=" | ".join(cells),
                    locator=DocumentLocator(
                        "sheet",
                        sheet=worksheet.title,
                        line_start=index,
                        line_end=index,
                    ),
                ))
    finally:
        workbook.close()
    return _finish(source, units)


_PARSERS: dict[str, Callable[[Path, str], ParsedDocument]] = {
    **{suffix: _parse_markdown for suffix in MARKDOWN_SUFFIXES},
    **{suffix: _parse_tabular for suffix in TABULAR_SUFFIXES},
    **{suffix: _parse_lines for suffix in TEXT_SUFFIXES},
    ".pdf": _parse_pdf,
    ".docx": _parse_docx,
    ".pptx": _parse_pptx,
    ".xlsx": _parse_xlsx,
}
