from __future__ import annotations

from drsai.content.documents import DocumentLocator, parse_document


def test_markdown_units_carry_heading_path_and_line_numbers(tmp_path) -> None:
    document = tmp_path / "runtime.md"
    document.write_text(
        "\n".join([
            "# OpenDrSai Runtime",
            "",
            "Runtime hosts Sessions and Runs.",
            "",
            "## Session",
            "",
            "A Session is one continuous user conversation.",
            "",
            "### Replay",
            "",
            "Replay always creates a new Run.",
        ]),
        encoding="utf-8",
    )

    parsed = parse_document(document)

    assert parsed.ok
    replay = next(unit for unit in parsed.units if unit.text.startswith("Replay always"))
    assert replay.locator.kind == "heading"
    assert replay.locator.heading_path == ("OpenDrSai Runtime", "Session", "Replay")
    assert replay.locator.line_start == 11
    assert replay.locator.label() == "OpenDrSai Runtime > Session > Replay (L11)"

    # A deeper heading must not leak into a later sibling section.
    session = next(unit for unit in parsed.units if unit.text.startswith("A Session is"))
    assert session.locator.heading_path == ("OpenDrSai Runtime", "Session")


def test_plain_text_units_carry_line_numbers(tmp_path) -> None:
    document = tmp_path / "notes.txt"
    document.write_text("alpha\n\nbeta\ngamma\n", encoding="utf-8")

    parsed = parse_document(document)

    assert [unit.text for unit in parsed.units] == ["alpha", "beta", "gamma"]
    assert [unit.locator.line_start for unit in parsed.units] == [1, 3, 4]
    assert parsed.units[0].locator.kind == "line"
    assert parsed.units[2].locator.label() == "L4"


def test_docx_units_carry_heading_path(tmp_path) -> None:
    from docx import Document

    document = Document()
    document.add_heading("OpenDrSai Runtime", level=1)
    document.add_paragraph("Runtime hosts Sessions and Runs.")
    document.add_heading("Replay", level=2)
    document.add_paragraph("Replay always creates a new Run.")
    target = tmp_path / "runtime.docx"
    document.save(str(target))

    parsed = parse_document(target)

    assert parsed.ok
    replay = next(unit for unit in parsed.units if unit.text.startswith("Replay always"))
    assert replay.locator.kind == "heading"
    assert replay.locator.heading_path == ("OpenDrSai Runtime", "Replay")

    intro = next(unit for unit in parsed.units if unit.text.startswith("Runtime hosts"))
    assert intro.locator.heading_path == ("OpenDrSai Runtime",)


def test_pptx_units_carry_slide_numbers(tmp_path) -> None:
    from pptx import Presentation

    presentation = Presentation()
    blank = presentation.slide_layouts[5]
    first = presentation.slides.add_slide(blank)
    first.shapes.title.text = "Runtime overview"
    second = presentation.slides.add_slide(blank)
    second.shapes.title.text = "Replay rules"
    target = tmp_path / "deck.pptx"
    presentation.save(str(target))

    parsed = parse_document(target)

    assert parsed.ok
    assert [unit.locator.slide for unit in parsed.units] == [1, 2]
    assert parsed.units[1].locator.label() == "slide 2"
    assert "Replay rules" in parsed.units[1].text


def test_xlsx_units_carry_sheet_and_row(tmp_path) -> None:
    from openpyxl import Workbook

    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Metrics"
    sheet.append(["metric", "value"])
    sheet.append(["bandwidth", "40 Tbps"])
    target = tmp_path / "metrics.xlsx"
    workbook.save(str(target))

    parsed = parse_document(target)

    assert parsed.ok
    bandwidth = next(unit for unit in parsed.units if "bandwidth" in unit.text)
    assert bandwidth.locator.kind == "sheet"
    assert bandwidth.locator.sheet == "Metrics"
    assert bandwidth.locator.line_start == 2
    assert bandwidth.locator.label() == "Metrics"


def test_oversized_document_is_reported_truncated_not_ok(tmp_path, monkeypatch) -> None:
    import drsai.content.documents as documents

    monkeypatch.setattr(documents, "_MAX_UNITS", 3)
    target = tmp_path / "long.txt"
    target.write_text("\n".join(f"line {index}" for index in range(10)), encoding="utf-8")

    parsed = documents.parse_document(target)

    # Reporting "ok" here would repeat the defect this module exists to
    # prevent: the tail is missing from the index while the corpus still
    # claims to be complete, so a question about the tail is answered
    # "not in the material".
    assert parsed.status == "truncated"
    assert not parsed.ok
    assert len(parsed.units) == 3
    assert "10_units" in parsed.detail


def test_unsupported_format_is_reported_not_silently_empty(tmp_path) -> None:
    document = tmp_path / "archive.zip"
    document.write_bytes(b"PK\x03\x04binary")

    parsed = parse_document(document)

    assert parsed.status == "unsupported_format"
    assert parsed.units == ()
    assert not parsed.ok


def test_missing_optional_parser_is_reported_as_unavailable(tmp_path, monkeypatch) -> None:
    document = tmp_path / "report.docx"
    document.write_bytes(b"PK\x03\x04")
    monkeypatch.setitem(__import__("sys").modules, "docx", None)

    parsed = parse_document(document)

    # `None` in sys.modules makes the import fail the way a missing package
    # does; the file must surface as unreadable rather than as empty content.
    assert parsed.status in {"parser_unavailable", "failed"}
    assert parsed.units == ()


def test_unreadable_document_does_not_raise(tmp_path) -> None:
    parsed = parse_document(tmp_path / "absent.md")

    assert parsed.status == "failed"
    assert parsed.detail == "FileNotFoundError"


def test_locator_merge_spans_adjacent_lines() -> None:
    first = DocumentLocator("line", line_start=4, line_end=4)
    second = DocumentLocator("line", line_start=5, line_end=7)

    merged = first.merged_with(second)

    assert merged.line_start == 4
    assert merged.line_end == 7
    assert merged.label() == "L4-7"


def test_locator_payload_is_serializable_and_typed() -> None:
    locator = DocumentLocator("page", page=4)

    assert locator.payload() == {"kind": "page", "label": "p.4", "page": 4}
